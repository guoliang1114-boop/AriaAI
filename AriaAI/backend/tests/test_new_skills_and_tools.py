"""Tests for edit_project_office_document tool."""
from __future__ import annotations

import asyncio
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from sqlmodel import Session, SQLModel

from app.models.db import Project, ProjectFile, ProjectFolder
from app.tools import registry
from app.tools import file_generators, office_documents
from app.tools.office_documents import EDIT_PROJECT_OFFICE_DOCUMENT_TOOL_NAME
from tests.test_database import create_test_engine, drop_all_tables


class EditOfficeDocumentToolTestCase(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.uploads_dir = Path(self.temp_dir.name) / "uploads"
        self.generated_dir = self.uploads_dir / "generated"
        self.generated_dir.mkdir(parents=True, exist_ok=True)
        self.engine = create_test_engine()
        drop_all_tables(self.engine)
        SQLModel.metadata.create_all(self.engine)

    def tearDown(self):
        self.engine.dispose()
        self.temp_dir.cleanup()

    def _create_project(self) -> int:
        with Session(self.engine) as session:
            project = Project(name="Test Project", client="Test Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            return project.id

    def _create_project_file(self, project_id: int, name: str, file_type: str, content_bytes: bytes) -> int:
        project_dir = self.uploads_dir / "projects" / str(project_id)
        project_dir.mkdir(parents=True, exist_ok=True)
        file_path = project_dir / name
        file_path.write_bytes(content_bytes)
        with Session(self.engine) as session:
            pf = ProjectFile(
                project_id=project_id,
                name=name,
                file_type=file_type,
                path=str(file_path.relative_to(self.uploads_dir)),
                size_bytes=len(content_bytes),
                summary="Test file",
            )
            session.add(pf)
            session.commit()
            session.refresh(pf)
            return pf.id

    def _generate_test_pptx(self) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches
        import io

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Title"
        body = slide.placeholders[1]
        body.text = "Point 1\nPoint 2\nPoint 3"

        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Second Slide"
        body2 = slide2.placeholders[1]
        body2.text = "Content A\nContent B"

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _generate_test_docx(self) -> bytes:
        from docx import Document
        import io

        doc = Document()
        doc.add_heading("Chapter 1", level=1)
        doc.add_paragraph("This is the first chapter content.")
        doc.add_heading("Chapter 2", level=1)
        doc.add_paragraph("This is the second chapter content.")

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def _generate_test_xlsx(self) -> bytes:
        from openpyxl import Workbook
        import io

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Item A"
        ws["B2"] = 100
        ws["A3"] = "Item B"
        ws["B3"] = 200

        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    def test_tool_is_registered(self):
        tool = registry.get(EDIT_PROJECT_OFFICE_DOCUMENT_TOOL_NAME)
        self.assertIsNotNone(tool)
        self.assertEqual(tool.name, EDIT_PROJECT_OFFICE_DOCUMENT_TOOL_NAME)

    def test_edit_pptx_update_slide(self):
        project_id = self._create_project()
        pptx_bytes = self._generate_test_pptx()
        file_id = self._create_project_file(project_id, "test.pptx", "pptx", pptx_bytes)

        with patch("app.tools.office_documents.engine", self.engine), \
             patch("app.tools.office_documents.UPLOADS_DIR", self.uploads_dir):
            result = asyncio.get_event_loop().run_until_complete(
                office_documents.edit_project_office_document(
                    project_id=project_id,
                    file_id=file_id,
                    edits=[{"action": "update_slide", "slide_index": 0, "title": "New Title"}],
                )
            )

        self.assertTrue(result.get("ok"))
        self.assertIn("changes", result)

    def test_edit_pptx_update_text(self):
        project_id = self._create_project()
        pptx_bytes = self._generate_test_pptx()
        file_id = self._create_project_file(project_id, "test.pptx", "pptx", pptx_bytes)

        with patch("app.tools.office_documents.engine", self.engine), \
             patch("app.tools.office_documents.UPLOADS_DIR", self.uploads_dir):
            result = asyncio.get_event_loop().run_until_complete(
                office_documents.edit_project_office_document(
                    project_id=project_id,
                    file_id=file_id,
                    edits=[{"action": "update_text", "slide_index": 0, "old_text": "Point 1", "new_text": "Updated Point"}],
                )
            )

        self.assertTrue(result.get("ok"))

    def test_edit_docx_update_text(self):
        project_id = self._create_project()
        docx_bytes = self._generate_test_docx()
        file_id = self._create_project_file(project_id, "test.docx", "docx", docx_bytes)

        with patch("app.tools.office_documents.engine", self.engine), \
             patch("app.tools.office_documents.UPLOADS_DIR", self.uploads_dir):
            result = asyncio.get_event_loop().run_until_complete(
                office_documents.edit_project_office_document(
                    project_id=project_id,
                    file_id=file_id,
                    edits=[{"action": "update_text", "old_text": "first chapter", "new_text": "updated chapter"}],
                )
            )

        self.assertTrue(result.get("ok"))

    def test_edit_xlsx_update_cell(self):
        project_id = self._create_project()
        xlsx_bytes = self._generate_test_xlsx()
        file_id = self._create_project_file(project_id, "test.xlsx", "xlsx", xlsx_bytes)

        with patch("app.tools.office_documents.engine", self.engine), \
             patch("app.tools.office_documents.UPLOADS_DIR", self.uploads_dir):
            result = asyncio.get_event_loop().run_until_complete(
                office_documents.edit_project_office_document(
                    project_id=project_id,
                    file_id=file_id,
                    edits=[{"action": "update_cell", "sheet": "Sheet1", "cell": "B2", "value": 999}],
                )
            )

        self.assertTrue(result.get("ok"))

    def test_edit_creates_copy_with_output_name(self):
        project_id = self._create_project()
        pptx_bytes = self._generate_test_pptx()
        file_id = self._create_project_file(project_id, "test.pptx", "pptx", pptx_bytes)

        with patch("app.tools.office_documents.engine", self.engine), \
             patch("app.tools.office_documents.UPLOADS_DIR", self.uploads_dir):
            result = asyncio.get_event_loop().run_until_complete(
                office_documents.edit_project_office_document(
                    project_id=project_id,
                    file_id=file_id,
                    edits=[{"action": "update_slide", "slide_index": 0, "title": "Copy Title"}],
                    output_name="modified.pptx",
                )
            )

        self.assertTrue(result.get("ok"))
        self.assertEqual(result.get("name"), "modified.pptx")

    def test_edit_rejects_invalid_file_type(self):
        project_id = self._create_project()
        file_id = self._create_project_file(project_id, "test.txt", "txt", b"hello")

        with patch("app.tools.office_documents.engine", self.engine), \
             patch("app.tools.office_documents.UPLOADS_DIR", self.uploads_dir):
            with self.assertRaises(Exception):
                asyncio.get_event_loop().run_until_complete(
                    office_documents.edit_project_office_document(
                        project_id=project_id,
                        file_id=file_id,
                        edits=[{"action": "update_text", "old_text": "a", "new_text": "b"}],
                    )
                )


class ManagePdfToolTestCase(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.uploads_dir = Path(self.temp_dir.name) / "uploads"
        self.generated_dir = self.uploads_dir / "generated"
        self.generated_dir.mkdir(parents=True, exist_ok=True)
        self.engine = create_test_engine()
        drop_all_tables(self.engine)
        SQLModel.metadata.create_all(self.engine)

    def tearDown(self):
        self.engine.dispose()
        self.temp_dir.cleanup()

    def _create_project(self) -> int:
        with Session(self.engine) as session:
            project = Project(name="Test Project", client="Test Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            return project.id

    def _create_project_file(self, project_id: int, name: str, file_type: str, content_bytes: bytes) -> int:
        project_dir = self.uploads_dir / "projects" / str(project_id)
        project_dir.mkdir(parents=True, exist_ok=True)
        file_path = project_dir / name
        file_path.write_bytes(content_bytes)
        with Session(self.engine) as session:
            pf = ProjectFile(
                project_id=project_id,
                name=name,
                file_type=file_type,
                path=str(file_path.relative_to(self.uploads_dir)),
                size_bytes=len(content_bytes),
                summary="Test file",
            )
            session.add(pf)
            session.commit()
            session.refresh(pf)
            return pf.id

    def _generate_test_pdf(self) -> bytes:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        import io

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=letter)
        c.drawString(100, 700, "Page 1 Content")
        c.showPage()
        c.drawString(100, 700, "Page 2 Content")
        c.showPage()
        c.drawString(100, 700, "Page 3 Content")
        c.showPage()
        c.save()
        return buf.getvalue()

    def test_tool_is_registered(self):
        from app.tools.pdf_tools import MANAGE_PDF_TOOL_NAME
        tool = registry.get(MANAGE_PDF_TOOL_NAME)
        self.assertIsNotNone(tool)
        self.assertEqual(tool.name, MANAGE_PDF_TOOL_NAME)

    def test_pdf_read(self):
        from app.tools import pdf_tools

        project_id = self._create_project()
        pdf_bytes = self._generate_test_pdf()
        file_id = self._create_project_file(project_id, "test.pdf", "pdf", pdf_bytes)

        with patch("app.tools.pdf_tools.engine", self.engine), \
             patch("app.tools.pdf_tools.UPLOADS_DIR", self.uploads_dir):
            result = asyncio.get_event_loop().run_until_complete(
                pdf_tools.manage_pdf(
                    project_id=project_id,
                    action="read",
                    file_id=file_id,
                )
            )

        self.assertTrue(result.get("ok"))
        self.assertEqual(result.get("total_pages"), 3)
        self.assertEqual(len(result.get("pages", [])), 3)

    def test_pdf_extract(self):
        from app.tools import pdf_tools

        project_id = self._create_project()
        pdf_bytes = self._generate_test_pdf()
        file_id = self._create_project_file(project_id, "test.pdf", "pdf", pdf_bytes)

        with patch("app.tools.pdf_tools.engine", self.engine), \
             patch("app.tools.pdf_tools.UPLOADS_DIR", self.uploads_dir):
            result = asyncio.get_event_loop().run_until_complete(
                pdf_tools.manage_pdf(
                    project_id=project_id,
                    action="extract",
                    file_id=file_id,
                    page_numbers=[1, 3],
                    output_name="extracted.pdf",
                )
            )

        self.assertTrue(result.get("ok"))
        self.assertEqual(len(result.get("extracted_pages", [])), 2)

    def test_pdf_merge(self):
        from app.tools import pdf_tools

        project_id = self._create_project()
        pdf_bytes = self._generate_test_pdf()
        file_id_1 = self._create_project_file(project_id, "part1.pdf", "pdf", pdf_bytes)
        file_id_2 = self._create_project_file(project_id, "part2.pdf", "pdf", pdf_bytes)

        with patch("app.tools.pdf_tools.engine", self.engine), \
             patch("app.tools.pdf_tools.UPLOADS_DIR", self.uploads_dir):
            result = asyncio.get_event_loop().run_until_complete(
                pdf_tools.manage_pdf(
                    project_id=project_id,
                    action="merge",
                    file_ids=[file_id_1, file_id_2],
                    output_name="merged.pdf",
                )
            )

        self.assertTrue(result.get("ok"))
        self.assertIn("merged", result.get("name", "").lower())

    def test_pdf_watermark(self):
        from app.tools import pdf_tools

        project_id = self._create_project()
        pdf_bytes = self._generate_test_pdf()
        file_id = self._create_project_file(project_id, "test.pdf", "pdf", pdf_bytes)

        with patch("app.tools.pdf_tools.engine", self.engine), \
             patch("app.tools.pdf_tools.UPLOADS_DIR", self.uploads_dir):
            result = asyncio.get_event_loop().run_until_complete(
                pdf_tools.manage_pdf(
                    project_id=project_id,
                    action="watermark",
                    file_id=file_id,
                    watermark_text="DRAFT",
                    output_name="watermarked.pdf",
                )
            )

        self.assertTrue(result.get("ok"))


class SkillLoadingTestCase(unittest.TestCase):
    """Test that file-backed skills load correctly."""

    def test_office_document_editor_skill_md_exists(self):
        skill_path = Path(__file__).resolve().parents[2] / "skills" / "office-document-editor" / "SKILL.md"
        self.assertTrue(skill_path.is_file(), f"Skill file not found: {skill_path}")

    def test_pdf_management_skill_md_exists(self):
        skill_path = Path(__file__).resolve().parents[2] / "skills" / "pdf-management" / "SKILL.md"
        self.assertTrue(skill_path.is_file(), f"Skill file not found: {skill_path}")

    def test_meeting_intelligence_skill_md_exists(self):
        skill_path = Path(__file__).resolve().parents[2] / "skills" / "meeting-intelligence" / "SKILL.md"
        self.assertTrue(skill_path.is_file(), f"Skill file not found: {skill_path}")

    def test_goal_definition_skill_md_exists(self):
        skill_path = Path(__file__).resolve().parents[2] / "skills" / "goal-definition" / "SKILL.md"
        self.assertTrue(skill_path.is_file(), f"Skill file not found: {skill_path}")

    def test_skill_packages_load_content(self):
        from app.routers.skills import _load_skill_package_prompt

        for package_name in ["office-document-editor", "pdf-management", "meeting-intelligence", "goal-definition"]:
            content = _load_skill_package_prompt(package_name)
            self.assertTrue(len(content) > 100, f"Skill '{package_name}' loaded empty or too short content")
            self.assertNotIn("---", content[:50], f"Skill '{package_name}' still has YAML frontmatter")

    def test_skill_constants_defined(self):
        from app.routers.skills import (
            OFFICE_DOCUMENT_EDITOR_SKILL_NAME,
            PDF_MANAGEMENT_SKILL_NAME,
            MEETING_INTELLIGENCE_SKILL_NAME,
            GOAL_DEFINITION_SKILL_NAME,
            OFFICE_DOCUMENT_EDITOR_TOOL_NAMES,
            PDF_MANAGEMENT_TOOL_NAMES,
            MEETING_INTELLIGENCE_TOOL_NAMES,
            GOAL_DEFINITION_TOOL_NAMES,
        )

        self.assertIn("edit_project_office_document", OFFICE_DOCUMENT_EDITOR_TOOL_NAMES)
        self.assertIn("manage_pdf", PDF_MANAGEMENT_TOOL_NAMES)
        self.assertIn("update_project_markdown_document", MEETING_INTELLIGENCE_TOOL_NAMES)
        self.assertIn("update_project_markdown_document", GOAL_DEFINITION_TOOL_NAMES)

    def test_skills_in_gstack_pro_skills(self):
        from app.routers.skills import GSTACK_PRO_SKILLS

        skill_names = {s["name"] for s in GSTACK_PRO_SKILLS}
        self.assertIn("Office 文档编辑", skill_names)
        self.assertIn("PDF 工具箱", skill_names)
        self.assertIn("会议纪要提取", skill_names)
        self.assertIn("目标定义", skill_names)


if __name__ == "__main__":
    unittest.main()
