from __future__ import annotations

import asyncio
import gc
import json
import os
import shutil
import tempfile
import time
import unittest
from datetime import datetime, timedelta
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import AsyncMock, patch

from fastapi import FastAPI
from fastapi.testclient import TestClient
from sqlmodel import Session, SQLModel, create_engine, select

from app.models.db import (
    ClientMemorySummary,
    ClientMemorySnapshot,
    ClientRecord,
    ClientStakeholder,
    Conversation,
    DocumentChunk,
    GeneratedFile,
    KnowledgeDocument,
    Message,
    Milestone,
    Project,
    ProjectFile,
    ProjectFolder,
    ProjectMemorySnapshot,
    ProjectMemorySummary,
    ProjectMember,
    ProjectPayment,
    ProjectTodo,
    Setting,
    Skill,
    ScheduledTask,
    SystemMessage,
    SystemMessageRead,
    TaskArtifact,
    TaskEvent,
    TaskRun,
    TaskStep,
    ToolCall,
    User,
)
from app import config as app_config
from app import database as database_module
from app.routers import auth as auth_router_module
from app.routers import chat as chat_router_module
from app.routers import clients as clients_router_module
from app.routers import memory_operations as memory_operations_router_module
from app.routers import messages as messages_router_module
from app.routers import projects as projects_router_module
from app.routers import projects_memory as projects_memory_module
from app.routers import projects_deps as projects_deps_module
from app.routers import projects_files as projects_files_module
from app.routers import skills as skills_router_module
from app.services.cache import conversations_cache, projects_cache
from app.services import chat_exports as chat_exports_module
from app.services import chat_streaming as chat_streaming_module
from app.services import client_contexts as client_contexts_module
from app.services import context_builder as context_builder_module
from app.services import document_text as document_text_module
from app.services import openai_compat as openai_compat_module
from app.services import project_ai as project_ai_module
from app.services import project_contexts as project_contexts_module
from app.services import provider_selector as provider_selector_module
from app.services import project_notes as project_notes_module
from app.services import rag as rag_module
from app.services import scheduler as scheduler_module
from contextlib import ExitStack, contextmanager

@contextmanager
def patch_project_llm(complete=None, stream=None):
    """Patch complete_with_selected_model and/or stream_with_selected_model across ALL project router modules."""
    with ExitStack() as stack:
        if complete is not None:
            for mod in (projects_router_module, projects_memory_module, projects_deps_module):
                stack.enter_context(patch.object(mod, "complete_with_selected_model", complete))
        if stream is not None:
            for mod in (projects_router_module, projects_memory_module, projects_deps_module):
                stack.enter_context(patch.object(mod, "stream_with_selected_model", stream))
        yield complete if complete is not None else stream
from app.services.chat_streaming import ChatRuntime, _build_slides_from_strategy_text, _route_ppt_tool_for_skill, stream_chat_events
from app.services.time_utils import utc_now_naive
from tests.test_database import create_test_engine, drop_all_tables


class SchedulerTimeTestCase(unittest.TestCase):
    def test_as_utc_aware_treats_naive_datetime_as_utc(self):
        naive = datetime(2026, 4, 27, 13, 19, 29)

        aware = scheduler_module._as_utc_aware(naive)

        self.assertIsNotNone(aware.tzinfo)
        self.assertEqual(aware.utcoffset(), timedelta(0))
        self.assertEqual(aware.hour, 13)

    def test_add_or_replace_date_job_uses_utc_aware_date_trigger(self):
        run_at = datetime(2026, 4, 27, 13, 19, 29)

        with patch.object(scheduler_module, "_scheduler") as mocked_scheduler:
            mocked_scheduler.running = True
            scheduler_module.add_or_replace_date_job("job", run_at, lambda: None)

        trigger = mocked_scheduler.add_job.call_args.kwargs["trigger"]
        self.assertEqual(trigger.run_date.utcoffset(), timedelta(0))
        self.assertEqual(trigger.run_date.hour, 13)
        self.assertEqual(mocked_scheduler.add_job.call_args.kwargs["misfire_grace_time"], 3600)


class OpenAICompatLoopResourceTestCase(unittest.TestCase):
    def test_kimi_complete_semaphore_is_event_loop_local(self):
        async def get_sem():
            return openai_compat_module._get_kimi_complete_sem()

        first = asyncio.run(get_sem())
        second = asyncio.run(get_sem())

        self.assertIsNot(first, second)

    def test_http_client_is_event_loop_local(self):
        async def get_client():
            return openai_compat_module._get_http_client()

        first = asyncio.run(get_client())
        second = asyncio.run(get_client())

        self.assertIsNot(first, second)


class FakeRetrievalContext:
    def __init__(self, text: str = "", results: list | None = None):
        self._text = text
        self.results = results or []

    def to_text(self) -> str:
        return self._text


class FakeStreamingLLM:
    def __init__(self, responses):
        self._responses = responses
        self.calls = 0

    async def stream_response(self, *args, **kwargs):
        current = self._responses[self.calls]
        self.calls += 1
        if isinstance(current, Exception):
            raise current
        for chunk in current:
            yield chunk

    async def complete(self, *args, **kwargs):
        return "ok"


def collect_async_generator(async_gen):
    async def _collect():
        items = []
        async for item in async_gen:
            items.append(item)
        return items

    return asyncio.run(_collect())


class ChatRouterTestCase(unittest.TestCase):
    def setUp(self):
        projects_cache.clear()
        conversations_cache.clear()
        self.engine = create_test_engine()
        drop_all_tables(self.engine)
        SQLModel.metadata.create_all(self.engine)

        def override_session():
            with Session(self.engine) as session:
                yield session

        app = FastAPI()
        app.include_router(auth_router_module.router)
        app.include_router(chat_router_module.router)
        app.dependency_overrides[auth_router_module.get_session] = override_session
        app.dependency_overrides[chat_router_module.get_session] = override_session
        self.client = TestClient(app)

    def tearDown(self):
        self.client.close()
        self.engine.dispose()

    def test_conversation_crud_and_markdown_export(self):
        create_resp = self.client.post(
            "/chat/conversations",
            json={"title": "Regression Chat"},
        )
        self.assertEqual(create_resp.status_code, 200)
        conv_id = create_resp.json()["id"]

        with Session(self.engine) as session:
            now = utc_now_naive()
            session.add(Message(conversation_id=conv_id, role="user", content="hello", created_at=now))
            session.add(
                Message(
                    conversation_id=conv_id,
                    role="assistant",
                    content="world",
                    created_at=now + timedelta(seconds=1),
                )
            )
            session.commit()

        list_resp = self.client.get("/chat/conversations")
        self.assertEqual(list_resp.status_code, 200)
        self.assertEqual(len(list_resp.json()), 1)

        messages_resp = self.client.get(f"/chat/conversations/{conv_id}/messages")
        self.assertEqual(messages_resp.status_code, 200)
        self.assertEqual([m["content"] for m in messages_resp.json()], ["hello", "world"])

        export_resp = self.client.post(
            f"/chat/conversations/{conv_id}/export",
            json={"format": "markdown"},
        )
        self.assertEqual(export_resp.status_code, 200)
        self.assertIn("# Regression Chat", export_resp.text)
        self.assertIn("**User**", export_resp.text)
        self.assertIn("**Assistant**", export_resp.text)

    def test_conversation_list_purges_items_inactive_for_more_than_7_days(self):
        now = utc_now_naive()
        with Session(self.engine) as session:
            project = Project(name="Retention Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)

            old_conv = Conversation(
                title="Old Chat",
                created_at=now - timedelta(days=10),
                updated_at=now - timedelta(days=8),
            )
            recent_conv = Conversation(
                title="Recent Chat",
                created_at=now - timedelta(days=20),
                updated_at=now - timedelta(days=1),
            )
            session.add(old_conv)
            session.add(recent_conv)
            session.commit()
            session.refresh(old_conv)
            session.refresh(recent_conv)

            session.add(Message(conversation_id=old_conv.id, role="user", content="old"))
            session.add(Message(conversation_id=recent_conv.id, role="user", content="recent"))
            task = TaskRun(
                project_id=project.id,
                conversation_id=old_conv.id,
                task_type="text_artifact",
                goal="old task",
            )
            artifact = GeneratedFile(
                conversation_id=old_conv.id,
                name="old.md",
                file_type="md",
                path="generated/old.md",
            )
            session.add(task)
            session.add(artifact)
            session.commit()
            old_conv_id = old_conv.id
            recent_conv_id = recent_conv.id
            task_id = task.id

        list_resp = self.client.get("/chat/conversations")
        self.assertEqual(list_resp.status_code, 200)
        self.assertEqual([item["id"] for item in list_resp.json()], [recent_conv_id])

        with Session(self.engine) as session:
            self.assertIsNone(session.get(Conversation, old_conv_id))
            self.assertIsNotNone(session.get(Conversation, recent_conv_id))
            old_messages = session.exec(select(Message).where(Message.conversation_id == old_conv_id)).all()
            old_artifacts = session.exec(select(GeneratedFile).where(GeneratedFile.conversation_id == old_conv_id)).all()
            self.assertEqual(old_messages, [])
            self.assertEqual(old_artifacts, [])
            task = session.get(TaskRun, task_id)
            self.assertIsNotNone(task)
            self.assertIsNone(task.conversation_id)

    def test_patch_conversation_updates_title(self):
        create_resp = self.client.post(
            "/chat/conversations",
            json={"title": "Old Title"},
        )
        self.assertEqual(create_resp.status_code, 200)
        conv_id = create_resp.json()["id"]

        patch_resp = self.client.patch(
            f"/chat/conversations/{conv_id}",
            json={"title": "New Title"},
        )
        self.assertEqual(patch_resp.status_code, 200)
        self.assertEqual(patch_resp.json()["title"], "New Title")

        get_resp = self.client.get(f"/chat/conversations/{conv_id}")
        self.assertEqual(get_resp.status_code, 200)
        self.assertEqual(get_resp.json()["title"], "New Title")

    def test_delete_conversation_unlinks_task_runs(self):
        with Session(self.engine) as session:
            project = Project(name="Task Project", client="Client", status="active")
            session.add(project)
            session.commit()
            session.refresh(project)
            conv = Conversation(title="Task Chat", project_id=project.id)
            session.add(conv)
            session.commit()
            session.refresh(conv)
            session.add(Message(conversation_id=conv.id, role="user", content="hello"))
            task = TaskRun(
                project_id=project.id,
                conversation_id=conv.id,
                task_type="generate_client_ppt",
                goal="prepare ppt",
            )
            session.add(task)
            session.commit()
            task_id = task.id
            conv_id = conv.id

        delete_resp = self.client.delete(f"/chat/conversations/{conv_id}")
        self.assertEqual(delete_resp.status_code, 200)

        with Session(self.engine) as session:
            self.assertIsNone(session.get(Conversation, conv_id))
            task = session.get(TaskRun, task_id)
            self.assertIsNotNone(task)
            self.assertIsNone(task.conversation_id)

    def test_send_route_streams_service_output(self):
        async def fake_stream(*args, **kwargs):
            yield 'data: {"type":"conversation_id","id":99}\n\n'
            yield 'data: {"type":"done"}\n\n'

        with patch.object(chat_router_module, "prepare_chat_runtime", return_value="runtime"), patch.object(
            chat_router_module,
            "stream_chat_events",
            side_effect=lambda runtime, req, bind: fake_stream(),
        ):
            resp = self.client.post("/chat/send", json={"content": "hello"})

        self.assertEqual(resp.status_code, 200)
        self.assertIn('"conversation_id"', resp.text)
        self.assertIn('"done"', resp.text)

    def test_login_rate_limit_blocks_repeated_failures(self):
        auth_router_module._LOGIN_ATTEMPTS.clear()

        with Session(self.engine) as session:
            session.add(
                User(
                    email="admin@example.com",
                    display_name="Admin",
                    password_hash=auth_router_module._hash("correct-password"),
                    is_admin=True,
                    is_active=True,
                )
            )
            session.commit()

        headers = {"x-forwarded-for": "203.0.113.5"}
        for _ in range(5):
            response = self.client.post(
                "/auth/login",
                json={"email": "admin@example.com", "password": "wrong-password"},
                headers=headers,
            )
            self.assertEqual(response.status_code, 401)

        blocked = self.client.post(
            "/auth/login",
            json={"email": "admin@example.com", "password": "wrong-password"},
            headers=headers,
        )
        self.assertEqual(blocked.status_code, 429)
        self.assertEqual(blocked.json()["detail"], "Too many login attempts. Please try again later.")
        self.assertIn("Retry-After", blocked.headers)


class ProviderSelectionTestCase(unittest.TestCase):
    def test_kimi_k26_resolves_to_kimi_provider(self):
        self.assertEqual(provider_selector_module.resolve_provider_from_model("kimi-k2.6"), "kimi")

    def test_deepseek_v4_resolves_to_deepseek_provider(self):
        self.assertEqual(provider_selector_module.resolve_provider_from_model("deepseek-v4-pro"), "deepseek")

    def test_mimo_models_resolve_to_mimo_provider(self):
        self.assertEqual(provider_selector_module.resolve_provider_from_model("xiaomi/mimo-v2.5-flash"), "mimo")
        self.assertEqual(provider_selector_module.resolve_provider_from_model("mimo-v2.5-pro"), "mimo")

    def test_xiaomi_provider_alias_normalizes_to_mimo(self):
        engine = create_test_engine()
        drop_all_tables(engine)
        SQLModel.metadata.create_all(engine)
        try:
            with Session(engine) as session:
                session.add(Setting(key="llm_provider", value="xiaomi"))
                session.commit()
                self.assertEqual(provider_selector_module.get_provider_name(session), "mimo")
        finally:
            engine.dispose()

    def test_kimi_k26_uses_k2_sampling_defaults(self):
        self.assertEqual(openai_compat_module._apply_moonshot_fixed_params("kimi-k2.6", 0.2), (1.0, 0.95))

    def test_mimo_complete_uses_openai_compatible_endpoint(self):
        class FakeResponse:
            status_code = 200
            text = ""

            def json(self):
                return {"choices": [{"message": {"content": "mimo ok", "tool_calls": None}}]}

        class FakeClient:
            def __init__(self):
                self.calls = []

            async def post(self, url, headers=None, json=None):
                self.calls.append({"url": url, "headers": headers, "json": json})
                return FakeResponse()

        fake_client = FakeClient()

        async def run():
            with patch.object(openai_compat_module, "get_mimo_api_key", return_value="mimo-key"), patch.object(
                openai_compat_module,
                "_get_http_client",
                return_value=fake_client,
            ):
                return await openai_compat_module.complete(
                    [{"role": "user", "content": "hi"}],
                    model="mimo-v2.5-flash",
                    max_tokens=20,
                )

        self.assertEqual(asyncio.run(run()), "mimo ok")
        self.assertEqual(fake_client.calls[0]["url"], f"{openai_compat_module.MIMO_BASE_URL}/chat/completions")
        self.assertEqual(fake_client.calls[0]["json"]["model"], "mimo-v2.5-flash")
        self.assertEqual(fake_client.calls[0]["headers"]["Authorization"], "Bearer mimo-key")

    def test_mimo_token_plan_key_uses_token_plan_endpoint(self):
        class FakeResponse:
            status_code = 200
            text = ""

            def json(self):
                return {"choices": [{"message": {"content": "token plan ok"}}]}

        class FakeClient:
            def __init__(self):
                self.calls = []

            async def post(self, url, headers=None, json=None):
                self.calls.append({"url": url, "headers": headers, "json": json})
                return FakeResponse()

        fake_client = FakeClient()

        async def run():
            with patch.object(openai_compat_module, "get_mimo_api_key", return_value="tp-test-key"), patch.object(
                openai_compat_module,
                "_get_http_client",
                return_value=fake_client,
            ):
                return await openai_compat_module.complete(
                    [{"role": "user", "content": "hi"}],
                    model="xiaomi/mimo-v2.5-pro",
                    max_tokens=20,
                )

        self.assertEqual(asyncio.run(run()), "token plan ok")
        self.assertEqual(fake_client.calls[0]["url"], f"{openai_compat_module.MIMO_TOKEN_PLAN_BASE_URL}/chat/completions")
        self.assertEqual(fake_client.calls[0]["json"]["model"], "mimo-v2.5-pro")

    def test_mimo_legacy_v2_model_aliases_to_v25(self):
        class FakeResponse:
            status_code = 200
            text = ""

            def json(self):
                return {"choices": [{"message": {"content": "alias ok"}}]}

        class FakeClient:
            def __init__(self):
                self.calls = []

            async def post(self, url, headers=None, json=None):
                self.calls.append({"url": url, "headers": headers, "json": json})
                return FakeResponse()

        fake_client = FakeClient()

        async def run():
            with patch.object(openai_compat_module, "get_mimo_api_key", return_value="tp-test-key"), patch.object(
                openai_compat_module,
                "_get_http_client",
                return_value=fake_client,
            ):
                return await openai_compat_module.complete(
                    [{"role": "user", "content": "hi"}],
                    model="mimo-v2-flash",
                    max_tokens=20,
                )

        self.assertEqual(asyncio.run(run()), "alias ok")
        self.assertEqual(fake_client.calls[0]["json"]["model"], "mimo-v2.5-flash")


class ProjectServiceHelperTestCase(unittest.TestCase):
    def test_validate_jwt_secret_rejects_default_secret(self):
        with patch.object(app_config, "ALLOW_INSECURE_JWT_SECRET", False):
            with self.assertRaises(RuntimeError):
                app_config.validate_jwt_secret(app_config.DEFAULT_JWT_SECRET)

    def test_validate_jwt_secret_accepts_strong_secret(self):
        with patch.object(app_config, "ALLOW_INSECURE_JWT_SECRET", False):
            app_config.validate_jwt_secret("this-is-a-very-strong-jwt-secret-2026")

    def test_build_project_context_prompt_emphasizes_current_project_only(self):
        prompt = project_contexts_module.build_project_context_prompt(
            "Project: Alpha\nClient: Client A\nStatus: active"
        )

        self.assertIn("current project as the only source of truth", prompt)
        self.assertIn("Do not blend in facts, progress, or risks from other projects", prompt)
        self.assertIn("Project data:\nProject: Alpha", prompt)

    def test_build_project_memory_view_prompt_supports_risk_summary(self):
        prompt = project_contexts_module.build_project_memory_view_prompt(
            {
                "project_brief": "Alpha rollout",
                "key_risks": ["Timeline risk"],
                "next_actions": ["Confirm scope"],
            },
            "Alpha",
            "risk",
        )

        self.assertIn("Summary type: risk", prompt)
        self.assertIn("focused on project risks", prompt)
        self.assertIn('"key_risks": ["Timeline risk"]', prompt)

    def test_build_project_summary_from_memory_prompt_respects_language(self):
        prompt = project_contexts_module.build_project_summary_from_memory_prompt(
            {
                "project_brief": "Alpha rollout",
                "current_stage": "delivery",
                "next_actions": ["Confirm scope"],
            },
            "Alpha",
            "zh-CN",
        )

        self.assertIn("Write the answer in Chinese", prompt)
        self.assertIn("Structured memory JSON", prompt)

    def test_build_project_memory_view_prompt_supports_financial_and_documents_summary(self):
        financial_prompt = project_contexts_module.build_project_memory_view_prompt(
            {
                "financial_status": "Large receivable pending",
                "key_risks": ["Collection delay"],
            },
            "Alpha",
            "financial",
        )
        documents_prompt = project_contexts_module.build_project_memory_view_prompt(
            {
                "important_documents": [{"name": "SOW", "reason": "Scope baseline"}],
            },
            "Alpha",
            "documents",
        )

        self.assertIn("Summary type: financial", financial_prompt)
        self.assertIn("focused on the project's financial picture", financial_prompt)
        self.assertIn("Summary type: documents", documents_prompt)
        self.assertIn("focused on project documents and knowledge signals", documents_prompt)

    def test_build_client_memory_summary_prompt_supports_relationship_and_delivery(self):
        relationship_prompt = clients_router_module.build_client_memory_summary_prompt(
            {
                "client_profile": "Global account with strong executive sponsorship",
                "decision_patterns": ["Prefers weekly steering rhythm"],
                "key_contacts": [{"name": "Jane", "role": "CFO", "note": "Final approver"}],
                "lessons_learned": ["Bring finance and IT together early"],
                "project_history": [{"project_name": "Alpha", "status": "won", "outcome": "Expanded"}],
                "sensitive_topics": ["Avoid overpromising delivery dates"],
            },
            "Acme",
            "relationship",
            "zh-CN",
        )
        delivery_prompt = clients_router_module.build_client_memory_summary_prompt(
            {
                "client_profile": "Global account with strong executive sponsorship",
                "decision_patterns": ["Prefers weekly steering rhythm"],
                "key_contacts": [{"name": "Jane", "role": "CFO", "note": "Final approver"}],
                "lessons_learned": ["Bring finance and IT together early"],
                "project_history": [{"project_name": "Alpha", "status": "won", "outcome": "Expanded"}],
                "sensitive_topics": ["Avoid overpromising delivery dates"],
            },
            "Acme",
            "delivery",
            "en-US",
        )

        self.assertIn("Summary type: relationship", relationship_prompt)
        self.assertIn("trust level", relationship_prompt)
        self.assertIn("Write the answer in Chinese", relationship_prompt)
        self.assertIn("Summary type: delivery", delivery_prompt)
        self.assertIn("delivery readiness", delivery_prompt)

    def test_database_migration_governance_detects_lightweight_mode(self):
        governance = database_module.get_database_migration_governance(
            tables=["project", "clientrecord"],
            current_revision=None,
            revisions=["001_init", "002_add_memory"],
        )

        self.assertEqual(governance["mode"], "lightweight")
        self.assertEqual(governance["pending_count"], 2)
        self.assertFalse(governance["up_to_date"])

    def test_database_migration_governance_detects_alembic_mode(self):
        governance = database_module.get_database_migration_governance(
            tables=["project", "alembic_version"],
            current_revision="001_init",
            revisions=["001_init", "002_add_memory"],
        )

        self.assertEqual(governance["mode"], "alembic")
        self.assertEqual(governance["pending_revisions"], ["002_add_memory"])

    def test_database_migration_governance_normalizes_short_revision_alias(self):
        governance = database_module.get_database_migration_governance(
            tables=["project", "alembic_version"],
            current_revision="005",
            revisions=["001_v1_1", "002_v1_2", "003_v1_3", "004_v1_4", "005_v1_5"],
        )

        self.assertEqual(governance["current_revision"], "005_v1_5")
        self.assertEqual(governance["latest_revision"], "005_v1_5")
        self.assertEqual(governance["pending_count"], 0)
        self.assertTrue(governance["up_to_date"])

    def test_database_migration_governance_keeps_unknown_revision_pending(self):
        governance = database_module.get_database_migration_governance(
            tables=["project", "alembic_version"],
            current_revision="999",
            revisions=["001_v1_1", "002_v1_2"],
        )

        self.assertEqual(governance["current_revision"], "999")
        self.assertEqual(governance["pending_revisions"], ["001_v1_1", "002_v1_2"])

    def test_get_project_memory_payload_flattens_pinned_slots(self):
        project = Project(
            name="Alpha",
            client="Client",
            context_memory_json=json.dumps(
                {
                    "key_risks": {
                        "ai": ["Timeline risk"],
                        "pinned": ["Decision-maker alignment risk"],
                    },
                    "stakeholder_notes": {
                        "ai": ["Finance team supports automation"],
                        "pinned": ["CFO is the final approver"],
                    },
                },
                ensure_ascii=False,
            ),
            memory_version=2,
            memory_stale=False,
        )

        payload = project_contexts_module.get_project_memory_payload(project)

        self.assertEqual(
            payload["key_risks"],
            ["Timeline risk", "Decision-maker alignment risk"],
        )
        self.assertEqual(
            payload["stakeholder_notes"],
            ["Finance team supports automation", "CFO is the final approver"],
        )
        self.assertEqual(
            payload["key_risks_detail"]["pinned"],
            ["Decision-maker alignment risk"],
        )

    def test_stream_llm_text_chunks_skips_tool_markers(self):
        async def fake_chunks():
            yield '{"type": "tool_use","id":"tool-1"}'
            yield "[TOOL_START:generate_docx]"
            yield "final text"

        chunks = collect_async_generator(project_contexts_module.stream_llm_text_chunks(fake_chunks()))
        self.assertEqual(chunks, ["final text"])

    def test_build_project_note_polish_messages_includes_project_context(self):
        project = Project(name="Alpha", client="Client A")
        messages = project_notes_module.build_project_note_polish_messages(project, "raw draft")

        self.assertEqual(len(messages), 2)
        self.assertEqual(messages[0]["role"], "system")
        self.assertIn("well-structured Markdown project notes", messages[0]["content"])
        self.assertEqual(messages[1]["role"], "user")
        self.assertIn("Project name: Alpha", messages[1]["content"])
        self.assertIn("Client: Client A", messages[1]["content"])
        self.assertIn("raw draft", messages[1]["content"])

    def test_build_project_ai_suggest_messages_include_client_context(self):
        messages = project_ai_module.build_project_ai_suggest_messages(
            "Need a China entry plan",
            client_name="Acme",
            client_industry="Retail",
        )

        self.assertEqual(len(messages), 1)
        self.assertIn('The project is for: Client: Acme (Retail)', messages[0]["content"])
        self.assertIn('The user described the project as: "Need a China entry plan"', messages[0]["content"])

    def test_parse_project_ai_suggestions_supports_fenced_json(self):
        raw = """```json
[
  {"name": "China Market Entry Strategy", "description": "Define market priorities."}
]
```"""
        suggestions = project_ai_module.parse_project_ai_suggestions(raw)
        self.assertEqual(
            suggestions,
            [{"name": "China Market Entry Strategy", "description": "Define market priorities."}],
        )

    def test_build_project_file_summary_prompt_contains_excerpt(self):
        prompt = project_ai_module.build_project_file_summary_prompt("Important project excerpt")
        self.assertIn("write a concise 2-3 sentence summary", prompt)
        self.assertIn("Document excerpt:\nImportant project excerpt", prompt)

    def test_extract_text_from_file_reads_markdown(self):
        fd, temp_path = tempfile.mkstemp(suffix=".md")
        os.close(fd)
        path = Path(temp_path)
        try:
            path.write_text("# Hello", encoding="utf-8")
            text = document_text_module.extract_text_from_file(path, "md", max_chars=4000)
        finally:
            path.unlink(missing_ok=True)
        self.assertEqual(text, "# Hello")

    def test_build_markdown_export_content_formats_messages(self):
        conversation = Conversation(
            title="Alpha Review",
            created_at=datetime(2026, 4, 16, 9, 30),
        )
        messages = [
            Message(role="user", content="hello", created_at=datetime(2026, 4, 16, 9, 31)),
            Message(role="assistant", content="world", created_at=datetime(2026, 4, 16, 9, 32)),
        ]

        content = chat_exports_module.build_markdown_export_content(conversation, messages)

        self.assertIn("# Alpha Review", content)
        self.assertIn("**User** *(09:31)*", content)
        self.assertIn("**Assistant** *(09:32)*", content)
        self.assertIn("hello", content)
        self.assertIn("world", content)


class ProjectConversationArchiveTestCase(unittest.TestCase):
    def setUp(self):
        self.engine = create_test_engine()
        drop_all_tables(self.engine)
        SQLModel.metadata.create_all(self.engine)

        def override_session():
            with Session(self.engine) as session:
                yield session

        with Session(self.engine) as session:
            current_user = User(
                email="current@example.com",
                display_name="Current User",
                password_hash="hashed",
            )
            session.add(current_user)
            session.commit()
            session.refresh(current_user)
            self.current_user_id = current_user.id

        def override_current_user():
            with Session(self.engine) as session:
                user = session.get(User, self.current_user_id)
                assert user is not None
                return user

        app = FastAPI()
        app.include_router(projects_router_module.router)
        from app.database import get_session as _gs
        from app.routers.auth import get_current_user as _gcu
        app.dependency_overrides[_gs] = override_session
        app.dependency_overrides[_gcu] = override_current_user
        self.client = TestClient(app)

        self.uploads_dir = Path(os.getcwd()) / "test_tmp_uploads"
        self.uploads_dir.mkdir(parents=True, exist_ok=True)
        self._upload_patches = []
        for mod in (projects_router_module, projects_deps_module, projects_files_module):
            p = patch.object(mod, "UPLOADS_DIR", self.uploads_dir)
            p.start()
            self._upload_patches.append(p)

    def tearDown(self):
        self.client.close()
        for p in self._upload_patches:
            p.stop()
        projects_cache.clear()
        shutil.rmtree(self.uploads_dir, ignore_errors=True)
        self.engine.dispose()

    def test_create_project_rejects_duplicate_active_project(self):
        first_resp = self.client.post(
            "/projects",
            json={
                "name": " 东阿阿胶新业务进入机会和策略 ",
                "client": "东阿阿胶股份有限公司",
                "status": "lead",
            },
        )
        self.assertEqual(first_resp.status_code, 201)
        first_project_id = first_resp.json()["id"]

        duplicate_resp = self.client.post(
            "/projects",
            json={
                "name": "东阿阿胶新业务进入机会和策略",
                "client": " 东阿阿胶股份有限公司 ",
                "status": "lead",
            },
        )

        self.assertEqual(duplicate_resp.status_code, 409)
        detail = duplicate_resp.json()["detail"]
        self.assertEqual(detail["code"], "duplicate_project")
        self.assertEqual(detail["project_id"], first_project_id)

        with Session(self.engine) as session:
            projects = session.exec(select(Project)).all()
        self.assertEqual(len(projects), 1)

    def test_save_project_conversation_as_markdown_file(self):
        with Session(self.engine) as session:
            project = Project(name="Alpha", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            folder = ProjectFolder(project_id=project.id, name="Deliverables", sort_order=2)
            session.add(folder)
            session.commit()
            session.refresh(folder)

            conv = Conversation(project_id=project.id, title="Weekly Sync")
            session.add(conv)
            session.commit()
            session.refresh(conv)

            session.add(Message(conversation_id=conv.id, role="user", content="Please summarize this week"))
            session.add(Message(conversation_id=conv.id, role="assistant", content="Here is the weekly summary"))
            session.commit()

            project_id = project.id
            folder_id = folder.id
            conv_id = conv.id

        resp = self.client.post(f"/projects/{project_id}/conversations/{conv_id}/save-markdown", json={})
        self.assertEqual(resp.status_code, 201)
        body = resp.json()
        self.assertTrue(body["ok"])
        self.assertEqual(body["action"], "new")
        self.assertIsNone(body["folder_id"])
        self.assertTrue(body["name"].endswith(".md"))

        with Session(self.engine) as session:
            saved_file = session.get(ProjectFile, body["id"])
            self.assertIsNotNone(saved_file)
            self.assertEqual(saved_file.project_id, project_id)
            self.assertEqual(saved_file.folder_id, None)
            saved_path = self.uploads_dir / saved_file.path

        self.assertTrue(saved_path.exists())
        content = saved_path.read_text(encoding="utf-8")
        self.assertIn("# Weekly Sync", content)
        self.assertIn("Please summarize this week", content)
        self.assertIn("Here is the weekly summary", content)

    @unittest.skip("Presales template bootstrap was removed; AI now owns project space structure.")
    def test_init_presales_template_creates_markdown_documents(self):
        with Session(self.engine) as session:
            project = Project(name="Beta", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

            legacy_folder = ProjectFolder(project_id=project_id, name="项目需求", sort_order=0)
            session.add(legacy_folder)
            session.commit()
            session.refresh(legacy_folder)

            legacy_path = Path("projects") / str(project_id) / "legacy_note.md"
            legacy_full_path = self.uploads_dir / legacy_path
            legacy_full_path.parent.mkdir(parents=True, exist_ok=True)
            legacy_full_path.write_text("# Legacy", encoding="utf-8")
            session.add(
                ProjectFile(
                    project_id=project_id,
                    folder_id=legacy_folder.id,
                    name="Legacy Note.md",
                    file_type="md",
                    path=str(legacy_path),
                    size_bytes=legacy_full_path.stat().st_size,
                )
            )
            session.commit()

        resp = self.client.post(f"/projects/{project_id}/notes/templates/presales", json={})
        self.assertEqual(resp.status_code, 201)
        body = resp.json()
        self.assertTrue(body["ok"])
        self.assertGreaterEqual(body["created_count"], 9)
        self.assertGreaterEqual(body["cleaned_folder_count"], 1)

        with Session(self.engine) as session:
            md_files = session.exec(
                select(ProjectFile).where(ProjectFile.project_id == project_id, ProjectFile.file_type == "md")
            ).all()
            self.assertGreaterEqual(len(md_files), 9)
            first_doc = next((f for f in md_files if f.name == "00_项目概览.md"), None)
            self.assertIsNotNone(first_doc)
            first_doc_path = self.uploads_dir / first_doc.path
            self.assertTrue(first_doc_path.exists())
            self.assertIn("## 客户想解决的核心问题", first_doc_path.read_text(encoding="utf-8"))
            self.assertFalse(any(f.name == "项目需求" for f in session.exec(select(ProjectFolder).where(ProjectFolder.project_id == project_id)).all()))
            migrated_file = next((f for f in md_files if f.name == "Legacy Note.md"), None)
            self.assertIsNotNone(migrated_file)
            new_folder = session.get(ProjectFolder, migrated_file.folder_id)
            self.assertIsNotNone(new_folder)
            self.assertEqual(new_folder.name, "02_需求与方案")

    def test_ai_suggest_project_parses_fenced_json_response(self):
        with patch_project_llm(complete=AsyncMock(
                return_value="""```json
[
  {"name": "China Market Entry Strategy", "description": "Assess demand, competitors, and go-to-market priorities."}
]
```"""
            ),
        ):
            resp = self.client.post(
                "/projects/ai-suggest",
                json={
                    "query": "Help expand into China",
                    "client_name": "Acme",
                    "client_industry": "Retail",
                },
            )

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(len(body), 1)
        self.assertEqual(body[0]["name"], "China Market Entry Strategy")
        self.assertIn("go-to-market", body[0]["description"])

    def test_project_memory_status_and_rebuild_flow(self):
        with Session(self.engine) as session:
            project = Project(name="Memory Project", client="Client", description="Alpha project")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        status_before = self.client.get(f"/projects/{project_id}/memory/status")
        self.assertEqual(status_before.status_code, 200)
        self.assertFalse(status_before.json()["has_memory"])
        self.assertTrue(status_before.json()["memory_stale"])

        with patch_project_llm(complete=AsyncMock(
                return_value=json.dumps(
                    {
                        "project_brief": "Alpha project brief",
                        "current_stage": "lead",
                        "recent_progress": ["Kickoff prepared"],
                        "key_risks": ["Timeline risk"],
                        "open_questions": [],
                        "next_actions": ["Confirm scope"],
                        "important_documents": [{"name": "brief.md", "reason": "Core scope"}],
                        "financial_status": "No contract signed yet",
                        "delivery_signals": [],
                        "stakeholder_notes": [],
                    },
                    ensure_ascii=False,
                )
            ),
        ):
            rebuild_resp = self.client.post(f"/projects/{project_id}/memory/rebuild")

        self.assertEqual(rebuild_resp.status_code, 200)
        rebuild_body = rebuild_resp.json()
        self.assertTrue(rebuild_body["ok"])
        self.assertEqual(rebuild_body["memory"]["project_brief"], "Alpha project brief")
        self.assertEqual(rebuild_body["memory_version"], 1)

        status_after = self.client.get(f"/projects/{project_id}/memory/status")
        self.assertEqual(status_after.status_code, 200)
        self.assertTrue(status_after.json()["has_memory"])
        self.assertFalse(status_after.json()["memory_stale"])
        self.assertEqual(status_after.json()["memory_version"], 1)

    def test_project_memory_batch_rebuild_only_updates_requested_projects(self):
        with Session(self.engine) as session:
            project_a = Project(name="Memory A", client="Client", description="Alpha project")
            project_b = Project(name="Memory B", client="Client", description="Beta project")
            project_c = Project(
                name="Memory C",
                client="Client",
                description="Gamma project",
                memory_stale=False,
                memory_version=2,
                context_memory_json=json.dumps({"project_brief": "ready"}, ensure_ascii=False),
            )
            session.add(project_a)
            session.add(project_b)
            session.add(project_c)
            session.commit()
            session.refresh(project_a)
            session.refresh(project_b)
            session.refresh(project_c)
            project_ids = [project_a.id, project_b.id, project_c.id]

        async def fake_complete(messages, max_tokens=4000):
            prompt = messages[0]["content"]
            if "Alpha project" in prompt:
                brief = "Alpha memory brief"
            elif "Beta project" in prompt:
                brief = "Beta memory brief"
            else:
                brief = "Unexpected"
            return json.dumps(
                {
                    "project_brief": brief,
                    "current_stage": "delivery",
                    "recent_progress": [],
                    "key_risks": [],
                    "open_questions": [],
                    "next_actions": ["Confirm scope"],
                    "important_documents": [],
                    "financial_status": "Healthy",
                    "delivery_signals": [],
                    "stakeholder_notes": [],
                },
                ensure_ascii=False,
            )

        with patch_project_llm(complete=AsyncMock(side_effect=fake_complete)):
            rebuild_resp = self.client.post(
                "/projects/memory/rebuild-batch",
                json={"project_ids": project_ids, "stale_only": True},
            )

        self.assertEqual(rebuild_resp.status_code, 200)
        body = rebuild_resp.json()
        self.assertTrue(body["ok"])
        self.assertEqual(body["requested_count"], 3)
        self.assertEqual(body["rebuilt_count"], 2)
        rebuilt_ids = {item["project_id"] for item in body["rebuilt"]}
        self.assertEqual(rebuilt_ids, {project_ids[0], project_ids[1]})
        self.assertEqual(body["skipped"], [{"project_id": project_ids[2], "reason": "not_stale"}])

        with Session(self.engine) as session:
            refreshed_a = session.get(Project, project_ids[0])
            refreshed_b = session.get(Project, project_ids[1])
            refreshed_c = session.get(Project, project_ids[2])

            self.assertFalse(refreshed_a.memory_stale)
            self.assertFalse(refreshed_b.memory_stale)
            self.assertEqual(refreshed_a.memory_version, 1)
            self.assertEqual(refreshed_b.memory_version, 1)
            self.assertEqual(json.loads(refreshed_a.context_memory_json)["project_brief"], "Alpha memory brief")
            self.assertEqual(json.loads(refreshed_b.context_memory_json)["project_brief"], "Beta memory brief")
            self.assertEqual(refreshed_c.memory_version, 2)
            self.assertFalse(refreshed_c.memory_stale)

    def test_project_memory_batch_rebuild_can_generate_missing_memory(self):
        with Session(self.engine) as session:
            project_a = Project(name="Missing Memory A", client="Client", description="Alpha project")
            project_b = Project(name="Missing Memory B", client="Client", description="Beta project")
            session.add(project_a)
            session.add(project_b)
            session.commit()
            session.refresh(project_a)
            session.refresh(project_b)
            project_ids = [project_a.id, project_b.id]

        async def fake_complete(messages, max_tokens=4000):
            prompt = messages[0]["content"]
            brief = "Alpha generated brief" if "Alpha project" in prompt else "Beta generated brief"
            return json.dumps(
                {
                    "project_brief": brief,
                    "current_stage": "lead",
                    "recent_progress": [],
                    "key_risks": [],
                    "open_questions": [],
                    "next_actions": ["Confirm scope"],
                    "important_documents": [],
                    "financial_status": "Not started",
                    "delivery_signals": [],
                    "stakeholder_notes": [],
                },
                ensure_ascii=False,
            )

        with patch_project_llm(complete=AsyncMock(side_effect=fake_complete)):
            rebuild_resp = self.client.post(
                "/projects/memory/rebuild-batch",
                json={"project_ids": project_ids, "stale_only": False},
            )

        self.assertEqual(rebuild_resp.status_code, 200)
        body = rebuild_resp.json()
        self.assertTrue(body["ok"])
        self.assertEqual(body["requested_count"], 2)
        self.assertEqual(body["rebuilt_count"], 2)
        self.assertEqual(body["skipped"], [])

        with Session(self.engine) as session:
            refreshed_a = session.get(Project, project_ids[0])
            refreshed_b = session.get(Project, project_ids[1])
            self.assertEqual(refreshed_a.memory_version, 1)
            self.assertEqual(refreshed_b.memory_version, 1)
            self.assertFalse(refreshed_a.memory_stale)
            self.assertFalse(refreshed_b.memory_stale)
            self.assertEqual(json.loads(refreshed_a.context_memory_json)["project_brief"], "Alpha generated brief")
            self.assertEqual(json.loads(refreshed_b.context_memory_json)["project_brief"], "Beta generated brief")

    def test_auto_summarize_file_persists_generated_summary(self):
        with Session(self.engine) as session:
            project = Project(name="Upload Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            pid = project.id

            file_path = Path("projects") / str(project.id) / "brief.md"
            full_path = self.uploads_dir / file_path
            full_path.parent.mkdir(parents=True, exist_ok=True)
            full_path.write_text("# Brief", encoding="utf-8")

            project_file = ProjectFile(
                project_id=project.id,
                name="brief.md",
                file_type="md",
                path=str(file_path),
                size_bytes=full_path.stat().st_size,
                origin="uploaded",
            )
            session.add(project_file)
            session.commit()
            session.refresh(project_file)
            file_id = project_file.id

        async def fake_summarize(file_id, *, file_path, file_type, extract_file_text, complete, session_factory):
            with session_factory() as s:
                pf = s.get(ProjectFile, file_id)
                if pf:
                    pf.summary = "Concise consultant summary"
                    s.add(pf)
                    s.commit()

        with patch("app.routers.projects.summarize_uploaded_project_file", side_effect=fake_summarize), patch("app.database.engine", self.engine):
            asyncio.run(projects_router_module._auto_summarize_file(file_id, str(full_path), "md", pid, None))

        with Session(self.engine) as session:
            refreshed = session.get(ProjectFile, file_id)
            self.assertEqual(refreshed.summary, "Concise consultant summary")

    def test_create_project_initializes_default_folders(self):
        resp = self.client.post(
            "/projects",
            json={
                "name": "New Project",
                "client": "Client",
                "description": "Desc",
            },
        )
        self.assertEqual(resp.status_code, 201)
        project_id = resp.json()["id"]

        folders_resp = self.client.get(f"/projects/{project_id}/folders")
        self.assertEqual(folders_resp.status_code, 200)
        folder_names = [folder["name"] for folder in folders_resp.json()]
        self.assertEqual(
            folder_names,
            ["项目需求", "方案和报价", "项目交付文档", "项目归档信息"],
        )

    def test_list_projects_filters_by_member_user_id(self):
        with Session(self.engine) as session:
            user = User(
                email="filter-member@example.com",
                display_name="Filter Member",
                password_hash="hashed",
            )
            member_project = Project(name="Member Project", client="Client")
            other_project = Project(name="Other Project", client="Client")
            session.add(user)
            session.add(member_project)
            session.add(other_project)
            session.commit()
            session.refresh(user)
            session.refresh(member_project)
            session.refresh(other_project)
            session.add(ProjectMember(project_id=member_project.id, user_id=user.id))
            session.commit()
            user_id = user.id

        resp = self.client.get(f"/projects?member_user_id={user_id}")
        self.assertEqual(resp.status_code, 200)
        projects = resp.json()
        self.assertEqual(len(projects), 1)
        self.assertEqual(projects[0]["name"], "Member Project")

    def test_project_members_crud_and_detail_payload(self):
        with Session(self.engine) as session:
            project = Project(name="Gamma", client="Client")
            user = User(
                email="member@example.com",
                display_name="Alice",
                password_hash="hashed",
            )
            session.add(project)
            session.add(user)
            session.commit()
            session.refresh(project)
            session.refresh(user)
            project_id = project.id
            user_id = user.id

        add_resp = self.client.post(
            f"/projects/{project_id}/members",
            json={"user_id": user_id},
        )
        self.assertEqual(add_resp.status_code, 201)
        self.assertEqual(add_resp.json()["user"]["display_name"], "Alice")

        list_resp = self.client.get(f"/projects/{project_id}/members")
        self.assertEqual(list_resp.status_code, 200)
        members = list_resp.json()
        self.assertEqual(len(members), 1)
        self.assertEqual(members[0]["user_id"], user_id)
        self.assertEqual(members[0]["user"]["display_name"], "Alice")

        detail_resp = self.client.get(f"/projects/{project_id}/detail")
        self.assertEqual(detail_resp.status_code, 200)
        detail_members = detail_resp.json()["members"]
        self.assertEqual(len(detail_members), 1)
        self.assertEqual(detail_members[0]["user_id"], user_id)
        self.assertEqual(detail_members[0]["user"]["display_name"], "Alice")

        remove_resp = self.client.delete(f"/projects/{project_id}/members/{user_id}")
        self.assertEqual(remove_resp.status_code, 200)
        self.assertEqual(remove_resp.json(), {"ok": True})

        detail_after_remove = self.client.get(f"/projects/{project_id}/detail")
        self.assertEqual(detail_after_remove.status_code, 200)
        self.assertEqual(detail_after_remove.json()["members"], [])

        with Session(self.engine) as session:
            self.assertEqual(
                session.exec(
                    select(ProjectMember).where(ProjectMember.project_id == project_id)
                ).all(),
                [],
            )

    def test_project_members_prevent_duplicate_member(self):
        with Session(self.engine) as session:
            project = Project(name="Delta", client="Client")
            user = User(
                email="duplicate@example.com",
                display_name="Bob",
                password_hash="hashed",
            )
            session.add(project)
            session.add(user)
            session.commit()
            session.refresh(project)
            session.refresh(user)
            session.add(ProjectMember(project_id=project.id, user_id=user.id))
            session.commit()
            project_id = project.id
            user_id = user.id

        resp = self.client.post(
            f"/projects/{project_id}/members",
            json={"user_id": user_id},
        )
        self.assertEqual(resp.status_code, 409)
        self.assertIn("already a member", resp.json()["detail"])

    def test_project_detail_smoke_includes_aggregate_sections(self):
        with Session(self.engine) as session:
            project = Project(
                name="Aggregate Project",
                client="Client",
                md_notes="# Overview",
                contract_amount=880000,
            )
            session.add(project)
            session.commit()
            session.refresh(project)

            folder = ProjectFolder(project_id=project.id, name="Archive", sort_order=3)
            member_user = User(
                email="aggregate-member@example.com",
                display_name="Eve",
                password_hash="hashed",
            )
            session.add(folder)
            session.add(member_user)
            session.commit()
            session.refresh(folder)
            session.refresh(member_user)

            session.add(ProjectMember(project_id=project.id, user_id=member_user.id))
            session.add(
                ProjectPayment(
                    project_id=project.id,
                    amount=120000,
                    payment_date="2026-04-10",
                    note="Kickoff invoice",
                    payment_type="invoiced",
                )
            )
            session.add(
                ProjectTodo(
                    project_id=project.id,
                    content="Prepare workshop",
                    due_date="2026-04-20",
                    assigned_to_user_id=member_user.id,
                )
            )
            session.commit()
            project_id = project.id

        resp = self.client.get(f"/projects/{project_id}/detail")
        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["project"]["name"], "Aggregate Project")
        self.assertEqual(body["md_notes"], "# Overview")
        self.assertEqual(len(body["folders"]), 1)
        self.assertEqual(body["folders"][0]["name"], "Archive")
        self.assertEqual(len(body["members"]), 1)
        self.assertEqual(body["members"][0]["user"]["display_name"], "Eve")
        self.assertEqual(len(body["todos"]), 1)
        self.assertEqual(body["todos"][0]["due_date"], "2026-04-20")
        self.assertEqual(body["financials"]["contract_amount"], 880000)
        self.assertEqual(body["financials"]["total_invoiced"], 120000)
        self.assertEqual(len(body["financials"]["payments"]), 1)

    def test_financials_endpoint_normalizes_expense_and_aggregates_totals(self):
        with Session(self.engine) as session:
            project = Project(name="Financial Project", client="Client", contract_amount=500000)
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        invoice_resp = self.client.post(
            f"/projects/{project_id}/financials",
            json={
                "amount": 120000,
                "payment_date": "2026-04-01",
                "note": "Invoice issued",
                "payment_type": "invoiced",
            },
        )
        self.assertEqual(invoice_resp.status_code, 201)
        self.assertEqual(invoice_resp.json()["amount"], 120000)

        expense_resp = self.client.post(
            f"/projects/{project_id}/financials",
            json={
                "amount": 3000,
                "payment_date": "2026-04-02",
                "note": "Travel cost",
                "payment_type": "expense",
            },
        )
        self.assertEqual(expense_resp.status_code, 201)
        self.assertEqual(expense_resp.json()["amount"], -3000)

        received_resp = self.client.post(
            f"/projects/{project_id}/financials",
            json={
                "amount": 50000,
                "payment_date": "2026-04-03",
                "note": "Advance payment",
                "payment_type": "received",
            },
        )
        self.assertEqual(received_resp.status_code, 201)

        detail_resp = self.client.get(f"/projects/{project_id}/detail")
        self.assertEqual(detail_resp.status_code, 200)
        financials = detail_resp.json()["financials"]
        self.assertEqual(financials["contract_amount"], 500000)
        self.assertEqual(financials["total_invoiced"], 120000)
        self.assertEqual(financials["total_expense"], 3000)
        self.assertEqual(financials["total_received"], 50000)
        self.assertEqual(financials["uncollected"], 70000)
        self.assertEqual(financials["remaining"], 450000)
        self.assertEqual(len(financials["payments"]), 3)

        payment_id = expense_resp.json()["id"]
        delete_resp = self.client.delete(f"/projects/{project_id}/financials/{payment_id}")
        self.assertEqual(delete_resp.status_code, 200)
        self.assertEqual(delete_resp.json(), {"ok": True})

        financials_resp = self.client.get(f"/projects/{project_id}/financials")
        self.assertEqual(financials_resp.status_code, 200)
        self.assertEqual(financials_resp.json()["total_expense"], 0)
        self.assertEqual(len(financials_resp.json()["payments"]), 2)

        detail_after_delete = self.client.get(f"/projects/{project_id}/detail")
        self.assertEqual(detail_after_delete.status_code, 200)
        self.assertEqual(detail_after_delete.json()["financials"]["total_expense"], 0)
        self.assertEqual(len(detail_after_delete.json()["financials"]["payments"]), 2)

    def test_list_my_todos_returns_only_current_user_pending_items_with_due_dates(self):
        with Session(self.engine) as session:
            other_user = User(
                email="other@example.com",
                display_name="Other User",
                password_hash="hashed",
            )
            project_a = Project(name="Alpha", client="Client")
            project_b = Project(name="Beta", client="Client")
            session.add(other_user)
            session.add(project_a)
            session.add(project_b)
            session.commit()
            session.refresh(other_user)
            session.refresh(project_a)
            session.refresh(project_b)

            session.add(
                ProjectTodo(
                    project_id=project_a.id,
                    content="Current user pending",
                    due_date="2026-05-01",
                    assigned_to_user_id=self.current_user_id,
                    updated_at=utc_now_naive() + timedelta(minutes=2),
                )
            )
            session.add(
                ProjectTodo(
                    project_id=project_b.id,
                    content="Current user latest",
                    due_date="2026-05-04",
                    assigned_to_user_id=self.current_user_id,
                    updated_at=utc_now_naive() + timedelta(minutes=3),
                )
            )
            session.add(
                ProjectTodo(
                    project_id=project_b.id,
                    content="Current user done",
                    due_date="2026-05-03",
                    assigned_to_user_id=self.current_user_id,
                    is_done=True,
                    updated_at=utc_now_naive() + timedelta(minutes=1),
                )
            )
            session.add(
                ProjectTodo(
                    project_id=project_b.id,
                    content="Other user pending",
                    due_date="2026-05-02",
                    assigned_to_user_id=other_user.id,
                    updated_at=utc_now_naive(),
                )
            )
            session.commit()

        resp = self.client.get("/projects/todos/my")
        self.assertEqual(resp.status_code, 200)
        todos = resp.json()
        self.assertEqual(len(todos), 2)
        self.assertEqual(todos[0]["content"], "Current user latest")
        self.assertEqual(todos[0]["project_name"], "Beta")
        self.assertEqual(todos[0]["due_date"], "2026-05-04")
        self.assertEqual(todos[0]["assigned_to_user_id"], self.current_user_id)
        self.assertIsNotNone(todos[0]["created_at"])
        self.assertIsNotNone(todos[0]["updated_at"])
        self.assertFalse(todos[0]["is_done"])
        self.assertEqual(todos[1]["content"], "Current user pending")
        self.assertEqual(todos[1]["project_name"], "Alpha")
        self.assertEqual(todos[1]["due_date"], "2026-05-01")

    def test_update_project_document_persists_content(self):
        with Session(self.engine) as session:
            project = Project(name="Gamma", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

            folder = ProjectFolder(project_id=project_id, name="Notes", sort_order=0)
            session.add(folder)
            session.commit()
            session.refresh(folder)

            path = Path("projects") / str(project_id) / "manual_doc.md"
            full_path = self.uploads_dir / path
            full_path.parent.mkdir(parents=True, exist_ok=True)
            full_path.write_text("# Initial", encoding="utf-8")

            project_file = ProjectFile(
                project_id=project_id,
                folder_id=folder.id,
                name="Manual Doc.md",
                file_type="md",
                path=str(path),
                size_bytes=full_path.stat().st_size,
            )
            session.add(project_file)
            session.commit()
            session.refresh(project_file)
            file_id = project_file.id

        get_resp = self.client.get(f"/projects/{project_id}/documents/{file_id}")
        self.assertEqual(get_resp.status_code, 200)
        self.assertEqual(get_resp.json()["content"], "# Initial")

        patch_resp = self.client.patch(
            f"/projects/{project_id}/documents/{file_id}",
            json={"content": "# Updated\n\nNew content"},
        )
        self.assertEqual(patch_resp.status_code, 200)

        saved_text = (self.uploads_dir / "projects" / str(project_id) / "manual_doc.md").read_text(encoding="utf-8")
        self.assertIn("# Updated", saved_text)
        self.assertIn("New content", saved_text)

        rename_resp = self.client.patch(
            f"/projects/{project_id}/documents/{file_id}",
            json={"name": "Renamed Consulting Note"},
        )
        self.assertEqual(rename_resp.status_code, 200)
        self.assertEqual(rename_resp.json()["name"], "Renamed_Consulting_Note.md")

    def test_create_project_document_sanitizes_name_and_keeps_folder(self):
        with Session(self.engine) as session:
            project = Project(name="Doc Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)

            folder = ProjectFolder(project_id=project.id, name="Knowledge", sort_order=1)
            session.add(folder)
            session.commit()
            session.refresh(folder)
            project_id = project.id
            folder_id = folder.id

        resp = self.client.post(
            f"/projects/{project_id}/documents",
            json={
                "name": "Client Kickoff Notes",
                "content": "# Kickoff\n\nAgenda",
                "folder_id": folder_id,
            },
        )
        self.assertEqual(resp.status_code, 201)
        body = resp.json()
        self.assertEqual(body["project_id"], project_id)
        self.assertEqual(body["folder_id"], folder_id)
        self.assertEqual(body["name"], "Client_Kickoff_Notes.md")
        self.assertEqual(body["file_type"], "md")

        saved_path = self.uploads_dir / body["path"]
        self.assertTrue(saved_path.exists())
        self.assertIn("Agenda", saved_path.read_text(encoding="utf-8"))

    def test_delete_folder_unlinks_files_instead_of_deleting_them(self):
        with Session(self.engine) as session:
            project = Project(name="Folder Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)

            folder = ProjectFolder(project_id=project.id, name="Archive", sort_order=2)
            session.add(folder)
            session.commit()
            session.refresh(folder)

            path = Path("projects") / str(project.id) / "folder_file.md"
            full_path = self.uploads_dir / path
            full_path.parent.mkdir(parents=True, exist_ok=True)
            full_path.write_text("# Folder file", encoding="utf-8")

            project_file = ProjectFile(
                project_id=project.id,
                folder_id=folder.id,
                name="Folder File.md",
                file_type="md",
                path=str(path),
                size_bytes=full_path.stat().st_size,
            )
            session.add(project_file)
            session.commit()
            session.refresh(project_file)
            project_id = project.id
            folder_id = folder.id
            file_id = project_file.id

        delete_resp = self.client.delete(f"/projects/{project_id}/folders/{folder_id}")
        self.assertEqual(delete_resp.status_code, 200)
        self.assertEqual(delete_resp.json(), {"ok": True})

        with Session(self.engine) as session:
            saved_file = session.get(ProjectFile, file_id)
            self.assertIsNotNone(saved_file)
            self.assertIsNone(saved_file.folder_id)
            self.assertIsNone(session.get(ProjectFolder, folder_id))

        detail_resp = self.client.get(f"/projects/{project_id}/detail")
        self.assertEqual(detail_resp.status_code, 200)
        self.assertEqual(len(detail_resp.json()["files"]), 1)

    def test_project_file_upload_download_and_delete_flow(self):
        with Session(self.engine) as session:
            project = Project(name="Upload Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        upload_resp = self.client.post(
            f"/projects/{project_id}/files",
            files={"file": ("brief.txt", b"hello project file", "text/plain")},
        )
        self.assertEqual(upload_resp.status_code, 201)
        body = upload_resp.json()
        file_id = body["id"]
        self.assertEqual(body["name"], "brief.txt")
        self.assertEqual(body["file_type"], "txt")

        with Session(self.engine) as session:
            saved_file = session.get(ProjectFile, file_id)
            self.assertIsNotNone(saved_file)
            saved_path = self.uploads_dir / saved_file.path
            self.assertTrue(saved_path.exists())
            self.assertEqual(saved_path.read_text(encoding="utf-8"), "hello project file")

        download_resp = self.client.get(f"/projects/{project_id}/files/{file_id}/download")
        self.assertEqual(download_resp.status_code, 200)
        self.assertEqual(download_resp.content, b"hello project file")

        delete_resp = self.client.delete(f"/projects/{project_id}/files/{file_id}")
        self.assertEqual(delete_resp.status_code, 200)
        self.assertEqual(delete_resp.json(), {"ok": True, "archived": True})

        with Session(self.engine) as session:
            archived_file = session.get(ProjectFile, file_id)
            self.assertIsNotNone(archived_file)
            self.assertIsNotNone(archived_file.deleted_at)

        self.assertTrue(saved_path.exists())

        download_after_delete = self.client.get(f"/projects/{project_id}/files/{file_id}/download")
        self.assertEqual(download_after_delete.status_code, 404)

        restore_resp = self.client.post(f"/projects/{project_id}/files/{file_id}/restore")
        self.assertEqual(restore_resp.status_code, 200)
        restored_download = self.client.get(f"/projects/{project_id}/files/{file_id}/download")
        self.assertEqual(restored_download.status_code, 200)
        self.assertEqual(restored_download.content, b"hello project file")

    def test_list_files_reflects_upload_and_delete(self):
        with Session(self.engine) as session:
            project = Project(name="List Files Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        upload_resp = self.client.post(
            f"/projects/{project_id}/files",
            files={"file": ("scope.txt", b"scope draft", "text/plain")},
        )
        self.assertEqual(upload_resp.status_code, 201)
        file_id = upload_resp.json()["id"]

        list_resp = self.client.get(f"/projects/{project_id}/files")
        self.assertEqual(list_resp.status_code, 200)
        self.assertEqual(len(list_resp.json()), 1)
        self.assertEqual(list_resp.json()[0]["id"], file_id)
        self.assertEqual(list_resp.json()[0]["name"], "scope.txt")

        delete_resp = self.client.delete(f"/projects/{project_id}/files/{file_id}")
        self.assertEqual(delete_resp.status_code, 200)
        self.assertEqual(delete_resp.json(), {"ok": True, "archived": True})

        list_after_delete_resp = self.client.get(f"/projects/{project_id}/files")
        self.assertEqual(list_after_delete_resp.status_code, 200)
        self.assertEqual(list_after_delete_resp.json(), [])

    def test_delete_project_file_cleans_task_artifact_and_generated_file_references(self):
        with Session(self.engine) as session:
            project = Project(name="Generated File Project", client="Client")
            conversation = Conversation(title="Project chat", project_id=None)
            session.add(project)
            session.add(conversation)
            session.commit()
            session.refresh(project)
            session.refresh(conversation)

            file_path = Path("projects") / str(project.id) / "generated.md"
            full_path = self.uploads_dir / file_path
            full_path.parent.mkdir(parents=True, exist_ok=True)
            full_path.write_text("# Generated", encoding="utf-8")

            project_file = ProjectFile(
                project_id=project.id,
                name="generated.md",
                file_type="md",
                path=str(file_path),
                size_bytes=full_path.stat().st_size,
                origin="ai_generated",
            )
            session.add(project_file)
            session.commit()
            session.refresh(project_file)

            derivative = ProjectFile(
                project_id=project.id,
                source_file_id=project_file.id,
                name="generated_extracted.md",
                file_type="md",
                path=str(Path("projects") / str(project.id) / "generated_extracted.md"),
                size_bytes=0,
                origin="markdown_derivative",
            )
            task = TaskRun(project_id=project.id, conversation_id=conversation.id, task_type="text", goal="write")
            generated_file = GeneratedFile(
                conversation_id=conversation.id,
                project_id=project.id,
                name=project_file.name,
                file_type=project_file.file_type,
                path=project_file.path,
                size_bytes=project_file.size_bytes,
            )
            session.add(derivative)
            session.add(task)
            session.add(generated_file)
            session.commit()
            session.refresh(derivative)
            session.refresh(task)

            artifact = TaskArtifact(
                task_run_id=task.id,
                project_file_id=project_file.id,
                name=project_file.name,
                file_type=project_file.file_type,
                path=project_file.path,
            )
            session.add(artifact)
            session.commit()
            session.refresh(artifact)

            project_id = project.id
            file_id = project_file.id
            derivative_id = derivative.id
            artifact_id = artifact.id

        delete_resp = self.client.delete(f"/projects/{project_id}/files/{file_id}")
        self.assertEqual(delete_resp.status_code, 200)
        self.assertEqual(delete_resp.json(), {"ok": True, "archived": True})

        with Session(self.engine) as session:
            archived_file = session.get(ProjectFile, file_id)
            self.assertIsNotNone(archived_file)
            self.assertIsNotNone(archived_file.deleted_at)
            saved_derivative = session.get(ProjectFile, derivative_id)
            self.assertIsNotNone(saved_derivative)
            self.assertIsNone(saved_derivative.source_file_id)
            saved_artifact = session.get(TaskArtifact, artifact_id)
            self.assertIsNotNone(saved_artifact)
            self.assertEqual(saved_artifact.project_file_id, file_id)
            generated_files = session.exec(
                select(GeneratedFile).where(
                    GeneratedFile.project_id == project_id,
                    GeneratedFile.path == str(file_path),
                )
            ).all()
            self.assertEqual(len(generated_files), 1)

        self.assertTrue(full_path.exists())

    def test_delete_project_cascades_related_records(self):
        with Session(self.engine) as session:
            project = Project(name="Cascade Project", client="Client")
            member_user = User(
                email="cascade-member@example.com",
                display_name="Cascade Member",
                password_hash="hashed",
            )
            session.add(project)
            session.add(member_user)
            session.commit()
            session.refresh(project)
            session.refresh(member_user)

            folder = ProjectFolder(project_id=project.id, name="Archive", sort_order=1)
            milestone = Milestone(project_id=project.id, title="Kickoff", due_date="2026-04-30")
            payment = ProjectPayment(project_id=project.id, amount=5000, payment_date="2026-04-20", payment_type="received")
            todo = ProjectTodo(project_id=project.id, content="Prepare deck")
            member = ProjectMember(project_id=project.id, user_id=member_user.id)
            conv = Conversation(project_id=project.id, title="Project Chat")
            session.add(folder)
            session.add(milestone)
            session.add(payment)
            session.add(todo)
            session.add(member)
            session.add(conv)
            session.commit()
            session.refresh(folder)
            session.refresh(conv)

            path = Path("projects") / str(project.id) / "cascade_file.md"
            full_path = self.uploads_dir / path
            full_path.parent.mkdir(parents=True, exist_ok=True)
            full_path.write_text("# Cascade", encoding="utf-8")

            project_file = ProjectFile(
                project_id=project.id,
                folder_id=folder.id,
                name="Cascade File.md",
                file_type="md",
                path=str(path),
                size_bytes=full_path.stat().st_size,
            )
            session.add(project_file)
            session.commit()
            session.refresh(project_file)

            message = Message(conversation_id=conv.id, role="assistant", content="hello")
            session.add(message)
            session.commit()
            session.refresh(message)

            generated_file = GeneratedFile(
                conversation_id=conv.id,
                project_id=project.id,
                name="Generated Deck.pptx",
                file_type="pptx",
                path="projects/generated-deck.pptx",
            )
            task = TaskRun(project_id=project.id, conversation_id=conv.id, task_type="generate_client_ppt", goal="deck")
            knowledge_document = KnowledgeDocument(
                name="Knowledge.md",
                file_type="md",
                path="projects/knowledge.md",
                project_id=project.id,
            )
            memory_summary = ProjectMemorySummary(
                project_id=project.id,
                summary_type="brief",
                language="zh",
                memory_version=1,
                content="summary",
            )
            memory_snapshot = ProjectMemorySnapshot(
                project_id=project.id,
                memory_version=1,
                trigger="test",
                memory_json="{}",
            )
            scheduled_task = ScheduledTask(name="Daily review", project_id=project.id, frequency="daily")
            session.add(generated_file)
            session.add(task)
            session.add(knowledge_document)
            session.add(memory_summary)
            session.add(memory_snapshot)
            session.add(scheduled_task)
            session.commit()
            session.refresh(generated_file)
            session.refresh(task)
            session.refresh(knowledge_document)

            tool_call = ToolCall(
                conversation_id=conv.id,
                message_id=message.id,
                tool_name="write_project_office_document",
                output_file_id=generated_file.id,
                status="success",
            )
            task_step = TaskStep(task_run_id=task.id, key="create", title="Create", step_type="write_project_office_document")
            document_chunk = DocumentChunk(document_id=knowledge_document.id, chunk_index=0, content="chunk")
            session.add(tool_call)
            session.add(task_step)
            session.add(document_chunk)
            session.commit()
            session.refresh(task_step)

            task_event = TaskEvent(task_run_id=task.id, step_id=task_step.id, event_type="step_completed", message="done")
            task_artifact = TaskArtifact(
                task_run_id=task.id,
                step_id=task_step.id,
                project_file_id=project_file.id,
                name="Text artifact",
                file_type="text",
                metadata_json='{"content":"hello"}',
            )
            session.add(task_event)
            session.add(task_artifact)
            session.commit()
            session.refresh(tool_call)
            session.refresh(task_event)
            session.refresh(task_artifact)
            session.refresh(document_chunk)
            session.refresh(memory_summary)
            session.refresh(memory_snapshot)
            session.refresh(scheduled_task)

            project_id = project.id
            folder_id = folder.id
            milestone_id = milestone.id
            payment_id = payment.id
            todo_id = todo.id
            member_id = member.id
            conv_id = conv.id
            message_id = message.id
            file_id = project_file.id
            generated_file_id = generated_file.id
            tool_call_id = tool_call.id
            task_id = task.id
            task_step_id = task_step.id
            task_event_id = task_event.id
            task_artifact_id = task_artifact.id
            knowledge_document_id = knowledge_document.id
            document_chunk_id = document_chunk.id
            memory_summary_id = memory_summary.id
            memory_snapshot_id = memory_snapshot.id
            scheduled_task_id = scheduled_task.id

        delete_resp = self.client.delete(f"/projects/{project_id}")
        self.assertEqual(delete_resp.status_code, 200)
        self.assertEqual(delete_resp.json(), {"ok": True})

        with Session(self.engine) as session:
            self.assertIsNone(session.get(Project, project_id))
            self.assertIsNone(session.get(ProjectFolder, folder_id))
            self.assertIsNone(session.get(Milestone, milestone_id))
            self.assertIsNone(session.get(ProjectPayment, payment_id))
            self.assertIsNone(session.get(ProjectTodo, todo_id))
            self.assertIsNone(session.get(ProjectMember, member_id))
            self.assertIsNone(session.get(Conversation, conv_id))
            self.assertIsNone(session.get(Message, message_id))
            self.assertIsNone(session.get(ProjectFile, file_id))
            self.assertIsNone(session.get(GeneratedFile, generated_file_id))
            self.assertIsNone(session.get(ToolCall, tool_call_id))
            self.assertIsNone(session.get(TaskRun, task_id))
            self.assertIsNone(session.get(TaskStep, task_step_id))
            self.assertIsNone(session.get(TaskEvent, task_event_id))
            self.assertIsNone(session.get(TaskArtifact, task_artifact_id))
            self.assertIsNone(session.get(KnowledgeDocument, knowledge_document_id))
            self.assertIsNone(session.get(DocumentChunk, document_chunk_id))
            self.assertIsNone(session.get(ProjectMemorySummary, memory_summary_id))
            self.assertIsNone(session.get(ProjectMemorySnapshot, memory_snapshot_id))
            self.assertIsNone(session.get(ScheduledTask, scheduled_task_id))

    def test_project_milestone_crud_flow(self):
        with Session(self.engine) as session:
            project = Project(name="Milestone Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        create_resp = self.client.post(
            f"/projects/{project_id}/milestones",
            json={"title": "Kickoff", "due_date": "2026-05-01"},
        )
        self.assertEqual(create_resp.status_code, 201)
        milestone_id = create_resp.json()["id"]
        self.assertEqual(create_resp.json()["title"], "Kickoff")
        self.assertEqual(create_resp.json()["due_date"], "2026-05-01")
        self.assertFalse(create_resp.json()["is_done"])

        list_resp = self.client.get(f"/projects/{project_id}/milestones")
        self.assertEqual(list_resp.status_code, 200)
        self.assertEqual(len(list_resp.json()), 1)
        self.assertEqual(list_resp.json()[0]["id"], milestone_id)

        update_resp = self.client.patch(
            f"/projects/{project_id}/milestones/{milestone_id}",
            json={"title": "Kickoff Complete", "is_done": True},
        )
        self.assertEqual(update_resp.status_code, 200)
        self.assertEqual(update_resp.json()["title"], "Kickoff Complete")
        self.assertTrue(update_resp.json()["is_done"])

        delete_resp = self.client.delete(f"/projects/{project_id}/milestones/{milestone_id}")
        self.assertEqual(delete_resp.status_code, 200)
        self.assertEqual(delete_resp.json(), {"ok": True})

        final_list_resp = self.client.get(f"/projects/{project_id}/milestones")
        self.assertEqual(final_list_resp.status_code, 200)
        self.assertEqual(final_list_resp.json(), [])

    def test_generate_project_context_uses_current_project_files_and_milestones_only(self):
        captured = {}

        async def fake_stream(messages, max_tokens=4000):
            captured["prompt"] = messages[0]["content"]
            yield "- **Current summary**"

        with Session(self.engine) as session:
            project = Project(
                name="Context Project",
                client="Client",
                status="active",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Current project file summary",
                        "current_stage": "active",
                        "recent_progress": ["Current Milestone"],
                        "key_risks": [],
                        "open_questions": [],
                        "next_actions": ["Review Current Note.md"],
                        "important_documents": [
                            {"name": "Current Note.md", "reason": "Latest project note"}
                        ],
                        "financial_status": "",
                        "delivery_signals": [],
                        "stakeholder_notes": [],
                    },
                    ensure_ascii=False,
                ),
                memory_version=1,
                memory_stale=False,
            )
            sibling = Project(name="Sibling Project", client="Client", status="lead")
            session.add(project)
            session.add(sibling)
            session.commit()
            session.refresh(project)
            session.refresh(sibling)

            session.add(Milestone(project_id=project.id, title="Current Milestone", priority="high"))
            session.add(Milestone(project_id=sibling.id, title="Sibling Milestone", priority="medium"))

            current_path = Path("projects") / str(project.id) / "current_note.md"
            sibling_path = Path("projects") / str(sibling.id) / "sibling_note.md"
            current_full_path = self.uploads_dir / current_path
            sibling_full_path = self.uploads_dir / sibling_path
            current_full_path.parent.mkdir(parents=True, exist_ok=True)
            sibling_full_path.parent.mkdir(parents=True, exist_ok=True)
            current_full_path.write_text("# Current", encoding="utf-8")
            sibling_full_path.write_text("# Sibling", encoding="utf-8")

            session.add(
                ProjectFile(
                    project_id=project.id,
                    name="Current Note.md",
                    file_type="md",
                    path=str(current_path),
                    size_bytes=current_full_path.stat().st_size,
                    summary="Current project file summary",
                )
            )
            session.add(
                ProjectFile(
                    project_id=sibling.id,
                    name="Sibling Note.md",
                    file_type="md",
                    path=str(sibling_path),
                    size_bytes=sibling_full_path.stat().st_size,
                    summary="Sibling project file summary",
                )
            )
            session.commit()
            project_id = project.id

        with patch_project_llm(stream=fake_stream), patch("app.database.engine", self.engine):
            resp = self.client.post(f"/projects/{project_id}/generate-context")

        self.assertEqual(resp.status_code, 200)
        self.assertIn("Current Milestone", captured["prompt"])
        self.assertIn("Current Note.md", captured["prompt"])
        self.assertNotIn("Sibling Milestone", captured["prompt"])
        self.assertNotIn("Sibling Note.md", captured["prompt"])
        self.assertNotIn("Sibling project file summary", captured["prompt"])

        with Session(self.engine) as session:
            saved_project = session.get(Project, project_id)
            self.assertEqual(saved_project.context_summary, "- **Current summary**")

    def test_generate_project_context_respects_requested_language(self):
        captured = {}

        async def fake_stream(messages, max_tokens=4000):
            captured["prompt"] = messages[0]["content"]
            yield "- 当前摘要"

        with Session(self.engine) as session:
            project = Project(
                name="Localized Project",
                client="Client",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Alpha brief",
                        "current_stage": "delivery",
                        "next_actions": ["Confirm scope"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=1,
                memory_stale=False,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        with patch_project_llm(stream=fake_stream), patch("app.database.engine", self.engine):
            resp = self.client.post(
                f"/projects/{project_id}/generate-context",
                json={"language": "zh-CN"},
            )

        self.assertEqual(resp.status_code, 200)
        self.assertIn("Write the answer in Chinese", captured["prompt"])
        self.assertIn("当前摘要", resp.text)

    def test_memory_summarize_streams_sse_events(self):
        captured = {}

        async def fake_stream(messages, max_tokens=4000):
            captured["prompt"] = messages[0]["content"]
            yield "- chunk 1"
            yield " chunk 2"

        with Session(self.engine) as session:
            project = Project(
                name="Memory Stream Project",
                client="Client",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Alpha brief",
                        "current_stage": "delivery",
                        "key_risks": ["Timeline risk"],
                        "next_actions": ["Confirm scope"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=2,
                memory_stale=False,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        with patch_project_llm(stream=fake_stream):
            resp = self.client.post(
                f"/projects/{project_id}/memory/summarize",
                json={
                    "summary_type": "risk",
                    "stream": True,
                    "language": "en-US",
                },
            )

        self.assertEqual(resp.status_code, 200)
        self.assertIn("data: ", resp.text)
        self.assertIn('"type": "text"', resp.text)
        self.assertIn('"type": "done"', resp.text)
        self.assertIn("- chunk 1 chunk 2", resp.text)
        self.assertIn("Write the answer in English", captured["prompt"])
        resp.close()

    def test_memory_summarize_uses_cached_summary_without_calling_llm(self):
        with Session(self.engine) as session:
            project = Project(
                name="Cached Memory Project",
                client="Client",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Alpha brief",
                        "current_stage": "delivery",
                        "key_risks": ["Timeline risk"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=3,
                memory_stale=False,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            session.add(
                ProjectMemorySummary(
                    project_id=project.id,
                    summary_type="risk",
                    language="zh",
                    memory_version=3,
                    content="- cached risk summary",
                )
            )
            session.commit()
            project_id = project.id

        with patch_project_llm(complete=AsyncMock(side_effect=AssertionError("should not call llm"))):
            resp = self.client.post(
                f"/projects/{project_id}/memory/summarize",
                json={
                    "summary_type": "risk",
                    "language": "zh-CN",
                },
            )

        self.assertEqual(resp.status_code, 200)
        self.assertEqual(resp.json()["content"], "- cached risk summary")
        self.assertTrue(resp.json()["cached"])
        resp.close()

    def test_memory_summary_cache_endpoint_does_not_generate(self):
        with Session(self.engine) as session:
            project = Project(
                name="Cached Summary Read Project",
                client="Client",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Alpha brief",
                        "current_stage": "delivery",
                        "key_risks": ["Timeline risk"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=7,
                memory_stale=False,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            session.add(
                ProjectMemorySummary(
                    project_id=project.id,
                    summary_type="delivery",
                    language="zh",
                    memory_version=7,
                    content="- cached delivery summary",
                )
            )
            session.commit()
            project_id = project.id

        with patch_project_llm(complete=AsyncMock(side_effect=AssertionError("should not call llm"))):
            resp = self.client.get(
                f"/projects/{project_id}/memory/summaries/delivery",
                params={"language": "zh-CN"},
            )

        self.assertEqual(resp.status_code, 200)
        self.assertEqual(resp.json()["content"], "- cached delivery summary")
        self.assertTrue(resp.json()["cached"])
        resp.close()

        missing_resp = self.client.get(
            f"/projects/{project_id}/memory/summaries/risk",
            params={"language": "zh-CN"},
        )
        self.assertEqual(missing_resp.status_code, 404)
        missing_resp.close()

    def test_memory_summaries_generate_creates_all_views_with_one_llm_call(self):
        calls = []

        async def fake_complete(messages, max_tokens=4000):
            calls.append(messages[0]["content"])
            return json.dumps(
                {
                    "overview": "- overview summary",
                    "risk": "- risk summary",
                    "delivery": "- delivery summary",
                    "stakeholder": "- stakeholder summary",
                    "client-facing": "- client-facing summary",
                    "financial": "- financial summary",
                    "documents": "- documents summary",
                },
                ensure_ascii=False,
            )

        with Session(self.engine) as session:
            project = Project(
                name="All Views Project",
                client="Client",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Alpha brief",
                        "current_stage": "delivery",
                        "key_risks": ["Timeline risk"],
                        "delivery_signals": ["Milestone moving"],
                        "financial_status": "Invoice pending",
                    },
                    ensure_ascii=False,
                ),
                memory_version=8,
                memory_stale=False,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        with patch_project_llm(complete=AsyncMock(side_effect=fake_complete)):
            resp = self.client.post(
                f"/projects/{project_id}/memory/summaries/generate",
                json={"language": "en-US", "force_refresh": True},
            )

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertFalse(body["cached"])
        self.assertEqual(len(calls), 1)
        self.assertEqual(body["summaries"]["risk"]["content"], "- risk summary")
        self.assertEqual(body["summaries"]["documents"]["content"], "- documents summary")
        resp.close()

        with Session(self.engine) as session:
            cached_count = len(
                session.exec(
                    select(ProjectMemorySummary).where(ProjectMemorySummary.project_id == project_id)
                ).all()
            )
            refreshed_project = session.get(Project, project_id)

        self.assertEqual(cached_count, 7)
        self.assertEqual(refreshed_project.context_summary, "- overview summary")

    def test_memory_summaries_generate_falls_back_for_missing_views(self):
        calls = []

        async def fake_complete(messages, max_tokens=4000):
            prompt = messages[0]["content"]
            calls.append(prompt)
            if "Return ONLY a valid JSON object" in prompt:
                return json.dumps(
                    {
                        "overview": "- overview summary",
                        "risk": "- risk summary",
                    },
                    ensure_ascii=False,
                )
            if "Summary type: delivery" in prompt:
                return "- delivery summary"
            if "Summary type: stakeholder" in prompt:
                return "- stakeholder summary"
            if "Summary type: client-facing" in prompt:
                return "- client-facing summary"
            if "Summary type: financial" in prompt:
                return "- financial summary"
            if "Summary type: documents" in prompt:
                return "- documents summary"
            raise AssertionError(f"Unexpected prompt: {prompt[:120]}")

        with Session(self.engine) as session:
            project = Project(
                name="Fallback Views Project",
                client="Client",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Alpha brief",
                        "current_stage": "delivery",
                        "key_risks": ["Timeline risk"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=9,
                memory_stale=False,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        with patch_project_llm(complete=AsyncMock(side_effect=fake_complete)):
            resp = self.client.post(
                f"/projects/{project_id}/memory/summaries/generate",
                json={"language": "en-US", "force_refresh": True},
            )

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["summaries"]["overview"]["content"], "- overview summary")
        self.assertEqual(body["summaries"]["documents"]["content"], "- documents summary")
        self.assertEqual(len(calls), 6)
        resp.close()

    def test_memory_summarize_force_refresh_bypasses_cache_and_updates_it(self):
        async def fake_complete(messages, max_tokens=4000):
            return "- fresh risk summary"

        with Session(self.engine) as session:
            project = Project(
                name="Forced Memory Project",
                client="Client",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Alpha brief",
                        "current_stage": "delivery",
                        "key_risks": ["Timeline risk"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=4,
                memory_stale=False,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            session.add(
                ProjectMemorySummary(
                    project_id=project.id,
                    summary_type="risk",
                    language="en",
                    memory_version=4,
                    content="- stale cached risk summary",
                )
            )
            session.commit()
            project_id = project.id

        with patch_project_llm(complete=AsyncMock(side_effect=fake_complete)):
            resp = self.client.post(
                f"/projects/{project_id}/memory/summarize",
                json={
                    "summary_type": "risk",
                    "language": "en-US",
                    "force_refresh": True,
                },
            )

        self.assertEqual(resp.status_code, 200)
        self.assertEqual(resp.json()["content"], "- fresh risk summary")
        self.assertFalse(resp.json()["cached"])
        resp.close()

        with Session(self.engine) as session:
            cached = session.exec(
                select(ProjectMemorySummary).where(
                    ProjectMemorySummary.project_id == project_id,
                    ProjectMemorySummary.summary_type == "risk",
                    ProjectMemorySummary.language == "en",
                    ProjectMemorySummary.memory_version == 4,
                )
            ).first()

        self.assertIsNotNone(cached)
        self.assertEqual(cached.content, "- fresh risk summary")

    def test_memory_warm_summaries_batch_generates_common_cached_views(self):
        async def fake_complete(messages, max_tokens=4000):
            prompt = messages[0]["content"]
            if "Summary type: overview" in prompt:
                return "- overview summary"
            if "Summary type: risk" in prompt:
                return "- risk summary"
            if "Summary type: stakeholder" in prompt:
                return "- stakeholder summary"
            return "- fallback summary"

        with Session(self.engine) as session:
            project = Project(
                name="Warm Summary Project",
                client="Client",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Alpha brief",
                        "current_stage": "delivery",
                        "key_risks": ["Timeline risk"],
                        "stakeholder_notes": ["Finance sponsor engaged"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=5,
                memory_stale=False,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        with patch_project_llm(complete=AsyncMock(side_effect=fake_complete)):
            resp = self.client.post(
                "/projects/memory/warm-summaries-batch",
                json={
                    "project_ids": [project_id],
                    "summary_types": ["overview", "risk", "stakeholder"],
                    "language": "en-US",
                },
            )

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertTrue(body["ok"])
        self.assertEqual(body["processed_count"], 1)
        self.assertEqual(body["warmed_count"], 3)

        with Session(self.engine) as session:
            cached = session.exec(
                select(ProjectMemorySummary).where(
                    ProjectMemorySummary.project_id == project_id,
                    ProjectMemorySummary.language == "en",
                    ProjectMemorySummary.memory_version == 5,
                )
            ).all()

        self.assertEqual(len(cached), 3)
        self.assertEqual(sorted(item.summary_type for item in cached), ["overview", "risk", "stakeholder"])

    def test_project_meeting_briefing_combines_memory_client_and_delivery_signals(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Acme",
                industry="Manufacturing",
                client_memory_json=json.dumps(
                    {
                        "decision_patterns": ["CEO needs quantified upside"],
                        "lessons_learned": ["Lock change requests before kickoff"],
                        "sensitive_topics": ["Avoid price details in steering meeting"],
                    },
                    ensure_ascii=False,
                ),
                client_memory_version=2,
                client_memory_stale=False,
            )
            session.add(client)
            session.commit()
            session.refresh(client)
            session.add(
                ClientStakeholder(
                    client_id=client.id,
                    name="Jane",
                    role="Sponsor",
                    influence_type="decision",
                    relationship_status="supportive",
                    concerns="Needs board-ready ROI story",
                    sensitivities="Do not surprise procurement",
                    last_action="Confirm security answer",
                )
            )
            project = Project(
                name="Meeting Project",
                client="Acme",
                status="delivering",
                md_notes="客户上次会议要求补充安全合规结论。",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Second phase rollout",
                        "current_objective": "Align on launch scope",
                        "key_risks": ["Security review is still open"],
                        "open_questions": ["Who signs off UAT?"],
                        "next_actions": ["Send revised timeline"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=3,
                memory_stale=False,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id
            session.add(Milestone(project_id=project_id, title="UAT signoff", priority="high", due_date="2026-04-25"))
            session.add(ProjectTodo(project_id=project_id, content="Prepare steering deck", due_date="2026-04-24"))
            conversation = Conversation(project_id=project_id, title="Steering prep")
            session.add(conversation)
            session.commit()
            session.refresh(conversation)
            session.add(
                Message(
                    conversation_id=conversation.id,
                    role="user",
                    content="李总监提醒业务侧还在等安全合规问题的书面答复。",
                )
            )
            session.commit()

        resp = self.client.get(f"/projects/{project_id}/briefing")

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["project"]["name"], "Meeting Project")
        self.assertEqual(body["client"]["name"], "Acme")
        self.assertIn("Align on launch scope", body["meeting_card"]["say"])
        self.assertIn("Who signs off UAT?", body["meeting_card"]["confirm"])
        self.assertIn("CEO needs quantified upside", body["meeting_card"]["experience"])
        self.assertEqual(body["stakeholders"][0]["name"], "Jane")
        self.assertEqual(body["signals"]["upcoming_milestones"][0]["title"], "UAT signoff")
        self.assertEqual(body["signals"]["communication_sources"][0]["type"], "markdown_note")
        self.assertEqual(body["signals"]["communication_sources"][0]["target"], "notes")
        self.assertIn("安全合规", body["signals"]["communication_sources"][0]["excerpt"])
        self.assertEqual(body["signals"]["communication_sources"][1]["type"], "chat")
        self.assertEqual(body["signals"]["communication_sources"][1]["target"], "chat")
        self.assertIsNotNone(body["signals"]["communication_sources"][1]["conversation_id"])

    def test_project_meeting_briefing_refine_uses_cache_for_same_sources(self):
        with Session(self.engine) as session:
            project = Project(
                name="Refine Project",
                client="Acme",
                status="delivering",
                md_notes="客户关心上线节奏和风险边界。",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Launch readiness",
                        "current_objective": "Align launch scope",
                        "key_risks": ["Procurement delay"],
                        "open_questions": ["Who confirms go-live?"],
                        "next_actions": ["Share launch checklist"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=4,
                memory_stale=False,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        mock_complete = AsyncMock(return_value="30 秒判断：主打上线范围和采购风险。\n会后行动：发送清单。")
        with patch_project_llm(complete=mock_complete):
            first = self.client.post(
                f"/projects/{project_id}/briefing/refine",
                json={"meeting_type": "risk", "language": "zh"},
            )
            second = self.client.post(
                f"/projects/{project_id}/briefing/refine",
                json={"meeting_type": "risk", "language": "zh"},
            )

        self.assertEqual(first.status_code, 200)
        self.assertEqual(second.status_code, 200)
        self.assertFalse(first.json()["cached"])
        self.assertTrue(second.json()["cached"])
        self.assertEqual(first.json()["content"], second.json()["content"])
        self.assertEqual(mock_complete.await_count, 1)

        with Session(self.engine) as session:
            cached = session.exec(
                select(ProjectMemorySummary).where(
                    ProjectMemorySummary.project_id == project_id,
                    ProjectMemorySummary.summary_type == "briefing:risk",
                    ProjectMemorySummary.language == "zh",
                )
            ).first()
        self.assertIsNotNone(cached)

    def test_project_memory_snapshots_are_listed_and_can_rollback(self):
        with Session(self.engine) as session:
            project = Project(name="Snapshot Project", client="Client", memory_version=0)
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

            project_contexts_module.save_project_memory(
                session,
                project_id,
                {"project_brief": "Version one", "key_risks": ["Risk A"]},
                trigger="first_build",
            )
            project_contexts_module.save_project_memory(
                session,
                project_id,
                {"project_brief": "Version two", "key_risks": ["Risk B"]},
                trigger="second_build",
            )

        snapshots_resp = self.client.get(f"/projects/{project_id}/memory/snapshots")
        self.assertEqual(snapshots_resp.status_code, 200)
        snapshots = snapshots_resp.json()
        self.assertEqual(len(snapshots), 2)
        older_snapshot = next(item for item in snapshots if item["memory_version"] == 1)

        detail_resp = self.client.get(f"/projects/{project_id}/memory/snapshots/{older_snapshot['id']}")
        self.assertEqual(detail_resp.status_code, 200)
        self.assertEqual(detail_resp.json()["memory"]["project_brief"], "Version one")

        diff_resp = self.client.get(f"/projects/{project_id}/memory/snapshots/{older_snapshot['id']}/diff")
        self.assertEqual(diff_resp.status_code, 200)
        diff_body = diff_resp.json()
        self.assertEqual(diff_body["scope"], "project")
        changed_fields = {item["field"]: item for item in diff_body["fields"]}
        self.assertEqual(changed_fields["project_brief"]["before"], "Version one")
        self.assertEqual(changed_fields["project_brief"]["after"], "Version two")
        self.assertEqual(changed_fields["key_risks"]["added"], ["Risk B"])
        self.assertEqual(changed_fields["key_risks"]["removed"], ["Risk A"])

        rollback_resp = self.client.post(f"/projects/{project_id}/memory/snapshots/{older_snapshot['id']}/rollback")
        self.assertEqual(rollback_resp.status_code, 200)
        self.assertEqual(rollback_resp.json()["memory"]["project_brief"], "Version one")
        self.assertEqual(rollback_resp.json()["memory_version"], 3)

        with Session(self.engine) as session:
            snapshot_count = session.exec(
                select(ProjectMemorySnapshot).where(ProjectMemorySnapshot.project_id == project_id)
            ).all()
        self.assertEqual(len(snapshot_count), 3)

    def test_project_chat_text_can_apply_stakeholder_candidates(self):
        with Session(self.engine) as session:
            client = ClientRecord(name="Acme", client_memory_stale=False)
            session.add(client)
            session.commit()
            session.refresh(client)
            project = Project(name="Stakeholder Capture", client="Acme")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id
            client_id = client.id

        preview = self.client.post(
            f"/projects/{project_id}/stakeholder-candidates",
            json={"text": "李总监提醒采购负责人还在关注预算审批，CFO 需要看到回款计划。"},
        )
        self.assertEqual(preview.status_code, 200)
        self.assertGreaterEqual(len(preview.json()["candidates"]), 2)

        resp = self.client.post(
            f"/projects/{project_id}/stakeholder-candidates/apply",
            json={"text": "李总监提醒采购负责人还在关注预算审批，CFO 需要看到回款计划。"},
        )

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertGreaterEqual(len(body["created"]), 2)
        created_names = {item["name"] for item in body["created"]}
        self.assertIn("李总监", created_names)
        self.assertIn("采购负责人", created_names)

        with Session(self.engine) as session:
            stakeholders = session.exec(
                select(ClientStakeholder).where(ClientStakeholder.client_id == client_id)
            ).all()
            client = session.get(ClientRecord, client_id)

        self.assertGreaterEqual(len(stakeholders), 2)
        self.assertTrue(client.client_memory_stale)

    def test_memory_jobs_list_exposes_rebuild_and_summary_warm_queue(self):
        with Session(self.engine) as session:
            project = Project(
                name="Queued Project",
                client="Client",
                memory_version=2,
                memory_stale=True,
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        fake_jobs = [
            SimpleNamespace(
                id=f"project_memory_rebuild_{project_id}",
                next_run_time=utc_now_naive(),
            ),
            SimpleNamespace(
                id=f"project_memory_summary_warm_{project_id}_zh",
                next_run_time=utc_now_naive(),
            ),
        ]

        with patch("app.services.scheduler.get_jobs", return_value=fake_jobs):
            resp = self.client.get("/projects/memory/jobs")

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["count"], 2)
        self.assertEqual({item["job_type"] for item in body["jobs"]}, {"rebuild", "summary_warm"})

    def test_memory_jobs_list_keeps_status_only_rebuild_when_only_summary_warm_is_queued(self):
        with Session(self.engine) as session:
            project = Project(
                name="Status Only Project",
                client="Client",
                memory_version=2,
                memory_stale=False,
                memory_rebuild_status="queued",
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        fake_jobs = [
            SimpleNamespace(
                id=f"project_memory_summary_warm_{project_id}_zh",
                next_run_time=utc_now_naive(),
            ),
        ]

        with patch("app.services.scheduler.get_jobs", return_value=fake_jobs):
            resp = self.client.get("/projects/memory/jobs")

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["count"], 2)
        jobs_by_source = {item["status_source"]: item for item in body["jobs"]}
        self.assertEqual(jobs_by_source["scheduler"]["job_type"], "summary_warm")
        self.assertEqual(jobs_by_source["project_status"]["job_type"], "rebuild")
        self.assertEqual(jobs_by_source["project_status"]["status_note"], "queued")

    def test_memory_jobs_list_includes_budget_retry_and_recent_failures(self):
        with Session(self.engine) as session:
            project = Project(
                name="Queued Project",
                client="Client",
                memory_version=3,
                memory_stale=True,
                context_memory_json=json.dumps(
                    {
                        "_last_failure": {
                            "stage": "summary_warm",
                            "message": "Rate limit exceeded",
                            "retry_count": 2,
                            "failed_at": "2026-04-19T10:00:00",
                        },
                        "rebuild_log": [
                            {"at": "2026-04-18T09:00:00", "trigger": "manual", "version": 3}
                        ],
                    }
                ),
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id
            session.add(
                ProjectMemorySummary(
                    project_id=project_id,
                    summary_type="overview",
                    language="zh",
                    memory_version=3,
                    content="cached",
                )
            )
            session.commit()

        fake_job = SimpleNamespace(
            id=f"project_memory_rebuild_{project_id}",
            next_run_time=utc_now_naive(),
        )
        with patch("app.services.scheduler.get_jobs", return_value=[fake_job]), patch(
            "app.services.scheduler.get_job_metadata",
            return_value={"retry_count": 1, "max_retries": 3, "trigger": "project_updated"},
        ):
            resp = self.client.get("/projects/memory/jobs")

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["budget"]["used"], 1)
        self.assertEqual(body["jobs"][0]["retry_count"], 1)
        self.assertEqual(body["jobs"][0]["trigger"], "project_updated")
        self.assertEqual(body["recent_failures"][0]["message"], "Rate limit exceeded")
        success_stages = {item["stage"] for item in body["recent_successes"]}
        self.assertIn("rebuild", success_stages)
        self.assertIn("summary_cache:overview", success_stages)
        rebuild_success = next(item for item in body["recent_successes"] if item["stage"] == "rebuild")
        self.assertEqual(rebuild_success["project_id"], project_id)
        self.assertEqual(rebuild_success["completed_at"], "2026-04-18T09:00:00")
        self.assertEqual(rebuild_success["version"], 3)

    def test_memory_jobs_cancel_removes_rebuild_and_summary_jobs(self):
        removed: list[str] = []

        with patch(
            "app.services.scheduler.remove_job",
            side_effect=lambda job_id: removed.append(job_id),
        ):
            resp = self.client.post("/projects/memory/jobs/42/cancel")

        self.assertEqual(resp.status_code, 200)
        self.assertIn("project_memory_rebuild_42", removed)
        self.assertIn("project_memory_summary_warm_42_zh", removed)
        self.assertIn("project_memory_summary_warm_42_en", removed)

    def test_memory_jobs_run_now_clears_status_only_queue_when_memory_is_ready(self):
        with Session(self.engine) as session:
            project = Project(
                name="Ready Queued Project",
                client="Client",
                memory_version=4,
                memory_stale=False,
                memory_rebuild_status="queued",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Ready memory",
                        "memory_version": 4,
                        "stale": False,
                    },
                    ensure_ascii=False,
                ),
            )
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        with patch(
            "app.routers.projects_deps._warm_project_memory_summary_caches",
            new=AsyncMock(return_value=["overview"]),
        ), patch("app.services.scheduler.remove_job"):
            resp = self.client.post(f"/projects/memory/jobs/{project_id}/run-now")

        self.assertEqual(resp.status_code, 200)
        self.assertEqual(resp.json()["action"], "summary_warm")

        with Session(self.engine) as session:
            refreshed = session.get(Project, project_id)
            self.assertIsNotNone(refreshed)
            self.assertEqual(refreshed.memory_rebuild_status, "idle")
            self.assertIsNone(refreshed.memory_rebuild_failed_at)

    def test_memory_rebuild_persists_rebuild_log_and_coverage(self):
        async def fake_complete(messages, max_tokens=4000):
            return json.dumps(
                {
                    "project_brief": "Alpha brief",
                    "current_stage": "delivery",
                    "key_risks": ["Timeline risk"],
                    "open_questions": ["Who signs off?"],
                    "stakeholder_notes": ["Finance is supportive"],
                },
                ensure_ascii=False,
            )

        with Session(self.engine) as session:
            project = Project(name="Rebuild Log Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            session.add(Milestone(project_id=project.id, title="Kickoff"))
            session.commit()
            project_id = project.id

        with patch_project_llm(complete=AsyncMock(side_effect=fake_complete)):
            resp = self.client.post(f"/projects/{project_id}/memory/rebuild")

        self.assertEqual(resp.status_code, 200)
        memory = resp.json()["memory"]
        self.assertEqual(memory["key_risks"], ["Timeline risk"])
        self.assertEqual(memory["open_questions"], ["Who signs off?"])
        self.assertIn("rebuild_log", memory)
        self.assertEqual(memory["rebuild_log"][-1]["trigger"], "manual")
        self.assertIn("_coverage", memory)
        self.assertIn("milestones_total", memory["_coverage"])

    def test_project_memory_stale_marks_matching_client_memory_stale(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Acme Corp",
                client_memory_json=json.dumps({"client_profile": "Existing memory"}, ensure_ascii=False),
                client_memory_version=1,
                client_memory_stale=False,
            )
            project = Project(name="Acme Project", client="Acme Corp", description="Before update")
            session.add(client)
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id
            client_id = client.id

        resp = self.client.patch(
            f"/projects/{project_id}",
            json={"description": "After update"},
        )

        self.assertEqual(resp.status_code, 200)
        with Session(self.engine) as session:
            refreshed_client = session.get(ClientRecord, client_id)
            self.assertIsNotNone(refreshed_client)
            self.assertTrue(refreshed_client.client_memory_stale)

    def test_archiving_project_auto_promotes_memory_to_client(self):
        with Session(self.engine) as session:
            client = ClientRecord(name="Acme Corp", notes="Strategic account")
            project = Project(
                name="Acme Delivery",
                client="Acme Corp",
                status="delivering",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Delivered finance workflow redesign",
                        "key_risks": ["Stakeholder alignment drift"],
                        "stakeholder_notes": ["CFO prefers phased rollout"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=2,
                memory_stale=False,
            )
            session.add(client)
            session.add(project)
            session.commit()
            session.refresh(client)
            session.refresh(project)
            client_id = client.id
            project_id = project.id

        with patch_project_llm(complete=AsyncMock(
                return_value=json.dumps(
                    {
                        "client_profile": "Strategic account with a finance modernization agenda",
                        "decision_patterns": ["Executive sponsors prefer phased delivery"],
                        "key_contacts": [{"name": "CFO", "role": "Executive Sponsor", "note": "Needs measurable wins"}],
                        "lessons_learned": ["Show value in stages before broader rollout"],
                        "project_history": [{"project_name": "Acme Delivery", "status": "archived", "outcome": "Delivered", "key_factor": "Phased wins"}],
                        "sensitive_topics": ["Avoid overpromising automation speed"],
                    },
                    ensure_ascii=False,
                )
            ),
        ) as mocked_complete:
            resp = self.client.patch(
                f"/projects/{project_id}",
                json={"status": "archived"},
            )

        self.assertEqual(resp.status_code, 200)
        self.assertEqual(mocked_complete.await_count, 2)

        with Session(self.engine) as session:
            refreshed_project = session.get(Project, project_id)
            refreshed_client = session.get(ClientRecord, client_id)
            self.assertIsNotNone(refreshed_project)
            self.assertEqual(refreshed_project.status, "archived")
            self.assertGreaterEqual(refreshed_project.memory_version, 3)
            self.assertIsNotNone(refreshed_client)
            self.assertEqual(refreshed_client.client_memory_version, 1)
            self.assertFalse(refreshed_client.client_memory_stale)
            payload = json.loads(refreshed_client.client_memory_json)
            self.assertIn(project_id, payload["source_project_ids"])
            self.assertEqual(payload["rebuild_log"][-1]["trigger"], "project_archived_auto_promoted")

    def test_save_conversation_markdown_merge_requires_file_id(self):
        with Session(self.engine) as session:
            project = Project(name="Merge Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)

            conv = Conversation(project_id=project.id, title="Weekly Sync")
            session.add(conv)
            session.commit()
            session.refresh(conv)

            session.add(Message(conversation_id=conv.id, role="user", content="hello"))
            session.commit()
            project_id = project.id
            conv_id = conv.id

        resp = self.client.post(
            f"/projects/{project_id}/conversations/{conv_id}/save-markdown",
            json={"action": "merge"},
        )
        self.assertEqual(resp.status_code, 400)
        self.assertIn("file_id is required", resp.json()["detail"])

    def test_save_conversation_markdown_merge_appends_existing_document(self):
        with Session(self.engine) as session:
            project = Project(name="Conversation Merge Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)

            conv = Conversation(project_id=project.id, title="Weekly Sync")
            session.add(conv)
            session.commit()
            session.refresh(conv)

            session.add(Message(conversation_id=conv.id, role="user", content="Please summarize this week"))
            session.add(Message(conversation_id=conv.id, role="assistant", content="Here is the weekly summary"))
            session.commit()

            path = Path("projects") / str(project.id) / "merge_target.md"
            full_path = self.uploads_dir / path
            full_path.parent.mkdir(parents=True, exist_ok=True)
            full_path.write_text("# Existing", encoding="utf-8")

            project_file = ProjectFile(
                project_id=project.id,
                name="Merge Target.md",
                file_type="md",
                path=str(path),
                size_bytes=full_path.stat().st_size,
            )
            session.add(project_file)
            session.commit()
            session.refresh(project_file)
            project_id = project.id
            conv_id = conv.id
            file_id = project_file.id

        resp = self.client.post(
            f"/projects/{project_id}/conversations/{conv_id}/save-markdown",
            json={"action": "merge", "file_id": file_id},
        )
        self.assertEqual(resp.status_code, 201)
        self.assertEqual(resp.json()["action"], "merge")

        merged_text = (self.uploads_dir / "projects" / str(project_id) / "merge_target.md").read_text(encoding="utf-8")
        self.assertIn("# Existing", merged_text)
        self.assertIn("From project conversation", merged_text)
        self.assertIn("# Weekly Sync", merged_text)
        self.assertIn("Here is the weekly summary", merged_text)

    def test_update_project_document_rejects_cross_project_file(self):
        with Session(self.engine) as session:
            project_a = Project(name="Project A", client="Client")
            project_b = Project(name="Project B", client="Client")
            session.add(project_a)
            session.add(project_b)
            session.commit()
            session.refresh(project_a)
            session.refresh(project_b)

            path = Path("projects") / str(project_b.id) / "foreign_doc.md"
            full_path = self.uploads_dir / path
            full_path.parent.mkdir(parents=True, exist_ok=True)
            full_path.write_text("# Foreign", encoding="utf-8")

            foreign_file = ProjectFile(
                project_id=project_b.id,
                name="Foreign Doc.md",
                file_type="md",
                path=str(path),
                size_bytes=full_path.stat().st_size,
            )
            session.add(foreign_file)
            session.commit()
            session.refresh(foreign_file)
            project_a_id = project_a.id
            foreign_file_id = foreign_file.id

        resp = self.client.patch(
            f"/projects/{project_a_id}/documents/{foreign_file_id}",
            json={"content": "# Should fail"},
        )
        self.assertEqual(resp.status_code, 404)
        self.assertIn("File not found", resp.json()["detail"])

    def test_save_message_to_document_merge_requires_file_id(self):
        with Session(self.engine) as session:
            project = Project(name="Message Merge Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)

            conv = Conversation(project_id=project.id, title="Discovery")
            session.add(conv)
            session.commit()
            session.refresh(conv)

            message = Message(conversation_id=conv.id, role="assistant", content="Important follow-up")
            session.add(message)
            session.commit()
            session.refresh(message)
            project_id = project.id
            message_id = message.id

        resp = self.client.post(
            f"/projects/{project_id}/messages/{message_id}/save-to-document",
            json={"action": "merge"},
        )
        self.assertEqual(resp.status_code, 400)
        self.assertIn("file_id is required", resp.json()["detail"])

    def test_save_message_to_document_merge_appends_with_header_when_requested(self):
        with Session(self.engine) as session:
            project = Project(name="Message Merge Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)

            conv = Conversation(project_id=project.id, title="Discovery")
            session.add(conv)
            session.commit()
            session.refresh(conv)

            message = Message(conversation_id=conv.id, role="assistant", content="Important follow-up")
            session.add(message)
            session.commit()
            session.refresh(message)

            path = Path("projects") / str(project.id) / "message_target.md"
            full_path = self.uploads_dir / path
            full_path.parent.mkdir(parents=True, exist_ok=True)
            full_path.write_text("# Existing message doc", encoding="utf-8")

            project_file = ProjectFile(
                project_id=project.id,
                name="Message Target.md",
                file_type="md",
                path=str(path),
                size_bytes=full_path.stat().st_size,
            )
            session.add(project_file)
            session.commit()
            session.refresh(project_file)
            project_id = project.id
            message_id = message.id
            file_id = project_file.id

        resp = self.client.post(
            f"/projects/{project_id}/messages/{message_id}/save-to-document",
            json={"action": "merge", "file_id": file_id, "prepend_header": True},
        )
        self.assertEqual(resp.status_code, 201)
        self.assertEqual(resp.json()["action"], "merge")

        merged_text = (self.uploads_dir / "projects" / str(project_id) / "message_target.md").read_text(encoding="utf-8")
        self.assertIn("# Existing message doc", merged_text)
        self.assertIn("From project conversation", merged_text)
        self.assertIn("Important follow-up", merged_text)

    def test_project_todo_due_date_round_trip(self):
        with Session(self.engine) as session:
            project = Project(name="Todo Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        create_resp = self.client.post(
            f"/projects/{project_id}/todos",
            json={
                "content": "Prepare proposal deck",
                "due_date": "2026-04-30",
            },
        )
        self.assertEqual(create_resp.status_code, 201)
        todo_id = create_resp.json()["id"]
        self.assertEqual(create_resp.json()["due_date"], "2026-04-30")

        list_resp = self.client.get(f"/projects/{project_id}/todos")
        self.assertEqual(list_resp.status_code, 200)
        self.assertEqual(len(list_resp.json()), 1)
        self.assertEqual(list_resp.json()[0]["due_date"], "2026-04-30")

        patch_resp = self.client.patch(
            f"/projects/{project_id}/todos/{todo_id}",
            json={
                "due_date": "2026-05-15",
                "is_done": True,
            },
        )
        self.assertEqual(patch_resp.status_code, 200)
        self.assertEqual(patch_resp.json()["due_date"], "2026-05-15")
        self.assertTrue(patch_resp.json()["is_done"])

        list_resp_after_patch = self.client.get(f"/projects/{project_id}/todos")
        self.assertEqual(list_resp_after_patch.status_code, 200)
        self.assertEqual(list_resp_after_patch.json()[0]["due_date"], "2026-05-15")
        self.assertTrue(list_resp_after_patch.json()[0]["is_done"])

    def test_project_detail_returns_todo_due_date(self):
        with Session(self.engine) as session:
            project = Project(name="Detail Project", client="Client")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        create_resp = self.client.post(
            f"/projects/{project_id}/todos",
            json={
                "content": "Confirm workshop agenda",
                "due_date": "2026-06-01",
            },
        )
        self.assertEqual(create_resp.status_code, 201)

        detail_resp = self.client.get(f"/projects/{project_id}/detail")
        self.assertEqual(detail_resp.status_code, 200)
        todos = detail_resp.json()["todos"]
        self.assertEqual(len(todos), 1)
        self.assertEqual(todos[0]["content"], "Confirm workshop agenda")
        self.assertEqual(todos[0]["due_date"], "2026-06-01")

    def test_save_project_note_appends_existing_notes_by_default(self):
        with Session(self.engine) as session:
            project = Project(name="Notes Project", client="Client", notes="Legacy notes")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        resp = self.client.post(
            f"/projects/{project_id}/notes",
            json={
                "content": "Fresh summary",
            },
        )
        self.assertEqual(resp.status_code, 200)
        notes = resp.json()["notes"]
        self.assertIn("Legacy notes", notes)
        self.assertIn("Fresh summary", notes)
        self.assertIn("---", notes)

        with Session(self.engine) as session:
            project = session.get(Project, project_id)
            self.assertIsNotNone(project)
            self.assertEqual(project.notes, notes)

    def test_save_project_note_overwrites_when_append_is_false(self):
        with Session(self.engine) as session:
            project = Project(name="Overwrite Notes Project", client="Client", notes="Legacy notes")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        resp = self.client.post(
            f"/projects/{project_id}/notes",
            json={
                "content": "Fresh summary",
                "append": False,
            },
        )
        self.assertEqual(resp.status_code, 200)
        self.assertEqual(resp.json()["notes"], "Fresh summary")

        with Session(self.engine) as session:
            project = session.get(Project, project_id)
            self.assertIsNotNone(project)
            self.assertEqual(project.notes, "Fresh summary")

    def test_save_project_note_appends_with_timestamp_separator(self):
        with Session(self.engine) as session:
            project = Project(name="Append Notes Project", client="Client", notes="Initial memo")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        resp = self.client.post(
            f"/projects/{project_id}/notes",
            json={
                "content": "Follow-up decision",
                "append": True,
            },
        )
        self.assertEqual(resp.status_code, 200)
        notes = resp.json()["notes"]
        self.assertIn("Initial memo", notes)
        self.assertIn("---", notes)
        self.assertIn("Follow-up decision", notes)
        self.assertRegex(notes, r"\[\d{4}-\d{2}-\d{2} \d{2}:\d{2}\]")

        with Session(self.engine) as session:
            project = session.get(Project, project_id)
            self.assertIsNotNone(project)
            self.assertEqual(project.notes, notes)


class ChatStreamingServiceTestCase(unittest.TestCase):
    def setUp(self):
        self.engine = create_test_engine()
        drop_all_tables(self.engine)
        SQLModel.metadata.create_all(self.engine)

    def tearDown(self):
        self.engine.dispose()

    def _create_conversation(self, title: str = "Existing Title") -> int:
        with Session(self.engine) as session:
            conv = Conversation(title=title)
            session.add(conv)
            session.commit()
            session.refresh(conv)
            return conv.id

    def test_digital_strategy_ppt_tool_is_forced_to_skill_template(self):
        runtime = SimpleNamespace(
            skill_name="数字化战略",
            system="",
        )

        routed_name, routed_input = _route_ppt_tool_for_skill(
            runtime,
            "generate_ppt",
            {
                "title": "Godiva Digital Strategy",
                "slides": [{"type": "content", "title": "Digital transformation roadmap", "content": "- Horizon 1"}],
            },
        )

        self.assertEqual(routed_name, "generate_ppt_from_skill")
        self.assertEqual(routed_input["skill_name"], "digital-strategy")

    def test_digital_strategy_tool_input_routes_even_with_legacy_runtime(self):
        runtime = SimpleNamespace(
            skill_name="",
            system="",
        )

        routed_name, routed_input = _route_ppt_tool_for_skill(
            runtime,
            "generate_ppt",
            {
                "title": "Digital transformation strategy",
                "slides": [{"type": "content", "title": "Three-Horizon Roadmap", "content": "- Scale"}],
            },
        )

        self.assertEqual(routed_name, "generate_ppt_from_skill")
        self.assertEqual(routed_input["skill_name"], "digital-strategy")

    def test_digital_strategy_auto_ppt_fallback_builds_rich_deck(self):
        _, slides = _build_slides_from_strategy_text("数字化战略方案\n\n只有一段战略说明。")

        self.assertGreaterEqual(len(slides), 16)
        titles = {slide["title"] for slide in slides}
        self.assertTrue("Use-Case Portfolio" in titles or "场景组合：平衡快赢、基础能力和战略差异化" in titles)
        self.assertTrue(
            any("90-Day Action Plan" == slide["title"] for slide in slides),
            "fallback deck should include an execution-oriented 90-day plan",
        )
        for slide in slides[-4:]:
            content = slide.get("content", "")
            self.assertGreaterEqual(content.count("\n- ") + int(content.startswith("- ")), 4)

    def test_prepare_chat_runtime_limits_history_window_for_standalone_chat(self):
        conv_id = self._create_conversation()
        with Session(self.engine) as session:
            for index in range(30):
                session.add(
                    Message(
                        conversation_id=conv_id,
                        role="user" if index % 2 == 0 else "assistant",
                        content=f"message-{index}",
                    )
                )
            session.commit()

            with patch.object(chat_streaming_module, "build_chat_context") as mocked_context, patch.object(
                chat_streaming_module,
                "_load_provider_module",
            ) as mocked_provider, patch.object(
                chat_streaming_module,
                "get_selected_model",
                return_value="kimi-k2.6",
            ):
                mocked_context.return_value = context_builder_module.ChatContext()
                mocked_provider.return_value = SimpleNamespace(
                    build_system_prompt=lambda skill_prompt, rag_context, project_context, **kwargs: "system"
                )

                runtime = chat_streaming_module.prepare_chat_runtime(
                    session,
                    chat_router_module.SendMessageRequest(
                        conversation_id=conv_id,
                        content="latest-message",
                    ),
                )

        self.assertEqual(len(runtime.api_messages), chat_streaming_module.CHAT_HISTORY_WINDOW)
        self.assertEqual(runtime.api_messages[0]["content"], "message-7")
        self.assertEqual(runtime.api_messages[-1]["content"], "latest-message")

    def test_prepare_chat_runtime_routes_standalone_short_chat_to_fast_model(self):
        conv_id = self._create_conversation()
        with Session(self.engine) as session:
            with patch.object(chat_streaming_module, "build_chat_context") as mocked_context, patch.object(
                chat_streaming_module,
                "_load_provider_module",
            ) as mocked_provider, patch.object(
                chat_streaming_module,
                "get_selected_model",
                return_value="kimi-k2.6",
            ):
                mocked_context.return_value = context_builder_module.ChatContext(max_tokens=8192)
                mocked_provider.return_value = SimpleNamespace(
                    build_system_prompt=lambda skill_prompt, rag_context, project_context, **kwargs: "system"
                )

                runtime = chat_streaming_module.prepare_chat_runtime(
                    session,
                    chat_router_module.SendMessageRequest(
                        conversation_id=conv_id,
                        content="帮我总结一下今天的重点事项",
                    ),
                )

        self.assertEqual(runtime.selected_model, chat_streaming_module.STANDALONE_FAST_PATH_MODEL)
        self.assertEqual(runtime.max_tokens, chat_streaming_module.STANDALONE_FAST_PATH_MAX_TOKENS)

    def test_prepare_chat_runtime_does_not_auto_apply_selected_skill_from_keywords(self):
        conv_id = self._create_conversation()
        with Session(self.engine) as session:
            skill = Skill(name="Strategy Skill", category="deep_task", system_prompt="skill system")
            session.add(skill)
            session.commit()
            session.refresh(skill)
            skill_id = skill.id

            with patch.object(chat_streaming_module, "build_chat_context") as mocked_context, patch.object(
                chat_streaming_module,
                "_load_provider_module",
            ) as mocked_provider, patch.object(
                chat_streaming_module,
                "get_selected_model",
                return_value="kimi-k2.6",
            ):
                mocked_context.return_value = context_builder_module.ChatContext(max_tokens=8192)
                mocked_provider.return_value = SimpleNamespace(
                    build_system_prompt=lambda skill_prompt, rag_context, project_context, **kwargs: f"system:{skill_prompt}"
                )

                runtime = chat_streaming_module.prepare_chat_runtime(
                    session,
                    chat_router_module.SendMessageRequest(
                        conversation_id=conv_id,
                        content="生成一份战略报告",
                        skill_id=skill_id,
                        force_skill=False,
                    ),
                )

        self.assertFalse(runtime.skill_name)
        self.assertEqual(runtime.prepare_metrics["skill_decision"], "selected_skill_not_armed")
        self.assertEqual(mocked_context.call_args.kwargs["skill_id"], None)

    def test_prepare_chat_runtime_applies_forced_skill_contract(self):
        conv_id = self._create_conversation()
        with Session(self.engine) as session:
            skill = Skill(name="Strategy Skill", category="deep_task", system_prompt="skill system")
            session.add(skill)
            session.commit()
            session.refresh(skill)
            skill_id = skill.id

            with patch.object(chat_streaming_module, "build_chat_context") as mocked_context, patch.object(
                chat_streaming_module,
                "_load_provider_module",
            ) as mocked_provider, patch.object(
                chat_streaming_module,
                "get_selected_model",
                return_value="kimi-k2.6",
            ):
                mocked_context.return_value = context_builder_module.ChatContext(max_tokens=8192)
                mocked_provider.return_value = SimpleNamespace(
                    build_system_prompt=lambda skill_prompt, rag_context, project_context, **kwargs: f"system:{skill_prompt}"
                )

                runtime = chat_streaming_module.prepare_chat_runtime(
                    session,
                    chat_router_module.SendMessageRequest(
                        conversation_id=conv_id,
                        content="生成一份战略报告",
                        skill_id=skill_id,
                        force_skill=True,
                    ),
                )

        self.assertEqual(runtime.skill_name, "Strategy Skill")
        self.assertEqual(runtime.prepare_metrics["skill_decision"], "forced_by_user")
        self.assertEqual(mocked_context.call_args.kwargs["skill_id"], skill_id)

    def test_prepare_chat_runtime_keeps_portfolio_query_on_selected_model(self):
        conv_id = self._create_conversation()
        with Session(self.engine) as session:
            with patch.object(context_builder_module, "build_chat_context") as mocked_context, patch.object(
                chat_streaming_module,
                "_load_provider_module",
            ) as mocked_provider, patch.object(
                chat_streaming_module,
                "get_selected_model",
                return_value="kimi-k2.6",
            ):
                mocked_context.return_value = context_builder_module.ChatContext(max_tokens=8192)
                mocked_provider.return_value = SimpleNamespace(
                    build_system_prompt=lambda skill_prompt, rag_context, project_context, **kwargs: "system"
                )

                runtime = chat_streaming_module.prepare_chat_runtime(
                    session,
                    chat_router_module.SendMessageRequest(
                        conversation_id=conv_id,
                        content="总结金科智慧服务集团股份有限公司的全部项目情况及风险。",
                    ),
                )

        self.assertEqual(runtime.selected_model, "kimi-k2.6")
        self.assertEqual(runtime.max_tokens, chat_streaming_module.CLIENT_PORTFOLIO_MAX_TOKENS)

    def test_prepare_chat_runtime_routes_portfolio_query_to_deepseek_flash_when_key_exists(self):
        conv_id = self._create_conversation()
        with Session(self.engine) as session:
            session.add(Setting(key="deepseek_api_key", value="test-key"))
            session.commit()

            with patch.object(context_builder_module, "build_chat_context") as mocked_context, patch.object(
                chat_streaming_module,
                "_load_provider_module",
            ) as mocked_provider, patch.object(
                chat_streaming_module,
                "get_selected_model",
                return_value="kimi-k2.6",
            ):
                mocked_context.return_value = context_builder_module.ChatContext(max_tokens=8192)
                mocked_provider.return_value = SimpleNamespace(
                    build_system_prompt=lambda skill_prompt, rag_context, project_context, **kwargs: "system"
                )

                runtime = chat_streaming_module.prepare_chat_runtime(
                    session,
                    chat_router_module.SendMessageRequest(
                        conversation_id=conv_id,
                        content="总结金科智慧服务集团股份有限公司的全部项目情况及风险。",
                    ),
                )

        self.assertEqual(runtime.selected_model, chat_streaming_module.CLIENT_PORTFOLIO_FAST_MODEL)
        self.assertEqual(runtime.max_tokens, chat_streaming_module.CLIENT_PORTFOLIO_MAX_TOKENS)

    def test_prepare_chat_runtime_keeps_recent_history_for_portfolio_query(self):
        conv_id = self._create_conversation()
        with Session(self.engine) as session:
            session.add(Message(conversation_id=conv_id, role="user", content="old question"))
            session.add(Message(conversation_id=conv_id, role="assistant", content="old incomplete answer"))
            session.commit()

            with patch.object(context_builder_module, "build_chat_context") as mocked_context, patch.object(
                chat_streaming_module,
                "_load_provider_module",
            ) as mocked_provider, patch.object(
                chat_streaming_module,
                "get_selected_model",
                return_value="kimi-k2.6",
            ):
                mocked_context.return_value = context_builder_module.ChatContext(max_tokens=8192)
                mocked_provider.return_value = SimpleNamespace(
                    build_system_prompt=lambda skill_prompt, rag_context, project_context, **kwargs: "system"
                )

                runtime = chat_streaming_module.prepare_chat_runtime(
                    session,
                    chat_router_module.SendMessageRequest(
                        conversation_id=conv_id,
                        content="总结金科智慧服务集团股份有限公司的全部项目情况及风险。",
                    ),
                )

        self.assertEqual(runtime.api_messages[-1], {"role": "user", "content": "总结金科智慧服务集团股份有限公司的全部项目情况及风险。"})
        self.assertEqual(runtime.prepare_metrics["history_message_count"], 3)

    def test_project_chat_context_disables_global_rag_auto_trigger(self):
        with Session(self.engine) as session:
            project = Project(name="Scoped Project", client="Acme")
            session.add(project)
            session.commit()
            session.refresh(project)

            with patch.object(
                context_builder_module,
                "retrieve_structured",
                return_value=rag_module.RetrievalContext([], "#doc summarize this"),
            ) as mocked_retrieve:
                ctx = context_builder_module.build_chat_context(
                    session=session,
                    project_id=project.id,
                    knowledge_scope="project",
                    content="#doc summarize this",
                )

        mocked_retrieve.assert_called_once()
        called_project_id = mocked_retrieve.call_args.kwargs.get("project_id")
        called_client_id = mocked_retrieve.call_args.kwargs.get("client_id")
        self.assertEqual(called_project_id, project.id)
        self.assertIsNone(called_client_id)
        self.assertEqual(ctx.rag_context, "")
        self.assertEqual(ctx.rag_sources, [])

    def test_project_chat_context_uses_client_scope_when_requested(self):
        with Session(self.engine) as session:
            client = ClientRecord(name="Acme", industry="Consulting")
            session.add(client)
            session.commit()
            session.refresh(client)

            project = Project(name="Scoped Project", client="Acme")
            session.add(project)
            session.commit()
            session.refresh(project)

            with patch.object(
                context_builder_module,
                "retrieve_structured",
                return_value=rag_module.RetrievalContext([], "#doc summarize this"),
            ) as mocked_retrieve:
                context_builder_module.build_chat_context(
                    session=session,
                    project_id=project.id,
                    knowledge_scope="client",
                    content="#doc summarize this",
                )

        mocked_retrieve.assert_called_once()
        called_project_id = mocked_retrieve.call_args.kwargs.get("project_id")
        called_client_id = mocked_retrieve.call_args.kwargs.get("client_id")
        self.assertIsNone(called_project_id)
        self.assertEqual(called_client_id, client.id)

    def test_project_chat_context_allows_explicit_rag_docs(self):
        with Session(self.engine) as session:
            project = Project(name="Scoped Project", client="Acme")
            session.add(project)
            session.commit()
            session.refresh(project)

            fake_result = rag_module.RetrievalContext(
                [
                    rag_module.RetrievalResult(
                        content="Relevant knowledge",
                        document_name="Selected Doc",
                        document_id=1,
                        chunk_index=0,
                        score=0.99,
                    )
                ],
                "#doc summarize this",
            )

            with patch.object(
                context_builder_module,
                "retrieve_structured",
                return_value=fake_result,
            ) as mocked_retrieve:
                ctx = context_builder_module.build_chat_context(
                    session=session,
                    project_id=project.id,
                    knowledge_scope="project",
                    rag_doc_ids=[1],
                    content="#doc summarize this",
                )

        mocked_retrieve.assert_called_once()
        self.assertIn("Relevant knowledge", ctx.rag_context)
        self.assertEqual(len(ctx.rag_sources), 1)

    def test_standalone_chat_context_uses_lightweight_workspace_brief(self):
        with Session(self.engine) as session:
            session.add(Project(name="Standalone Context Project", client="Acme", status="active"))
            session.commit()

            ctx = context_builder_module.build_chat_context(
                session=session,
                project_id=None,
                skill_id=None,
                knowledge_scope="project",
                content="hello",
            )

        self.assertIn("Workspace Brief", ctx.project_context)
        self.assertIn("Standalone Context Project", ctx.project_context)
        self.assertNotIn("## Project:", ctx.project_context)

    def test_standalone_client_portfolio_query_includes_all_matching_projects(self):
        client_name = "金科智慧服务集团股份有限公司"
        with Session(self.engine) as session:
            for index in range(7):
                project = Project(
                    name=f"金科项目 {index + 1}",
                    client=client_name,
                    status="archived" if index == 6 else "delivering",
                    context_memory_json=json.dumps(
                        {
                            "project_brief": f"项目 {index + 1} 摘要",
                            "key_risks": [f"风险 {index + 1}"],
                            "next_actions": [f"行动 {index + 1}"],
                        },
                        ensure_ascii=False,
                    ),
                    memory_version=1,
                )
                session.add(project)
            session.add(Project(name="Other Client Project", client="Other Client", status="delivering"))
            session.commit()

            ctx = context_builder_module.build_chat_context(
                session=session,
                project_id=None,
                skill_id=None,
                content="总结金科智慧服务集团股份有限公司的全部项目情况及风险。",
            )

        self.assertIn("Client Project Portfolio Context", ctx.project_context)
        self.assertIn("Matched projects: 7", ctx.project_context)
        for index in range(7):
            self.assertIn(f"金科项目 {index + 1}", ctx.project_context)
            self.assertIn(f"风险 {index + 1}", ctx.project_context)
        self.assertNotIn("Other Client Project", ctx.project_context)

    def test_project_client_portfolio_query_respects_current_project_scope(self):
        client_name = "金科智慧服务集团股份有限公司"
        with Session(self.engine) as session:
            current_project = Project(
                name="金科当前项目",
                client=client_name,
                status="delivering",
                context_memory_json=json.dumps(
                    {"project_brief": "当前项目摘要", "key_risks": ["当前风险"]},
                    ensure_ascii=False,
                ),
                memory_version=1,
            )
            session.add(current_project)
            for index in range(6):
                session.add(
                    Project(
                        name=f"金科其他项目 {index + 1}",
                        client=client_name,
                        status="archived" if index == 5 else "delivering",
                        context_memory_json=json.dumps(
                            {
                                "project_brief": f"其他项目 {index + 1} 摘要",
                                "key_risks": [f"其他风险 {index + 1}"],
                            },
                            ensure_ascii=False,
                        ),
                        memory_version=1,
                    )
                )
            session.add(Project(name="Other Client Project", client="Other Client", status="delivering"))
            session.commit()
            session.refresh(current_project)

            ctx = context_builder_module.build_chat_context(
                session=session,
                project_id=current_project.id,
                skill_id=None,
                knowledge_scope="project",
                content="总结金科智慧服务集团股份有限公司的全部项目情况及风险。",
            )

        self.assertIn("Scope Guard", ctx.project_context)
        self.assertIn("金科当前项目", ctx.project_context)
        self.assertNotIn("Client Project Portfolio Context", ctx.project_context)
        for index in range(6):
            self.assertNotIn(f"金科其他项目 {index + 1}", ctx.project_context)
            self.assertNotIn(f"其他风险 {index + 1}", ctx.project_context)
        self.assertNotIn("Other Client Project", ctx.project_context)

    def test_context_mode_can_force_client_portfolio_from_current_project(self):
        client_name = "金科智慧服务集团股份有限公司"
        with Session(self.engine) as session:
            current_project = Project(name="金科当前项目", client=client_name, status="delivering")
            session.add(current_project)
            session.add(Project(name="金科同客户项目", client=client_name, status="delivering"))
            session.add(Project(name="Other Client Project", client="Other Client", status="delivering"))
            session.commit()
            session.refresh(current_project)

            ctx = context_builder_module.build_chat_context(
                session=session,
                project_id=current_project.id,
                skill_id=None,
                knowledge_scope="project",
                content="看一下这个客户的组合情况",
                context_mode="client_portfolio",
            )

        self.assertIn("Client Project Portfolio Context", ctx.project_context)
        self.assertIn("Matched projects: 2", ctx.project_context)
        self.assertIn("金科同客户项目", ctx.project_context)
        self.assertNotIn("Other Client Project", ctx.project_context)

    def test_context_mode_project_prevents_legacy_portfolio_auto_upgrade(self):
        client_name = "金科智慧服务集团股份有限公司"
        with Session(self.engine) as session:
            current_project = Project(name="金科当前项目", client=client_name, status="delivering")
            session.add(current_project)
            session.add(Project(name="金科同客户项目", client=client_name, status="delivering"))
            session.commit()
            session.refresh(current_project)

            ctx = context_builder_module.build_chat_context(
                session=session,
                project_id=current_project.id,
                skill_id=None,
                knowledge_scope="global",
                content="总结金科智慧服务集团股份有限公司的全部项目情况及风险。",
                context_mode="project",
            )

        self.assertIn("Scope Guard", ctx.project_context)
        self.assertIn("金科当前项目", ctx.project_context)
        self.assertNotIn("Client Project Portfolio Context", ctx.project_context)
        self.assertNotIn("金科同客户项目", ctx.project_context)

    def test_project_client_portfolio_query_can_use_current_project_client(self):
        client_name = "金科智慧服务集团股份有限公司"
        with Session(self.engine) as session:
            current_project = Project(name="金科当前项目", client=client_name, status="delivering")
            session.add(current_project)
            for index in range(2):
                session.add(Project(name=f"金科同客户项目 {index + 1}", client=client_name, status="delivering"))
            session.commit()
            session.refresh(current_project)

            ctx = context_builder_module.build_chat_context(
                session=session,
                project_id=current_project.id,
                knowledge_scope="client",
                content="总结这个客户的全部项目情况及风险。",
            )

        self.assertIn("Client Project Portfolio Context", ctx.project_context)
        self.assertIn("Matched projects: 3", ctx.project_context)
        self.assertIn("金科同客户项目 1", ctx.project_context)

    def test_project_chat_context_global_scope_does_not_force_scope_filters(self):
        with Session(self.engine) as session:
            project = Project(name="Scoped Project", client="Acme")
            session.add(project)
            session.commit()
            session.refresh(project)

            with patch.object(
                context_builder_module,
                "retrieve_structured",
                return_value=rag_module.RetrievalContext([], "#doc summarize this"),
            ) as mocked_retrieve:
                context_builder_module.build_chat_context(
                    session=session,
                    project_id=project.id,
                    knowledge_scope="global",
                    content="#doc summarize this",
                )

        mocked_retrieve.assert_called_once()
        called_project_id = mocked_retrieve.call_args.kwargs.get("project_id")
        called_client_id = mocked_retrieve.call_args.kwargs.get("client_id")
        self.assertIsNone(called_project_id)
        self.assertIsNone(called_client_id)

    def test_explicit_rag_doc_ids_bypass_scope_filters(self):
        with Session(self.engine) as session:
            client = ClientRecord(name="Acme", industry="Consulting")
            session.add(client)
            session.commit()
            session.refresh(client)

            project = Project(name="Scoped Project", client="Acme")
            session.add(project)
            session.commit()
            session.refresh(project)

            with patch.object(
                context_builder_module,
                "retrieve_structured",
                return_value=rag_module.RetrievalContext([], "#doc summarize this"),
            ) as mocked_retrieve:
                context_builder_module.build_chat_context(
                    session=session,
                    project_id=project.id,
                    knowledge_scope="client",
                    rag_doc_ids=[99],
                    content="#doc summarize this",
                )

        mocked_retrieve.assert_called_once()
        self.assertEqual(mocked_retrieve.call_args.args[2], [99])
        called_project_id = mocked_retrieve.call_args.kwargs.get("project_id")
        called_client_id = mocked_retrieve.call_args.kwargs.get("client_id")
        self.assertIsNone(called_project_id)
        self.assertIsNone(called_client_id)

    def test_build_project_context_excludes_sibling_project_content(self):
        with Session(self.engine) as session:
            primary_project = Project(
                name="Primary Project",
                client="Acme",
                description="Current project description",
                notes="Current project note",
                context_summary="Current summary",
            )
            sibling_project = Project(
                name="Sibling Project",
                client="Acme",
                description="Sibling secret description",
                notes="Sibling secret note",
                context_summary="Sibling secret summary",
            )
            session.add(primary_project)
            session.add(sibling_project)
            session.commit()
            session.refresh(primary_project)

            project_context = context_builder_module.build_project_context(
                session=session,
                project_id=primary_project.id,
            )

        self.assertIn("Primary Project", project_context)
        self.assertIn("Current project description", project_context)
        self.assertIn("Current project note", project_context)
        self.assertIn("Current summary", project_context)
        self.assertNotIn("Sibling Project", project_context)
        self.assertNotIn("Sibling secret description", project_context)
        self.assertNotIn("Sibling secret note", project_context)
        self.assertNotIn("Sibling secret summary", project_context)

    def test_build_project_context_includes_structured_client_stakeholders(self):
        with Session(self.engine) as session:
            client = ClientRecord(name="Acme", industry="Manufacturing")
            project = Project(
                name="Current Delivery",
                client="Acme",
                description="Current project description",
            )
            session.add(client)
            session.add(project)
            session.commit()
            session.refresh(client)
            session.refresh(project)
            session.add(
                ClientStakeholder(
                    client_id=client.id,
                    name="Jane Sponsor",
                    role="Executive sponsor",
                    influence_type="decision",
                    relationship_status="supportive",
                    concerns="Needs executive rollout clarity",
                )
            )
            session.commit()

            project_context = context_builder_module.build_project_context(
                session=session,
                project_id=project.id,
            )

        self.assertIn("Structured Client Stakeholders", project_context)
        self.assertIn("Jane Sponsor", project_context)
        self.assertIn("Executive sponsor", project_context)
        self.assertIn("Needs executive rollout clarity", project_context)

    def test_project_memory_data_and_summary_payload_include_structured_stakeholders(self):
        with Session(self.engine) as session:
            client = ClientRecord(name="Acme", industry="Manufacturing")
            project = Project(name="Memory Project", client="Acme")
            session.add(client)
            session.add(project)
            session.commit()
            session.refresh(client)
            session.refresh(project)
            session.add(
                ClientStakeholder(
                    client_id=client.id,
                    name="Finance Lead",
                    role="Finance",
                    influence_type="budget",
                    relationship_status="neutral",
                    concerns="Needs payment schedule confirmation",
                )
            )
            session.commit()

            _, project_data, coverage = project_contexts_module.build_project_memory_data(session, project.id)

        self.assertIn("Structured client stakeholders", project_data)
        self.assertIn("Finance Lead", project_data)
        self.assertEqual(coverage["client_stakeholders_total"], 1)
        payload = project_contexts_module.build_project_memory_summary_payload(
            {"client_stakeholders": coverage["client_stakeholders"]},
            "stakeholder",
        )
        self.assertEqual(payload["client_stakeholders"][0]["name"], "Finance Lead")

    def test_client_memory_data_and_summary_payload_include_structured_stakeholders(self):
        with Session(self.engine) as session:
            client = ClientRecord(name="Acme", industry="Manufacturing")
            session.add(client)
            session.commit()
            session.refresh(client)
            session.add(
                ClientStakeholder(
                    client_id=client.id,
                    name="Procurement Owner",
                    role="Procurement",
                    influence_type="purchase",
                    relationship_status="supportive",
                    concerns="Needs procurement timeline visibility",
                )
            )
            session.commit()

            _, client_data, _ = client_contexts_module.build_client_memory_data(session, client.id)
            memory = client_contexts_module.save_client_memory(
                session,
                client.id,
                {"client_profile": "Acme profile"},
                trigger="test",
            )

        self.assertIn("Structured client stakeholders", client_data)
        self.assertIn("Procurement Owner", client_data)
        self.assertEqual(memory["structured_stakeholders"][0]["name"], "Procurement Owner")
        payload = client_contexts_module.build_client_memory_summary_payload(memory, "stakeholder")
        self.assertEqual(payload["structured_stakeholders"][0]["name"], "Procurement Owner")

    def test_build_chat_context_project_scope_keeps_current_project_only(self):
        with Session(self.engine) as session:
            primary_project = Project(
                name="Current Delivery",
                client="Acme",
                description="Current delivery workstream",
                notes="Only current project notes",
            )
            sibling_project = Project(
                name="Other Delivery",
                client="Acme",
                description="Other project should stay isolated",
                notes="Other project confidential notes",
            )
            session.add(primary_project)
            session.add(sibling_project)
            session.commit()
            session.refresh(primary_project)

            with patch.object(
                context_builder_module,
                "retrieve_structured",
                return_value=FakeRetrievalContext(""),
            ):
                ctx = context_builder_module.build_chat_context(
                    session=session,
                    project_id=primary_project.id,
                    knowledge_scope="project",
                    content="Summarize status",
                )

        self.assertIn("Current Delivery", ctx.project_context)
        self.assertIn("Current delivery workstream", ctx.project_context)
        self.assertIn("Only current project notes", ctx.project_context)
        self.assertNotIn("Other Delivery", ctx.project_context)
        self.assertNotIn("Other project should stay isolated", ctx.project_context)
        self.assertNotIn("Other project confidential notes", ctx.project_context)

    def test_build_chat_context_client_scope_without_matching_client_record_falls_back_to_project_scope(self):
        with Session(self.engine) as session:
            project = Project(
                name="Unmatched Client Project",
                client="Acme Corp",
                description="Current project description",
            )
            session.add(project)
            session.commit()
            session.refresh(project)

            with patch.object(
                context_builder_module,
                "retrieve_structured",
                return_value=rag_module.RetrievalContext([], "#doc summarize this"),
            ) as mocked_retrieve:
                ctx = context_builder_module.build_chat_context(
                    session=session,
                    project_id=project.id,
                    knowledge_scope="client",
                    content="#doc summarize this",
                )

        mocked_retrieve.assert_called_once()
        self.assertEqual(mocked_retrieve.call_args.kwargs.get("project_id"), project.id)
        self.assertIsNone(mocked_retrieve.call_args.kwargs.get("client_id"))
        self.assertEqual(ctx.rag_context, "")
        self.assertEqual(ctx.rag_sources, [])

    def test_build_chat_context_client_scope_injects_client_memory_when_available(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Acme Corp",
                client_memory_json=json.dumps(
                    {
                        "client_profile": "Large enterprise account with multi-step approvals",
                        "decision_patterns": ["Budget closes in Q4"],
                        "lessons_learned": ["Technical team prefers phased rollouts"],
                    },
                    ensure_ascii=False,
                ),
                client_memory_version=1,
                client_memory_stale=False,
            )
            project = Project(
                name="Client Memory Project",
                client="Acme Corp",
                description="Current project description",
            )
            session.add(client)
            session.add(project)
            session.commit()
            session.refresh(project)

            with patch.object(
                context_builder_module,
                "retrieve_structured",
                return_value=FakeRetrievalContext(""),
            ):
                ctx = context_builder_module.build_chat_context(
                    session=session,
                    project_id=project.id,
                    knowledge_scope="client",
                    content="Summarize current relationship",
                )

        self.assertIn("Structured Client Memory", ctx.project_context)
        self.assertIn("Budget closes in Q4", ctx.project_context)
        self.assertIn("Current project description", ctx.project_context)

    def test_build_chat_context_file_ids_do_not_attach_sibling_project_files(self):
        with Session(self.engine) as session:
            primary_project = Project(name="Primary Project", client="Acme")
            sibling_project = Project(name="Sibling Project", client="Acme")
            session.add(primary_project)
            session.add(sibling_project)
            session.commit()
            session.refresh(primary_project)
            session.refresh(sibling_project)

            sibling_file = ProjectFile(
                project_id=sibling_project.id,
                name="sibling-note.md",
                file_type="md",
                path="projects/sibling-note.md",
                size_bytes=10,
            )
            session.add(sibling_file)
            session.commit()
            session.refresh(sibling_file)

            with patch.object(
                context_builder_module,
                "extract_file_text",
                return_value="Sibling confidential file",
            ):
                ctx = context_builder_module.build_chat_context(
                    session=session,
                    project_id=primary_project.id,
                    knowledge_scope="project",
                    file_ids=[sibling_file.id],
                    content="summarize",
                )

        self.assertIn("Primary Project", ctx.project_context)
        self.assertNotIn("Sibling confidential file", ctx.project_context)
        self.assertNotIn("sibling-note.md", ctx.project_context)

    def test_stream_chat_events_persists_successful_response(self):
        conv_id = self._create_conversation()
        runtime = ChatRuntime(
            conv_id=conv_id,
            selected_model="claude-sonnet-4-6",
            llm=FakeStreamingLLM([["hello ", "world"]]),
            system="system",
            api_messages=[{"role": "user", "content": "hello"}],
            rag_sources=[],
            tools=None,
            max_tokens=1024,
            temperature=0.7,
        )
        req = chat_router_module.SendMessageRequest(content="hello")

        events = collect_async_generator(stream_chat_events(runtime, req, self.engine))
        joined = "".join(events)
        self.assertIn('"conversation_id"', joined)
        self.assertIn('"done"', joined)
        self.assertIn('"stage_timings"', joined)
        self.assertNotIn('"step_index"', joined)

        with Session(self.engine) as session:
            assistant_messages = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).all()
            self.assertEqual(len(assistant_messages), 1)
            self.assertEqual(assistant_messages[0].content, "hello world")
            self.assertIn("stage_timings", assistant_messages[0].metadata_json)

    def test_stream_chat_events_repairs_project_office_tool_in_p3_follow_up(self):
        with Session(self.engine) as session:
            project = Project(name="Office Repair Project", client="Client", status="active")
            session.add(project)
            session.commit()
            session.refresh(project)
            project_id = project.id

        conv_id = self._create_conversation()
        prompt = "访谈的excel参数补齐"
        llm = FakeStreamingLLM(
            [
                ['{"type":"tool_use","id":"tool-read","name":"read_project_markdown_document","input":{"file_id":1}}'],
                [
                    '我会创建访谈 Excel。 '
                    '{"type":"tool_use","id":"tool-office","name":"write_project_office_document","input":{"title":"访谈计划"}}'
                ],
                ["访谈 Excel 已保存。"],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            project_id=project_id,
            selected_model="claude-sonnet-4-6",
            llm=llm,
            system="system",
            api_messages=[{"role": "user", "content": prompt}],
            rag_sources=[],
            tools=[{"name": "read_project_markdown_document"}, {"name": "write_project_office_document"}],
            max_tokens=1024,
            temperature=0.7,
        )
        req = chat_router_module.SendMessageRequest(content=prompt, project_id=project_id)

        async def mock_execute(name, input_data):
            if name == "read_project_markdown_document":
                return {"type": "tool_result", "tool_name": name, "status": "success", "output": {"content": "# Context"}}
            if name == "write_project_office_document":
                return {
                    "type": "tool_result",
                    "tool_name": name,
                    "status": "success",
                    "output": {
                        "success": True,
                        "project_file_id": 99,
                        "file_name": input_data["file_name"],
                        "file_type": input_data["file_type"],
                        "file_path": f"projects/{project_id}/{input_data['file_name']}",
                    },
                }
            raise AssertionError(f"Unexpected tool call: {name}")

        execute_mock = AsyncMock(side_effect=mock_execute)
        with patch("app.services.chat_streaming.registry.execute", new=execute_mock):
            events = collect_async_generator(stream_chat_events(runtime, req, self.engine))

        office_call = [call for call in execute_mock.await_args_list if call.args[0] == "write_project_office_document"][0]
        read_call = [call for call in execute_mock.await_args_list if call.args[0] == "read_project_markdown_document"][0]
        self.assertEqual(read_call.args[1]["action"], "read")
        office_input = office_call.args[1]
        self.assertEqual(office_input["project_id"], project_id)
        self.assertEqual(office_input["file_type"], "xlsx")
        self.assertTrue(office_input["file_name"].endswith(".xlsx"))
        self.assertEqual(office_input["sheets"][0]["name"], "访谈计划")
        joined = "".join(events)
        self.assertIn('"step_index"', joined)
        self.assertIn("访谈 Excel 已保存。", joined)

    def test_stream_chat_events_emits_and_persists_references(self):
        conv_id = self._create_conversation()
        runtime = ChatRuntime(
            conv_id=conv_id,
            selected_model="claude-sonnet-4-6",
            llm=FakeStreamingLLM([["answer with source"]]),
            system="system",
            api_messages=[{"role": "user", "content": "hello"}],
            rag_sources=[{"type": "file", "id": 7, "title": "project-brief.md"}],
            tools=None,
            max_tokens=1024,
            temperature=0.7,
        )
        req = chat_router_module.SendMessageRequest(content="hello", project_id=12)

        events = collect_async_generator(stream_chat_events(runtime, req, self.engine))
        joined = "".join(events)
        self.assertIn('"references"', joined)
        self.assertIn('"project-brief.md"', joined)

        with Session(self.engine) as session:
            assistant_messages = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).all()
            self.assertEqual(len(assistant_messages), 1)
            metadata = json.loads(assistant_messages[0].metadata_json)
            self.assertEqual(metadata["references"][0]["title"], "project-brief.md")
            self.assertEqual(metadata["project_id"], 12)

    def test_stream_chat_events_continues_truncated_plain_response(self):
        conv_id = self._create_conversation()
        llm = FakeStreamingLLM(
            [
                ["first half ", "[OUTPUT_TRUNCATED]"],
                ["second half"],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            selected_model="kimi-k2.6",
            llm=llm,
            system="system",
            api_messages=[{"role": "user", "content": "long answer"}],
            rag_sources=[],
            tools=None,
            max_tokens=1024,
            temperature=0.7,
        )
        req = chat_router_module.SendMessageRequest(content="long answer")

        events = collect_async_generator(stream_chat_events(runtime, req, self.engine))
        joined = "".join(events)

        self.assertNotIn("[OUTPUT_TRUNCATED]", joined)
        self.assertIn("正在尝试继续生成", joined)
        self.assertIn("first half ", joined)
        self.assertIn("second half", joined)
        self.assertEqual(llm.calls, 2)

        with Session(self.engine) as session:
            assistant_messages = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).all()
            self.assertEqual(assistant_messages[0].content, "first half second half")

    def test_stream_chat_events_continues_truncated_tool_response_before_execution(self):
        conv_id = self._create_conversation()
        llm = FakeStreamingLLM(
            [
                [
                    "# 数字化战略方案\n",
                    "## 执行摘要\n",
                    "- 前半段内容",
                    "[OUTPUT_TRUNCATED]",
                    '{"type":"tool_use","id":"tool-1","name":"generate_ppt_from_skill","input":{"skill_name":"digital-strategy","title":"Deck","slides":[{"type":"content","title":"Thin","content":"- A"}]}}',
                ],
                ["\n## 路线图\n- 后半段补齐"],
                ["done"],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            selected_model="deepseek-v4",
            llm=llm,
            system="digital-strategy",
            api_messages=[{"role": "user", "content": "生成数字化战略 PPT"}],
            rag_sources=[],
            tools=[{"name": "generate_ppt_from_skill"}],
            max_tokens=1024,
            temperature=0.7,
            skill_name="digital-strategy",
        )
        req = chat_router_module.SendMessageRequest(content="生成数字化战略 PPT", skill_id=24)

        execute_mock = AsyncMock(
            return_value={
                "type": "tool_result",
                "tool_name": "generate_ppt_from_skill",
                "status": "success",
                "output": {
                    "success": True,
                    "file_name": "continued_deck.pptx",
                    "file_type": "pptx",
                    "file_path": "generated/continued_deck.pptx",
                },
            }
        )
        with patch("app.services.chat_streaming.registry.execute", new=execute_mock):
            events = collect_async_generator(stream_chat_events(runtime, req, self.engine))

        joined = "".join(events)
        self.assertIn("后半段补齐", joined)
        self.assertEqual(llm.calls, 3)
        _, repaired_input = execute_mock.await_args.args
        self.assertGreaterEqual(len(repaired_input["slides"]), 2)
        self.assertNotEqual(repaired_input["slides"][0]["title"], "Thin")

    def test_stream_chat_events_handles_tool_follow_up(self):
        conv_id = self._create_conversation()
        llm = FakeStreamingLLM(
            [
                [
                    '{"type":"tool_use","id":"tool-1","name":"generate_docx","input":{"title":"Spec"}}',
                ],
                ["follow-up answer"],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            selected_model="claude-sonnet-4-6",
            llm=llm,
            system="system",
            api_messages=[{"role": "user", "content": "make doc"}],
            rag_sources=[],
            tools=[{"name": "generate_docx"}],
            max_tokens=1024,
            temperature=0.7,
        )
        req = chat_router_module.SendMessageRequest(content="make doc")

        with patch(
            "app.services.chat_streaming.registry.execute",
            new=AsyncMock(
                return_value={
                    "status": "ok",
                    "tool_name": "generate_docx",
                    "success": True,
                    "file_name": "spec.docx",
                    "file_type": "docx",
                    "file_path": "generated/spec.docx",
                    "output": {"file": "spec.docx"},
                }
            ),
        ):
            events = collect_async_generator(stream_chat_events(runtime, req, self.engine))

        joined = "".join(events)
        self.assertIn('"tool_executing"', joined)
        self.assertIn('"tool_result"', joined)
        self.assertIn("follow-up answer", joined)
        self.assertEqual(llm.calls, 2)

        with Session(self.engine) as session:
            assistant_messages = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).all()
            self.assertIn("follow-up answer", assistant_messages[0].content)
            self.assertIn("已生成附件：spec.docx", assistant_messages[0].content)
            metadata = json.loads(assistant_messages[0].metadata_json)
            self.assertEqual(metadata["tool_calls"][0]["tool_name"], "generate_docx")
            self.assertEqual(metadata["tool_calls"][0]["status"], "completed")
            self.assertEqual(metadata["artifacts"][0]["name"], "spec.docx")
            self.assertEqual(metadata["artifacts"][0]["path"], "generated/spec.docx")

    def test_stream_chat_events_emits_done_artifacts_and_notice_for_nested_tool_output(self):
        conv_id = self._create_conversation()
        llm = FakeStreamingLLM(
            [
                [
                    '{"type":"tool_use","id":"tool-1","name":"generate_ppt_from_skill","input":{"skill_name":"digital-strategy","title":"Deck","slides":[{"type":"content","title":"One","content":"- A"}]}}',
                ],
                ["strategy summary"],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            selected_model="claude-sonnet-4-6",
            llm=llm,
            system="system",
            api_messages=[{"role": "user", "content": "make deck"}],
            rag_sources=[],
            tools=[{"name": "generate_ppt_from_skill"}],
            max_tokens=1024,
            temperature=0.7,
            skill_name="digital-strategy",
        )
        req = chat_router_module.SendMessageRequest(content="make deck", skill_id=24)

        with patch(
            "app.services.chat_streaming.registry.execute",
            new=AsyncMock(
                return_value={
                    "type": "tool_result",
                    "tool_name": "generate_ppt_from_skill",
                    "status": "success",
                    "output": {
                        "success": True,
                        "file_name": "generated_test_deck.pptx",
                        "file_type": "pptx",
                        "file_path": "generated/generated_test_deck.pptx",
                        "template_applied": True,
                        "template_mode": "cloned_prototype_slides",
                    },
                }
            ),
        ):
            events = collect_async_generator(stream_chat_events(runtime, req, self.engine))

        done_event = next(
            json.loads(event.replace("data: ", "").strip())
            for event in events
            if json.loads(event.replace("data: ", "").strip()).get("type") == "done"
        )
        self.assertEqual(done_event["artifacts"][0]["name"], "generated_test_deck.pptx")
        self.assertEqual(done_event["artifacts"][0]["path"], "generated/generated_test_deck.pptx")

        with Session(self.engine) as session:
            assistant_messages = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).all()
            self.assertEqual(len(assistant_messages), 1)
            self.assertIn("已生成附件：generated_test_deck.pptx", assistant_messages[0].content)
            metadata = json.loads(assistant_messages[0].metadata_json)
            self.assertEqual(metadata["artifacts"][0]["name"], "generated_test_deck.pptx")
            self.assertEqual(metadata["artifacts"][0]["path"], "generated/generated_test_deck.pptx")
            self.assertIsNotNone(metadata["artifacts"][0].get("id"))

            saved_artifact = session.exec(
                select(GeneratedFile).where(GeneratedFile.conversation_id == conv_id)
            ).first()
            self.assertIsNotNone(saved_artifact)
            self.assertEqual(saved_artifact.name, "generated_test_deck.pptx")
            self.assertEqual(saved_artifact.path, "generated/generated_test_deck.pptx")

    def test_stream_chat_events_repairs_empty_digital_strategy_ppt_tool_input(self):
        conv_id = self._create_conversation()
        llm = FakeStreamingLLM(
            [
                [
                    "# Digital Strategy Plan\n",
                    "## Executive Summary\n",
                    "- Build a data foundation\n",
                    "- Prioritize operating efficiency\n",
                    '{"type":"tool_use","id":"tool-1","name":"generate_ppt_from_skill","input":{}}',
                ],
                ["deck created"],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            selected_model="deepseek-v4",
            llm=llm,
            system="Use digital-strategy to create a PPT.",
            api_messages=[{"role": "user", "content": "make deck"}],
            rag_sources=[],
            tools=[{"name": "generate_ppt_from_skill"}],
            max_tokens=1024,
            temperature=0.7,
            skill_name="digital-strategy",
        )
        req = chat_router_module.SendMessageRequest(content="make deck", skill_id=24)

        execute_mock = AsyncMock(
            return_value={
                "type": "tool_result",
                "tool_name": "generate_ppt_from_skill",
                "status": "success",
                "output": {
                    "success": True,
                    "file_name": "repaired_deck.pptx",
                    "file_type": "pptx",
                    "file_path": "generated/repaired_deck.pptx",
                },
            }
        )
        with patch("app.services.chat_streaming.registry.execute", new=execute_mock):
            events = collect_async_generator(stream_chat_events(runtime, req, self.engine))

        self.assertTrue(any('"tool_result"' in event for event in events))
        _, repaired_input = execute_mock.await_args.args
        self.assertEqual(repaired_input["skill_name"], "digital-strategy")
        self.assertTrue(repaired_input["title"])
        self.assertGreaterEqual(len(repaired_input["slides"]), 1)

    def test_stream_chat_events_does_not_auto_generate_ppt_after_skill_followup(self):
        conv_id = self._create_conversation()
        llm = FakeStreamingLLM(
            [
                [
                    '{"type":"tool_use","id":"tool-1","name":"save_json","input":{"filename":"金科服务_数字化战略_核心数据","data":{"maturity":"L2"}}}',
                ],
                ["digital strategy final answer"],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            selected_model="claude-sonnet-4-6",
            llm=llm,
            system="digital-strategy 工作流\n请生成数字化转型战略方案",
            api_messages=[{"role": "user", "content": "生成数字化战略方案"}],
            rag_sources=[],
            tools=[{"name": "generate_ppt_from_skill"}, {"name": "save_json"}],
            max_tokens=1024,
            temperature=0.7,
            skill_name="digital-strategy",
        )
        req = chat_router_module.SendMessageRequest(content="生成数字化战略方案", skill_id=24)

        async def execute_side_effect(name, input_data):
            if name == "save_json":
                return {
                    "type": "tool_result",
                    "tool_name": "save_json",
                    "status": "success",
                    "output": {
                        "success": True,
                        "file_name": "金科服务_数字化战略_核心数据.json",
                        "file_type": "json",
                        "file_path": "generated/金科服务_数字化战略_核心数据.json",
                    },
                }
            if name == "generate_ppt_from_skill":
                return {
                    "type": "tool_result",
                    "tool_name": "generate_ppt_from_skill",
                    "status": "success",
                    "output": {
                        "success": True,
                        "file_name": "金科服务_数字化战略_汇报材料.pptx",
                        "file_type": "pptx",
                        "file_path": "generated/金科服务_数字化战略_汇报材料.pptx",
                    },
                }
            raise AssertionError(f"Unexpected tool call: {name}")

        execute_mock = AsyncMock(side_effect=execute_side_effect)
        with patch("app.services.chat_streaming.registry.execute", new=execute_mock):
            events = collect_async_generator(stream_chat_events(runtime, req, self.engine))

        done_event = next(
            json.loads(event.replace("data: ", "").strip())
            for event in events
            if json.loads(event.replace("data: ", "").strip()).get("type") == "done"
        )
        artifact_names = {item["name"] for item in done_event["artifacts"]}
        self.assertIn("金科服务_数字化战略_核心数据.json", artifact_names)
        self.assertNotIn("金科服务_数字化战略_汇报材料.pptx", artifact_names)
        called_tools = [call.args[0] for call in execute_mock.await_args_list]
        self.assertEqual(called_tools, ["save_json"])

        with Session(self.engine) as session:
            assistant_message = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).one()
            metadata = json.loads(assistant_message.metadata_json)
            saved_names = {item["name"] for item in metadata["artifacts"]}
            self.assertIn("金科服务_数字化战略_核心数据.json", saved_names)
            self.assertNotIn("金科服务_数字化战略_汇报材料.pptx", saved_names)

    def test_stream_chat_events_surfaces_friendly_errors(self):
        conv_id = self._create_conversation()
        runtime = ChatRuntime(
            conv_id=conv_id,
            selected_model="claude-sonnet-4-6",
            llm=FakeStreamingLLM([Exception("429 engine_overloaded")]),
            system="system",
            api_messages=[{"role": "user", "content": "hello"}],
            rag_sources=[],
            tools=None,
            max_tokens=1024,
            temperature=0.7,
        )
        req = chat_router_module.SendMessageRequest(content="hello")

        events = collect_async_generator(stream_chat_events(runtime, req, self.engine))
        error_event = next(
            json.loads(event.replace("data: ", "").strip())
            for event in events
            if json.loads(event.replace("data: ", "").strip()).get("type") == "error"
        )
        self.assertEqual(error_event["type"], "error")
        self.assertIn("AI 服务当前繁忙", error_event["message"])

        with Session(self.engine) as session:
            assistant_messages = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).all()
            self.assertEqual(len(assistant_messages), 0)

    def test_stream_chat_events_detects_pseudo_tool_use_json_in_p3_follow_up(self):
        """If Claude outputs tool_use JSON as plain text in P3 follow-up, we detect and execute it."""
        conv_id = self._create_conversation()
        llm = FakeStreamingLLM(
            [
                # P1: read tool_use
                ['{"type":"tool_use","id":"tool-read","name":"read_project_markdown_document","input":{"action":"read","file_id":1}}'],
                # P3: text containing pseudo tool_use JSON (not real function calling)
                ['I will update now. {"type":"tool_use","id":"tool-write","name":"update_project_markdown_document","input":{"mode":"replace","file_id":1,"content":"# Updated"}}'],
                # P3 re-follow-up after tool execution
                ['Done updating the file.'],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            project_id=1,
            selected_model="claude-sonnet-4-6",
            llm=llm,
            system="system",
            api_messages=[{"role": "user", "content": "read and update the project document"}],
            rag_sources=[],
            tools=[{"name": "read_project_markdown_document"}, {"name": "update_project_markdown_document"}],
            max_tokens=1024,
            temperature=0.7,
        )
        req = chat_router_module.SendMessageRequest(content="read and update the project document")

        async def mock_execute(name, input_data):
            if name == "read_project_markdown_document":
                return {"type": "tool_result", "tool_name": name, "status": "success", "output": {"content": "# Old"}}
            if name == "update_project_markdown_document":
                return {"type": "tool_result", "tool_name": name, "status": "success", "output": {"ok": True, "action": "updated"}}
            return {"type": "tool_result", "tool_name": name, "status": "success", "output": {}}

        with patch("app.services.chat_streaming.registry.execute", new=AsyncMock(side_effect=mock_execute)):
            events = collect_async_generator(stream_chat_events(runtime, req, self.engine))

        joined = "".join(events)
        self.assertIn("Done updating the file.", joined)
        self.assertNotIn('"type":"tool_use"', joined)
        self.assertEqual(llm.calls, 3)  # P1, P3, P3 re-follow-up

        with Session(self.engine) as session:
            assistant_messages = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).all()
            self.assertEqual(len(assistant_messages), 1)
            self.assertIn("Done updating the file.", assistant_messages[0].content)
            metadata = json.loads(assistant_messages[0].metadata_json)
            tool_names = [tc["tool_name"] for tc in metadata["tool_calls"]]
            self.assertIn("read_project_markdown_document", tool_names)
            self.assertIn("update_project_markdown_document", tool_names)
            self.assertNotIn('"type":"tool_use"', assistant_messages[0].content)

    def test_stream_chat_events_suppresses_mixed_p1_tool_use_json(self):
        """If tool_use JSON leaks into P1 text, execute it without showing raw JSON to the user."""
        conv_id = self._create_conversation()
        create_tool_json = (
            '{"type":"tool_use","id":"tool-write","name":"update_project_markdown_document",'
            '"input":{"mode":"create","file_name":"first-meeting.md","content":"# First meeting"}}'
        )
        llm = FakeStreamingLLM(
            [
                [
                    '文档状态：初稿，待第一次沟通后更新\n'
                    '下一步行动：准备第一次沟通的详细内容和 PPT\n\n'
                    '{"type":"tool_use","id":"tool-empty","name":"update_project_markdown_document","input":{}}'
                    + create_tool_json
                ],
                ['已创建第一次沟通会详细内容文档。'],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            project_id=27,
            selected_model="claude-sonnet-4-6",
            llm=llm,
            system="system",
            api_messages=[{"role": "user", "content": "创建第一次沟通文档"}],
            rag_sources=[],
            tools=[{"name": "update_project_markdown_document"}],
            max_tokens=1024,
            temperature=0.7,
        )
        req = chat_router_module.SendMessageRequest(content="创建第一次沟通文档")

        async def mock_execute(name, input_data):
            return {
                "type": "tool_result",
                "tool_name": name,
                "status": "success",
                "output": {"ok": True, "action": "created", "id": 123, "name": input_data.get("file_name")},
            }

        with patch("app.services.chat_streaming.registry.execute", new=AsyncMock(side_effect=mock_execute)):
            events = collect_async_generator(stream_chat_events(runtime, req, self.engine))

        joined = "".join(events)
        self.assertIn("文档状态：初稿", joined)
        self.assertIn("已创建第一次沟通会详细内容文档。", joined)
        self.assertNotIn('{"type":"tool_use"', joined)
        self.assertEqual(llm.calls, 2)

        with Session(self.engine) as session:
            assistant_messages = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).all()
            self.assertEqual(len(assistant_messages), 1)
            self.assertIn("文档状态：初稿", assistant_messages[0].content)
            self.assertIn("已创建第一次沟通会详细内容文档。", assistant_messages[0].content)
            self.assertNotIn('{"type":"tool_use"', assistant_messages[0].content)

    def test_stream_chat_events_ignores_unrequested_markdown_write_for_direct_risk_answer(self):
        """Direct advisory answers must not become project document writes."""
        conv_id = self._create_conversation()
        llm = FakeStreamingLLM(
            [
                [
                    "基于结构化记忆，主要风险是决策信息不完整，建议本周确认客户决策链。\n\n"
                    "[TOOL_START:update_project_markdown_document]\n"
                    '{"type":"tool_use","id":"tool-write","name":"update_project_markdown_document",'
                    '"input":{"file_name":"项目风险清单.md","content":"# 风险清单"}}'
                ],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            project_id=26,
            selected_model="glm-5.1",
            llm=llm,
            system="system",
            api_messages=[{"role": "user", "content": "识别项目风险"}],
            rag_sources=[],
            tools=[{"name": "update_project_markdown_document"}],
            max_tokens=1024,
            temperature=0.7,
            action_policy="read_only_tool",
        )
        req = chat_router_module.SendMessageRequest(
            content="请基于当前项目的结构化记忆，识别最重要的项目风险和阻塞点，并给出建议的缓解动作。",
            project_id=26,
        )

        execute_mock = AsyncMock()
        with patch("app.services.chat_streaming.registry.execute", new=execute_mock):
            events = collect_async_generator(stream_chat_events(runtime, req, self.engine))

        joined = "".join(events)
        execute_mock.assert_not_awaited()
        self.assertNotIn("TOOL_START", joined)
        self.assertNotIn('{"type":"tool_use"', joined)

        with Session(self.engine) as session:
            assistant_messages = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).all()
            self.assertEqual(len(assistant_messages), 1)
            self.assertIn("主要风险", assistant_messages[0].content)
            self.assertNotIn("TOOL_START", assistant_messages[0].content)
            metadata = json.loads(assistant_messages[0].metadata_json)
            self.assertNotIn("tool_calls", metadata)
            self.assertNotIn("artifacts", metadata)

    def test_stream_chat_events_skips_invalid_tool_in_p2_without_breaking_pairing(self):
        """Invalid tool in P2 must still emit a tool_result so tool_use/tool_result pairing is preserved."""
        conv_id = self._create_conversation()
        llm = FakeStreamingLLM(
            [
                # P1: valid tool_use + invalid tool_use
                [
                    '{"type":"tool_use","id":"tool-1","name":"generate_docx","input":{"title":"OK"}}',
                    '{"type":"tool_use","id":"tool-2","name":"","input":null}',
                ],
                # P3: follow-up
                ["follow-up text"],
            ]
        )
        runtime = ChatRuntime(
            conv_id=conv_id,
            selected_model="claude-sonnet-4-6",
            llm=llm,
            system="system",
            api_messages=[{"role": "user", "content": "make doc"}],
            rag_sources=[],
            tools=[{"name": "generate_docx"}],
            max_tokens=1024,
            temperature=0.7,
        )
        req = chat_router_module.SendMessageRequest(content="make doc")

        with patch(
            "app.services.chat_streaming.registry.execute",
            new=AsyncMock(return_value={"type": "tool_result", "tool_name": "generate_docx", "status": "success", "output": {"file": "ok"}}),
        ):
            events = collect_async_generator(stream_chat_events(runtime, req, self.engine))

        # Should complete without API errors and persist message
        with Session(self.engine) as session:
            assistant_messages = session.exec(
                select(Message).where(Message.conversation_id == conv_id, Message.role == "assistant")
            ).all()
            self.assertEqual(len(assistant_messages), 1)


class ClientMemoryRouterTestCase(unittest.TestCase):
    def setUp(self):
        self.engine = create_test_engine()
        drop_all_tables(self.engine)
        SQLModel.metadata.create_all(self.engine)

        def override_session():
            with Session(self.engine) as session:
                yield session

        app = FastAPI()
        app.include_router(clients_router_module.router)
        app.dependency_overrides[clients_router_module.get_session] = override_session
        self.client = TestClient(app)

    def tearDown(self):
        self.client.close()
        self.engine.dispose()

    def test_memory_operations_summary_aggregates_project_and_client_failures(self):
        with Session(self.engine) as session:
            project = Project(
                name="Ops Project",
                client="Acme",
                memory_version=2,
                context_memory_json=json.dumps(
                    {
                        "_last_failure": {
                            "stage": "summary_warm",
                            "message": "Rate limit exceeded",
                            "retry_count": 1,
                            "failed_at": "2026-04-20T10:00:00",
                        },
                        "rebuild_log": [{"at": "2026-04-19T09:00:00", "trigger": "manual", "version": 2}],
                    },
                    ensure_ascii=False,
                ),
            )
            client = ClientRecord(
                name="Ops Client",
                client_memory_version=3,
                client_memory_json=json.dumps(
                    {
                        "_last_failure": {
                            "stage": "rebuild",
                            "message": "Client not found in source data",
                            "retry_count": 0,
                            "failed_at": "2026-04-20T11:00:00",
                        },
                        "rebuild_log": [{"at": "2026-04-19T12:00:00", "trigger": "manual", "version": 3}],
                    },
                    ensure_ascii=False,
                ),
            )
            session.add(project)
            session.add(client)
            session.commit()
            session.refresh(project)
            session.refresh(client)
            session.add(
                ProjectMemorySummary(
                    project_id=project.id,
                    summary_type="overview",
                    language="zh",
                    memory_version=2,
                    content="cached",
                )
            )
            session.add(
                ClientMemorySummary(
                    client_id=client.id,
                    summary_type="overview",
                    language="zh",
                    memory_version=3,
                    content="cached",
                )
            )
            session.commit()

        app = FastAPI()

        def override_session():
            with Session(self.engine) as session:
                yield session

        app.include_router(memory_operations_router_module.router)
        app.dependency_overrides[memory_operations_router_module.get_session] = override_session
        client = TestClient(app)
        try:
            with patch("app.services.scheduler.get_jobs", return_value=[]):
                resp = client.get("/memory/operations/summary")
        finally:
            client.close()

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["counts"]["recent_failures"], 2)
        self.assertEqual(body["counts"]["manual_attention"], 1)
        self.assertEqual(body["failure_summary"]["category_counts"]["rate_limit"], 1)
        self.assertEqual(body["failure_summary"]["category_counts"]["data"], 1)
        self.assertTrue(body["recent_failures"][0]["manual_attention"])
        self.assertEqual(body["budget"]["project"]["used"], 1)
        self.assertEqual(body["budget"]["client"]["used"], 1)

    def test_client_memory_snapshots_are_listed_and_can_rollback(self):
        with Session(self.engine) as session:
            client = ClientRecord(name="Snapshot Client", industry="Tech")
            session.add(client)
            session.commit()
            session.refresh(client)
            client_id = client.id

            client_contexts_module.save_client_memory(
                session,
                client_id,
                {"client_profile": "Version one", "decision_patterns": ["Pattern A"]},
                trigger="first_build",
            )
            client_contexts_module.save_client_memory(
                session,
                client_id,
                {"client_profile": "Version two", "decision_patterns": ["Pattern B"]},
                trigger="second_build",
            )

        snapshots_resp = self.client.get(f"/clients/{client_id}/memory/snapshots")
        self.assertEqual(snapshots_resp.status_code, 200)
        snapshots = snapshots_resp.json()
        self.assertEqual(len(snapshots), 2)
        older_snapshot = next(item for item in snapshots if item["memory_version"] == 1)

        detail_resp = self.client.get(f"/clients/{client_id}/memory/snapshots/{older_snapshot['id']}")
        self.assertEqual(detail_resp.status_code, 200)
        self.assertEqual(detail_resp.json()["memory"]["client_profile"], "Version one")

        diff_resp = self.client.get(f"/clients/{client_id}/memory/snapshots/{older_snapshot['id']}/diff")
        self.assertEqual(diff_resp.status_code, 200)
        diff_body = diff_resp.json()
        self.assertEqual(diff_body["scope"], "client")
        changed_fields = {item["field"]: item for item in diff_body["fields"]}
        self.assertEqual(changed_fields["client_profile"]["before"], "Version one")
        self.assertEqual(changed_fields["client_profile"]["after"], "Version two")
        self.assertEqual(changed_fields["decision_patterns"]["added"], ["Pattern B"])
        self.assertEqual(changed_fields["decision_patterns"]["removed"], ["Pattern A"])

        rollback_resp = self.client.post(f"/clients/{client_id}/memory/snapshots/{older_snapshot['id']}/rollback")
        self.assertEqual(rollback_resp.status_code, 200)
        self.assertEqual(rollback_resp.json()["memory"]["client_profile"], "Version one")
        self.assertEqual(rollback_resp.json()["memory_version"], 3)

        with Session(self.engine) as session:
            snapshots = session.exec(
                select(ClientMemorySnapshot).where(ClientMemorySnapshot.client_id == client_id)
            ).all()
        self.assertEqual(len(snapshots), 3)

    def test_client_memory_rebuild_flow(self):
        with Session(self.engine) as session:
            client = ClientRecord(name="Acme Corp", notes="Large enterprise client")
            project = Project(
                name="Alpha Delivery",
                client="Acme Corp",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Finance transformation project",
                        "key_risks": ["Slow approvals"],
                        "next_actions": ["Confirm sponsor"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=1,
                memory_stale=False,
            )
            session.add(client)
            session.add(project)
            session.commit()
            session.refresh(client)
            client_id = client.id

        with patch.object(
            clients_router_module,
            "complete_with_selected_model",
            new=AsyncMock(
                return_value=json.dumps(
                    {
                        "client_profile": "Large enterprise account with long procurement cycles",
                        "decision_patterns": ["Budget closes in Q4"],
                        "key_contacts": [{"name": "Jane", "role": "CFO", "note": "Economic buyer"}],
                        "lessons_learned": ["Senior sponsors want phased delivery"],
                        "project_history": [{"project_name": "Alpha Delivery", "status": "delivering", "outcome": "", "key_factor": "Finance transformation"}],
                        "sensitive_topics": ["Avoid aggressive automation claims"],
                    },
                    ensure_ascii=False,
                )
            ),
        ):
            resp = self.client.post(f"/clients/{client_id}/memory/rebuild")

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["memory_version"], 1)
        self.assertFalse(body["memory_stale"])
        self.assertIn("Budget closes in Q4", body["memory"]["decision_patterns"])

    def test_promote_project_memory_to_client(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Acme Corp",
                client_memory_json=json.dumps({"client_profile": "Existing client profile"}, ensure_ascii=False),
                client_memory_version=1,
                client_memory_stale=False,
            )
            project = Project(
                name="Alpha Delivery",
                client="Acme Corp",
                context_memory_json=json.dumps(
                    {
                        "project_brief": "Finance transformation project",
                        "key_risks": ["Slow approvals"],
                        "stakeholder_notes": ["CFO is supportive"],
                    },
                    ensure_ascii=False,
                ),
                memory_version=2,
                memory_stale=False,
            )
            session.add(client)
            session.add(project)
            session.commit()
            session.refresh(client)
            session.refresh(project)
            client_id = client.id
            project_id = project.id

        with patch.object(
            clients_router_module,
            "complete_with_selected_model",
            new=AsyncMock(
                return_value=json.dumps(
                    {
                        "client_profile": "Existing client profile",
                        "decision_patterns": ["Budget closes in Q4"],
                        "key_contacts": [],
                        "lessons_learned": ["Finance transformation needs executive sponsorship"],
                        "project_history": [{"project_name": "Alpha Delivery", "status": "delivering", "outcome": "", "key_factor": "Executive sponsorship"}],
                        "sensitive_topics": [],
                    },
                    ensure_ascii=False,
                )
            ),
        ):
            resp = self.client.post(
                f"/clients/{client_id}/memory/promote-project",
                json={"project_id": project_id},
            )

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["memory_version"], 2)
        self.assertIn(project_id, body["memory"]["source_project_ids"])

    def test_client_memory_batch_rebuild_updates_requested_clients(self):
        with Session(self.engine) as session:
            first = ClientRecord(name="Acme Corp", notes="Manufacturing group")
            second = ClientRecord(
                name="Beta Labs",
                client_memory_json=json.dumps({"client_profile": "Existing"}, ensure_ascii=False),
                client_memory_version=1,
                client_memory_stale=False,
            )
            session.add(first)
            session.add(second)
            session.add(Project(name="Alpha Delivery", client="Acme Corp", status="delivering"))
            session.commit()
            session.refresh(first)
            session.refresh(second)
            first_id = first.id
            second_id = second.id

        with patch.object(
            clients_router_module,
            "complete_with_selected_model",
            new=AsyncMock(
                return_value=json.dumps(
                    {
                        "client_profile": "Manufacturing account",
                        "decision_patterns": ["Needs phased rollout"],
                        "key_contacts": [],
                        "lessons_learned": ["Delivery proof matters"],
                        "project_history": [],
                        "sensitive_topics": [],
                    },
                    ensure_ascii=False,
                )
            ),
        ) as mocked_complete:
            resp = self.client.post(
                "/clients/memory/rebuild-batch",
                json={"client_ids": [first_id, second_id], "stale_only": True},
            )

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["requested_count"], 2)
        self.assertEqual(body["rebuilt_count"], 1)
        self.assertEqual(body["rebuilt"][0]["client_id"], first_id)
        self.assertEqual(body["skipped"][0]["client_id"], second_id)
        self.assertEqual(body["skipped"][0]["reason"], "not_stale")
        self.assertEqual(mocked_complete.await_count, 1)

    def test_list_client_memory_jobs_returns_queued_rebuilds(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Acme Corp",
                industry="Manufacturing",
                client_memory_version=2,
                client_memory_stale=True,
            )
            session.add(client)
            session.commit()
            session.refresh(client)
            client_id = client.id

        fake_job = SimpleNamespace(
            id=f"client_memory_rebuild_{client_id}",
            next_run_time=utc_now_naive() + timedelta(minutes=5),
        )

        with patch.object(clients_router_module.scheduler_service, "get_jobs", return_value=[fake_job]):
            resp = self.client.get("/clients/memory/jobs")

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["count"], 1)
        self.assertEqual(body["jobs"][0]["client_id"], client_id)
        self.assertEqual(body["jobs"][0]["job_type"], "rebuild")
        self.assertEqual(body["jobs"][0]["client_name"], "Acme Corp")
        self.assertEqual(body["jobs"][0]["job_type"], "rebuild")

    def test_list_client_memory_jobs_includes_summary_warm_tasks(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Acme Corp",
                industry="Manufacturing",
                client_memory_version=2,
                client_memory_stale=False,
            )
            session.add(client)
            session.commit()
            session.refresh(client)
            client_id = client.id

        fake_job = SimpleNamespace(
            id=f"client_memory_summary_warm_{client_id}_zh",
            next_run_time=utc_now_naive() + timedelta(minutes=3),
        )

        with patch.object(clients_router_module.scheduler_service, "get_jobs", return_value=[fake_job]):
            resp = self.client.get("/clients/memory/jobs")

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["count"], 1)
        self.assertEqual(body["jobs"][0]["job_type"], "summary_warm")
        self.assertEqual(body["jobs"][0]["language"], "zh")

    def test_list_client_memory_jobs_includes_budget_retry_and_recent_failures(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Acme Corp",
                industry="Manufacturing",
                client_memory_version=4,
                client_memory_stale=False,
                client_memory_json=json.dumps(
                    {
                        "_last_failure": {
                            "stage": "rebuild",
                            "message": "Temporary model timeout",
                            "retry_count": 1,
                            "failed_at": "2026-04-19T11:00:00",
                        },
                        "rebuild_log": [
                            {"at": "2026-04-18T12:30:00", "trigger": "client_manual", "version": 4}
                        ],
                    }
                ),
            )
            session.add(client)
            session.commit()
            session.refresh(client)
            client_id = client.id
            session.add(
                ClientMemorySummary(
                    client_id=client_id,
                    summary_type="overview",
                    language="zh",
                    memory_version=4,
                    content="cached",
                )
            )
            session.commit()

        fake_job = SimpleNamespace(
            id=f"client_memory_summary_warm_{client_id}_zh",
            next_run_time=utc_now_naive(),
        )
        with patch.object(clients_router_module.scheduler_service, "get_jobs", return_value=[fake_job]), patch.object(
            clients_router_module.scheduler_service,
            "get_job_metadata",
            return_value={
                "retry_count": 2,
                "max_retries": 3,
                "trigger": "rebuild_completed",
                "summary_types": ["overview", "risk"],
            },
        ):
            resp = self.client.get("/clients/memory/jobs")

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["budget"]["used"], 1)
        self.assertEqual(body["jobs"][0]["retry_count"], 2)
        self.assertEqual(body["jobs"][0]["summary_types"], ["overview", "risk"])
        self.assertEqual(body["recent_failures"][0]["stage"], "rebuild")
        self.assertEqual(body["recent_successes"][0]["client_id"], client_id)
        self.assertEqual(body["recent_successes"][0]["completed_at"], "2026-04-18T12:30:00")
        self.assertEqual(body["recent_successes"][0]["version"], 4)

    def test_cancel_client_memory_jobs_removes_rebuild_job(self):
        removed: list[str] = []

        with patch.object(
            clients_router_module.scheduler_service,
            "remove_job",
            side_effect=lambda job_id: removed.append(job_id),
        ):
            resp = self.client.post("/clients/memory/jobs/42/cancel")

        self.assertEqual(resp.status_code, 200)
        self.assertIn("client_memory_rebuild_42", removed)
        self.assertIn("client_memory_summary_warm_42_zh", removed)
        self.assertIn("client_memory_summary_warm_42_en", removed)
        self.assertIn("client_memory_summary_warm_42_default", removed)

    def test_run_client_memory_jobs_now_rebuilds_client(self):
        with Session(self.engine) as session:
            client = ClientRecord(name="Acme Corp", notes="Strategic account")
            session.add(client)
            session.add(Project(name="Alpha Delivery", client="Acme Corp", status="delivering"))
            session.commit()
            session.refresh(client)
            client_id = client.id

        with patch.object(
            clients_router_module,
            "complete_with_selected_model",
            new=AsyncMock(
                return_value=json.dumps(
                    {
                        "client_profile": "Strategic account with long buying cycles",
                        "decision_patterns": ["Needs internal champion"],
                        "key_contacts": [],
                        "lessons_learned": ["Pilot success matters"],
                        "project_history": [],
                        "sensitive_topics": [],
                    },
                    ensure_ascii=False,
                )
            ),
        ), patch.object(clients_router_module.scheduler_service, "remove_job") as mocked_remove_job:
            resp = self.client.post(f"/clients/memory/jobs/{client_id}/run-now")

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertTrue(body["ok"])
        self.assertEqual(body["action"], "rebuild")
        self.assertGreaterEqual(body["memory_version"], 1)
        removed_job_ids = [call.args[0] for call in mocked_remove_job.call_args_list]
        self.assertIn(f"client_memory_rebuild_{client_id}", removed_job_ids)

    def test_client_memory_warm_summaries_batch_generates_common_cached_views(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Acme Corp",
                client_memory_json=json.dumps(
                    {
                        "client_profile": "Strategic manufacturing account",
                        "decision_patterns": ["Prefers phased rollout"],
                        "key_contacts": [{"name": "Jane", "role": "CFO", "note": "Executive sponsor"}],
                        "lessons_learned": ["Value proof matters before scale"],
                        "project_history": [],
                        "sensitive_topics": ["Avoid promising fixed dates too early"],
                    },
                    ensure_ascii=False,
                ),
                client_memory_version=3,
                client_memory_stale=False,
            )
            session.add(client)
            session.commit()
            session.refresh(client)
            client_id = client.id

        with patch.object(
            clients_router_module,
            "complete_with_selected_model",
            new=AsyncMock(return_value="- Cached summary"),
        ) as mocked_complete:
            resp = self.client.post(
                "/clients/memory/warm-summaries-batch",
                json={"client_ids": [client_id], "summary_types": ["overview", "stakeholder"], "language": "zh-CN"},
            )

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["processed_count"], 1)
        self.assertGreaterEqual(body["warmed_count"], 2)
        self.assertEqual(mocked_complete.await_count, 2)

        with Session(self.engine) as session:
            cached = session.exec(
                select(ClientMemorySummary).where(ClientMemorySummary.client_id == client_id)
            ).all()
            self.assertEqual(len(cached), 2)

    def test_client_memory_warm_summaries_batch_supports_relationship_and_delivery(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Acme Corp",
                client_memory_json=json.dumps(
                    {
                        "client_profile": "Strategic manufacturing account",
                        "decision_patterns": ["Prefers phased rollout"],
                        "key_contacts": [{"name": "Jane", "role": "CFO", "note": "Executive sponsor"}],
                        "lessons_learned": ["Value proof matters before scale"],
                        "project_history": [{"project_name": "Alpha", "status": "won", "outcome": "Expanded"}],
                        "sensitive_topics": ["Avoid promising fixed dates too early"],
                    },
                    ensure_ascii=False,
                ),
                client_memory_version=3,
                client_memory_stale=False,
            )
            session.add(client)
            session.commit()
            session.refresh(client)
            client_id = client.id

        with patch.object(
            clients_router_module,
            "complete_with_selected_model",
            new=AsyncMock(return_value="- Cached summary"),
        ) as mocked_complete:
            resp = self.client.post(
                "/clients/memory/warm-summaries-batch",
                json={
                    "client_ids": [client_id],
                    "summary_types": ["relationship", "delivery"],
                    "language": "en-US",
                },
            )

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["processed_count"], 1)
        self.assertGreaterEqual(body["warmed_count"], 2)
        self.assertEqual(mocked_complete.await_count, 2)

        with Session(self.engine) as session:
            cached = session.exec(
                select(ClientMemorySummary).where(ClientMemorySummary.client_id == client_id)
            ).all()
            self.assertEqual(sorted(item.summary_type for item in cached), ["delivery", "relationship"])

    def test_client_memory_summary_uses_cache_by_memory_version(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Acme Corp",
                client_memory_json=json.dumps(
                    {
                        "client_profile": "Strategic manufacturing account",
                        "decision_patterns": ["Prefers phased rollout"],
                        "key_contacts": [{"name": "Jane", "role": "CFO", "note": "Executive sponsor"}],
                        "lessons_learned": ["Value proof matters before scale"],
                        "project_history": [],
                        "sensitive_topics": ["Avoid promising fixed dates too early"],
                    },
                    ensure_ascii=False,
                ),
                client_memory_version=3,
                client_memory_stale=False,
            )
            session.add(client)
            session.commit()
            session.refresh(client)
            session.add(
                ClientMemorySummary(
                    client_id=client.id,
                    summary_type="overview",
                    language="zh",
                    memory_version=3,
                    content="- 已缓存客户摘要",
                )
            )
            session.commit()
            client_id = client.id

        with patch.object(
            clients_router_module,
            "complete_with_selected_model",
            new=AsyncMock(return_value="- 新摘要"),
        ) as mocked_complete:
            cached_resp = self.client.post(
                f"/clients/{client_id}/memory/summarize",
                json={"language": "zh-CN", "summary_type": "overview"},
            )

        self.assertEqual(cached_resp.status_code, 200)
        self.assertTrue(cached_resp.json()["cached"])
        self.assertEqual(cached_resp.json()["content"], "- 已缓存客户摘要")
        self.assertEqual(mocked_complete.await_count, 0)

        with patch.object(
            clients_router_module,
            "complete_with_selected_model",
            new=AsyncMock(return_value="- 新摘要"),
        ) as mocked_complete:
            fresh_resp = self.client.post(
                f"/clients/{client_id}/memory/summarize",
                json={"language": "zh-CN", "summary_type": "lessons", "force_refresh": True},
            )

        self.assertEqual(fresh_resp.status_code, 200)
        self.assertFalse(fresh_resp.json()["cached"])
        self.assertEqual(fresh_resp.json()["summary_type"], "lessons")
        self.assertEqual(mocked_complete.await_count, 1)

    def test_client_memory_summary_cache_endpoint_does_not_generate(self):
        with Session(self.engine) as session:
            client = ClientRecord(
                name="Cached Client",
                client_memory_json=json.dumps(
                    {
                        "client_profile": "Strategic account",
                        "decision_patterns": ["Needs executive proof"],
                        "project_history": [],
                    },
                    ensure_ascii=False,
                ),
                client_memory_version=4,
                client_memory_stale=False,
            )
            session.add(client)
            session.commit()
            session.refresh(client)
            session.add(
                ClientMemorySummary(
                    client_id=client.id,
                    summary_type="risk",
                    language="en",
                    memory_version=4,
                    content="- cached client risk summary",
                )
            )
            session.commit()
            client_id = client.id

        with patch.object(
            clients_router_module,
            "complete_with_selected_model",
            new=AsyncMock(side_effect=AssertionError("should not call llm")),
        ):
            cached_resp = self.client.get(
                f"/clients/{client_id}/memory/summaries/risk",
                params={"language": "en-US"},
            )

        self.assertEqual(cached_resp.status_code, 200)
        self.assertTrue(cached_resp.json()["cached"])
        self.assertEqual(cached_resp.json()["content"], "- cached client risk summary")

        missing_resp = self.client.get(
            f"/clients/{client_id}/memory/summaries/opportunity",
            params={"language": "en-US"},
        )
        self.assertEqual(missing_resp.status_code, 404)


class BuiltinSkillsTestCase(unittest.TestCase):
    def setUp(self):
        self.engine = create_test_engine()
        drop_all_tables(self.engine)
        SQLModel.metadata.create_all(self.engine)

    def tearDown(self):
        self.engine.dispose()

    def test_builtin_digital_skills_are_seeded_idempotently(self):
        expected_digital_skills = {
            "AI 用例优先级矩阵",
            "数字化战略设计",
            "数字化成熟度评估",
            "企业架构蓝图设计",
            "数据治理咨询方案",
            "流程数字化改造",
            "数字技术路线图",
            "数字化组织变革",
            "数字化 ROI 商业案例",
            "行业数字化蓝图",
        }

        with Session(self.engine) as session:
            first_count = skills_router_module.ensure_builtin_pro_skills(session)
            second_count = skills_router_module.ensure_builtin_pro_skills(session)
            digital_skills = session.exec(
                select(Skill).where(Skill.category == "数字化与技术")
            ).all()

        self.assertGreaterEqual(first_count, len(expected_digital_skills))
        self.assertEqual(second_count, 0)
        names = {skill.name for skill in digital_skills}
        self.assertTrue(expected_digital_skills <= names)
        strategy = next(skill for skill in digital_skills if skill.name == "数字化战略设计")
        self.assertIn("数字化战略方案", strategy.user_template)
        self.assertIn("业务战略", strategy.system_prompt)
        self.assertIn("Huawei 5-See 3-Define", strategy.system_prompt)
        self.assertIn("Foundation", strategy.system_prompt)
        self.assertIn("generate_ppt_from_skill", strategy.system_prompt)
        self.assertEqual(strategy.max_tokens, 32768)
        self.assertIn("generate_ppt_from_skill", strategy.tools)
        self.assertNotIn("save_json", strategy.tools)
        tool_defs = json.loads(strategy.tools_definition_json)
        tool_def_names = {tool.get("name") for tool in tool_defs}
        self.assertIn("generate_ppt_from_skill", tool_def_names)
        self.assertNotIn("save_json", tool_def_names)

    def test_presentation_builder_skill_is_seeded_with_ppt_tooling(self):
        with Session(self.engine) as session:
            skills_router_module.ensure_builtin_pro_skills(session)
            skill = session.exec(
                select(Skill).where(Skill.name == skills_router_module.PRESENTATION_BUILDER_SKILL_NAME)
            ).one()

        self.assertEqual(skill.category, "提案与项目交付")
        self.assertIn("presentation-builder workflow", skill.system_prompt)
        self.assertIn("deck_type", skill.system_prompt)
        self.assertIn("generate_ppt_from_skill", skill.system_prompt)
        self.assertEqual(skill.tools, ["generate_ppt_from_skill"])
        self.assertEqual(skill.max_tokens, 24576)
        tool_defs = json.loads(skill.tools_definition_json)
        tool_def_names = {tool.get("name") for tool in tool_defs}
        self.assertEqual(tool_def_names, {"generate_ppt_from_skill"})

    def test_existing_digital_strategy_skill_tools_are_upgraded(self):
        old_skill = Skill(
            name="数字化战略设计",
            category="数字化与技术",
            description="old",
            system_prompt="digital-strategy 工作流\n旧提示",
            user_template="旧模板",
            estimated_time="~25 min",
            max_tokens=4096,
            tools_definition_json=json.dumps(
                [
                    {"name": "strategy", "type": "legacy"},
                    {"name": "roadmap", "type": "legacy"},
                    {"name": "report", "type": "legacy"},
                ]
            ),
        )
        old_skill.tools = ["strategy", "roadmap", "report"]

        with Session(self.engine) as session:
            session.add(old_skill)
            session.commit()

            changed = skills_router_module.ensure_builtin_pro_skills(session)
            strategy = session.exec(
                select(Skill).where(Skill.name == "数字化战略设计")
            ).one()
            second_count = skills_router_module.ensure_builtin_pro_skills(session)

        self.assertGreaterEqual(changed, 1)
        self.assertEqual(second_count, 0)
        self.assertEqual(strategy.tools, ["generate_ppt_from_skill"])
        self.assertEqual(strategy.max_tokens, 32768)
        tool_defs = json.loads(strategy.tools_definition_json)
        tool_def_names = {tool.get("name") for tool in tool_defs}
        self.assertIn("generate_ppt_from_skill", tool_def_names)
        self.assertNotIn("save_json", tool_def_names)
        self.assertNotIn("strategy", tool_def_names)

    def test_digital_strategy_file_skill_assets_exist(self):
        repo_root = Path(__file__).resolve().parents[2]
        skill_dir = repo_root / "skills" / "digital-strategy"

        self.assertTrue((skill_dir / "SKILL.md").exists())
        self.assertTrue((skill_dir / "references" / "frameworks.md").exists())
        self.assertTrue((skill_dir / "references" / "industry-notes.md").exists())

        skill_text = (skill_dir / "SKILL.md").read_text(encoding="utf-8")
        framework_text = (skill_dir / "references" / "frameworks.md").read_text(encoding="utf-8")
        industry_text = (skill_dir / "references" / "industry-notes.md").read_text(encoding="utf-8")
        self.assertIn("name: digital-strategy", skill_text)
        self.assertIn("generate_ppt_from_skill", skill_text)
        self.assertNotIn("save_json", skill_text)
        self.assertIn("Huawei 5-See 3-Define", framework_text)
        self.assertIn("Manufacturing", industry_text)

    def test_digital_strategy_ppt_normalizer_expands_sparse_deck(self):
        from app.tools.file_generators import _normalize_digital_strategy_slides, _wants_visual_slide

        slides = _normalize_digital_strategy_slides(
            [{"type": "content", "title": "Executive Summary", "content": "- Initial thesis"}]
        )

        self.assertGreaterEqual(len(slides), 20)
        titles = {slide["title"] for slide in slides}
        self.assertTrue("Use-Case Portfolio" in titles or "场景组合：平衡快赢、基础能力和战略差异化" in titles)
        self.assertTrue("Executive Alignment" in titles or "高层共识" in titles)
        self.assertTrue("Target Blueprint" in titles or "目标蓝图" in titles)
        self.assertIn("数据治理优先级：先解决会阻断多个场景复制的共性问题", titles)
        self.assertIn("90 天行动计划：把共识转成可执行启动方案", titles)
        investment_slide = next(
            slide for slide in slides
            if slide["title"] in {"Investment, KPI and Risk Controls", "投资、KPI 与风险控制：建立可追踪的价值闭环"}
        )
        self.assertTrue(
            "technology, data, talent, change" in investment_slide["content"]
            or "技术、数据、人才、变革" in investment_slide["content"]
        )
        self.assertGreaterEqual(investment_slide["content"].count("\n- ") + 1, 5)
        self.assertTrue(_wants_visual_slide("Maturity Heatmap: Strengths vs Constraints"))
        self.assertTrue(_wants_visual_slide("Investment, KPI and Risk Controls"))
        self.assertTrue(_wants_visual_slide("Risk Register and Mitigation Plan"))

    def test_digital_strategy_ppt_normalizer_replaces_thin_twenty_page_deck(self):
        from app.tools.file_generators import _normalize_digital_strategy_slides

        thin_slides = [
            {"type": "content", "title": f"Thin slide {idx}", "content": "- A"}
            for idx in range(1, 21)
        ]

        slides = _normalize_digital_strategy_slides(thin_slides)
        titles = {slide["title"] for slide in slides}

        self.assertGreaterEqual(len(slides), 22)
        self.assertIn("执行摘要：把数字化作为业务价值组合来管理", titles)
        self.assertIn("三阶段路线图：夯实基础、规模复制、领先优化", titles)
        self.assertIn("行业标杆与差距启示：用外部参照校准转型目标", titles)
        self.assertIn("附录：评估与访谈指南", titles)
        self.assertGreaterEqual(
            len({slide.get("layout_key") for slide in slides if slide.get("layout_key")}),
            8,
        )

    def test_digital_strategy_generation_renders_richer_template_deck(self):
        from app.tools.file_generators import generate_ppt_from_skill

        result = asyncio.run(
            generate_ppt_from_skill(
                "digital-strategy",
                "Digital Strategy Rich Deck Test",
                [{"type": "content", "title": "Executive Summary", "content": "- Initial thesis"}],
                "Local validation run",
            )
        )

        self.assertTrue(result["success"], result)
        self.assertEqual(result["template_name"], "Template.pptx")
        self.assertEqual(result["template_mode"], "cloned_prototype_slides")
        self.assertGreaterEqual(result["slide_count"], 22)
        generated_path = Path(result["full_path"])
        self.assertTrue(generated_path.exists())
        generated_path.unlink(missing_ok=True)

    def test_presentation_builder_normalizer_supports_deck_presets(self):
        from app.tools.file_generators import _normalize_presentation_builder_slides

        proposal_slides = _normalize_presentation_builder_slides(
            [{"type": "content", "title": "Opening", "content": "- Initial point"}],
            "proposal",
        )
        update_slides = _normalize_presentation_builder_slides([], "project-update")

        self.assertGreaterEqual(len(proposal_slides), 12)
        self.assertIn("建议方法：诊断、设计、对齐与动员", {slide["title"] for slide in proposal_slides})
        self.assertIn("商务、风险与下一步：明确启动条件", {slide["title"] for slide in proposal_slides})
        self.assertIn("关键假设：说明材料当前依据和待验证项", {slide["title"] for slide in proposal_slides})
        self.assertGreaterEqual(len(update_slides), 12)
        self.assertIn("计划与实际：识别偏差、依赖和纠偏动作", {slide["title"] for slide in update_slides})
        self.assertIn("风险、问题与决策：聚焦本周期必须处理的事项", {slide["title"] for slide in update_slides})

    def test_presentation_builder_formats_rich_slide_types_for_ppt(self):
        from app.tools.file_generators import (
            _normalize_presentation_builder_slide_format,
        )

        slides = _normalize_presentation_builder_slide_format(
            [
                {
                    "type": "roadmap",
                    "title": "Roadmap",
                    "left_content": "Phase 1\n- Align",
                    "content": "Phase 2\n- Scale",
                    "right_content": "Phase 3\n- Optimize",
                },
                {
                    "type": "matrix",
                    "title": "Prioritization",
                    "labels": ["Quick wins", "Foundations"],
                    "content": "- Use value and feasibility to sequence work",
                },
                {
                    "type": "content",
                    "title": "Dense Page",
                    "content": "- One\n- Two\n- Three\n- Four\n- Five\n- Six\n- Seven",
                },
            ]
        )

        self.assertEqual(slides[0]["type"], "roadmap")
        self.assertIn("Phase 1", slides[0]["content"])
        self.assertIn("Phase 3", slides[0]["content"])
        self.assertEqual(slides[1]["type"], "matrix")
        self.assertIn("Quick wins", slides[1]["content"])
        self.assertLessEqual(slides[2]["content"].count("\n- ") + 1, 6)

    def test_ppt_text_cleanup_removes_markdown_formatting(self):
        from app.tools.file_generators import _split_bullets

        bullets = _split_bullets(
            "### **关键发现**\n"
            "- **价值机会**：减少人工处理\n"
            "| 维度 | 结论 |\n"
            "| --- | --- |\n"
            "| 数据 | 需要统一口径 |\n",
            limit=5,
        )

        self.assertIn("关键发现", bullets[0])
        self.assertIn("价值机会：减少人工处理", bullets[1])
        self.assertIn("数据；需要统一口径", bullets[-1])
        self.assertNotIn("**", "\n".join(bullets))

    def test_presentation_builder_uses_digital_strategy_template_fallback(self):
        from app.tools.file_generators import generate_ppt_from_skill

        result = asyncio.run(
            generate_ppt_from_skill(
                "presentation-builder",
                "Presentation Builder Smoke Test",
                [],
                "Local validation run",
                deck_type="proposal",
            )
        )

        self.assertTrue(result["success"], result)
        self.assertEqual(result["file_type"], "pptx")
        self.assertEqual(result["template_name"], "Template.pptx")
        self.assertTrue(result["template_applied"])
        self.assertGreaterEqual(result["slide_count"], 14)
        generated_path = Path(result["full_path"])
        self.assertTrue(generated_path.exists())
        generated_path.unlink(missing_ok=True)

    def test_presentation_builder_file_skill_assets_exist(self):
        repo_root = Path(__file__).resolve().parents[2]
        skill_dir = repo_root / "skills" / "presentation-builder"

        self.assertTrue((skill_dir / "SKILL.md").exists())
        skill_text = (skill_dir / "SKILL.md").read_text(encoding="utf-8")
        self.assertIn("name: presentation-builder", skill_text)
        self.assertIn("deck_type", skill_text)
        self.assertIn("generate_ppt_from_skill", skill_text)

    def test_ppt_generation_clones_template_slide_artwork(self):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        from app.tools.file_generators import generate_ppt

        with tempfile.TemporaryDirectory() as tmpdir:
            template_path = Path(tmpdir) / "Template.pptx"
            prs = Presentation()
            blank = prs.slide_layouts[6]
            while len(prs.slides) > 0:
                slide_id = list(prs.slides._sldIdLst)[0]
                prs.slides._sldIdLst.remove(slide_id)

            cover = prs.slides.add_slide(blank)
            title = cover.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            title.name = "aria_cover_title"
            subtitle = cover.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
            subtitle.name = "aria_cover_subtitle"

            content = prs.slides.add_slide(blank)
            content.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1), Inches(7)).name = "brand_artifact"
            content.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).name = "aria_slide_title"
            content.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3)).name = "aria_slide_body"

            two_col = prs.slides.add_slide(blank)
            two_col.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).name = "aria_slide_title"
            two_col.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(3)).name = "aria_left_body"
            two_col.shapes.add_textbox(Inches(6), Inches(2), Inches(4), Inches(3)).name = "aria_right_body"

            visual = prs.slides.add_slide(blank)
            visual.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).name = "aria_slide_title"
            visual.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(3)).name = "aria_slide_body"
            visual.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2), Inches(3), Inches(3)).name = "aria_visual_area"

            prs.slides.add_slide(blank)
            prs.save(template_path)

            result = asyncio.run(generate_ppt(
                "Template clone check",
                [
                    {"type": "content", "title": "One", "content": "- A"},
                    {"type": "content", "title": "Two", "content": "- B"},
                    {"type": "content", "title": "Three", "content": "- C"},
                ],
                "Subtitle",
                str(template_path),
            ))

            self.assertTrue(result["success"])
            generated_path = Path(result["full_path"])
            generated = Presentation(generated_path)
            body_slides = list(generated.slides)[1:4]
            self.assertEqual(len(body_slides), 3)
            for slide in body_slides:
                self.assertIn("brand_artifact", {shape.name for shape in slide.shapes})
            generated_path.unlink(missing_ok=True)

    def test_ppt_generation_preserves_named_title_style_and_fits_one_line(self):
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        from app.tools.file_generators import generate_ppt

        with tempfile.TemporaryDirectory() as tmpdir:
            template_path = Path(tmpdir) / "Template.pptx"
            prs = Presentation()
            blank = prs.slide_layouts[6]

            cover = prs.slides.add_slide(blank)
            cover.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).name = "aria_cover_title"
            cover.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1)).name = "aria_cover_subtitle"

            section = prs.slides.add_slide(blank)
            section_title = section.shapes.add_textbox(Inches(1), Inches(2), Inches(6.5), Inches(0.9))
            section_title.name = "aria_section_title"
            section_title.text_frame.text = "Section"
            section_title.text_frame.paragraphs[0].font.size = Pt(40)
            section_title.text_frame.paragraphs[0].font.bold = True
            section_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
            section.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1), Inches(0.5)).name = "aria_section_number"

            content = prs.slides.add_slide(blank)
            content.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1), Inches(7)).name = "brand_artifact"
            content_title = content.shapes.add_textbox(Inches(1), Inches(0.5), Inches(7), Inches(0.6))
            content_title.name = "aria_slide_title"
            content_title.text_frame.text = "Title"
            content_title.text_frame.paragraphs[0].font.size = Pt(30)
            content.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(3)).name = "aria_slide_body"

            two_col = prs.slides.add_slide(blank)
            two_col.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.6)).name = "aria_slide_title"
            two_col.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(3)).name = "aria_left_body"
            two_col.shapes.add_textbox(Inches(6), Inches(1.5), Inches(4), Inches(3)).name = "aria_right_body"

            visual = prs.slides.add_slide(blank)
            visual.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.6)).name = "aria_slide_title"
            visual.shapes.add_textbox(Inches(1), Inches(1.3), Inches(4), Inches(3)).name = "aria_slide_body"
            visual.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2), Inches(3), Inches(3)).name = "aria_visual_area"

            prs.slides.add_slide(blank)
            prs.save(template_path)

            long_section_title = "Executive Alignment and Governance Model for a Multi-Year Digital Transformation"
            long_content_title = "A very long content page title should be intelligently resized to remain on one line"
            result = asyncio.run(generate_ppt(
                "Template title fit check",
                [
                    {"type": "section", "title": long_section_title},
                    {"type": "content", "title": long_content_title, "content": "- First point\n- Second point"},
                ],
                "Subtitle",
                str(template_path),
            ))

            self.assertTrue(result["success"], result)
            generated_path = Path(result["full_path"])
            generated = Presentation(generated_path)

            section_slide = generated.slides[1]
            generated_section_title = next(shape for shape in section_slide.shapes if shape.name == "aria_section_title")
            self.assertEqual(generated_section_title.text, long_section_title)
            section_font = generated_section_title.text_frame.paragraphs[0].font
            self.assertEqual(section_font.color.rgb, RGBColor(0xC0, 0x00, 0x00))
            self.assertLessEqual(section_font.size.pt, 40)
            self.assertFalse(generated_section_title.text_frame.word_wrap)

            content_slide = generated.slides[2]
            generated_content_title = next(shape for shape in content_slide.shapes if shape.name == "aria_slide_title")
            content_font = generated_content_title.text_frame.paragraphs[0].font
            self.assertLessEqual(content_font.size.pt, 30)
            self.assertFalse(generated_content_title.text_frame.word_wrap)

            generated_path.unlink(missing_ok=True)


    def test_get_skill_detail_by_id(self):
        with Session(self.engine) as session:
            skill = Skill(
                name="Detail Skill",
                category="Digital",
                description="Detail page skill",
                system_prompt="System prompt",
                user_template="User template",
                estimated_time="~5 min",
                tools_definition_json="[]",
            )
            session.add(skill)
            session.commit()
            session.refresh(skill)
            skill_id = skill.id

        def override_session():
            with Session(self.engine) as session:
                yield session

        app = FastAPI()
        app.include_router(skills_router_module.router)
        app.dependency_overrides[skills_router_module.get_session] = override_session
        client = TestClient(app)
        try:
            resp = client.get(f"/skills/{skill_id}")
            missing_resp = client.get("/skills/999999")
        finally:
            client.close()

        self.assertEqual(resp.status_code, 200)
        body = resp.json()
        self.assertEqual(body["id"], skill_id)
        self.assertEqual(body["name"], "Detail Skill")
        self.assertEqual(body["system_prompt"], "System prompt")
        self.assertEqual(body["user_template"], "User template")
        self.assertEqual(missing_resp.status_code, 404)


class MessageRouterTestCase(unittest.TestCase):
    def setUp(self):
        self.engine = create_test_engine()
        drop_all_tables(self.engine)
        SQLModel.metadata.create_all(self.engine)

        def override_session():
            with Session(self.engine) as session:
                yield session

        app = FastAPI()
        app.include_router(auth_router_module.router)
        app.include_router(messages_router_module.router)
        app.dependency_overrides[auth_router_module.get_session] = override_session
        app.dependency_overrides[messages_router_module.get_session] = override_session
        self.client = TestClient(app)

        with Session(self.engine) as session:
            admin = User(
                email="admin@example.com",
                display_name="Admin",
                password_hash=auth_router_module._hash("password123"),
                is_admin=True,
                is_active=True,
            )
            member = User(
                email="user@example.com",
                display_name="Member",
                password_hash=auth_router_module._hash("password123"),
                is_admin=False,
                is_active=True,
            )
            session.add(admin)
            session.add(member)
            session.commit()

        admin_login = self.client.post(
            "/auth/login",
            json={"email": "admin@example.com", "password": "password123"},
        )
        user_login = self.client.post(
            "/auth/login",
            json={"email": "user@example.com", "password": "password123"},
        )
        self.admin_headers = {"X-Auth-Token": admin_login.json()["token"]}
        self.user_headers = {"X-Auth-Token": user_login.json()["token"]}

    def tearDown(self):
        self.client.close()
        self.engine.dispose()

    def test_admin_can_publish_and_user_can_mark_read(self):
        create_resp = self.client.post(
            "/messages/admin",
            headers=self.admin_headers,
            json={
                "title": "System maintenance",
                "content": "We will update the server tonight.",
                "level": "warning",
                "link": "/settings/server",
                "is_published": True,
            },
        )
        self.assertEqual(create_resp.status_code, 200)
        message_id = create_resp.json()["id"]

        list_resp = self.client.get("/messages", headers=self.user_headers)
        self.assertEqual(list_resp.status_code, 200)
        self.assertEqual(list_resp.json()["unread_count"], 1)
        self.assertEqual(list_resp.json()["items"][0]["title"], "System maintenance")

        mark_resp = self.client.post(f"/messages/{message_id}/read", headers=self.user_headers)
        self.assertEqual(mark_resp.status_code, 200)

        count_resp = self.client.get("/messages/unread-count", headers=self.user_headers)
        self.assertEqual(count_resp.status_code, 200)
        self.assertEqual(count_resp.json()["unread_count"], 0)


if __name__ == "__main__":
    unittest.main()
