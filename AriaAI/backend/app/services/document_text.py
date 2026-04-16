from __future__ import annotations

from pathlib import Path

try:
    import pdfplumber as _pdfplumber

    _HAS_PDF = True
except ImportError:
    _HAS_PDF = False

try:
    from docx import Document as _DocxDocument

    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

try:
    from pptx import Presentation as _Presentation

    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

try:
    import openpyxl as _openpyxl

    _HAS_XLSX = True
except ImportError:
    _HAS_XLSX = False


def extract_text_from_file(
    path: Path,
    file_type: str,
    *,
    max_chars: int = 4000,
    empty_placeholder: str = "",
    unsupported_placeholder: str = "",
    error_prefix: str = "",
) -> str:
    if not path.exists():
        return "[File not found]"

    try:
        ft = file_type.lower()
        if ft == "pdf" and _HAS_PDF:
            with _pdfplumber.open(path) as pdf:
                pages = [page.extract_text() or "" for page in pdf.pages[:15]]
            text = "\n".join(pages)
        elif ft == "docx" and _HAS_DOCX:
            doc = _DocxDocument(str(path))
            text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
        elif ft == "pptx" and _HAS_PPTX:
            parts = []
            presentation = _Presentation(str(path))
            for index, slide in enumerate(presentation.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    parts.append(f"[Slide {index + 1}]\n" + "\n".join(slide_texts))
            text = "\n\n".join(parts)
        elif ft in ("xlsx", "xls") and _HAS_XLSX:
            workbook = _openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            parts = []
            for sheet in workbook.worksheets:
                rows = []
                for row in sheet.iter_rows(max_row=200, values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in cells):
                        rows.append("\t".join(cells))
                if rows:
                    parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
            workbook.close()
            text = "\n\n".join(parts)
        elif ft in ("txt", "md", "csv", "json"):
            text = path.read_text(encoding="utf-8", errors="replace")
        else:
            return unsupported_placeholder

        text = text.strip()
        if not text:
            return empty_placeholder
        if len(text) > max_chars:
            return text[:max_chars] + "\n…[truncated]"
        return text
    except Exception as exc:
        return f"{error_prefix}{exc}" if error_prefix else ""
