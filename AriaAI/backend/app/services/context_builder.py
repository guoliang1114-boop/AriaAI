"""Context builder — assemble AI context from projects, files, RAG, and skills."""
from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Optional

from sqlmodel import Session, select

from app.config import UPLOADS_DIR
from app.models.db import (
    ClientRecord,
    Milestone,
    Project,
    ProjectFile,
    ProjectPayment,
    Skill,
)
from app.services.rag import retrieve_structured
from app.services.tool_executor import format_tools_for_claude

try:
    import pdfplumber as _pdfplumber
    _HAS_PDF = True
except ImportError:
    _HAS_PDF = False

try:
    from docx import Document as _DocxDocument
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

try:
    from pptx import Presentation as _Presentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

try:
    import openpyxl as _openpyxl
    _HAS_XLSX = True
except ImportError:
    _HAS_XLSX = False

MAX_FILE_CONTENT_CHARS = 40000  # cap total injected content to ~10k tokens
MAX_SINGLE_FILE_CHARS = 8000


def extract_file_text(path: Path, file_type: str, max_chars: int = 4000) -> str:
    """Extract plain text from a project file for AI context injection."""
    if not path.exists():
        return "[File not found]"
    try:
        ft = file_type.lower()
        if ft == "pdf" and _HAS_PDF:
            with _pdfplumber.open(path) as pdf:
                pages = [p.extract_text() or "" for p in pdf.pages[:15]]
            text = "\n".join(pages)
        elif ft == "docx" and _HAS_DOCX:
            doc = _DocxDocument(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ft == "pptx" and _HAS_PPTX:
            prs = _Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
            text = "\n\n".join(parts)
        elif ft in ("xlsx", "xls") and _HAS_XLSX:
            wb = _openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                rows = []
                for row in sheet.iter_rows(max_row=200, values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        rows.append("\t".join(cells))
                if rows:
                    parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
            wb.close()
            text = "\n\n".join(parts)
        elif ft in ("txt", "md", "csv", "json"):
            text = path.read_text(encoding="utf-8", errors="replace")
        else:
            return "[Binary file — text extraction not supported for this format]"
        text = text.strip()
        if len(text) > max_chars:
            return text[:max_chars] + "\n…[truncated]"
        return text if text else "[File appears to be empty]"
    except Exception as exc:
        return f"[Could not extract text: {exc}]"


class SkillContext:
    """Context for a skill including prompt, tools, and max_tokens."""
    def __init__(
        self,
        skill_prompt: str = "",
        tools: Optional[list] = None,
        max_tokens: Optional[int] = None,
    ):
        self.skill_prompt = skill_prompt
        self.tools = tools
        self.max_tokens = max_tokens


def build_skill_context(
    session: Session,
    skill_id: Optional[int],
    default_max_tokens: int = 4096,
) -> SkillContext:
    """Build context from a skill definition."""
    skill_prompt = ""
    tools = None
    max_tokens = default_max_tokens
    
    if skill_id:
        skill = session.get(Skill, skill_id)
        if skill:
            skill_prompt = skill.system_prompt
            max_tokens = skill.max_tokens or max_tokens
            # Load tools from skill
            if skill.tools_definition_json and skill.tools_definition_json.strip():
                try:
                    tools = format_tools_for_claude(
                        __import__("json").loads(skill.tools_definition_json)
                    )
                except Exception:
                    pass
    
    return SkillContext(
        skill_prompt=skill_prompt,
        tools=tools,
        max_tokens=max_tokens,
    )


def build_global_workspace_context(session: Session) -> str:
    """Build global workspace context with all non-archived projects."""
    all_projects = session.exec(
        select(Project).where(Project.status != "archived").order_by(Project.updated_at.desc())
    ).all()
    
    if not all_projects:
        return ""
    
    today_str = datetime.utcnow().strftime("%Y-%m-%d")
    ws_lines = [
        f"# 工作台全局数据（截至 {today_str}）",
        f"当前共有 {len(all_projects)} 个活跃项目。以下是每个项目的详细信息：\n",
    ]
    
    for p in all_projects:
        ws_lines.append(f"## 项目：{p.name}")
        ws_lines.append(f"- 客户：{p.client}")
        ws_lines.append(f"- 阶段：{p.status}")
        if p.contract_amount:
            ws_lines.append(f"- 合同金额：¥{p.contract_amount:,.0f}")
        if p.description:
            ws_lines.append(f"- 简介：{p.description[:200]}")
        if p.context_summary:
            ws_lines.append(f"- AI 摘要：{p.context_summary[:300]}")
        if p.notes:
            ws_lines.append(f"- 项目笔记：{p.notes[:200]}")
        
        # Milestones
        milestones = session.exec(
            select(Milestone).where(Milestone.project_id == p.id)
        ).all()
        if milestones:
            done = sum(1 for m in milestones if m.is_done)
            overdue = [
                m for m in milestones
                if not m.is_done and m.due_date and m.due_date < today_str
            ]
            ws_lines.append(
                f"- 里程碑：{done}/{len(milestones)} 已完成" +
                (f"，{len(overdue)} 个已逾期" if overdue else "")
            )
            for m in milestones:
                status_icon = "✓" if m.is_done else (
                    "⚠ 逾期" if m.due_date and m.due_date < today_str else "○"
                )
                due = f"（截止 {m.due_date}）" if m.due_date else ""
                priority = "【高优先级】" if m.priority == "high" else ""
                ws_lines.append(f"  {status_icon} {m.title}{priority}{due}")
        
        # Financials
        payments = session.exec(
            select(ProjectPayment).where(ProjectPayment.project_id == p.id)
        ).all()
        if payments:
            received = sum(pay.amount for pay in payments if pay.payment_type == "received")
            expense = sum(pay.amount for pay in payments if pay.payment_type == "expense")
            if received or expense:
                ws_lines.append(f"- 财务：已收款 ¥{received:,.0f}，支出 ¥{abs(expense):,.0f}")
        ws_lines.append("")
    
    return "\n".join(ws_lines)


def build_project_context(
    session: Session,
    project_id: int,
    file_ids: Optional[list[int]] = None,
) -> str:
    """Build context for a specific project including files, milestones, financials."""
    project = session.get(Project, project_id)
    if not project:
        return ""
    
    lines = [
        "## Scope Guard",
        "Use only the current project's context as the primary source of truth.",
        "Do not assume facts from other projects under the same client unless the user explicitly asks for cross-project comparison.",
        "If outside context is mentioned, label it as a hypothesis or reference rather than current-project fact.",
        "",
        f"**Project Name:** {project.name}",
        f"**Client:** {project.client}",
        f"**Status:** {project.status}",
    ]
    if project.description:
        lines.append(f"**Description:** {project.description}")
    if project.contract_amount:
        lines.append(f"**Contract Amount:** ¥{project.contract_amount:,.0f}")
    
    # AI-generated context summary
    if project.context_summary:
        lines.append(f"\n**Project Context Summary:**\n{project.context_summary}")
    
    # Accumulated project notes
    if project.notes:
        lines.append(f"\n**Project Notes:**\n{project.notes}")
    
    # Milestones
    milestones = session.exec(
        select(Milestone).where(Milestone.project_id == project.id)
    ).all()
    if milestones:
        lines.append("\n**Milestones:**")
        for m in milestones:
            status_icon = "✓" if m.is_done else "○"
            due = f" (due: {m.due_date})" if m.due_date else ""
            priority = f" [{m.priority} priority]" if m.priority == "high" else ""
            lines.append(f"  {status_icon} {m.title}{priority}{due}")
    
    # Files
    files = session.exec(
        select(ProjectFile).where(ProjectFile.project_id == project.id)
    ).all()
    file_content_sections = []
    if files:
        lines.append("\n**Uploaded Documents (full content auto-injected below):**")
        total_chars = 0
        for f in files:
            summary_hint = f" — {f.summary[:80]}" if f.summary else ""
            lines.append(f"  - {f.name} ({f.file_type.upper()}){summary_hint}")
            # Auto-inject readable file content
            if (
                f.file_type.lower() in ("pdf", "docx", "pptx", "xlsx", "xls", "txt", "md", "csv", "json")
                and total_chars < MAX_FILE_CONTENT_CHARS
            ):
                full_path = UPLOADS_DIR / f.path
                text = extract_file_text(
                    full_path,
                    f.file_type,
                    max_chars=min(MAX_SINGLE_FILE_CHARS, MAX_FILE_CONTENT_CHARS - total_chars)
                )
                if text and not text.startswith("["):
                    file_content_sections.append(f"### {f.name}\n{text}")
                    total_chars += len(text)
    
    # Financials
    payments = session.exec(
        select(ProjectPayment).where(ProjectPayment.project_id == project.id)
    ).all()
    if payments:
        received = sum(p.amount for p in payments if p.payment_type == "received")
        expense = sum(p.amount for p in payments if p.payment_type == "expense")
        if received or expense:
            lines.append(f"\n**Financials:** Received ¥{received:,.0f} | Expenses ¥{abs(expense):,.0f}")
    
    project_context = "\n".join(lines)
    
    # Append auto-injected file contents
    if file_content_sections:
        project_context += "\n\n## Project File Contents\n" + "\n\n---\n\n".join(file_content_sections)
    
    # Inject selected project file contents
    if file_ids:
        file_sections = []
        for fid in file_ids:
            pf = session.get(ProjectFile, fid)
            if pf:
                full_path = UPLOADS_DIR / pf.path
                text = extract_file_text(full_path, pf.file_type)
                file_sections.append(f"### {pf.name}\n{text}")
        if file_sections:
            attachment_block = "\n\n---\n\n".join(file_sections)
            project_context += "\n\n## Attached Files\n" + attachment_block
    
    return project_context


def build_rag_context(
    session: Session,
    query: str,
    rag_doc_ids: Optional[list[int]] = None,
    project_id: Optional[int] = None,
    knowledge_scope: str = "global",
    auto_trigger: bool = True,
) -> dict:
    """
    Build RAG context from knowledge documents.
    
    Returns structured dict with both text for LLM and sources for citations.
    """
    # Check if RAG should trigger
    should_retrieve = bool(rag_doc_ids) or (auto_trigger and "#doc" in query)
    
    if not should_retrieve:
        return {"text": "", "sources": []}
    
    # Perform structured retrieval
    client_id = None
    effective_project_id = None
    if not rag_doc_ids and knowledge_scope == "project" and project_id is not None:
        effective_project_id = project_id
    elif not rag_doc_ids and knowledge_scope == "client" and project_id is not None:
        project = session.get(Project, project_id)
        if project and project.client.strip():
            client = session.exec(
                select(ClientRecord).where(ClientRecord.name.ilike(project.client.strip()))
            ).first()
            client_id = client.id if client else None

    ctx = retrieve_structured(
        query,
        session,
        rag_doc_ids,
        project_id=effective_project_id,
        client_id=client_id,
    )
    
    return {
        "text": ctx.to_text(),
        "sources": [r.to_dict() for r in ctx.results],
        "query": ctx.query,
    }


class ChatContext:
    """Complete context for a chat session."""
    def __init__(
        self,
        skill_prompt: str = "",
        project_context: str = "",
        rag_context: str = "",
        rag_sources: Optional[list] = None,
        tools: Optional[list] = None,
        max_tokens: int = 4096,
    ):
        self.skill_prompt = skill_prompt
        self.project_context = project_context
        self.rag_context = rag_context
        self.rag_sources = rag_sources or []
        self.tools = tools
        self.max_tokens = max_tokens


def build_chat_context(
    session: Session,
    skill_id: Optional[int] = None,
    project_id: Optional[int] = None,
    knowledge_scope: str = "global",
    rag_doc_ids: Optional[list[int]] = None,
    file_ids: Optional[list[int]] = None,
    content: str = "",
    default_max_tokens: int = 4096,
) -> ChatContext:
    """Build complete chat context including skill, project, and RAG."""
    # Build skill context
    skill_ctx = build_skill_context(session, skill_id, default_max_tokens)
    
    # Build project context
    if project_id:
        project_context = build_project_context(session, project_id, file_ids)
    else:
        project_context = build_global_workspace_context(session)
    
    # Build RAG context
    rag_data = build_rag_context(
        session,
        content,
        rag_doc_ids,
        project_id=project_id,
        knowledge_scope=knowledge_scope,
        auto_trigger=True,
    )
    
    return ChatContext(
        skill_prompt=skill_ctx.skill_prompt,
        project_context=project_context,
        rag_context=rag_data["text"],
        rag_sources=rag_data["sources"],
        tools=skill_ctx.tools,
        max_tokens=skill_ctx.max_tokens or default_max_tokens,
    )
