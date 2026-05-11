"""File generation tools for Claude Function Calling.

Provides tools to generate PPT, Word, Excel, and PDF files.
"""
from __future__ import annotations

import base64
import json
import mimetypes
import os
import re
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import Any

from app.config import UPLOADS_DIR
from app.tools import registry

# Ensure output directory exists
GENERATED_DIR = UPLOADS_DIR / "generated"
GENERATED_DIR.mkdir(parents=True, exist_ok=True)

# Skills directory for templates
# file is at backend/app/tools/ → need 4 parents to reach project root
SKILLS_DIR = Path(__file__).parent.parent.parent.parent / "skills"

GRAPHIC_LIBRARY_H5_THEME = """
:root {
  --gl-bg: #f5f7fb;
  --gl-surface: #ffffff;
  --gl-ink: #17191d;
  --gl-body: #333333;
  --gl-muted: #64748b;
  --gl-brand: #003294;
  --gl-brand-2: #00338d;
  --gl-brand-soft: #e8eefb;
  --gl-line: #d9dee8;
  --gl-shadow: 0 18px 50px rgba(15, 23, 42, 0.08);
  --gl-radius: 24px;
  --gl-font-display: "KPMG Bold", "Segoe UI", "PingFang SC", "Microsoft YaHei", sans-serif;
  --gl-font-body: "Aptos", "Segoe UI", "PingFang SC", "Microsoft YaHei", sans-serif;
}
* { box-sizing: border-box; }
html, body {
  margin: 0;
  padding: 0;
  background:
    radial-gradient(circle at top right, rgba(0, 50, 148, 0.1), transparent 28%),
    linear-gradient(180deg, #eef3fb 0%, #f7f9fc 100%);
  color: var(--gl-ink);
  font-family: var(--gl-font-body);
}
body { min-height: 100vh; }
.deck-shell { width: min(100vw, 1480px); margin: 0 auto; padding: 28px 24px 60px; }
.deck-header { display:flex; justify-content:space-between; align-items:center; margin-bottom:18px; color:var(--gl-muted); font-size:14px; }
.deck-title { font-family:var(--gl-font-display); color:var(--gl-brand); font-size:15px; letter-spacing:.08em; text-transform:uppercase; }
.deck-nav { display:flex; gap:8px; flex-wrap:wrap; }
.deck-nav a { color:var(--gl-muted); text-decoration:none; padding:6px 10px; border-radius:999px; background:rgba(255,255,255,.7); border:1px solid rgba(0,50,148,.08); }
.deck { display:grid; gap:24px; }
.slide { position:relative; width:100%; aspect-ratio:16/9; min-height:540px; background:var(--gl-surface); border-radius:28px; box-shadow:var(--gl-shadow); overflow:hidden; }
.slide-inner { position:absolute; inset:0; padding:56px 72px; }
.eyebrow { display:inline-flex; align-items:center; gap:10px; color:var(--gl-brand); font-size:14px; font-weight:700; letter-spacing:.06em; text-transform:uppercase; }
.eyebrow::before { content:""; width:44px; height:4px; background:var(--gl-brand); border-radius:999px; }
.slide-title { margin:18px 0 0; max-width:980px; font-family:var(--gl-font-display); font-size:clamp(30px,2.4vw,42px); line-height:1.08; letter-spacing:-.02em; color:var(--gl-ink); }
.slide-subtitle { margin-top:14px; max-width:820px; font-size:18px; line-height:1.6; color:var(--gl-muted); }
.slide-no { position:absolute; right:42px; bottom:26px; font-size:13px; color:rgba(23,25,29,.36); }
.cover { background:linear-gradient(135deg, rgba(0, 50, 148, 0.94), rgba(0, 20, 56, 0.96)), #0b1c47; color:#fff; }
.cover::before { content:""; position:absolute; inset:auto -10% -18% 48%; height:72%; background:linear-gradient(135deg, rgba(255,255,255,.08), rgba(255,255,255,.02)); transform:skewX(-26deg); }
.cover::after { content:""; position:absolute; inset:0 0 auto auto; width:240px; height:240px; background:radial-gradient(circle, rgba(172,234,255,.28) 0%, rgba(172,234,255,0) 72%); }
.cover .slide-inner { display:grid; grid-template-columns:1.1fr .9fr; align-items:end; }
.cover .cover-kicker { font-size:15px; text-transform:uppercase; letter-spacing:.1em; color:rgba(255,255,255,.72); }
.cover .slide-title { color:#fff; max-width:620px; font-size:clamp(42px,3.6vw,62px); }
.cover .slide-subtitle { color:rgba(255,255,255,.78); max-width:560px; }
.cover-visual { position:relative; justify-self:end; width:100%; max-width:430px; height:100%; min-height:460px; }
.cover-visual .frame { position:absolute; inset:14% 0 6% 14%; border:1px solid rgba(255,255,255,.18); border-radius:26px; background:linear-gradient(180deg, rgba(255,255,255,.08), rgba(255,255,255,.02)); }
.cover-visual .frame::before, .cover-visual .frame::after { content:""; position:absolute; border-radius:18px; background:rgba(255,255,255,.12); }
.cover-visual .frame::before { width:68%; height:18%; left:10%; top:14%; }
.cover-visual .frame::after { width:52%; height:10%; right:12%; bottom:16%; }
.section-slide { background:linear-gradient(135deg, rgba(0, 50, 148, 0.96), rgba(0, 28, 82, 0.98)), #0d2458; }
.section-slide .slide-inner { display:grid; align-content:center; }
.section-no { margin-bottom:18px; color:rgba(255,255,255,.75); font-family:var(--gl-font-display); font-size:28px; letter-spacing:.12em; }
.section-slide .slide-title { color:#fff; max-width:760px; font-size:clamp(38px,3vw,56px); }
.content-layout, .two-col-layout, .roadmap-layout, .matrix-layout, .cards-layout { display:grid; gap:24px; margin-top:42px; }
.content-layout { grid-template-columns:.95fr .85fr; }
.lead-note { padding:18px 22px; border-radius:20px; background:var(--gl-brand-soft); color:var(--gl-brand); font-weight:700; line-height:1.55; }
.bullet-list { display:grid; gap:14px; margin:0; padding:0; list-style:none; }
.bullet-list li { position:relative; padding-left:18px; font-size:18px; line-height:1.65; color:var(--gl-body); }
.bullet-list li::before { content:""; position:absolute; left:0; top:12px; width:8px; height:8px; border-radius:50%; background:var(--gl-brand); }
.visual-card { min-height:420px; border-radius:24px; border:1px solid var(--gl-line); background:linear-gradient(180deg,#fff 0%,#f9fbfd 100%); position:relative; overflow:hidden; }
.visual-card .visual-ribbon { position:absolute; inset:24px 24px auto auto; padding:8px 12px; border-radius:999px; font-size:12px; color:var(--gl-brand); background:rgba(0,50,148,.08); }
.visual-card .visual-stack { position:absolute; inset:90px 28px 28px; display:grid; gap:14px; }
.visual-card .visual-stack span { display:block; border-radius:18px; background:linear-gradient(90deg, rgba(0,50,148,.1), rgba(0,50,148,.02)); }
.visual-card .visual-stack span:nth-child(1){height:78px;width:76%}.visual-card .visual-stack span:nth-child(2){height:110px;width:100%}.visual-card .visual-stack span:nth-child(3){height:90px;width:84%;justify-self:end}
.two-col-layout { grid-template-columns:repeat(2,minmax(0,1fr)); }
.compare-card { border-radius:24px; border:1px solid var(--gl-line); background:#fff; padding:28px 28px 24px; }
.compare-card.current { box-shadow: inset 0 4px 0 #94a3b8; }
.compare-card.target { box-shadow: inset 0 4px 0 var(--gl-brand); }
.compare-card h3 { margin:0 0 18px; font-family:var(--gl-font-display); font-size:24px; color:var(--gl-ink); }
.roadmap-layout { grid-template-columns:repeat(3,minmax(0,1fr)); position:relative; margin-top:56px; }
.roadmap-layout::before { content:""; position:absolute; left:8%; right:8%; top:50px; height:2px; background:var(--gl-line); }
.roadmap-phase { position:relative; padding-top:78px; }
.roadmap-phase::before { content:""; position:absolute; top:40px; left:50%; width:18px; height:18px; margin-left:-9px; border-radius:50%; background:var(--gl-brand); box-shadow:0 0 0 10px rgba(0,50,148,.08); }
.roadmap-phase h3 { margin:0 0 16px; font-size:24px; color:var(--gl-brand); font-family:var(--gl-font-display); text-align:center; }
.roadmap-phase .phase-card { border-radius:22px; border:1px solid var(--gl-line); background:#fff; padding:22px 22px 18px; min-height:280px; }
.matrix-layout { grid-template-columns:1fr 320px; align-items:start; }
.matrix-board { position:relative; height:430px; border-radius:24px; border:1px solid var(--gl-line); background:linear-gradient(180deg,#fff 0%,#f8fafc 100%); }
.matrix-board::before, .matrix-board::after { content:""; position:absolute; background:var(--gl-line); }
.matrix-board::before { left:50%; top:20px; bottom:20px; width:1px; }
.matrix-board::after { top:50%; left:20px; right:20px; height:1px; }
.matrix-pill { position:absolute; padding:8px 12px; border-radius:999px; background:rgba(0,50,148,.1); color:var(--gl-brand); font-size:13px; font-weight:700; }
.matrix-note { border-radius:24px; border:1px solid var(--gl-line); background:#fff; padding:24px; }
.cards-layout { grid-template-columns:repeat(4,minmax(0,1fr)); }
.metric-card { border-radius:22px; border:1px solid var(--gl-line); background:#fff; padding:24px 22px 20px; min-height:170px; position:relative; }
.metric-card::before { content:""; position:absolute; left:0; top:0; width:100%; height:6px; border-radius:22px 22px 0 0; background:var(--accent, var(--gl-brand)); }
.metric-label { font-size:14px; color:var(--gl-muted); }
.metric-value { margin-top:18px; font-family:var(--gl-font-display); font-size:34px; color:var(--gl-ink); }
.metric-detail { margin-top:12px; font-size:15px; line-height:1.55; color:var(--gl-body); }
.steps { display:grid; gap:18px; margin-top:42px; }
.step { display:grid; grid-template-columns:54px 1fr; gap:18px; align-items:start; padding:18px 20px; border-radius:22px; border:1px solid var(--gl-line); background:#fff; }
.step-no { width:42px; height:42px; display:grid; place-items:center; border-radius:50%; background:var(--gl-brand); color:#fff; font-weight:700; }
.step-text { font-size:18px; line-height:1.6; color:var(--gl-body); }
.closing { background:linear-gradient(135deg, rgba(0, 50, 148, 0.97), rgba(8, 20, 44, 0.98)), #0f172a; }
.closing .slide-title,.closing .slide-subtitle,.closing .slide-no { color:#fff; }
@media (max-width: 1100px) {
  .cover .slide-inner, .content-layout, .two-col-layout, .roadmap-layout, .matrix-layout, .cards-layout { grid-template-columns:1fr; }
  .slide { min-height:760px; aspect-ratio:auto; }
  .cover-visual { max-width:none; min-height:280px; }
}
"""


def _generate_filename(extension: str) -> str:
    """Generate a unique filename with timestamp."""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"generated_{timestamp}.{extension}"


def _slugify_filename(value: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9._-]+", "-", str(value or "").strip()).strip("-").lower()
    return slug or "deck"


def _html_escape(value: Any) -> str:
    text = str(value or "")
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _html_bullets(content: str) -> str:
    bullets = _split_bullets(content, limit=8)
    if not bullets:
        return ""
    items = "".join(f"<li>{_html_escape(item)}</li>" for item in bullets)
    return f'<ul class="bullet-list">{items}</ul>'


def _normalize_graphic_library_slide_type(slide_type: Any) -> str:
    value = str(slide_type or "content").strip().lower().replace("-", "_").replace(" ", "_")
    aliases = {
        "title": "section",
        "section_divider": "section",
        "divider": "section",
        "executive_summary": "summary",
        "summary": "summary",
        "card_grid": "summary",
        "numbered_cards": "summary",
        "insight": "content",
        "insight_content": "content",
        "compare": "two_column",
        "current_target": "two_column",
        "three_part_framework": "framework",
        "three_column": "framework",
        "framework_3": "framework",
        "kpi_cards": "kpi",
        "risks": "risk",
        "risk_mitigation": "risk",
        "risk_mitigations": "risk",
        "risks_and_mitigations": "risk",
        "stakeholders": "stakeholder",
        "stakeholder_map": "stakeholder",
        "nextsteps": "next_steps",
        "next_step": "next_steps",
    }
    return aliases.get(value, value)


def _normalize_graphic_library_style_key(style_key: Any) -> str:
    value = str(style_key or "").strip().lower().replace("-", "_").replace(" ", "_")
    aliases = {
        "summary_cards": "executive_summary",
        "executive_summary_cards": "executive_summary",
        "summary_stack": "executive_summary",
        "text": "text_focus",
        "text_only": "text_focus",
        "insight_visual": "insight_visual",
        "visual": "insight_visual",
        "two_column": "compare",
        "three_part": "framework",
        "three_column": "framework",
        "card_grid": "quad_grid",
        "tile_grid": "quad_grid",
        "four_up": "quad_grid",
        "grid_4": "quad_grid",
        "numbered_grid": "quad_numbered",
        "grid_numbered": "quad_numbered",
        "three_tiles": "triple_tiles",
        "triple_tile": "triple_tiles",
        "tile_triptych": "triple_tiles",
        "detail_grid": "detail_grid",
        "grid_6": "six_tiles",
        "six_grid": "six_tiles",
        "risk_mitigation": "risk",
        "risks_and_mitigations": "risk",
        "visual_note": "visual_note",
        "visual_caption": "visual_note",
        "nextsteps": "next_steps",
    }
    return aliases.get(value, value)


def _graphic_library_style_to_template_slide(style_key: Any) -> int | None:
    normalized = _normalize_graphic_library_style_key(style_key)
    mapping = {
        "cover": 1,
        "agenda": 2,
        "section": 3,
        "executive_summary": 4,
        "text_focus": 5,
        "insight_visual": 6,
        "compare": 7,
        "framework": 8,
        "roadmap": 9,
        "matrix": 10,
        "kpi": 11,
        "risk": 12,
        "visual_note": 13,
        "stakeholder": 13,
        "next_steps": 14,
        "executive_summary_alt": 22,
        "quad_grid": 23,
        "quad_numbered": 24,
        "triple_tiles": 25,
        "detail_grid": 26,
        "six_tiles": 27,
        "closing": 15,
    }
    return mapping.get(normalized)


def _graphic_library_template_slide_to_style_key(template_slide: Any) -> str | None:
    try:
        slide_no = int(template_slide)
    except (TypeError, ValueError):
        return None
    mapping = {
        1: "cover",
        2: "agenda",
        3: "section",
        4: "executive_summary",
        5: "text_focus",
        6: "insight_visual",
        7: "compare",
        8: "framework",
        9: "roadmap",
        10: "matrix",
        11: "kpi",
        12: "risk",
        13: "visual_note",
        14: "next_steps",
        15: "closing",
        21: "executive_summary",
        22: "executive_summary_alt",
        23: "quad_grid",
        24: "quad_numbered",
        25: "triple_tiles",
        26: "detail_grid",
        27: "six_tiles",
    }
    return mapping.get(slide_no)


def _infer_graphic_library_style_key(slide: dict[str, Any], index: int) -> str | None:
    explicit_style = _normalize_graphic_library_style_key(slide.get("style_key"))
    if explicit_style:
        return explicit_style

    template_style = _graphic_library_template_slide_to_style_key(slide.get("template_slide"))
    if template_style:
        if template_style == "visual_note":
            normalized_type = _normalize_graphic_library_slide_type(slide.get("type"))
            if normalized_type == "stakeholder":
                return "stakeholder"
        return template_style

    normalized_type = _normalize_graphic_library_slide_type(slide.get("type"))
    bullets = _split_bullets(str(slide.get("content") or ""), limit=8)

    if normalized_type == "section":
        return "section"
    if normalized_type == "framework":
        return "framework"
    if normalized_type == "two_column":
        return "compare"
    if normalized_type == "roadmap":
        return "roadmap"
    if normalized_type == "matrix":
        return "matrix"
    if normalized_type == "kpi":
        return "kpi"
    if normalized_type == "risk":
        return "risk"
    if normalized_type == "stakeholder":
        return "stakeholder"
    if normalized_type in {"next_steps", "steps"}:
        return "next_steps"
    if normalized_type == "summary":
        return "executive_summary"
    if normalized_type == "content":
        if slide.get("visual_caption") or slide.get("caption"):
            return "insight_visual"
        labels = slide.get("labels") or slide.get("card_titles") or []
        if len(labels) >= 6:
            return "six_tiles"
        if len(labels) == 4:
            return "quad_grid"
        if len(labels) == 3 and len(bullets) >= 6:
            return "executive_summary_alt"
        if slide.get("card_titles") or slide.get("labels"):
            return "executive_summary"
        if index == 0 and len(bullets) >= 3:
            return "executive_summary"
        if 5 <= len(bullets) <= 6:
            return "detail_grid"
        if len(bullets) == 4:
            return "quad_grid"
        if len(bullets) <= 4:
            return "text_focus"
        return "insight_visual"
    return None


def _prepare_graphic_library_h5_slides(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    prepared: list[dict[str, Any]] = []
    seen_business_page = False
    for index, slide in enumerate(slides):
        current = deepcopy(slide)
        explicit_style = _normalize_graphic_library_style_key(current.get("style_key"))
        template_style = _graphic_library_template_slide_to_style_key(current.get("template_slide"))
        normalized_type = _normalize_graphic_library_slide_type(current.get("type"))
        bullets = _split_bullets(str(current.get("content") or ""), limit=8)
        if (
            not seen_business_page
            and not explicit_style
            and not template_style
            and normalized_type == "content"
            and len(bullets) >= 3
        ):
            inferred_style = "executive_summary"
        else:
            inferred_style = _infer_graphic_library_style_key(current, index)
        if inferred_style and not current.get("style_key"):
            current["style_key"] = inferred_style
        if current.get("template_slide") is None and inferred_style:
            template_slide = _graphic_library_style_to_template_slide(inferred_style)
            if template_slide is not None and (3 <= template_slide <= 14 or 21 <= template_slide <= 27):
                current["template_slide"] = template_slide
        prepared.append(current)
        resolved_style = _normalize_graphic_library_style_key(current.get("style_key"))
        if resolved_style and resolved_style not in {"cover", "agenda", "section", "closing"}:
            seen_business_page = True
    return prepared


def _data_uri_for_asset(path: Path) -> str:
    if not path.is_file():
        return ""
    mime_type, _ = mimetypes.guess_type(path.name)
    mime_type = mime_type or "application/octet-stream"
    encoded = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime_type};base64,{encoded}"


def _graphic_library_source_slide_for_template(
    template_slide_no: int | None, style_key: str | None = None
) -> int | None:
    normalized_style = _normalize_graphic_library_style_key(style_key)
    if normalized_style == "executive_summary":
        return 21
    if normalized_style == "executive_summary_alt":
        return 22
    if template_slide_no is None:
        return None
    if 3 <= template_slide_no <= 15 or 21 <= template_slide_no <= 27:
        return template_slide_no
    return None


def _build_graphic_library_h5_html(title: str, subtitle: str, slides: list[dict]) -> str:
    nav_items = []
    sections = []

    cover_id = "cover"
    nav_items.append(f'<a href="#{cover_id}">封面</a>')
    sections.append(
        f"""
        <section class="slide cover" id="{cover_id}">
          <div class="slide-inner">
            <div class="cover-meta">
              <div class="cover-kicker">Graphic Library Inspired</div>
              <h1 class="slide-title">{_html_escape(title)}</h1>
              <p class="slide-subtitle">{_html_escape(subtitle or "Consulting-style H5 deck generated from the company visual system.")}</p>
            </div>
            <div class="cover-visual"><div class="frame"></div></div>
            <div class="slide-no">01</div>
          </div>
        </section>
        """
    )

    for index, slide in enumerate(slides, start=2):
        slide_type = str(slide.get("type") or "content").strip().lower()
        slide_title = _html_escape(slide.get("title") or f"Slide {index}")
        slide_id = f"slide-{index:02d}"
        nav_items.append(f'<a href="#{slide_id}">{slide_title}</a>')

        if slide_type in {"title", "section"}:
            body = f"""
            <section class="slide section-slide" id="{slide_id}">
              <div class="slide-inner">
                <div class="section-no">{index:02d}</div>
                <h2 class="slide-title">{slide_title}</h2>
                <div class="slide-no">{index:02d}</div>
              </div>
            </section>
            """
        elif slide_type == "two_column":
            left_content = _html_bullets(str(slide.get("left_content") or ""))
            right_content = _html_bullets(str(slide.get("right_content") or ""))
            body = f"""
            <section class="slide" id="{slide_id}">
              <div class="slide-inner">
                <div class="eyebrow">Compare</div>
                <h2 class="slide-title">{slide_title}</h2>
                <div class="two-col-layout">
                  <article class="compare-card current">
                    <h3>Current / Foundation</h3>
                    {left_content}
                  </article>
                  <article class="compare-card target">
                    <h3>Target / Direction</h3>
                    {right_content}
                  </article>
                </div>
                <div class="slide-no">{index:02d}</div>
              </div>
            </section>
            """
        elif slide_type == "roadmap":
            columns = [
                ("Phase 1", _html_bullets(str(slide.get("left_content") or slide.get("phase_1") or ""))),
                ("Phase 2", _html_bullets(str(slide.get("content") or slide.get("phase_2") or ""))),
                ("Phase 3", _html_bullets(str(slide.get("right_content") or slide.get("phase_3") or ""))),
            ]
            cards = "".join(
                f'<div class="roadmap-phase"><h3>{label}</h3><div class="phase-card">{content}</div></div>'
                for label, content in columns
            )
            body = f"""
            <section class="slide" id="{slide_id}">
              <div class="slide-inner">
                <div class="eyebrow">Roadmap</div>
                <h2 class="slide-title">{slide_title}</h2>
                <div class="roadmap-layout">{cards}</div>
                <div class="slide-no">{index:02d}</div>
              </div>
            </section>
            """
        elif slide_type == "matrix":
            body_text = _html_bullets(str(slide.get("content") or ""))
            labels = slide.get("labels") or ["Priority A", "Priority B", "Priority C", "Priority D"]
            pills = "".join(
                f'<div class="matrix-pill" style="left:{left}%; top:{top}%;">{_html_escape(label)}</div>'
                for (left, top), label in zip(((16, 24), (28, 64), (60, 32), (72, 58)), labels)
            )
            body = f"""
            <section class="slide" id="{slide_id}">
              <div class="slide-inner">
                <div class="eyebrow">Matrix</div>
                <h2 class="slide-title">{slide_title}</h2>
                <div class="matrix-layout">
                  <div class="matrix-board">{pills}</div>
                  <aside class="matrix-note">
                    <div class="lead-note">Decision note</div>
                    {body_text}
                  </aside>
                </div>
                <div class="slide-no">{index:02d}</div>
              </div>
            </section>
            """
        elif slide_type == "kpi":
            metrics = slide.get("metrics") or [
                {"label": "Metric 1", "value": "XX%", "detail": "Key reading", "accent": "#003294"},
                {"label": "Metric 2", "value": "XX%", "detail": "Key reading", "accent": "#16A34A"},
                {"label": "Metric 3", "value": "XX%", "detail": "Key reading", "accent": "#D97706"},
                {"label": "Metric 4", "value": "XX%", "detail": "Key reading", "accent": "#DC2626"},
            ]
            cards = "".join(
                f"""
                <div class="metric-card" style="--accent:{_html_escape(item.get('accent') or '#003294')}">
                  <div class="metric-label">{_html_escape(item.get('label') or '')}</div>
                  <div class="metric-value">{_html_escape(item.get('value') or '')}</div>
                  <div class="metric-detail">{_html_escape(item.get('detail') or '')}</div>
                </div>
                """
                for item in metrics[:4]
            )
            body = f"""
            <section class="slide" id="{slide_id}">
              <div class="slide-inner">
                <div class="eyebrow">KPI Cards</div>
                <h2 class="slide-title">{slide_title}</h2>
                <div class="cards-layout">{cards}</div>
                <div class="slide-no">{index:02d}</div>
              </div>
            </section>
            """
        elif slide_type in {"next_steps", "steps"}:
            steps = _split_bullets(str(slide.get("content") or ""), limit=5)
            step_html = "".join(
                f'<div class="step"><div class="step-no">{idx}</div><div class="step-text">{_html_escape(item)}</div></div>'
                for idx, item in enumerate(steps or ["Action item"], start=1)
            )
            body = f"""
            <section class="slide" id="{slide_id}">
              <div class="slide-inner">
                <div class="eyebrow">Next Steps</div>
                <h2 class="slide-title">{slide_title}</h2>
                <div class="steps">{step_html}</div>
                <div class="slide-no">{index:02d}</div>
              </div>
            </section>
            """
        else:
            body_html = _html_bullets(str(slide.get("content") or ""))
            lead = _split_bullets(str(slide.get("content") or ""), limit=1)
            lead_note = (
                f'<div class="lead-note">{_html_escape(lead[0])}</div>'
                if lead
                else '<div class="lead-note">Key takeaway</div>'
            )
            body = f"""
            <section class="slide" id="{slide_id}">
              <div class="slide-inner">
                <div class="eyebrow">Consulting Page</div>
                <h2 class="slide-title">{slide_title}</h2>
                <div class="content-layout">
                  <div class="content-body">
                    {lead_note}
                    {body_html}
                  </div>
                  <div class="visual-card">
                    <div class="visual-ribbon">Visual area</div>
                    <div class="visual-stack"><span></span><span></span><span></span></div>
                  </div>
                </div>
                <div class="slide-no">{index:02d}</div>
              </div>
            </section>
            """

        sections.append(body)

    return f"""<!doctype html>
<html lang="zh-CN">
  <head>
    <meta charset="UTF-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1.0" />
    <title>{_html_escape(title)}</title>
    <style>{GRAPHIC_LIBRARY_H5_THEME}</style>
  </head>
  <body>
    <div class="deck-shell">
      <div class="deck-header">
        <div class="deck-title">Graphic Library H5</div>
        <div class="deck-nav">{''.join(nav_items)}</div>
      </div>
      <main class="deck" id="deck-track">
        {''.join(sections)}
      </main>
      <div class="deck-counter" id="deck-counter">1 / {len(sections)}</div>
      <div class="deck-controls">
        <button type="button" id="deck-prev">上一页</button>
        <button type="button" id="deck-next">下一页</button>
      </div>
    </div>
    <script>
      (() => {{
        const track = document.getElementById('deck-track');
        const slides = Array.from(track.querySelectorAll('.slide'));
        const prev = document.getElementById('deck-prev');
        const next = document.getElementById('deck-next');
        const counter = document.getElementById('deck-counter');
        let current = 0;

        const update = () => {{
          const width = window.innerWidth || 1;
          current = Math.round(track.scrollLeft / width);
          current = Math.max(0, Math.min(current, slides.length - 1));
          counter.textContent = `${{current + 1}} / ${{slides.length}}`;
          prev.disabled = current === 0;
          next.disabled = current === slides.length - 1;
        }};

        const goTo = (index) => {{
          const width = window.innerWidth || 1;
          const target = Math.max(0, Math.min(index, slides.length - 1));
          track.scrollTo({{ left: target * width, behavior: 'smooth' }});
        }};

        prev.addEventListener('click', () => goTo(current - 1));
        next.addEventListener('click', () => goTo(current + 1));
        track.addEventListener('scroll', () => requestAnimationFrame(update), {{ passive: true }});
        window.addEventListener('resize', update);
        document.addEventListener('keydown', (event) => {{
          if (event.key === 'ArrowRight' || event.key === 'PageDown' || event.key === ' ') {{
            event.preventDefault();
            goTo(current + 1);
          }}
          if (event.key === 'ArrowLeft' || event.key === 'PageUp') {{
            event.preventDefault();
            goTo(current - 1);
          }}
        }});

        document.querySelectorAll('.deck-nav a').forEach((link, index) => {{
          link.addEventListener('click', (event) => {{
            event.preventDefault();
            goTo(index);
          }});
        }});

        update();
      }})();
    </script>
  </body>
</html>"""






def _build_graphic_library_h5_html_v2(title: str, subtitle: str, slides: list[dict]) -> str:
    skill_dir = SKILLS_DIR / "graphic-library-h5"
    template_map_path = skill_dir / "references" / "automation-template-map.json"
    if not template_map_path.is_file():
        raise FileNotFoundError(
            f"graphic-library-h5 template map not found: {template_map_path}. "
            "Ensure the skill directory and references/automation-template-map.json exist."
        )
    template_map = json.loads(template_map_path.read_text(encoding="utf-8"))
    templates = {item["slide"]: item for item in template_map}
    source_media_dir = skill_dir / "assets" / "source-media"
    source_svg_dir = skill_dir / "assets" / "source-svg"
    cover_image = _data_uri_for_asset(source_media_dir / "image2.png")

    def get_entry(slide_no: int, name: str) -> dict[str, Any]:
        for entry in templates[slide_no]["placeholders"]:
            if entry["name"] == name:
                return entry
        raise KeyError(f"Missing placeholder {name} on slide {slide_no}")

    def box_style(entry: dict[str, Any], extra: str = "") -> str:
        return f"left:{entry['x']}%;top:{entry['y']}%;width:{entry['w']}%;height:{entry['h']}%;{extra}"

    def render_text(entry: dict[str, Any], value: str, cls: str, extra: str = "") -> str:
        safe = _html_escape(value)
        return f'<div class="{cls}" style="{box_style(entry, extra)}">{safe}</div>'

    def render_box(entry: dict[str, Any], cls: str, extra: str = "") -> str:
        return f'<div class="{cls}" style="{box_style(entry, extra)}"></div>'

    def render_bullets(items: list[str], cls: str, extra: str = "") -> str:
        if not items:
            items = ["-"]
        lis = "".join(f"<li>{_html_escape(item)}</li>" for item in items)
        return f'<ul class="{cls}" style="{extra}">{lis}</ul>'

    def compact_copy(value: Any, limit: int = 72) -> str:
        text = re.sub(r"\s+", " ", str(value or "")).strip()
        if len(text) <= limit:
            return text
        clipped = text[: limit - 1].rstrip(" ,;:.")
        return f"{clipped}…"

    def first_items(slide: dict[str, Any], count: int, fallback_prefix: str) -> list[str]:
        labels = slide.get("labels") or slide.get("card_titles") or []
        items = [str(value) for value in labels if str(value).strip()]
        if len(items) < count:
            items.extend(_split_bullets(str(slide.get("content") or ""), limit=count * 2))
        items = [item for item in items if item.strip()]
        while len(items) < count:
            items.append(f"{fallback_prefix} {len(items) + 1}")
        return items[:count]

    def render_source_background(template_slide_no: int | None, style_key: Any = None) -> str:
        source_slide_no = _graphic_library_source_slide_for_template(template_slide_no, str(style_key or ""))
        if source_slide_no is None:
            return ""
        source_asset = source_svg_dir / f"source-slide-{source_slide_no}.svg"
        source_uri = _data_uri_for_asset(source_asset)
        if not source_uri:
            return ""
        normalized_style = _normalize_graphic_library_style_key(style_key)
        wash = "rgba(255,255,255,.84)"
        texture = "rgba(255,255,255,.18)"
        blur = "0.2px"
        image_opacity = ".26"
        if normalized_style in {"framework", "roadmap", "matrix", "stakeholder"}:
            wash = "rgba(255,255,255,.72)"
            texture = "rgba(255,255,255,.14)"
            blur = "0px"
            image_opacity = ".34"
        elif normalized_style in {"compare", "risk", "next_steps"}:
            wash = "rgba(255,255,255,.78)"
            texture = "rgba(255,255,255,.16)"
            image_opacity = ".30"
        elif normalized_style in {"text_focus"}:
            wash = "rgba(255,255,255,.92)"
            texture = "rgba(255,255,255,.12)"
            blur = "0.5px"
            image_opacity = ".18"
        elif normalized_style in {"executive_summary", "executive_summary_alt", "quad_grid", "quad_numbered", "triple_tiles", "detail_grid", "six_tiles"}:
            wash = "rgba(255,255,255,.90)"
            texture = "rgba(255,255,255,.12)"
            image_opacity = ".18"

        mask_specs: list[tuple[float, float, float, float, str]] = []
        if template_slide_no in {4, 21, 22}:
            mask_specs.extend([
                (8.0, 9.8, 50.0, 13.8, "rgba(255,255,255,.96)"),
                (8.0, 54.6, 84.0, 15.2, "rgba(255,255,255,.92)"),
            ])
        elif template_slide_no in {5, 6}:
            mask_specs.extend([
                (8.0, 10.0, 48.0, 13.0, "rgba(255,255,255,.96)"),
                (8.0, 24.0, 38.0, 42.0, "rgba(255,255,255,.88)"),
            ])
        elif template_slide_no in {23, 24, 25, 26, 27}:
            mask_specs.extend([
                (8.0, 10.0, 52.0, 13.0, "rgba(255,255,255,.96)"),
                (8.0, 50.0, 84.0, 18.0, "rgba(255,255,255,.90)") if template_slide_no in {23, 24, 25} else (0, 0, 0, 0, ""),
            ])
            mask_specs = [spec for spec in mask_specs if spec[2] > 0]
        mask_html = "".join(
            f'<div class="ppt-source-mask" style="left:{x}%;top:{y}%;width:{w}%;height:{h}%;background:{color};"></div>'
            for x, y, w, h, color in mask_specs
        )
        return (
            '<div class="ppt-source-bg" aria-hidden="true">'
            f'<img class="ppt-source-svg" src="{source_uri}" alt="" style="opacity:{image_opacity};filter:blur({blur}) saturate(.88) contrast(.94);" />'
            f'<div class="ppt-source-wash" style="background:{wash};"></div>'
            f'<div class="ppt-source-texture" style="background:{texture};"></div>'
            f'{mask_html}'
            "</div>"
        )

    css = """
:root {
  --brand: #00338d;
  --brand-2: #005eb8;
  --ink: #17191d;
  --muted: #5b6472;
  --line: rgba(0, 51, 141, 0.12);
  --surface: rgba(255,255,255,0.96);
  --surface-soft: #f5f7fb;
  --canvas: #eef2f7;
  --display: "KPMG Bold", Arial, "Microsoft YaHei", "PingFang SC", sans-serif;
  --body: Arial, "Microsoft YaHei", "PingFang SC", "Helvetica Neue", sans-serif;
  --shadow: 0 6px 18px rgba(15, 23, 42, 0.05);
}
* { box-sizing: border-box; }
html, body { margin:0; width:100%; height:100%; overflow:hidden; background:var(--canvas); font-family:var(--body); }
body { color: var(--ink); }
.deck-shell { position:relative; width:100vw; height:100vh; overflow:hidden; }
.deck { display:flex; width:100vw; height:100vh; overflow-x:auto; overflow-y:hidden; scroll-snap-type:x mandatory; scroll-behavior:smooth; scrollbar-width:none; }
.deck::-webkit-scrollbar { display:none; }
.slide { position:relative; flex:0 0 100vw; width:100vw; height:100vh; overflow:hidden; scroll-snap-align:start; background:#ffffff; }
.slide-inner { position:absolute; inset:0; }
.ppt-source-bg { position:absolute; inset:0; overflow:hidden; pointer-events:none; }
.ppt-source-svg { position:absolute; inset:0; width:100%; height:100%; object-fit:cover; object-position:center; transform:scale(1.008); }
.ppt-source-wash { position:absolute; inset:0; }
.ppt-source-texture { position:absolute; inset:0; backdrop-filter:blur(1px); }
.ppt-source-mask { position:absolute; border-radius:14px; box-shadow:0 4px 10px rgba(15,23,42,.03); }
.slide-no { position:absolute; right:2.8vw; bottom:2.5vh; font-size:12px; letter-spacing:.06em; color:rgba(23,25,29,.4); font-family:var(--body); }
.ppt-title, .ppt-big-title, .ppt-section-title, .ppt-closing-title {
  position:absolute; white-space:pre-wrap; font-family:var(--display); letter-spacing:-0.03em; font-weight:700;
  -webkit-font-smoothing:antialiased; text-rendering:geometricPrecision;
}
.ppt-title { font-size:30px; line-height:1.06; color:var(--brand); }
.ppt-big-title { font-size:44px; line-height:1.0; color:var(--brand); }
.ppt-section-title { font-size:50px; line-height:1.0; color:#ffffff; }
.ppt-closing-title { font-size:46px; line-height:1.02; color:#ffffff; }
.ppt-subtitle, .ppt-body, .ppt-caption, .ppt-axis, .ppt-note, .ppt-closing-body {
  position:absolute; white-space:pre-wrap; color:#334155; font-family:var(--body); font-weight:400;
  -webkit-font-smoothing:antialiased; text-rendering:geometricPrecision;
}
.ppt-subtitle { font-size:14px; line-height:1.36; color:#4b5563; }
.ppt-body { font-size:14px; line-height:1.42; color:#2f3947; }
.ppt-caption { font-size:11px; line-height:1.3; color:#526072; }
.ppt-note { font-size:13px; line-height:1.38; color:#334155; }
.ppt-axis { font-size:12px; font-weight:700; color:#334155; }
.ppt-closing-body { font-size:16px; line-height:1.36; color:rgba(255,255,255,.82); }
.ppt-chip-number {
  position:absolute; display:grid; place-items:center; font-family:var(--display); font-size:38px; line-height:0.96; letter-spacing:-0.04em; font-weight:800; color:#ffffff;
}
.ppt-box, .ppt-card, .ppt-panel, .ppt-visual, .ppt-summary-card, .ppt-kpi-card, .ppt-step-badge, .ppt-matrix-q, .ppt-copy-panel {
  position:absolute; border-radius:10px; background:var(--surface); border:1px solid var(--line); box-shadow:var(--shadow);
}
.ppt-copy-panel { border-radius:14px; background:rgba(255,255,255,.92); backdrop-filter:blur(8px); }
.ppt-summary-card { border-radius:10px; background:#ffffff; }
.summary-chip {
  position:absolute; display:grid; place-items:center; border-radius:999px;
  background:rgba(0, 51, 141, 0.08); color:var(--brand); font-family:var(--display);
  font-size:13px; font-weight:700; line-height:1;
}
.ppt-panel { background:#ffffff; }
.ppt-visual { background:linear-gradient(180deg,rgba(0,51,141,.05),rgba(0,94,184,.02)); }
.ppt-cover-plate { position:absolute; background:rgba(255,255,255,.96); }
.ppt-cover-image { position:absolute; object-fit:cover; }
.ppt-cover-marker { position:absolute; background:var(--brand); clip-path: polygon(0 0, 100% 0, 78% 100%, 0 100%); }
.cover-slide { background:#ffffff; }
.contents-slide { background:#ffffff; }
.section-slide { background:#00338d; }
.closing-slide { background:#00338d; }
.contents-bar { position:absolute; background:var(--brand); border-radius:999px; }
.contents-frame { position:absolute; border:1px solid rgba(0,51,141,.12); border-radius:12px; background:#ffffff; box-shadow:none; }
.agenda-list { position:absolute; margin:0; padding:0; list-style:none; display:grid; gap:14px; }
.agenda-list li { display:grid; grid-template-columns:48px 1fr; gap:14px; padding:0; border-top:1px solid rgba(0,51,141,.12); padding-top:10px; color:#2f3947; font-size:14px; line-height:1.34; }
.agenda-list li:last-child { border-bottom:1px solid rgba(0,51,141,.12); padding-bottom:10px; }
.agenda-index { font-family:var(--display); color:var(--brand); font-size:12px; letter-spacing:.03em; }
.summary-card-text { position:absolute; padding:14px 14px 0; font-family:var(--display); font-size:16px; line-height:1.08; color:var(--brand); }
.summary-supporting-note { position:absolute; padding:16px 18px; border-radius:12px; background:rgba(255,255,255,.88); border:1px solid rgba(0,51,141,.10); box-shadow:var(--shadow); color:#334155; font-size:13px; line-height:1.45; }
.ppt-grid-card { position:absolute; border-radius:14px; background:rgba(255,255,255,.90); border:1px solid rgba(0,51,141,.10); box-shadow:var(--shadow); backdrop-filter:blur(8px); }
.ppt-grid-card-title { position:absolute; font-family:var(--display); font-size:15px; line-height:1.08; color:var(--brand); }
.ppt-grid-card-body { position:absolute; font-size:12px; line-height:1.36; color:#334155; }
.ppt-grid-badge { position:absolute; display:grid; place-items:center; border-radius:999px; background:rgba(0,51,141,.10); color:var(--brand); font-family:var(--display); font-size:14px; line-height:1; }
.body-list, .two-col-list, .timeline-list, .matrix-list, .stakeholder-list, .step-list, .kpi-body-list, .framework-list, .risk-list { margin:0; padding:0; list-style:none; display:grid; gap:12px; }
.body-list li, .two-col-list li, .timeline-list li, .matrix-list li, .stakeholder-list li, .step-list li, .kpi-body-list li, .framework-list li, .risk-list li {
  position:relative; padding-left:16px; font-size:14px; line-height:1.44; color:#2f3947;
}
.body-list li::before, .two-col-list li::before, .timeline-list li::before, .matrix-list li::before, .stakeholder-list li::before, .step-list li::before, .kpi-body-list li::before, .framework-list li::before, .risk-list li::before {
  content:""; position:absolute; left:0; top:8px; width:6px; height:6px; border-radius:999px; background:var(--brand);
}
.kpi-accent { position:absolute; border-radius:10px 10px 0 0; }
.framework-panel { position:absolute; border-radius:14px; background:var(--surface); border:1px solid var(--line); box-shadow:var(--shadow); }
.framework-panel-title { position:absolute; font-family:var(--display); font-size:16px; line-height:1.1; color:var(--brand); }
.framework-topbar { position:absolute; border-radius:10px 10px 0 0; background:var(--brand); }
.framework-badge { position:absolute; display:grid; place-items:center; border-radius:999px; background:rgba(0,51,141,.08); color:var(--brand); font-family:var(--display); font-size:14px; }
.risk-panel { position:absolute; border-radius:14px; background:var(--surface); border:1px solid var(--line); box-shadow:var(--shadow); }
.risk-head { position:absolute; font-family:var(--display); font-size:16px; color:var(--brand); }
.risk-ribbon { position:absolute; border-radius:10px 10px 0 0; background:var(--brand); }
.timeline-line { position:absolute; background:rgba(0,51,141,.18); }
.timeline-dot { position:absolute; border-radius:999px; background:var(--brand); box-shadow:0 0 0 10px rgba(0,51,141,.08); }
.timeline-label { position:absolute; font-family:var(--display); font-size:17px; color:var(--brand); text-align:center; }
.matrix-line-v, .matrix-line-h { position:absolute; background:rgba(0,51,141,.16); }
.matrix-pill {
  position:absolute; display:grid; place-items:center; text-align:center; padding:9px 10px; border-radius:14px; background:var(--surface); border:1px solid var(--line); box-shadow:var(--shadow); color:var(--brand); font-weight:700; font-size:12px;
}
.matrix-body { position:absolute; white-space:pre-wrap; color:#334155; font-size:13px; line-height:1.36; }
.step-row { position:absolute; border-radius:18px; background:var(--surface); border:1px solid var(--line); box-shadow:var(--shadow); }
.step-row-text { position:absolute; font-size:15px; line-height:1.4; color:#2f3947; }
.stakeholder-grid { position:absolute; inset:0; }
.stakeholder-node { position:absolute; display:grid; place-items:center; border-radius:999px; background:linear-gradient(180deg, rgba(0,51,141,.96), rgba(0,94,184,.72)); color:#ffffff; font-family:var(--display); font-size:12px; line-height:1.15; text-align:center; padding:8px; box-shadow:var(--shadow); }
.stakeholder-ring { position:absolute; border-radius:50%; border:1px dashed rgba(0,51,141,.18); }
.ppt-floating-label { position:absolute; padding:6px 12px; border-radius:999px; background:rgba(255,255,255,.86); color:var(--brand); font-family:var(--display); font-size:12px; line-height:1; border:1px solid rgba(0,51,141,.10); box-shadow:var(--shadow); }
.deck-pager { position:absolute; left:50%; bottom:18px; transform:translateX(-50%); z-index:30; display:inline-flex; align-items:center; gap:14px; padding:8px 14px; border-radius:999px; background:rgba(255,255,255,.92); color:#334155; border:1px solid rgba(0,51,141,.12); box-shadow:0 4px 12px rgba(15,23,42,.05); }
.deck-pager button { border:0; background:transparent; color:inherit; font-size:18px; cursor:pointer; padding:0 4px; }
.deck-pager button:disabled { opacity:.35; cursor:not-allowed; }
.deck-status { min-width:72px; text-align:center; font-size:12px; letter-spacing:.06em; }
@media (max-width: 960px) {
  .ppt-big-title, .ppt-section-title, .ppt-closing-title { width:78% !important; }
}
"""

    sections: list[str] = []

    cover_title = get_entry(1, "aria_cover_title")
    cover_subtitle = get_entry(1, "aria_cover_subtitle")
    sections.append(
        f"""
        <section class="slide cover-slide" id="cover">
          <div class="slide-inner">
            <div class="ppt-cover-marker" style="left:6%;top:7.62%;width:7.48%;height:5.41%;"></div>
            <div class="ppt-cover-plate" style="left:49.06%;top:21.24%;width:44.94%;height:55.92%;"></div>
            <img class="ppt-cover-image" src="{cover_image}" alt="cover visual" style="left:37.38%;top:23.71%;width:56.71%;height:71.75%;" />
            {render_text(cover_title, title, 'ppt-big-title')}
            {render_text(cover_subtitle, subtitle or 'Subtitle | Client | Date', 'ppt-subtitle')}
            <div class="slide-no">01</div>
          </div>
        </section>
        """
    )

    agenda_title = get_entry(2, "aria_slide_title")
    agenda_bar = get_entry(2, "aria_accent_bar")
    agenda_items_entry = get_entry(2, "aria_agenda_items")
    agenda_frame = get_entry(2, "aria_visual_frame")
    agenda_caption = get_entry(2, "aria_visual_caption")
    agenda_items = []
    for idx, slide in enumerate(slides, start=3):
        agenda_items.append(f'<li><span class="agenda-index">{idx:02d}</span><span>{_html_escape(slide.get("title") or f"Slide {idx}")}</span></li>')
    sections.append(
        f"""
        <section class="slide contents-slide" id="contents">
          <div class="slide-inner">
            {render_box(agenda_bar, 'contents-bar')}
            {render_text(agenda_title, 'Agenda', 'ppt-title')}
            <ol class="agenda-list" style="{box_style(agenda_items_entry)}">{''.join(agenda_items)}</ol>
            {render_box(agenda_frame, 'contents-frame')}
            {render_text(agenda_caption, 'Overview visual area', 'ppt-caption', 'display:grid;place-items:center;')}
            <div class="slide-no">02</div>
          </div>
        </section>
        """
    )

    def section_slide(page_no: int, heading: str) -> str:
        return f"""
        <section class="slide section-slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_text(get_entry(3, 'aria_section_number'), f'{page_no:02d}', 'ppt-chip-number')}
            {render_text(get_entry(3, 'aria_section_title'), heading, 'ppt-section-title')}
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def summary_slide(page_no: int, slide: dict[str, Any]) -> str:
        bullets = _split_bullets(str(slide.get('content') or ''), limit=6)
        labels = slide.get("labels") or slide.get("card_titles") or []
        cards = labels[:3] or bullets[:3] or ['Insight 1', 'Insight 2', 'Insight 3']
        details = bullets[3:] or bullets[:3] or ['Key evidence', 'Strategic implication', 'Immediate action']
        supporting_text = slide.get("supporting_text") or slide.get("caption") or "Use this page to land the three most important messages before moving into detail."
        card_html = []
        for i, name in enumerate(['aria_summary_card_1', 'aria_summary_card_2', 'aria_summary_card_3'], start=1):
            entry = get_entry(4, name)
            card_html.append(render_box(entry, 'ppt-summary-card'))
            card_html.append(
                f'<div class="summary-chip" style="left:{entry["x"] + 1.5}%;top:{entry["y"] + 2.2}%;width:3.8%;height:5.8%;">{i:02d}</div>'
            )
        text_html = []
        for i, name in enumerate(['aria_kpi_1', 'aria_kpi_2', 'aria_kpi_3'], start=0):
            value = cards[i] if i < len(cards) else f'Insight {i+1}'
            text_html.append(render_text(get_entry(4, name), value, 'summary-card-text'))
        body_entry = get_entry(4, 'aria_slide_body')
        visual_entry = get_entry(4, 'aria_visual_area')
        caption_entry = get_entry(4, 'aria_visual_caption')
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(4, slide.get('style_key') or 'executive_summary')}
            {render_text(get_entry(4, 'aria_slide_title'), slide.get('title') or f'Slide {page_no}', 'ppt-title')}
            {''.join(card_html)}
            {''.join(text_html)}
            <div class="summary-supporting-note" style="{box_style(visual_entry)}">{_html_escape(str(supporting_text))}</div>
            <div class="ppt-copy-panel" style="{box_style(body_entry)}">
              <div class="ppt-body" style="position:absolute;inset:18px 20px 18px 20px;">{render_bullets(details, 'body-list')}</div>
            </div>
            {render_text(caption_entry, 'Executive summary | original slide grammar preserved', 'ppt-caption', 'display:grid;place-items:center;')}
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def insight_slide(page_no: int, slide: dict[str, Any]) -> str:
        body_entry = get_entry(6, 'aria_slide_body')
        visual_entry = get_entry(6, 'aria_visual_area')
        caption_entry = get_entry(6, 'aria_visual_caption')
        footer_entry = get_entry(6, 'aria_footer_note')
        bullets = _split_bullets(str(slide.get('content') or ''), limit=6)
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(6, slide.get('style_key') or 'insight_visual')}
            {render_text(get_entry(6, 'aria_slide_title'), slide.get('title') or f'Slide {page_no}', 'ppt-title')}
            <div class="ppt-copy-panel" style="{box_style(body_entry)}">
              <div class="ppt-body" style="position:absolute;inset:20px 22px;">{render_bullets(bullets, 'body-list')}</div>
            </div>
            {render_box(visual_entry, 'ppt-visual', 'background:rgba(255,255,255,.22);border-color:rgba(0,51,141,.08);')}
            {render_text(caption_entry, slide.get('visual_caption') or 'Visual area', 'ppt-caption', 'display:grid;place-items:center;background:rgba(255,255,255,.72);border-radius:999px;')}
            {render_text(footer_entry, f'{page_no:02d}', 'ppt-caption', 'text-align:right;color:#64748b;')}
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def compare_slide(page_no: int, slide: dict[str, Any]) -> str:
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(7, slide.get('style_key') or 'compare')}
            {render_text(get_entry(7, 'aria_slide_title'), slide.get('title') or f'Slide {page_no}', 'ppt-title')}
            {render_box(get_entry(7, 'aria_left_panel'), 'ppt-copy-panel')}
            {render_box(get_entry(7, 'aria_right_panel'), 'ppt-copy-panel')}
            <div class="ppt-floating-label" style="left:11%;top:24.5%;">Option A</div>
            <div class="ppt-floating-label" style="left:54.3%;top:24.5%;">Option B</div>
            <div class="ppt-body" style="{box_style(get_entry(7, 'aria_left_body'), 'padding:14px 18px 0;')}">{render_bullets(_split_bullets(str(slide.get('left_content') or ''), limit=6), 'two-col-list')}</div>
            <div class="ppt-body" style="{box_style(get_entry(7, 'aria_right_body'), 'padding:14px 18px 0;')}">{render_bullets(_split_bullets(str(slide.get('right_content') or ''), limit=6), 'two-col-list')}</div>
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def framework_slide(page_no: int, slide: dict[str, Any]) -> str:
        columns = [
            ("left", str(slide.get("left_content") or "")),
            ("center", str(slide.get("content") or "")),
            ("right", str(slide.get("right_content") or "")),
        ]
        labels = ["Dimension 1", "Dimension 2", "Dimension 3"]
        parts = []
        for index, (key, content) in enumerate(columns):
            panel = get_entry(8, f"aria_{key}_panel")
            body = get_entry(8, f"aria_{key}_body")
            lines = _split_bullets(content, limit=4)
            heading = lines[0] if lines else labels[index]
            bullets = lines[1:] or ["Key point", "Key point"]
            parts.append(render_box(panel, "framework-panel"))
            parts.append(render_box(panel, "framework-topbar", "height:6%;"))
            parts.append(f'<div class="framework-badge" style="left:{panel["x"] + 1.5}%;top:{panel["y"] + 3.2}%;width:3.2%;height:5.8%;">{index + 1}</div>')
            parts.append(render_text(body, heading, "framework-panel-title"))
            parts.append(f'<div class="ppt-body" style="{box_style(body, "padding-top:34px;")}">{render_bullets(bullets, "framework-list")}</div>')
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(8, slide.get('style_key') or 'framework')}
            {render_text(get_entry(8, 'aria_slide_title'), slide.get('title') or 'Three-part framework', 'ppt-title')}
            {''.join(parts)}
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def text_focus_slide(page_no: int, slide: dict[str, Any]) -> str:
        body_entry = get_entry(5, 'aria_slide_body')
        footer_entry = get_entry(5, 'aria_footer_note')
        bullets = _split_bullets(str(slide.get('content') or ''), limit=7)
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(5, slide.get('style_key') or 'text_focus')}
            {render_text(get_entry(5, 'aria_slide_title'), slide.get('title') or f'Slide {page_no}', 'ppt-title')}
            <div class="ppt-copy-panel" style="{box_style(body_entry)}">
              <div class="ppt-body" style="position:absolute;inset:20px 22px;">{render_bullets(bullets, 'body-list')}</div>
            </div>
            {render_text(footer_entry, f'{page_no:02d}', 'ppt-caption', 'text-align:right;color:#64748b;')}
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def roadmap_slide(page_no: int, slide: dict[str, Any]) -> str:
        blocks = [
            ('left', str(slide.get('left_content') or slide.get('phase_1') or '')),
            ('center', str(slide.get('content') or slide.get('phase_2') or '')),
            ('right', str(slide.get('right_content') or slide.get('phase_3') or '')),
        ]
        body_parts = []
        for key, content in blocks:
            body_parts.append(render_box(get_entry(9, f'aria_{key}_dot'), 'timeline-dot'))
            body_parts.append(render_text(get_entry(9, f'aria_{key}_label'), f'Phase {1 if key=="left" else 2 if key=="center" else 3}', 'timeline-label'))
            body_parts.append(f'<div class="ppt-body" style="{box_style(get_entry(9, f"aria_{key}_body"))}">{render_bullets(_split_bullets(content, limit=4), "timeline-list")}</div>')
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(9, slide.get('style_key') or 'roadmap')}
            {render_text(get_entry(9, 'aria_slide_title'), slide.get('title') or 'Roadmap', 'ppt-title')}
            {render_box(get_entry(9, 'aria_timeline'), 'timeline-line')}
            {''.join(body_parts)}
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def matrix_slide(page_no: int, slide: dict[str, Any]) -> str:
        labels = slide.get('labels') or ['Priority A', 'Priority B', 'Priority C', 'Priority D']
        q_names = ['aria_matrix_q1', 'aria_matrix_q2', 'aria_matrix_q3', 'aria_matrix_q4']
        pills = ''.join(
            f'<div class="matrix-pill" style="{box_style(get_entry(10, q_names[i]))}">{_html_escape(labels[i] if i < len(labels) else f"Priority {i+1}")}</div>'
            for i in range(4)
        )
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(10, slide.get('style_key') or 'matrix')}
            {render_text(get_entry(10, 'aria_slide_title'), slide.get('title') or 'Prioritization matrix', 'ppt-title')}
            {render_box(get_entry(10, 'aria_matrix_v'), 'matrix-line-v')}
            {render_box(get_entry(10, 'aria_matrix_h'), 'matrix-line-h')}
            {pills}
            <div class="ppt-copy-panel" style="{box_style(get_entry(10, 'aria_slide_body'))}">
              <div class="matrix-body" style="position:absolute;inset:16px 18px;">{render_bullets(_split_bullets(str(slide.get('content') or ''), limit=5), 'matrix-list')}</div>
            </div>
            {render_text(get_entry(10, 'aria_y_axis'), 'Impact', 'ppt-axis')}
            {render_text(get_entry(10, 'aria_x_axis'), 'Effort', 'ppt-axis')}
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def kpi_slide(page_no: int, slide: dict[str, Any]) -> str:
        metrics = slide.get('metrics') or []
        cards = []
        for i in range(4):
            metric = metrics[i] if i < len(metrics) else {'label': f'Metric {i+1}', 'value': 'XX%', 'detail': 'Detail', 'accent': '#00338d'}
            cards.append(render_box(get_entry(11, f'aria_kpi_{i+1}_card'), 'ppt-kpi-card'))
            cards.append(render_box(get_entry(11, f'aria_kpi_{i+1}_accent'), 'kpi-accent', f'background:{metric.get("accent") or "#00338d"};'))
            metric_text = f'{metric.get("label") or "Metric"}\n{metric.get("value") or "XX%"}'
            cards.append(render_text(get_entry(11, f'aria_kpi_{i+1}'), metric_text, 'ppt-note', 'font-family:var(--display);font-size:24px;line-height:1.15;color:#17191d;'))
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(11, slide.get('style_key') or 'kpi')}
            {render_text(get_entry(11, 'aria_slide_title'), slide.get('title') or 'KPI dashboard', 'ppt-title')}
            {''.join(cards)}
            <div class="ppt-copy-panel" style="{box_style(get_entry(11, 'aria_slide_body'))}">
              <div class="ppt-body" style="position:absolute;inset:18px 20px;">{render_bullets(_split_bullets(' '.join(str((m.get('detail') or '')) for m in metrics), limit=4), 'kpi-body-list')}</div>
            </div>
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def risk_slide(page_no: int, slide: dict[str, Any]) -> str:
        left_panel = get_entry(12, "aria_left_panel")
        right_panel = get_entry(12, "aria_right_panel")
        left_body = get_entry(12, "aria_left_body")
        right_body = get_entry(12, "aria_right_body")
        risk_points = _split_bullets(str(slide.get("left_content") or slide.get("content") or ""), limit=5)
        mitigation_points = _split_bullets(str(slide.get("right_content") or ""), limit=5)
        if not mitigation_points:
            mitigation_points = ["Assign owner and response cadence", "Define trigger thresholds", "Review progress monthly"]
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(12, slide.get('style_key') or 'risk')}
            {render_text(get_entry(12, 'aria_slide_title'), slide.get('title') or 'Risks and mitigations', 'ppt-title')}
            {render_box(left_panel, 'risk-panel', 'background:rgba(255,255,255,.92);backdrop-filter:blur(8px);')}
            {render_box(right_panel, 'risk-panel', 'background:rgba(255,255,255,.92);backdrop-filter:blur(8px);')}
            {render_box(left_panel, 'risk-ribbon', 'height:5.2%;')}
            {render_box(right_panel, 'risk-ribbon', 'height:5.2%;')}
            <div class="risk-head" style="left:{left_panel['x'] + 2.0}%;top:{left_panel['y'] + 2.2}%;width:14%;height:4%;">Risks</div>
            <div class="risk-head" style="left:{right_panel['x'] + 2.0}%;top:{right_panel['y'] + 2.2}%;width:20%;height:4%;">Mitigations</div>
            <div class="ppt-body" style="{box_style(left_body)}">{render_bullets(risk_points, 'risk-list')}</div>
            <div class="ppt-body" style="{box_style(right_body)}">{render_bullets(mitigation_points, 'risk-list')}</div>
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def stakeholder_slide(page_no: int, slide: dict[str, Any]) -> str:
        visual = get_entry(13, "aria_visual_area")
        node_specs = [
            (visual["x"] + 6.0, visual["y"] + 8.0, 10.0, 10.0, "Sponsor"),
            (visual["x"] + 22.0, visual["y"] + 20.0, 11.0, 11.0, "Business"),
            (visual["x"] + 8.0, visual["y"] + 38.0, 10.0, 10.0, "IT"),
        ]
        nodes = ''.join(
            f'<div class="stakeholder-node" style="left:{x}%;top:{y}%;width:{w}%;height:{h}%;">{_html_escape(label)}</div>'
            for x, y, w, h, label in node_specs
        )
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(13, slide.get('style_key') or 'stakeholder')}
            {render_text(get_entry(13, 'aria_slide_title'), slide.get('title') or 'Stakeholder map', 'ppt-title')}
            {render_box(visual, 'ppt-visual')}
            <div class="stakeholder-ring" style="left:{visual['x'] + 7.5}%;top:{visual['y'] + 7.5}%;width:24%;height:24%;"></div>
            <div class="stakeholder-ring" style="left:{visual['x'] + 2.5}%;top:{visual['y'] + 2.5}%;width:34%;height:34%;"></div>
            <div class="stakeholder-grid">{nodes}</div>
            {render_text(get_entry(13, 'aria_visual_caption'), 'Stakeholder influence / alignment map', 'ppt-caption', 'display:grid;place-items:center;')}
            <div class="ppt-body" style="{box_style(get_entry(13, 'aria_slide_body'))}">{render_bullets(_split_bullets(str(slide.get('content') or ''), limit=5), 'stakeholder-list')}</div>
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def visual_note_slide(page_no: int, slide: dict[str, Any]) -> str:
        visual = get_entry(13, "aria_visual_area")
        caption = slide.get('visual_caption') or slide.get('caption') or 'Visual area'
        body = _split_bullets(str(slide.get('content') or ''), limit=5)
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(13, slide.get('style_key') or 'visual_note')}
            {render_text(get_entry(13, 'aria_slide_title'), slide.get('title') or f'Slide {page_no}', 'ppt-title')}
            {render_box(visual, 'ppt-visual')}
            {render_text(get_entry(13, 'aria_visual_caption'), str(caption), 'ppt-caption', 'display:grid;place-items:center;background:rgba(255,255,255,.74);border-radius:999px;')}
            <div class="ppt-copy-panel" style="{box_style(get_entry(13, 'aria_slide_body'))}">
              <div class="ppt-body" style="position:absolute;inset:18px 20px;">{render_bullets(body, 'stakeholder-list')}</div>
            </div>
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def executive_summary_alt_slide(page_no: int, slide: dict[str, Any]) -> str:
        cards = first_items(slide, 3, "Priority")
        notes = _split_bullets(str(slide.get('content') or ''), limit=6)
        supporting = notes[3:] or notes[:3] or ['Supporting text here']
        card_specs = [
            (10.4, 28.0, 23.5, 23.0),
            (38.3, 28.0, 23.5, 23.0),
            (66.2, 28.0, 23.5, 23.0),
        ]
        body_html = []
        for idx, (x, y, w, h) in enumerate(card_specs, start=1):
            title = cards[idx - 1]
            detail = compact_copy(supporting[idx - 1] if idx - 1 < len(supporting) else supporting[-1], 78)
            body_html.append(
                f'<div class="ppt-grid-card" style="left:{x}%;top:{y}%;width:{w}%;height:{h}%;"></div>'
                f'<div class="ppt-grid-badge" style="left:{x + 1.6}%;top:{y + 2.1}%;width:4.2%;height:6.2%;">{idx:02d}</div>'
                f'<div class="ppt-grid-card-title" style="left:{x + 2.0}%;top:{y + 10.0}%;width:{w - 4.0}%;height:7%;">{_html_escape(title)}</div>'
                f'<div class="ppt-grid-card-body" style="left:{x + 2.0}%;top:{y + 17.0}%;width:{w - 4.0}%;height:8%;">{_html_escape(detail)}</div>'
            )
        footer_note = slide.get("supporting_text") or slide.get("caption") or "Supporting text here"
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(22, slide.get('style_key') or 'executive_summary_alt')}
            <div class="ppt-copy-panel" style="left:8.2%;top:10.5%;width:48%;height:12.5%;"></div>
            <div class="ppt-copy-panel" style="left:8.2%;top:55.5%;width:82%;height:12.5%;"></div>
            {render_text(get_entry(4, 'aria_slide_title'), slide.get('title') or f'Slide {page_no}', 'ppt-title')}
            <div class="ppt-body" style="left:9.8%;top:57.8%;width:79%;height:8.4%;">{render_bullets(notes[:3] or notes, 'body-list')}</div>
            {''.join(body_html)}
            <div class="ppt-caption" style="left:9.8%;top:67.6%;width:77%;height:4.6%;">{_html_escape(str(footer_note))}</div>
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def quad_grid_slide(page_no: int, slide: dict[str, Any], numbered: bool = False, source_slide_no: int = 23) -> str:
        items = first_items(slide, 4, "Theme")
        details = _split_bullets(str(slide.get('content') or ''), limit=8)
        card_specs = [
            (9.2, 25.0, 17.8, 20.5),
            (31.0, 25.0, 17.8, 20.5),
            (52.8, 25.0, 17.8, 20.5),
            (74.6, 25.0, 17.8, 20.5),
        ]
        cards = []
        for idx, (x, y, w, h) in enumerate(card_specs, start=1):
            title = items[idx - 1]
            body = compact_copy(details[idx - 1] if idx - 1 < len(details) else "Theme color makes PPT more convenient to change.", 62)
            badge = (
                f'<div class="ppt-grid-badge" style="left:{x + 1.6}%;top:{y + 2.0}%;width:3.4%;height:5.2%;">{idx}</div>'
                if numbered else ""
            )
            cards.append(
                f'<div class="ppt-grid-card" style="left:{x}%;top:{y}%;width:{w}%;height:{h}%;"></div>'
                f'{badge}'
                f'<div class="ppt-grid-card-title" style="left:{x + 1.8}%;top:{y + 7.0}%;width:{w - 3.6}%;height:6.2%;">{_html_escape(title)}</div>'
                f'<div class="ppt-grid-card-body" style="left:{x + 1.8}%;top:{y + 13.2}%;width:{w - 3.6}%;height:6.8%;">{_html_escape(body)}</div>'
            )
        footer = slide.get("supporting_text") or slide.get("caption") or "Supporting text here."
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(source_slide_no, slide.get('style_key') or ('quad_numbered' if numbered else 'quad_grid'))}
            <div class="ppt-copy-panel" style="left:8.3%;top:10.8%;width:44%;height:10.8%;"></div>
            <div class="ppt-copy-panel" style="left:8.3%;top:50.5%;width:84%;height:14.5%;"></div>
            {render_text(get_entry(5, 'aria_slide_title'), slide.get('title') or f'Slide {page_no}', 'ppt-title')}
            {''.join(cards)}
            <div class="ppt-body" style="left:10.1%;top:52.9%;width:81%;height:9.6%;">{render_bullets(details[4:] or details[:4], 'body-list')}</div>
            <div class="ppt-caption" style="left:10.1%;top:63.0%;width:79%;height:4.0%;">{_html_escape(str(footer))}</div>
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def triple_tiles_slide(page_no: int, slide: dict[str, Any]) -> str:
        items = first_items(slide, 3, "Pillar")
        details = _split_bullets(str(slide.get('content') or ''), limit=6)
        card_specs = [
            (12.5, 28.5, 20.5, 24.5),
            (39.8, 28.5, 20.5, 24.5),
            (67.1, 28.5, 20.5, 24.5),
        ]
        cards = []
        for idx, (x, y, w, h) in enumerate(card_specs, start=1):
            cards.append(
                f'<div class="ppt-grid-card" style="left:{x}%;top:{y}%;width:{w}%;height:{h}%;"></div>'
                f'<div class="ppt-grid-card-title" style="left:{x + 1.8}%;top:{y + 10.4}%;width:{w - 3.6}%;height:6.2%;">{_html_escape(items[idx - 1])}</div>'
                f'<div class="ppt-grid-card-body" style="left:{x + 1.8}%;top:{y + 16.7}%;width:{w - 3.6}%;height:6.2%;">{_html_escape(compact_copy(details[idx - 1] if idx - 1 < len(details) else "Copy paste fonts. Choose the only option to retain text.", 58))}</div>'
            )
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(25, slide.get('style_key') or 'triple_tiles')}
            <div class="ppt-copy-panel" style="left:8.3%;top:10.8%;width:48%;height:11.5%;"></div>
            <div class="ppt-copy-panel" style="left:8.3%;top:54.5%;width:84%;height:12.0%;"></div>
            {render_text(get_entry(5, 'aria_slide_title'), slide.get('title') or f'Slide {page_no}', 'ppt-title')}
            {''.join(cards)}
            <div class="ppt-body" style="left:10.0%;top:57.0%;width:81%;height:7.6%;">{render_bullets(details[3:] or details, 'body-list')}</div>
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def detail_grid_slide(page_no: int, slide: dict[str, Any]) -> str:
        items = first_items(slide, 4, "Detail")
        details = _split_bullets(str(slide.get('content') or ''), limit=8)
        card_specs = [
            (10.2, 28.5, 18.5, 20.8),
            (31.0, 28.5, 18.5, 20.8),
            (51.8, 28.5, 18.5, 20.8),
            (72.6, 28.5, 18.5, 20.8),
        ]
        cards = []
        for idx, (x, y, w, h) in enumerate(card_specs, start=1):
            supporting = compact_copy(details[idx - 1] if idx - 1 < len(details) else "Supporting text here.", 60)
            cards.append(
                f'<div class="ppt-grid-card" style="left:{x}%;top:{y}%;width:{w}%;height:{h}%;"></div>'
                f'<div class="ppt-grid-card-title" style="left:{x + 1.6}%;top:{y + 6.5}%;width:{w - 3.2}%;height:5.5%;">{_html_escape(items[idx - 1])}</div>'
                f'<div class="ppt-grid-card-body" style="left:{x + 1.6}%;top:{y + 12.6}%;width:{w - 3.2}%;height:8.6%;">{_html_escape(supporting)}</div>'
            )
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(26, slide.get('style_key') or 'detail_grid')}
            <div class="ppt-copy-panel" style="left:8.3%;top:10.8%;width:45%;height:11.0%;"></div>
            {render_text(get_entry(5, 'aria_slide_title'), slide.get('title') or f'Slide {page_no}', 'ppt-title')}
            {''.join(cards)}
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def six_tiles_slide(page_no: int, slide: dict[str, Any]) -> str:
        items = first_items(slide, 6, "Item")
        details = _split_bullets(str(slide.get('content') or ''), limit=12)
        card_specs = [
            (10.0, 23.8, 24.0, 15.5),
            (38.0, 23.8, 24.0, 15.5),
            (66.0, 23.8, 24.0, 15.5),
            (10.0, 43.2, 24.0, 15.5),
            (38.0, 43.2, 24.0, 15.5),
            (66.0, 43.2, 24.0, 15.5),
        ]
        cards = []
        for idx, (x, y, w, h) in enumerate(card_specs, start=1):
            supporting = compact_copy(details[idx - 1] if idx - 1 < len(details) else "Copy paste fonts. Choose the only option to retain text.", 54)
            cards.append(
                f'<div class="ppt-grid-card" style="left:{x}%;top:{y}%;width:{w}%;height:{h}%;"></div>'
                f'<div class="ppt-grid-card-title" style="left:{x + 1.6}%;top:{y + 2.8}%;width:{w - 3.2}%;height:4.8%;">{_html_escape(items[idx - 1])}</div>'
                f'<div class="ppt-grid-card-body" style="left:{x + 1.6}%;top:{y + 7.6}%;width:{w - 3.2}%;height:6.0%;">{_html_escape(supporting)}</div>'
            )
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(27, slide.get('style_key') or 'six_tiles')}
            <div class="ppt-copy-panel" style="left:8.3%;top:10.8%;width:45%;height:11.0%;"></div>
            {render_text(get_entry(5, 'aria_slide_title'), slide.get('title') or f'Slide {page_no}', 'ppt-title')}
            {''.join(cards)}
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    def render_named_template_slide(page_no: int, slide: dict[str, Any], template_slide_no: int) -> str:
        if template_slide_no == 3:
            return section_slide(page_no, slide.get('title') or f'Section {page_no}')
        if template_slide_no == 4:
            return summary_slide(page_no, slide)
        if template_slide_no == 5:
            return text_focus_slide(page_no, slide)
        if template_slide_no == 6:
            return insight_slide(page_no, slide)
        if template_slide_no == 7:
            return compare_slide(page_no, slide)
        if template_slide_no == 8:
            return framework_slide(page_no, slide)
        if template_slide_no == 9:
            return roadmap_slide(page_no, slide)
        if template_slide_no == 10:
            return matrix_slide(page_no, slide)
        if template_slide_no == 11:
            return kpi_slide(page_no, slide)
        if template_slide_no == 12:
            return risk_slide(page_no, slide)
        if template_slide_no == 13:
            style_key = _normalize_graphic_library_style_key(slide.get('style_key'))
            if style_key == 'stakeholder':
                return stakeholder_slide(page_no, slide)
            return visual_note_slide(page_no, slide)
        if template_slide_no == 14:
            return steps_slide(page_no, slide)
        if template_slide_no == 21:
            return summary_slide(page_no, slide)
        if template_slide_no == 22:
            return executive_summary_alt_slide(page_no, slide)
        if template_slide_no == 23:
            return quad_grid_slide(page_no, slide, numbered=False, source_slide_no=23)
        if template_slide_no == 24:
            return quad_grid_slide(page_no, slide, numbered=True, source_slide_no=24)
        if template_slide_no == 25:
            return triple_tiles_slide(page_no, slide)
        if template_slide_no == 26:
            return detail_grid_slide(page_no, slide)
        if template_slide_no == 27:
            return six_tiles_slide(page_no, slide)
        return insight_slide(page_no, slide)

    def steps_slide(page_no: int, slide: dict[str, Any]) -> str:
        items = _split_bullets(str(slide.get('content') or ''), limit=4)
        rows = []
        for i in range(4):
            rows.append(render_box(get_entry(14, f'aria_step_{i+1}_badge'), 'ppt-step-badge', 'background:#00338d;border:none;box-shadow:none;'))
            rows.append(render_text(get_entry(14, f'aria_step_{i+1}_num'), str(i+1), 'ppt-chip-number', 'font-size:22px;'))
            rows.append(render_text(get_entry(14, f'aria_step_{i+1}_body'), items[i] if i < len(items) else f'Action item {i+1}', 'ppt-note', 'font-size:17px;line-height:1.5;color:#2f3947;'))
        return f"""
        <section class="slide" id="slide-{page_no:02d}">
          <div class="slide-inner">
            {render_source_background(14, slide.get('style_key') or 'next_steps')}
            {render_text(get_entry(14, 'aria_slide_title'), slide.get('title') or 'Next steps', 'ppt-title')}
            {''.join(rows)}
            <div class="slide-no">{page_no:02d}</div>
          </div>
        </section>
        """

    current_page = 3
    for idx, slide in enumerate(slides):
        template_slide_no = None
        if slide.get('template_slide') is not None:
            try:
                template_slide_no = int(slide.get('template_slide'))
            except (TypeError, ValueError):
                template_slide_no = None
        if template_slide_no is None:
            template_slide_no = _graphic_library_style_to_template_slide(slide.get('style_key'))
        slide_type = _normalize_graphic_library_slide_type(slide.get('type'))
        is_first_body_page = idx == 0
        if template_slide_no in {3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 21, 22, 23, 24, 25, 26, 27}:
            sections.append(render_named_template_slide(current_page, slide, template_slide_no))
        elif is_first_body_page and slide_type in {'content', 'summary'}:
            sections.append(summary_slide(current_page, slide))
        elif slide_type == 'section':
            sections.append(section_slide(current_page, slide.get('title') or f'Section {current_page}'))
        elif slide_type == 'framework':
            sections.append(framework_slide(current_page, slide))
        elif slide_type == 'two_column':
            sections.append(compare_slide(current_page, slide))
        elif slide_type == 'roadmap':
            sections.append(roadmap_slide(current_page, slide))
        elif slide_type == 'matrix':
            sections.append(matrix_slide(current_page, slide))
        elif slide_type == 'kpi':
            sections.append(kpi_slide(current_page, slide))
        elif slide_type in {'risk', 'risks', 'risk_mitigation', 'risk-mitigation'}:
            sections.append(risk_slide(current_page, slide))
        elif slide_type in {'next_steps', 'steps'}:
            sections.append(steps_slide(current_page, slide))
        elif slide_type in {'stakeholder', 'stakeholders'}:
            sections.append(stakeholder_slide(current_page, slide))
        else:
            sections.append(insight_slide(current_page, slide))
        current_page += 1

    sections.append(
        f"""
        <section class="slide closing-slide" id="closing">
          <div class="slide-inner">
            {render_box(get_entry(15, 'aria_closing_bg'), 'ppt-box', 'background:transparent;border:none;box-shadow:none;')}
            {render_text(get_entry(15, 'aria_slide_title'), 'Thank you', 'ppt-closing-title')}
            {render_text(get_entry(15, 'aria_slide_body'), 'Discussion / Q&A / Next discussion', 'ppt-closing-body')}
            <div class="slide-no">{current_page:02d}</div>
          </div>
        </section>
        """
    )

    total_count = len(sections)
    return f"""<!doctype html>
<html lang="zh-CN">
  <head>
    <meta charset="UTF-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1.0" />
    <title>{_html_escape(title)}</title>
    <style>{css}</style>
  </head>
  <body>
    <div class="deck-shell">
      <main class="deck" id="deck-track">{''.join(sections)}</main>
      <div class="deck-pager" aria-live="polite">
        <button type="button" id="deck-prev" aria-label="Previous slide">&#8592;</button>
        <div class="deck-status" id="deck-status">01 / {total_count:02d}</div>
        <button type="button" id="deck-next" aria-label="Next slide">&#8594;</button>
      </div>
    </div>
    <script>
      (() => {{
        const track = document.getElementById('deck-track');
        const slides = Array.from(track.querySelectorAll('.slide'));
        const prev = document.getElementById('deck-prev');
        const next = document.getElementById('deck-next');
        const status = document.getElementById('deck-status');
        let current = 0;
        const pad = (value) => String(value).padStart(2, '0');
        const update = () => {{
          const width = window.innerWidth || 1;
          current = Math.round(track.scrollLeft / width);
          current = Math.max(0, Math.min(current, slides.length - 1));
          status.textContent = `${{pad(current + 1)}} / ${{pad(slides.length)}}`;
          prev.disabled = current === 0;
          next.disabled = current === slides.length - 1;
        }};
        const goTo = (index) => {{
          const width = window.innerWidth || 1;
          const target = Math.max(0, Math.min(index, slides.length - 1));
          track.scrollTo({{ left: target * width, behavior: 'smooth' }});
        }};
        prev.addEventListener('click', () => goTo(current - 1));
        next.addEventListener('click', () => goTo(current + 1));
        track.addEventListener('scroll', () => requestAnimationFrame(update), {{ passive: true }});
        window.addEventListener('resize', update);
        document.addEventListener('keydown', (event) => {{
          if (event.key === 'ArrowRight' || event.key === 'PageDown' || event.key === ' ') {{
            event.preventDefault();
            goTo(current + 1);
          }}
          if (event.key === 'ArrowLeft' || event.key === 'PageUp') {{
            event.preventDefault();
            goTo(current - 1);
          }}
        }});
        update();
      }})();
    </script>
  </body>
</html>"""


@registry.register(
    name="generate_html_deck_from_skill",
    description="Generate a consulting-style HTML/H5 deck using a Skill's approved visual system.",
    input_schema={
        "type": "object",
        "properties": {
            "skill_name": {
                "type": "string",
                "description": "Skill folder name (e.g., 'graphic-library-h5')"
            },
            "title": {
                "type": "string",
                "description": "Deck title"
            },
            "subtitle": {
                "type": "string",
                "description": "Optional deck subtitle"
            },
            "slides": {
                "type": "array",
                "description": "Structured slide/page definitions for the H5 deck",
                "items": {
                    "type": "object",
                    "properties": {
                        "type": {
                            "type": "string",
                            "enum": [
                                "title",
                                "section",
                                "section_divider",
                                "executive_summary",
                                "summary",
                                "content",
                                "insight",
                                "insight_content",
                                "two_column",
                                "compare",
                                "three_part",
                                "three-part",
                                "three_column",
                                "framework",
                                "roadmap",
                                "matrix",
                                "kpi",
                                "kpi_cards",
                                "risk",
                                "risks",
                                "risks_and_mitigations",
                                "stakeholder",
                                "stakeholder_map",
                                "next_steps",
                                "steps"
                            ]
                        },
                        "style_key": {
                            "type": "string",
                            "description": "Optional explicit Graphic Library style key. Examples include executive_summary, executive_summary_alt, text_focus, insight_visual, compare, framework, roadmap, matrix, kpi, risk, visual_note, stakeholder, next_steps, quad_grid, quad_numbered, triple_tiles, detail_grid, and six_tiles. If omitted, the generator will infer the closest existing page style automatically."
                        },
                        "template_slide": {
                            "type": "integer",
                            "description": "Optional explicit template slide number from the Graphic Library H5 automation map. Standard body pages use 3-14; source-near extended styles also support 21-27. If omitted, it can be inferred from style_key."
                        },
                        "title": {"type": "string"},
                        "content": {"type": "string"},
                        "left_content": {"type": "string"},
                        "right_content": {"type": "string"},
                        "visual_caption": {"type": "string"},
                        "caption": {"type": "string"},
                        "phase_1": {"type": "string"},
                        "phase_2": {"type": "string"},
                        "phase_3": {"type": "string"},
                        "labels": {
                            "type": "array",
                            "items": {"type": "string"},
                        },
                        "card_titles": {
                            "type": "array",
                            "items": {"type": "string"},
                        },
                        "metrics": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "label": {"type": "string"},
                                    "value": {"type": "string"},
                                    "detail": {"type": "string"},
                                    "accent": {"type": "string"},
                                },
                            },
                        },
                    },
                    "required": ["type", "title"]
                }
            }
        },
        "required": ["skill_name", "title", "slides"]
    }
)
async def generate_html_deck_from_skill(
    skill_name: str,
    title: str,
    slides: list[dict],
    subtitle: str = "",
) -> dict[str, Any]:
    """Generate a branded HTML/H5 deck from a skill-specific visual grammar."""
    if skill_name != "graphic-library-h5":
        return {
            "success": False,
            "error": f"{skill_name} does not support HTML/H5 deck generation.",
        }

    if not slides:
        return {
            "success": False,
            "error": "slides cannot be empty for HTML deck generation.",
        }

    skill_dir = SKILLS_DIR / skill_name
    required_paths = [
        skill_dir / "SKILL.md",
        skill_dir / "references" / "style-system.md",
        skill_dir / "references" / "page-families.md",
        skill_dir / "references" / "deck-structure.md",
    ]
    missing_paths = [str(path) for path in required_paths if not path.exists()]
    if missing_paths:
        return {
            "success": False,
            "error": f"{skill_name} references are incomplete; refusing to generate an ungoverned deck.",
            "missing_paths": missing_paths,
        }

    normalized_slides = _prepare_graphic_library_h5_slides(slides)
    html = _build_graphic_library_h5_html_v2(title, subtitle, normalized_slides)
    filename = f"{_slugify_filename(title)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
    filepath = GENERATED_DIR / filename
    filepath.write_text(html, encoding="utf-8")

    return {
        "success": True,
        "file_type": "html",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
        "template_applied": True,
        "template_name": "graphic-library-h5",
        "template_mode": "graphic_library_h5",
        "slide_count": html.count('<section class="slide'),
        "resolved_styles": [
            {
                "title": str(slide.get("title") or ""),
                "type": str(slide.get("type") or ""),
                "style_key": str(slide.get("style_key") or ""),
                "template_slide": slide.get("template_slide"),
            }
            for slide in normalized_slides
        ],
    }


@registry.register(
    name="generate_ppt",
    description="Generate a PowerPoint (.pptx) presentation with structured slides. Supports title slides, content slides, and two-column layouts.",
    input_schema={
        "type": "object",
        "properties": {
            "title": {
                "type": "string",
                "description": "Presentation title"
            },
            "subtitle": {
                "type": "string",
                "description": "Optional subtitle for title slide"
            },
            "slides": {
                "type": "array",
                "description": "Array of slide objects",
                "items": {
                    "type": "object",
                    "properties": {
                        "type": {
                            "type": "string",
                            "enum": ["title", "content", "two_column"],
                            "description": "Slide layout type"
                        },
                        "title": {
                            "type": "string",
                            "description": "Slide title"
                        },
                        "content": {
                            "type": "string",
                            "description": "Main content (bullet points supported with - or *)"
                        },
                        "left_content": {
                            "type": "string",
                            "description": "Left column content (for two_column type)"
                        },
                        "right_content": {
                            "type": "string",
                            "description": "Right column content (for two_column type)"
                        }
                    },
                    "required": ["type", "title"]
                }
            },
            "template_path": {
                "type": "string",
                "description": "Optional path to a .pptx template file. If not provided, uses default blank template."
            }
        },
        "required": ["title", "slides"]
    }
)
def _safe_layout(prs, preferred_idx: int):
    """Return slide_layouts[preferred_idx] if it exists, otherwise the last available layout."""
    layouts = prs.slide_layouts
    if preferred_idx < len(layouts):
        return layouts[preferred_idx]
    return layouts[-1]


def _find_body_placeholder(slide):
    """Find the first non-title body/text placeholder on a slide.

    Templates can use non-standard idx values (e.g. 10, 11 instead of 1).
    Try common indices first, then fall back to any non-title placeholder.
    """
    for idx in (1, 10, 11, 2, 12, 13, 14, 15):
        try:
            return slide.placeholders[idx]
        except KeyError:
            continue
    # Last resort: first placeholder that is not the title (idx != 0)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:
            return ph
    return None


def _layout_placeholder_idx(slide, placeholder_name: str) -> int | None:
    for placeholder in slide.slide_layout.placeholders:
        if placeholder.name == placeholder_name:
            return placeholder.placeholder_format.idx
    return None


def _placeholder_by_layout_name(slide, placeholder_name: str):
    idx = _layout_placeholder_idx(slide, placeholder_name)
    if idx is None:
        return None
    try:
        return slide.placeholders[idx]
    except KeyError:
        return None


def _shape_by_name(slide, shape_name: str):
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape
    return None


def _write_text_preserving_style(frame, text: str) -> None:
    frame.word_wrap = True
    paragraphs = list(frame.paragraphs)
    if paragraphs:
        first_paragraph = paragraphs[0]
        if first_paragraph.runs:
            first_paragraph.runs[0].text = text
            for run in first_paragraph.runs[1:]:
                run.text = ""
        else:
            first_paragraph.text = text
        for paragraph in paragraphs[1:]:
            paragraph.text = ""
        return
    frame.text = text


def _text_display_units(text: str) -> float:
    units = 0.0
    for char in str(text or ""):
        if char.isspace():
            units += 0.35
        elif ord(char) > 127:
            units += 1.05
        elif char.isupper():
            units += 0.68
        else:
            units += 0.55
    return max(units, 1.0)


def _fit_title_text_frame_one_line(frame, text: str, width, *, max_size: int | None = None, min_size: int = 14) -> None:
    """Shrink title text to stay on one line while preserving template styling."""
    from pptx.util import Pt

    clean_text = " ".join(str(text or "").split())
    if not clean_text:
        return

    frame.word_wrap = False
    margin_pt = sum(
        (getattr(frame, attr, 0) or 0) / 12700
        for attr in ("margin_left", "margin_right")
    )
    available_pt = max(width / 12700 - margin_pt - 2, 24)
    estimated_size = int(available_pt / (_text_display_units(clean_text) * 0.56))

    current_sizes: list[int] = []
    for paragraph in frame.paragraphs:
        if paragraph.font.size:
            current_sizes.append(int(paragraph.font.size.pt))
        for run in paragraph.runs:
            if run.font.size:
                current_sizes.append(int(run.font.size.pt))
    cap = max_size or (max(current_sizes) if current_sizes else 28)
    fitted_size = max(min_size, min(cap, estimated_size))

    for paragraph in frame.paragraphs:
        paragraph.font.size = Pt(fitted_size)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.size = Pt(fitted_size)


def _write_title_preserving_style(shape, text: str, *, min_size: int = 14, max_size: int | None = None) -> None:
    clean_text = " ".join(str(text or "").split())
    _write_text_preserving_style(shape.text_frame, clean_text)
    _fit_title_text_frame_one_line(shape.text_frame, clean_text, shape.width, min_size=min_size, max_size=max_size)


def _set_placeholder_text(slide, placeholder_name: str, text: str) -> bool:
    placeholder = _placeholder_by_layout_name(slide, placeholder_name)
    if placeholder is None or not placeholder.has_text_frame:
        return False
    _write_text_preserving_style(placeholder.text_frame, text)
    return True


def _set_named_or_placeholder_text(slide, shape_name: str, text: str) -> bool:
    """Write text by normal shape name first, then by layout placeholder name."""
    shape = _shape_by_name(slide, shape_name)
    if shape is not None and getattr(shape, "has_text_frame", False):
        _write_text_preserving_style(shape.text_frame, text)
        return True
    return _set_placeholder_text(slide, shape_name, text)


def _set_title_named_or_placeholder_text(slide, shape_name: str, text: str, *, min_size: int = 14, max_size: int | None = None) -> bool:
    """Write a named title and shrink it to one line without replacing template style."""
    shape = _shape_by_name(slide, shape_name)
    if shape is not None and getattr(shape, "has_text_frame", False):
        _write_title_preserving_style(shape, text, min_size=min_size, max_size=max_size)
        return True

    placeholder = _placeholder_by_layout_name(slide, shape_name)
    if placeholder is None or not placeholder.has_text_frame:
        return False
    _write_title_preserving_style(placeholder, text, min_size=min_size, max_size=max_size)
    return True


def _has_usable_text_bounds(shape, *, min_width_inches: float = 1.0, min_height_inches: float = 0.25) -> bool:
    from pptx.util import Inches

    return bool(
        shape is not None
        and getattr(shape, "width", 0) >= Inches(min_width_inches)
        and getattr(shape, "height", 0) >= Inches(min_height_inches)
    )


def _bump_text_frame_font_size(frame, *, default_size: int = 14, delta: int = 2, max_size: int = 17, min_size: int = 12) -> None:
    from pptx.util import Pt

    for paragraph in frame.paragraphs:
        current = paragraph.font.size.pt if paragraph.font.size else default_size
        paragraph.font.size = Pt(max(min_size, min(max_size, current + delta)))
        for run in paragraph.runs:
            run_current = run.font.size.pt if run.font.size else current
            run.font.size = Pt(max(min_size, min(max_size, run_current + delta)))


def _ensure_min_text_frame_font_size(frame, *, min_size: int = 12, exclude_empty: bool = True) -> None:
    from pptx.util import Pt

    for paragraph in frame.paragraphs:
        if exclude_empty and not paragraph.text.strip():
            continue
        current = paragraph.font.size.pt if paragraph.font.size else min_size
        paragraph.font.size = Pt(max(min_size, current))
        for run in paragraph.runs:
            if exclude_empty and not run.text.strip():
                continue
            run_current = run.font.size.pt if run.font.size else current
            run.font.size = Pt(max(min_size, run_current))


def _ensure_body_min_font_sizes(prs, *, min_size: int = 12) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            shape_name = str(getattr(shape, "name", "") or "").lower()
            is_body_shape = any(token in shape_name for token in ("body", "content", "lead", "metric", "note", "caption"))
            if is_body_shape:
                _ensure_min_text_frame_font_size(shape.text_frame, min_size=min_size)


def _set_body_named_or_placeholder_text(slide, shape_name: str, text: str, *, default_size: int = 14) -> bool:
    shape = _shape_by_name(slide, shape_name)
    if shape is not None and getattr(shape, "has_text_frame", False):
        if not _has_usable_text_bounds(shape):
            return False
        _write_text_preserving_style(shape.text_frame, text)
        _bump_text_frame_font_size(shape.text_frame, default_size=default_size)
        return True

    placeholder = _placeholder_by_layout_name(slide, shape_name)
    if placeholder is None or not placeholder.has_text_frame:
        return False
    if not _has_usable_text_bounds(placeholder):
        return False
    _write_text_preserving_style(placeholder.text_frame, text)
    _bump_text_frame_font_size(placeholder.text_frame, default_size=default_size)
    return True


def _push_body_below_lead(slide, body_names: tuple[str, ...] = ("aria_slide_body",), gap_inches: float = 0.16) -> None:
    from pptx.util import Inches

    lead = _shape_by_name(slide, "aria_generated_lead") or _shape_by_name_or_placeholder(slide, "aria_slide_lead")
    if lead is None:
        return
    lead_bottom = lead.top + lead.height + Inches(gap_inches)
    slide_bottom = Inches(7.2)
    min_height = Inches(1.2)

    for name in body_names:
        body = _shape_by_name_or_placeholder(slide, name)
        if body is None or body.top >= lead_bottom:
            continue
        delta = lead_bottom - body.top
        body.top = lead_bottom
        body.height = max(min_height, min(body.height - delta, slide_bottom - body.top))


def _placeholder_bounds(slide, placeholder_name: str):
    placeholder = _placeholder_by_layout_name(slide, placeholder_name)
    if placeholder is None:
        return None
    return placeholder.left, placeholder.top, placeholder.width, placeholder.height


def _bounds_by_name_or_placeholder(slide, shape_name: str):
    shape = _shape_by_name(slide, shape_name)
    if shape is not None:
        return shape.left, shape.top, shape.width, shape.height
    return _placeholder_bounds(slide, shape_name)


def _shape_by_name_or_placeholder(slide, shape_name: str):
    return _shape_by_name(slide, shape_name) or _placeholder_by_layout_name(slide, shape_name)


def _resize_named_or_placeholder(slide, shape_name: str, x, y, w, h) -> bool:
    shape = _shape_by_name_or_placeholder(slide, shape_name)
    if shape is None:
        return False
    shape.left = x
    shape.top = y
    shape.width = w
    shape.height = h
    return True


def _style_text_frame(frame, role: str):
    """Apply a clean consulting-style baseline to template text frames."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    role = role.lower()
    is_title = "title" in role
    is_cover = "cover" in role
    is_section = "section" in role
    for paragraph in frame.paragraphs:
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(30 if is_cover and is_title else 27 if is_section and is_title else 21 if is_title else 12)
        paragraph.font.bold = is_title
        if is_cover:
            paragraph.font.color.rgb = RGBColor.from_string("FFFFFF")
        else:
            paragraph.font.color.rgb = RGBColor.from_string("111827" if is_title else "334155")
        paragraph.space_after = Pt(4 if is_title else 3)
        paragraph.line_spacing = 1.08 if is_title else 1.18
        if paragraph.text.strip().startswith("-"):
            paragraph.level = 0


def _set_content_slide_text(slide, title: str, content: str):
    """Write visible title/body text even when templates use unusual placeholders."""
    from pptx.util import Inches, Pt

    if slide.shapes.title:
        _write_title_preserving_style(slide.shapes.title, title)
    else:
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.6))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.text_frame.paragraphs[0].font.bold = True
        _fit_title_text_frame_one_line(title_box.text_frame, title, title_box.width, min_size=13, max_size=24)

    body = _find_body_placeholder(slide)
    if body is not None:
        body.text_frame.text = content
        body.text_frame.word_wrap = True
    else:
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(11.6), Inches(5.7))
        body_box.text_frame.word_wrap = True
        body_box.text_frame.text = content


def _set_template_cover_text(slide, title: str, subtitle: str):
    """Best-effort cover replacement for branded templates."""
    replaced_title = _set_named_or_placeholder_text(slide, "aria_cover_title", title)
    replaced_subtitle = bool(subtitle) and _set_named_or_placeholder_text(slide, "aria_cover_subtitle", subtitle)
    if replaced_title and (replaced_subtitle or not subtitle):
        return

    replaced_title = False
    replaced_subtitle = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        current = (shape.text or "").strip()
        if "客户名称" in current or "方案建议名称" in current or shape.name.lower() in {"cover_title", "title 3"}:
            shape.text_frame.text = f"{title}\n方案建议书"
            replaced_title = True
        elif subtitle and ("KPMG China" in current or shape.name.lower() in {"cover_subtitle", "subtitle 4"}):
            if not replaced_subtitle:
                shape.text_frame.text = subtitle
                replaced_subtitle = True
            elif shape.name.lower() in {"cover_subtitle", "subtitle 4"}:
                shape.text_frame.text = ""

    if not replaced_title:
        from pptx.util import Inches, Pt

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(10.8), Inches(1.2))
        title_box.text_frame.text = f"{title}\n方案建议书"
        for paragraph in title_box.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
    if subtitle and not replaced_subtitle:
        from pptx.util import Inches, Pt

        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(10.8), Inches(0.5))
        subtitle_box.text_frame.text = subtitle
        subtitle_box.text_frame.paragraphs[0].font.size = Pt(14)


def _remove_slide(prs, index: int):
    """Remove a template slide by index."""
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    if 0 <= index < len(slides):
        slide_id = slides[index]
        rel_id = slide_id.rId
        slide_id_list.remove(slide_id)
        prs.part.drop_rel(rel_id)


def _slide_ref(prs, index: int):
    refs = list(prs.slides._sldIdLst)
    if 0 <= index < len(refs):
        return refs[index]
    return None


def _remove_slide_ref(prs, slide_ref):
    if slide_ref is None:
        return
    slide_id_list = prs.slides._sldIdLst
    if slide_ref not in slide_id_list:
        return
    rel_id = slide_ref.rId
    slide_id_list.remove(slide_ref)
    prs.part.drop_rel(rel_id)


def _move_slide_ref_to_end(prs, slide_ref):
    if slide_ref is None:
        return
    slide_id_list = prs.slides._sldIdLst
    if slide_ref not in slide_id_list:
        return
    slide_id_list.remove(slide_ref)
    slide_id_list.append(slide_ref)


def _replace_relationship_ids(element, rel_id_map: dict[str, str]) -> None:
    if not rel_id_map:
        return
    for node in element.iter():
        for attr_name, attr_value in list(node.attrib.items()):
            if attr_value in rel_id_map:
                node.attrib[attr_name] = rel_id_map[attr_value]


def _clone_slide_from_prototype(prs, source_slide):
    """Duplicate a normal template slide so custom slide-level artwork survives.

    python-pptx's add_slide(layout) only applies the master/layout. It does not
    copy shapes that designers place on a normal slide, which made generated
    decks look like empty layouts. This clones the prototype slide's XML shapes
    and non-layout relationships before the content renderer writes text.
    """
    slide = prs.slides.add_slide(source_slide.slide_layout)
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape.element)

    cloned_elements = [deepcopy(shape.element) for shape in source_slide.shapes]
    rel_id_map: dict[str, str] = {}
    for rel in source_slide.part.rels.values():
        if rel.reltype.endswith("/slideLayout") or rel.reltype.endswith("/notesSlide"):
            continue
        new_rel_id = slide.part.rels._add_relationship(
            rel.reltype,
            rel._target,
            getattr(rel, "is_external", False),
        )
        rel_id_map[rel.rId] = new_rel_id

    for element in cloned_elements:
        _replace_relationship_ids(element, rel_id_map)
        sp_tree.insert_element_before(element, "p:extLst")
    return slide


def _clear_text_shapes(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.clear()


def _clear_generated_text_shapes(slide):
    for shape in slide.shapes:
        shape_name = str(getattr(shape, "name", "") or "")
        if shape_name.startswith("aria_"):
            continue
        if shape.has_text_frame and not shape.is_placeholder:
            shape.text_frame.clear()


def _remove_shape_by_name(slide, shape_name: str) -> bool:
    removed = False
    for shape in list(slide.shapes):
        if getattr(shape, "name", "") != shape_name:
            continue
        shape.element.getparent().remove(shape.element)
        removed = True
    return removed


def _slide_has_named_shapes(slide, required_names: set[str]) -> bool:
    shape_names = {getattr(shape, "name", "") for shape in slide.shapes}
    return required_names.issubset(shape_names)


def _find_named_prototype_slide(prs, required_names: set[str], preferred_layout_name: str | None = None):
    if preferred_layout_name:
        for slide in prs.slides:
            if slide.slide_layout.name == preferred_layout_name and _slide_has_named_shapes(slide, required_names):
                return slide
    for slide in prs.slides:
        if _slide_has_named_shapes(slide, required_names):
            return slide
    return None


def _shorten_slide_lead(text: str, limit: int = 96) -> str:
    text = " ".join(_clean_ppt_text(text).split())
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip("，,。.;； ") + "…"


def _slide_lead_text(slide_data: dict | None, title: str = "", content: str = "") -> str:
    slide_data = slide_data or {}
    for key in ("lead", "insight", "subtitle", "summary", "headline"):
        value = str(slide_data.get(key) or "").strip()
        if value:
            return _shorten_slide_lead(value)

    bullets = _split_bullets(_combined_slide_content(slide_data) or content, limit=1)
    if bullets:
        return _shorten_slide_lead(bullets[0])

    title = str(title or "").strip()
    if title:
        return _shorten_slide_lead(f"本页围绕“{title}”展开关键判断、管理含义与下一步动作。")
    return ""


def _add_generated_slide_lead(slide, text: str, *, role: str = "content") -> bool:
    from pptx.util import Inches

    lead = _shorten_slide_lead(text)
    if not lead:
        return False
    if role == "section":
        lead = _shorten_slide_lead(lead, limit=72)

    _remove_shape_by_name(slide, "aria_generated_lead")
    placeholder_name = "aria_section_lead" if role == "section" else "aria_slide_lead"
    if _set_named_or_placeholder_text(slide, placeholder_name, lead):
        shape = _shape_by_name_or_placeholder(slide, placeholder_name)
        if role == "section" and shape is not None and getattr(shape, "has_text_frame", False):
            from pptx.dml.color import RGBColor

            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor.from_string("FFFFFF")
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor.from_string("FFFFFF")
        return True

    title_shape = _shape_by_name_or_placeholder(
        slide,
        "aria_section_title" if role == "section" else "aria_slide_title",
    )
    if title_shape is None and slide.shapes.title is not None:
        title_shape = slide.shapes.title

    if title_shape is not None:
        left = title_shape.left
        width = title_shape.width
        top = title_shape.top + title_shape.height + Inches(0.08)
        if role == "section" and width < Inches(5.0):
            width = Inches(4.35)
    else:
        left = Inches(0.82)
        width = Inches(10.9)
        top = Inches(0.95 if role != "section" else 3.05)
        if role == "section":
            left = Inches(1.53)
            width = Inches(4.35)

    height = Inches(0.34 if role != "section" else 1.02)
    lead_box = _add_textbox(
        slide,
        left,
        top,
        width,
        height,
        lead,
        size=12 if role != "section" else 11,
        color="64748B" if role != "section" else "FFFFFF",
    )
    lead_box.name = "aria_generated_lead"
    return True


def _render_section_slide(slide, title: str, section_number: int, slide_data: dict | None = None):
    _clear_generated_text_shapes(slide)
    if not _set_title_named_or_placeholder_text(slide, "aria_section_title", title, min_size=18):
        if slide.shapes.title is not None:
            _write_title_preserving_style(slide.shapes.title, title, min_size=18)
        else:
            title_shape = None
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and not getattr(shape, "is_placeholder", False):
                    title_shape = shape
                    break
            if title_shape is not None:
                title_shape.text_frame.clear()
                title_shape.text_frame.text = title
                title_shape.text_frame.word_wrap = True
                _style_text_frame(title_shape.text_frame, "aria_section_title")
                _fit_title_text_frame_one_line(title_shape.text_frame, title, title_shape.width, min_size=18)
    if not _set_named_or_placeholder_text(slide, "aria_section_number", f"{section_number:02d}"):
        body = _find_body_placeholder(slide)
        if body is not None:
            body.text = f"{section_number:02d}"
    _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title), role="section")


def _render_graphic_library_title_slide(slide, title: str, slide_data: dict | None = None):
    _clear_generated_text_shapes(slide)
    if not _set_title_named_or_placeholder_text(slide, "aria_slide_title", title) and slide.shapes.title is not None:
        _write_title_preserving_style(slide.shapes.title, title)
    _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title), role="section")


def _render_graphic_library_content_slide(slide, title: str, content: str, slide_number: int, slide_data: dict | None = None):
    from pptx.util import Inches

    _clear_generated_text_shapes(slide)
    bullets = _split_bullets(content, limit=6)
    body_text = "\n".join(f"- {bullet}" for bullet in bullets) if bullets else ""

    title_shape = None
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and "Title" in getattr(shape, "name", ""):
            title_shape = shape
            break
    if title_shape is not None:
        _write_title_preserving_style(title_shape, title)
    _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, content), role="content")

    placed = False
    body = _find_body_placeholder(slide)
    if body is not None:
        body.text_frame.clear()
        body.text_frame.text = body_text
        body.text_frame.word_wrap = True
        _bump_text_frame_font_size(body.text_frame, default_size=14)
        placed = True

    if not placed:
        _add_textbox(slide, Inches(0.9), Inches(1.55), Inches(5.8), Inches(4.8), body_text, size=13, color="334155")

    _add_textbox(slide, Inches(11.55), Inches(6.75), Inches(0.7), Inches(0.3), f"{slide_number:02d}", size=8, color="94A3B8")


def _add_textbox(slide, x, y, w, h, text: str, *, size: int = 14, bold: bool = False, color: str = "1F2937"):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    text_value = str(text or "").strip()
    if 7 <= size <= 8:
        size += 2
    elif 9 <= size <= 13:
        size += 1
    if len(text_value) >= 8 or sum(1 for char in text_value if ord(char) > 127) >= 4:
        size = max(size, 12)

    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(4)
    frame.margin_bottom = Pt(4)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor.from_string(color)
    paragraph.line_spacing = 1.08
    return box


def _add_card(slide, x, y, w, h, *, fill: str = "F8FAFC", line: str = "E2E8F0"):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = Pt(1)
    return shape


def _add_shape(slide, shape_type, x, y, w, h, *, fill: str = "E0F2FE", line: str = "BAE6FD"):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = Pt(1)
    return shape


def _clean_ppt_text(value: Any) -> str:
    text = str(value or "").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"```.*?```", "", text, flags=re.S)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"^\s{0,3}#{1,6}\s*", "", text, flags=re.M)
    text = re.sub(r"^\s{0,3}>\s*", "", text, flags=re.M)
    text = re.sub(r"[*_]{1,3}([^*_]+)[*_]{1,3}", r"\1", text)
    text = re.sub(r"<[^>]+>", "", text)
    cleaned_lines: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("|") and stripped.endswith("|"):
            cells = [cell.strip() for cell in stripped.strip("|").split("|") if cell.strip()]
            if cells and all(set(cell) <= {"-", ":", " "} for cell in cells):
                continue
            stripped = "；".join(cells)
        cleaned_lines.append(stripped)
    return "\n".join(cleaned_lines).strip()


def _split_bullets(content: str, limit: int = 6) -> list[str]:
    bullets: list[str] = []
    for raw in _clean_ppt_text(content).splitlines():
        item = raw.strip()
        item = re.sub(r"^[-*•·]\s*", "", item)
        item = re.sub(r"^\d{1,2}[.、)）]\s*", "", item)
        item = _clean_ppt_text(item)
        if item:
            bullets.append(item)
    if not bullets and _clean_ppt_text(content):
        bullets = [_clean_ppt_text(content)]
    return bullets[:limit]


def _prepare_ppt_slide_text(slides: list[dict]) -> list[dict]:
    prepared: list[dict] = []
    for slide in slides:
        item = dict(slide)
        for key in ("title", "subtitle", "content", "left_content", "right_content"):
            if key in item:
                item[key] = _clean_ppt_text(item.get(key))
        prepared.append(item)
    return prepared


def _add_slide_header(slide, title: str, slide_number: int):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string("2563EB")
    accent.line.fill.background()

    title_box = _add_textbox(slide, Inches(0.65), Inches(0.32), Inches(10.8), Inches(0.62), " ".join(str(title or "").split()), size=22, bold=True, color="111827")
    _fit_title_text_frame_one_line(title_box.text_frame, title, title_box.width, min_size=13, max_size=22)
    _add_textbox(slide, Inches(11.65), Inches(0.42), Inches(1.1), Inches(0.35), f"{slide_number:02d}", size=11, bold=True, color="94A3B8")


def _add_slide_footer(slide, label: str = "AriaAI generated deck"):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(7.02), Inches(12), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string("E5E7EB")
    line.line.fill.background()
    _add_textbox(slide, Inches(0.65), Inches(7.08), Inches(5.5), Inches(0.25), label, size=8, color="94A3B8")


def _add_template_frame(slide, title: str, slide_number: int, *, section: str = "Digital strategy"):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    top_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(0.18), Inches(12.18), Inches(0.025))
    top_rule.fill.solid()
    top_rule.fill.fore_color.rgb = RGBColor.from_string("D7DEE8")
    top_rule.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(0.18), Inches(1.45), Inches(0.025))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string("1D4ED8")
    accent.line.fill.background()

    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.35), Inches(0.05), Inches(5.3))
    rail.fill.solid()
    rail.fill.fore_color.rgb = RGBColor.from_string("E2E8F0")
    rail.line.fill.background()

    marker_top = Inches(1.35 + ((slide_number - 1) % 5) * 0.56)
    marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.51), marker_top, Inches(0.19), Inches(0.5))
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor.from_string("1D4ED8")
    marker.line.fill.background()

    _add_textbox(slide, Inches(0.86), Inches(6.92), Inches(4.2), Inches(0.22), section.upper(), size=7, bold=True, color="94A3B8")
    _add_textbox(slide, Inches(11.85), Inches(6.9), Inches(0.8), Inches(0.24), f"{slide_number:02d}", size=8, bold=True, color="94A3B8")


def _add_consulting_visual_panel(slide, x, y, w, h, bullets: list[str]):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    _add_card(slide, x, y, w, h, fill="FAFBFC", line="D7DEE8")
    _add_textbox(slide, x + Inches(0.24), y + Inches(0.16), w - Inches(0.48), Inches(0.28), "管理层视角", size=11, bold=True, color="0F172A")
    _add_textbox(slide, x + Inches(0.24), y + Inches(0.48), w - Inches(0.48), Inches(0.24), "价值驱动因素", size=8, color="64748B")

    colors = ["1D4ED8", "047857", "B45309"]
    labels = ["价值", "采用", "规模化"]
    for idx, label in enumerate(labels):
        top = y + Inches(0.88 + idx * 0.64)
        _add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.28), top, Inches(0.30), Inches(0.30), fill=colors[idx], line=colors[idx])
        _add_textbox(slide, x + Inches(0.72), top - Inches(0.03), w - Inches(1.0), Inches(0.26), label, size=9, bold=True, color="334155")
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.72), top + Inches(0.28), Inches(1.35 + idx * 0.34), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(colors[idx])
        bar.line.fill.background()

    insight = bullets[0] if bullets else "将数字化投入与可衡量的业务结果绑定。"
    matrix_x = x + Inches(0.26)
    matrix_y = y + Inches(2.95)
    matrix_w = w - Inches(0.52)
    matrix_h = Inches(0.86)
    _add_card(slide, matrix_x, matrix_y, matrix_w, matrix_h, fill="EEF2FF", line="C7D2FE")
    _add_textbox(slide, matrix_x + Inches(0.16), matrix_y + Inches(0.13), matrix_w - Inches(0.32), Inches(0.46), insight[:110], size=9, bold=True, color="1E1B4B")
    _add_textbox(slide, x + Inches(0.28), y + Inches(3.98), w - Inches(0.56), Inches(0.22), "决策重点：优先推进高确定性、高影响力举措", size=7, color="64748B")


def _add_value_chain_visual(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    labels = ["诊断", "设计", "动员", "规模化"]
    colors = ["EFF6FF", "ECFDF5", "FFF7ED", "F8FAFC"]
    accents = ["1D4ED8", "047857", "C2410C", "475569"]
    step_w = w / 4 - Inches(0.08)
    for idx, label in enumerate(labels):
        left = x + idx * (step_w + Inches(0.1))
        _add_card(slide, left, y, step_w, h, fill=colors[idx], line="D8DEE9")
        _add_shape(slide, MSO_SHAPE.OVAL, left + Inches(0.12), y + Inches(0.12), Inches(0.34), Inches(0.34), fill=accents[idx], line=accents[idx])
        _add_textbox(slide, left + Inches(0.22), y + Inches(0.19), Inches(0.14), Inches(0.12), str(idx + 1), size=7, bold=True, color="FFFFFF")
        _add_textbox(slide, left + Inches(0.12), y + Inches(0.55), step_w - Inches(0.24), Inches(0.24), label, size=9, bold=True, color=accents[idx])
        if idx < len(bullets):
            _add_textbox(slide, left + Inches(0.12), y + Inches(0.88), step_w - Inches(0.24), h - Inches(1.0), bullets[idx][:80], size=7, color="475569")


def _add_roadmap_visual(slide, title: str, bullets: list[str], slide_number: int) -> bool:
    if not any(token in title.lower() for token in ("roadmap", "horizon", "phase", "路线", "阶段")):
        return False

    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    _clear_text_shapes(slide)
    _add_slide_header(slide, title, slide_number)
    _add_generated_slide_lead(slide, _slide_lead_text(None, title, "\n".join(bullets)), role="content")
    phases = bullets[:3] or ["阶段一：夯实基础", "阶段二：规模复制", "阶段三：领先优化"]
    colors = [("DBEAFE", "2563EB"), ("DCFCE7", "16A34A"), ("FEF3C7", "D97706")]
    for idx, phase in enumerate(phases):
        x = Inches(0.85 + idx * 4.05)
        fill, accent = colors[idx % len(colors)]
        _add_card(slide, x, Inches(1.55), Inches(3.55), Inches(4.55), fill=fill, line=accent)
        _add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.25), Inches(1.85), Inches(0.58), Inches(0.58), fill=accent, line=accent)
        _add_textbox(slide, x + Inches(0.42), Inches(1.97), Inches(0.28), Inches(0.22), str(idx + 1), size=11, bold=True, color="FFFFFF")
        label = ["夯实基础", "规模复制", "领先优化"][idx] if idx < 3 else f"阶段 {idx + 1}"
        _add_textbox(slide, x + Inches(0.95), Inches(1.88), Inches(2.2), Inches(0.35), label, size=15, bold=True, color=accent)
        _add_textbox(slide, x + Inches(0.35), Inches(2.55), Inches(2.95), Inches(2.45), phase[:260], size=12, color="334155")
        if idx < len(phases) - 1:
            _add_shape(slide, MSO_SHAPE.CHEVRON, x + Inches(3.38), Inches(3.42), Inches(0.48), Inches(0.44), fill="E5E7EB", line="E5E7EB")
    _add_slide_footer(slide)
    return True


def _wants_visual_slide(title: str, content: str = "") -> bool:
    text = f"{title}\n{content}".lower()
    keywords = (
        "roadmap", "horizon", "phase", "blueprint", "capability", "maturity",
        "heatmap", "root cause", "operating model", "use-case", "use case",
        "portfolio", "prioritization", "investment", "funding", "kpi",
        "risk", "mitigation", "90-day", "90 day", "action plan",
        "路线", "路线图", "阶段", "里程碑", "蓝图", "能力", "成熟度", "路径", "规划",
        "组合", "优先级", "投资", "资金", "指标", "风险", "缓释", "行动计划", "治理",
    )
    return any(keyword in text for keyword in keywords)


def _render_visual_slide(slide, title: str, content: str, slide_number: int, visual_kind: str = "", slide_data: dict | None = None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    _clear_generated_text_shapes(slide)
    bullets = _split_bullets(content, limit=6)
    used_template = (
        _set_title_named_or_placeholder_text(slide, "aria_slide_title", title)
        and _set_body_named_or_placeholder_text(slide, "aria_slide_body", "\n".join(f"- {bullet}" for bullet in bullets[:4]))
    )
    _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, content), role="content")
    _push_body_below_lead(slide)
    visual_bounds = _bounds_by_name_or_placeholder(slide, "aria_visual_area")
    if not used_template or visual_bounds is None:
        template_frame_used = _set_title_named_or_placeholder_text(slide, "aria_slide_title", title)
        if not template_frame_used:
            _clear_text_shapes(slide)
            _add_slide_header(slide, title, slide_number)
        used_template = used_template or template_frame_used
        _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, content), role="content")
        visual_bounds = (Inches(7.4), Inches(1.45), Inches(4.7), Inches(4.8))
        _add_textbox(slide, Inches(0.85), Inches(1.45), Inches(5.95), Inches(4.9), "\n".join(f"- {bullet}" for bullet in bullets), size=13, color="334155")

    x, y, w, h = visual_bounds
    lower_title = f"{visual_kind} {title}".lower()

    if any(token in lower_title for token in ("matrix", "heatmap", "maturity", "prioritization", "portfolio", "矩阵", "成熟度", "优先级", "组合")):
        labels = ["高价值快赢", "基础能力", "规模复制", "暂缓推进"]
        positions = [
            (0.06, 0.08, "DBEAFE", "1D4ED8"),
            (0.52, 0.08, "ECFDF5", "047857"),
            (0.06, 0.52, "FFF7ED", "C2410C"),
            (0.52, 0.52, "F8FAFC", "475569"),
        ]
        cell_w = w * 0.42
        cell_h = h * 0.35
        for idx, (left_ratio, top_ratio, fill, accent) in enumerate(positions):
            left = x + w * left_ratio
            top = y + h * top_ratio
            _add_card(slide, left, top, cell_w, cell_h, fill=fill, line=accent)
            _add_textbox(slide, left + Inches(0.14), top + Inches(0.12), cell_w - Inches(0.28), Inches(0.24), labels[idx], size=9, bold=True, color=accent)
            if idx < len(bullets):
                _add_textbox(slide, left + Inches(0.14), top + Inches(0.42), cell_w - Inches(0.28), cell_h - Inches(0.52), bullets[idx][:120], size=7, color="334155")
        if not used_template:
            _add_slide_footer(slide)
        return

    if any(token in lower_title for token in ("kpi", "investment", "funding", "business case", "指标", "投资", "资金", "商业价值")):
        metrics = [
            ("价值", "收入 / 利润 / 成本影响", "1D4ED8"),
            ("采用", "用户 / 流程覆盖率", "047857"),
            ("交付", "里程碑 / 依赖项", "C2410C"),
            ("风险", "控制点 / 缓释动作", "7C3AED"),
        ]
        card_w = w / 2 - Inches(0.12)
        card_h = h / 2 - Inches(0.14)
        for idx, (label, default, accent) in enumerate(metrics):
            left = x + (idx % 2) * (card_w + Inches(0.24))
            top = y + (idx // 2) * (card_h + Inches(0.28))
            _add_card(slide, left, top, card_w, card_h, fill="FFFFFF", line="D7DEE8")
            _add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, card_w, Inches(0.08), fill=accent, line=accent)
            _add_textbox(slide, left + Inches(0.16), top + Inches(0.22), card_w - Inches(0.32), Inches(0.26), label, size=10, bold=True, color=accent)
            metric_text = bullets[idx] if idx < len(bullets) else default
            _add_textbox(slide, left + Inches(0.16), top + Inches(0.58), card_w - Inches(0.32), card_h - Inches(0.68), metric_text[:135], size=8, color="334155")
        if not used_template:
            _add_slide_footer(slide)
        return

    if any(token in lower_title for token in ("risk", "mitigation", "风险", "缓释")):
        colors = [("FEF2F2", "DC2626"), ("FFF7ED", "EA580C"), ("F8FAFC", "475569")]
        for idx, (fill, accent) in enumerate(colors):
            top = y + idx * (h / 3)
            _add_card(slide, x, top, w, h / 3 - Inches(0.16), fill=fill, line=accent)
            _add_textbox(slide, x + Inches(0.16), top + Inches(0.12), Inches(0.6), Inches(0.28), f"R{idx + 1}", size=9, bold=True, color=accent)
            risk_text = bullets[idx] if idx < len(bullets) else "明确责任人、缓释动作和监控节奏。"
            _add_textbox(slide, x + Inches(0.86), top + Inches(0.1), w - Inches(1.05), h / 3 - Inches(0.34), risk_text[:170], size=8, color="334155")
        if not used_template:
            _add_slide_footer(slide)
        return

    phases = bullets[:3] or ["阶段一：夯实基础", "阶段二：规模复制", "阶段三：领先优化"]
    colors = [("EFF6FF", "1D4ED8"), ("ECFDF5", "047857"), ("FFF7ED", "C2410C")]
    card_h = h / 3 - Inches(0.16)
    for idx, phase in enumerate(phases[:3]):
        fill, accent = colors[idx]
        top = y + idx * (card_h + Inches(0.22))
        _add_card(slide, x, top, w, card_h, fill=fill, line=accent)
        _add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.18), top + Inches(0.16), Inches(0.38), Inches(0.38), fill=accent, line=accent)
        _add_textbox(slide, x + Inches(0.31), top + Inches(0.23), Inches(0.14), Inches(0.12), str(idx + 1), size=8, bold=True, color="FFFFFF")
        label = ["夯实基础", "规模复制", "领先优化"][idx]
        _add_textbox(slide, x + Inches(0.68), top + Inches(0.12), w - Inches(0.9), Inches(0.26), label, size=11, bold=True, color=accent)
        _add_textbox(slide, x + Inches(0.68), top + Inches(0.45), w - Inches(0.9), card_h - Inches(0.58), phase[:180], size=9, color="334155")

    if not used_template:
        _add_slide_footer(slide)


def _digital_strategy_layout_key(title: str, slide_type: str = "") -> str:
    text = f"{slide_type} {title}".lower()
    rules = [
        ("executive_summary", ("执行摘要", "executive summary", "executive answer")),
        ("maturity_heatmap", ("成熟度", "热力图", "maturity", "heatmap", "diagnosis")),
        ("root_cause", ("根因", "痛点", "root cause", "pain point")),
        ("target_blueprint", ("目标蓝图", "愿景", "target state", "vision")),
        ("capability_blueprint", ("能力蓝图", "capability blueprint")),
        ("operating_model", ("运营模式", "operating model", "治理与运营")),
        ("portfolio_matrix", ("场景组合", "用例组合", "use-case", "use case", "portfolio")),
        ("prioritization_matrix", ("优先级", "prioritization", "priority", "matrix")),
        ("roadmap", ("路线图", "三阶段", "roadmap", "horizon", "milestone")),
        ("investment_kpi", ("投资", "资金", "kpi", "指标", "investment", "funding", "business case")),
        ("risk_register", ("风险", "缓释", "risk", "mitigation")),
        ("action_plan", ("90", "行动计划", "下一步", "next step", "action plan")),
    ]
    for key, tokens in rules:
        if any(token in text for token in tokens):
            return key
    if slide_type in {"matrix", "kpi", "risk", "roadmap", "next_steps"}:
        return {
            "matrix": "prioritization_matrix",
            "kpi": "investment_kpi",
            "risk": "risk_register",
            "roadmap": "roadmap",
            "next_steps": "action_plan",
        }[slide_type]
    return ""


def _resolve_digital_strategy_layout(slide_data: dict) -> str:
    explicit = str(
        slide_data.get("layout_key")
        or slide_data.get("visualization_type")
        or slide_data.get("visualization")
        or ""
    ).strip().lower()
    aliases = {
        "executive": "executive_summary",
        "summary": "executive_summary",
        "maturity": "maturity_heatmap",
        "heatmap": "maturity_heatmap",
        "maturity_radar": "maturity_heatmap",
        "gap_matrix": "prioritization_matrix",
        "matrix": "prioritization_matrix",
        "portfolio": "portfolio_matrix",
        "use_case_portfolio": "portfolio_matrix",
        "capability": "capability_blueprint",
        "blueprint": "capability_blueprint",
        "tom": "operating_model",
        "operating": "operating_model",
        "investment": "investment_kpi",
        "kpi_dashboard": "investment_kpi",
        "risk": "risk_register",
        "risk_heatmap": "risk_register",
        "next_steps": "action_plan",
        "action": "action_plan",
        "current_target": "current_target",
        "current_vs_target": "current_target",
        "as_is_to_be": "current_target",
    }
    if explicit:
        return aliases.get(explicit, explicit)
    resolved = _digital_strategy_layout_key(str(slide_data.get("title") or ""), str(slide_data.get("type") or ""))
    if resolved:
        return resolved
    combined_text = " ".join(
        str(slide_data.get(key) or "")
        for key in ("title", "content", "left_content", "right_content")
    ).lower()
    if (
        ("current state" in combined_text and "target state" in combined_text)
        or ("as-is" in combined_text and "to-be" in combined_text)
        or ("现状" in combined_text and "目标" in combined_text)
    ):
        return "current_target"
    if any(token in combined_text for token in ("why now", "market shift", "competitive shift", "为什么现在", "市场变化", "竞争变化", "机会窗口")):
        return "strategic_context"
    if any(token in combined_text for token in ("customer journey", "journey", "客户旅程", "旅程", "获客", "留存")):
        return "customer_journey"
    if any(token in combined_text for token in ("initiative", "milestone", "year 1", "year 2", "year 3", "举措", "里程碑", "第一年", "第二年", "第三年")):
        return "initiative_milestones"
    return ""


def _combined_slide_content(slide_data: dict) -> str:
    parts = [
        str(slide_data.get("insight") or ""),
        str(slide_data.get("content") or ""),
        str(slide_data.get("left_content") or ""),
        str(slide_data.get("right_content") or ""),
        "\n".join(str(item) for item in slide_data.get("data_points", []) if str(item).strip())
        if isinstance(slide_data.get("data_points"), list)
        else "",
        "\n".join(str(item) for item in slide_data.get("management_implications", []) if str(item).strip())
        if isinstance(slide_data.get("management_implications"), list)
        else "",
    ]
    return "\n".join(part for part in parts if part.strip())


def _prepare_strategy_canvas(slide, title: str, body: str, slide_number: int, full_canvas: bool = False, slide_data: dict | None = None):
    from pptx.util import Inches

    _clear_generated_text_shapes(slide)
    bullets = _split_bullets(body, limit=8)
    title_set = _set_title_named_or_placeholder_text(slide, "aria_slide_title", title)
    _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, body), role="content")
    if full_canvas:
        _remove_shape_by_name(slide, "aria_visual_area")
        _set_body_named_or_placeholder_text(slide, "aria_slide_body", "")
        used_template = title_set
        visual_bounds = (Inches(1.02), Inches(1.38), Inches(11.28), Inches(5.22))
    else:
        used_template = (
            title_set
            and _set_body_named_or_placeholder_text(slide, "aria_slide_body", "\n".join(f"- {bullet}" for bullet in bullets[:4]))
        )
        _push_body_below_lead(slide)
        visual_bounds = _bounds_by_name_or_placeholder(slide, "aria_visual_area")
    if not used_template or visual_bounds is None:
        _clear_text_shapes(slide)
        _add_slide_header(slide, title, slide_number)
        _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, body), role="content")
        if full_canvas:
            visual_bounds = (Inches(0.92), Inches(1.58), Inches(11.48), Inches(4.98))
        else:
            visual_bounds = (Inches(6.9), Inches(1.55), Inches(5.6), Inches(5.05))
            _add_textbox(slide, Inches(0.82), Inches(1.55), Inches(5.65), Inches(4.9), "\n".join(f"- {bullet}" for bullet in bullets[:6]), size=13, color="334155")
    return bullets, visual_bounds, used_template


def _render_strategy_heatmap(slide, x, y, w, h, bullets: list[str]):
    from pptx.util import Inches

    dimensions = ["战略", "客户", "运营", "组织", "数据", "技术"]
    levels = ["L1", "L2", "L3", "L4", "L5"]
    colors = ["FEE2E2", "FFEDD5", "FEF3C7", "DCFCE7", "DBEAFE"]
    _add_textbox(slide, x, y - Inches(0.55), w, Inches(0.22), "成熟度等级：从机会式试点到可管理、可复制的企业能力", size=8, bold=True, color="1E3A8A")
    cell_w = w / 5
    row_h = (h - Inches(0.55)) / 6
    for row, dim in enumerate(dimensions):
        top = y + Inches(0.28) + row * row_h
        _add_textbox(slide, x - Inches(0.92), top + Inches(0.08), Inches(0.82), Inches(0.25), dim, size=8, bold=True, color="334155")
        active = [2, 1, 2, 1, 1, 2][row]
        for col, level in enumerate(levels):
            left = x + col * cell_w
            fill = colors[col] if col <= active else "F8FAFC"
            _add_card(slide, left, top, cell_w - Inches(0.06), row_h - Inches(0.07), fill=fill, line="E2E8F0")
            if row == 0:
                _add_textbox(slide, left + Inches(0.05), y - Inches(0.18), cell_w - Inches(0.1), Inches(0.18), level, size=7, bold=True, color="64748B")
            if col == active:
                _add_textbox(slide, left + Inches(0.05), top + Inches(0.24), cell_w - Inches(0.1), Inches(0.16), "当前", size=6, bold=True, color="0F172A")
    if bullets:
        _add_card(slide, x, y + h - Inches(0.2), w, Inches(0.52), fill="EFF6FF", line="BFDBFE")
        _add_textbox(slide, x + Inches(0.15), y + h - Inches(0.08), w - Inches(0.3), Inches(0.24), bullets[0][:110], size=8, bold=True, color="1E3A8A")


def _render_strategy_executive_summary(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    labels = ["战略判断", "价值目标", "优先动作", "高层决策"]
    accents = ["1D4ED8", "047857", "C2410C", "7C3AED"]
    for idx, label in enumerate(labels):
        left = x + (idx % 2) * (w * 0.52)
        top = y + (idx // 2) * (h * 0.52)
        _add_card(slide, left, top, w * 0.46, h * 0.42, fill="FFFFFF", line="D7DEE8")
        _add_shape(slide, MSO_SHAPE.OVAL, left + Inches(0.16), top + Inches(0.16), Inches(0.34), Inches(0.34), fill=accents[idx], line=accents[idx])
        _add_textbox(slide, left + Inches(0.6), top + Inches(0.16), w * 0.34, Inches(0.25), label, size=9, bold=True, color=accents[idx])
        text = bullets[idx] if idx < len(bullets) else "明确管理层需要确认的取舍、范围和投入节奏"
        _add_textbox(slide, left + Inches(0.16), top + Inches(0.58), w * 0.40, h * 0.20, text[:125], size=7, color="334155")


def _render_strategy_context(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    labels = ["市场压力", "竞争压力", "内部约束", "机会窗口"]
    colors = [("EFF6FF", "1D4ED8"), ("ECFDF5", "047857"), ("FFF7ED", "C2410C"), ("F5F3FF", "7C3AED")]
    center_x = x + w * 0.5
    center_y = y + h * 0.46
    _add_shape(slide, MSO_SHAPE.OVAL, center_x - Inches(0.84), center_y - Inches(0.62), Inches(1.68), Inches(1.24), fill="0F172A", line="0F172A")
    _add_textbox(slide, center_x - Inches(0.62), center_y - Inches(0.22), Inches(1.24), Inches(0.32), "现在必须\n加速转型", size=10, bold=True, color="FFFFFF")
    positions = [(0.02, 0.04), (0.62, 0.04), (0.02, 0.66), (0.62, 0.66)]
    for idx, label in enumerate(labels):
        fill, accent = colors[idx]
        left = x + w * positions[idx][0]
        top = y + h * positions[idx][1]
        card_w = w * 0.34
        _add_card(slide, left, top, card_w, Inches(1.22), fill=fill, line=accent)
        _add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, card_w, Inches(0.07), fill=accent, line=accent)
        _add_textbox(slide, left + Inches(0.16), top + Inches(0.18), card_w - Inches(0.32), Inches(0.28), label, size=10, bold=True, color=accent)
        text = bullets[idx] if idx < len(bullets) else "明确变化对增长、效率、风险和客户体验的管理影响"
        _add_textbox(slide, left + Inches(0.16), top + Inches(0.55), card_w - Inches(0.32), Inches(0.46), text[:120], size=7, color="334155")
        _add_shape(slide, MSO_SHAPE.LINE_INVERSE, left + card_w * 0.5, top + Inches(1.22), center_x, center_y, fill="CBD5E1", line="CBD5E1")


def _render_strategy_root_cause(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    labels = ["业务症状", "结构根因", "管理动作"]
    fills = ["EFF6FF", "FFF7ED", "ECFDF5"]
    accents = ["2563EB", "C2410C", "047857"]
    for idx, label in enumerate(labels):
        top = y + idx * (h / 3)
        width = w - Inches(0.58 * idx)
        left = x + Inches(0.29 * idx)
        _add_card(slide, left, top, width, h / 3 - Inches(0.2), fill=fills[idx], line=accents[idx])
        _add_textbox(slide, left + Inches(0.14), top + Inches(0.12), Inches(0.34), Inches(0.28), f"{idx + 1}", size=10, bold=True, color=accents[idx])
        _add_textbox(slide, left + Inches(0.52), top + Inches(0.12), Inches(1.0), Inches(0.26), label, size=9, bold=True, color=accents[idx])
        text = bullets[idx] if idx < len(bullets) else "补充访谈证据，明确责任边界和治理动作"
        _add_textbox(slide, left + Inches(1.55), top + Inches(0.1), width - Inches(1.75), h / 3 - Inches(0.36), text[:165], size=8, color="334155")
        if idx < 2:
            _add_shape(slide, MSO_SHAPE.DOWN_ARROW, x + w * 0.48, top + h / 3 - Inches(0.25), Inches(0.34), Inches(0.34), fill="CBD5E1", line="CBD5E1")


def _render_strategy_capability_blueprint(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    center_x = x + w * 0.5
    center_y = y + h * 0.45
    _add_shape(slide, MSO_SHAPE.OVAL, center_x - Inches(0.82), center_y - Inches(0.62), Inches(1.64), Inches(1.24), fill="1D4ED8", line="1D4ED8")
    _add_textbox(slide, center_x - Inches(0.62), center_y - Inches(0.28), Inches(1.24), Inches(0.42), "数字能力\n操作系统", size=9, bold=True, color="FFFFFF")
    labels = ["客户智能", "数字运营", "数据基础", "AI 决策", "平台架构"]
    positions = [(0.02, 0.04), (0.60, 0.04), (0.02, 0.70), (0.60, 0.70), (0.31, 0.38)]
    for idx, label in enumerate(labels):
        left = x + w * positions[idx][0]
        top = y + h * positions[idx][1]
        _add_card(slide, left, top, w * 0.38, Inches(0.96), fill="FFFFFF", line="C7D2FE")
        _add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, w * 0.38, Inches(0.06), fill="1D4ED8", line="1D4ED8")
        _add_textbox(slide, left + Inches(0.12), top + Inches(0.14), w * 0.34, Inches(0.22), label, size=8, bold=True, color="1D4ED8")
        detail = bullets[idx] if idx < len(bullets) else "定义能力、场景、数据和指标"
        _add_textbox(slide, left + Inches(0.12), top + Inches(0.42), w * 0.34, Inches(0.38), detail[:92], size=6, color="475569")


def _render_strategy_customer_journey(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    stages = ["获客", "销售", "上线", "服务", "留存"]
    accents = ["1D4ED8", "047857", "C2410C", "7C3AED", "0F766E"]
    step_w = w / len(stages) - Inches(0.08)
    top = y + Inches(0.34)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, x + Inches(0.3), top + Inches(0.45), w - Inches(0.6), Inches(0.06), fill="CBD5E1", line="CBD5E1")
    for idx, stage in enumerate(stages):
        left = x + idx * (step_w + Inches(0.1))
        accent = accents[idx]
        _add_shape(slide, MSO_SHAPE.OVAL, left + step_w / 2 - Inches(0.24), top + Inches(0.24), Inches(0.48), Inches(0.48), fill=accent, line=accent)
        _add_textbox(slide, left + step_w / 2 - Inches(0.08), top + Inches(0.35), Inches(0.16), Inches(0.12), str(idx + 1), size=7, bold=True, color="FFFFFF")
        _add_card(slide, left, top + Inches(1.0), step_w, h - Inches(1.2), fill="FFFFFF", line="D7DEE8")
        _add_textbox(slide, left + Inches(0.14), top + Inches(1.18), step_w - Inches(0.28), Inches(0.28), stage, size=11, bold=True, color=accent)
        text = bullets[idx] if idx < len(bullets) else "识别断点、数据需求和业务负责人"
        _add_textbox(slide, left + Inches(0.14), top + Inches(1.58), step_w - Inches(0.28), Inches(0.92), text[:115], size=7, color="334155")
    insight = bullets[5] if len(bullets) > 5 else "用旅程断点连接增长场景、数据资产、流程变化和责任人机制"
    _add_card(slide, x + Inches(0.25), y + h - Inches(0.48), w - Inches(0.5), Inches(0.42), fill="EFF6FF", line="BFDBFE")
    _add_textbox(slide, x + Inches(0.45), y + h - Inches(0.37), w - Inches(0.9), Inches(0.18), insight[:130], size=8, bold=True, color="1E3A8A")


def _render_strategy_operating_model(slide, x, y, w, h, bullets: list[str]):
    from pptx.util import Inches

    _add_card(slide, x + w * 0.18, y, w * 0.64, Inches(0.75), fill="EFF6FF", line="2563EB")
    _add_textbox(slide, x + w * 0.22, y + Inches(0.18), w * 0.56, Inches(0.28), "指导委员会：范围、资金与跨部门取舍", size=9, bold=True, color="1E3A8A")
    roles = ["业务负责人", "数据负责人", "技术平台", "转型 PMO"]
    for idx, role in enumerate(roles):
        left = x + (idx % 2) * (w * 0.52)
        top = y + Inches(1.15) + (idx // 2) * Inches(1.55)
        _add_card(slide, left, top, w * 0.46, Inches(1.12), fill="FFFFFF", line="D7DEE8")
        _add_textbox(slide, left + Inches(0.15), top + Inches(0.14), w * 0.42, Inches(0.24), role, size=9, bold=True, color="0F172A")
        detail = bullets[idx] if idx < len(bullets) else "明确价值、采用率、交付节奏和升级路径"
        _add_textbox(slide, left + Inches(0.15), top + Inches(0.45), w * 0.4, Inches(0.42), detail[:96], size=7, color="475569")


def _render_strategy_portfolio_matrix(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    _add_textbox(slide, x + w * 0.34, y - Inches(0.36), w * 0.32, Inches(0.18), "业务价值", size=7, bold=True, color="64748B")
    _add_textbox(slide, x - Inches(0.55), y + h * 0.42, Inches(0.45), Inches(0.32), "可行性", size=7, bold=True, color="64748B")
    _add_shape(slide, MSO_SHAPE.RECTANGLE, x, y + h * 0.5, w, Inches(0.01), fill="CBD5E1", line="CBD5E1")
    _add_shape(slide, MSO_SHAPE.RECTANGLE, x + w * 0.5, y, Inches(0.01), h, fill="CBD5E1", line="CBD5E1")
    labels = ["快赢场景", "基础能力", "差异化能力", "暂缓事项"]
    fills = ["DBEAFE", "ECFDF5", "FFF7ED", "F8FAFC"]
    accents = ["1D4ED8", "047857", "C2410C", "475569"]
    for idx, label in enumerate(labels):
        left = x + (idx % 2) * (w * 0.52)
        top = y + (idx // 2) * (h * 0.52)
        _add_card(slide, left, top, w * 0.46, h * 0.42, fill=fills[idx], line=accents[idx])
        _add_textbox(slide, left + Inches(0.16), top + Inches(0.12), w * 0.42, Inches(0.26), label, size=9, bold=True, color=accents[idx])
        text = bullets[idx] if idx < len(bullets) else "按价值、可行性、依赖和变革准备度排序"
        _add_textbox(slide, left + Inches(0.16), top + Inches(0.48), w * 0.4, h * 0.24, text[:120], size=7, color="334155")


def _render_strategy_roadmap(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    phases = ["夯实基础", "规模复制", "领先优化"]
    gates = ["价值样板", "平台复用", "AI 原生"]
    fills = ["EFF6FF", "ECFDF5", "FFF7ED"]
    accents = ["2563EB", "047857", "C2410C"]
    card_w = w / 3 - Inches(0.12)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.21), w - Inches(0.4), Inches(0.04), fill="CBD5E1", line="CBD5E1")
    for idx, phase in enumerate(phases):
        left = x + idx * (card_w + Inches(0.18))
        _add_card(slide, left, y + Inches(0.52), card_w, h - Inches(0.7), fill=fills[idx], line=accents[idx])
        _add_shape(slide, MSO_SHAPE.OVAL, left + card_w / 2 - Inches(0.2), y, Inches(0.4), Inches(0.4), fill=accents[idx], line=accents[idx])
        _add_textbox(slide, left + card_w / 2 - Inches(0.08), y + Inches(0.08), Inches(0.16), Inches(0.12), str(idx + 1), size=7, bold=True, color="FFFFFF")
        _add_textbox(slide, left + Inches(0.2), y + Inches(0.78), card_w - Inches(0.4), Inches(0.3), phase, size=10, bold=True, color=accents[idx])
        _add_textbox(slide, left + Inches(0.2), y + Inches(1.08), card_w - Inches(0.4), Inches(0.2), f"阶段门：{gates[idx]}", size=7, bold=True, color="64748B")
        text = bullets[idx] if idx < len(bullets) else "明确里程碑、依赖和管理闸口"
        _add_textbox(slide, left + Inches(0.2), y + Inches(1.42), card_w - Inches(0.4), h - Inches(1.75), text[:190], size=8, color="334155")
        if idx < 2:
            _add_shape(slide, MSO_SHAPE.CHEVRON, left + card_w + Inches(0.03), y + h * 0.48, Inches(0.28), Inches(0.28), fill="CBD5E1", line="CBD5E1")


def _render_strategy_initiative_milestones(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    rows = ["Year 1", "Year 2", "Year 3"]
    row_labels = ["基线与试点", "平台与复制", "AI 运营与优化"]
    accents = ["1D4ED8", "047857", "C2410C"]
    row_h = (h - Inches(0.34)) / 3
    for idx, year in enumerate(rows):
        top = y + idx * (row_h + Inches(0.1))
        accent = accents[idx]
        _add_card(slide, x, top, w, row_h, fill="FFFFFF", line="D7DEE8")
        _add_shape(slide, MSO_SHAPE.RECTANGLE, x, top, Inches(1.16), row_h, fill=accent, line=accent)
        _add_textbox(slide, x + Inches(0.16), top + Inches(0.22), Inches(0.84), Inches(0.28), year, size=11, bold=True, color="FFFFFF")
        _add_textbox(slide, x + Inches(1.38), top + Inches(0.18), Inches(2.0), Inches(0.28), row_labels[idx], size=10, bold=True, color=accent)
        source = bullets[idx] if idx < len(bullets) else "明确负责人、价值 KPI、用户群体、数据依赖和里程碑"
        chunks = [part.strip() for part in source.replace("；", ";").replace("，", ";").split(";") if part.strip()] or [source]
        lane_w = (w - Inches(3.45)) / 3
        for lane_idx in range(3):
            left = x + Inches(3.25) + lane_idx * (lane_w + Inches(0.12))
            _add_card(slide, left, top + Inches(0.18), lane_w, row_h - Inches(0.36), fill=["EFF6FF", "ECFDF5", "FFF7ED"][lane_idx], line="E2E8F0")
            text = chunks[lane_idx] if lane_idx < len(chunks) else source[:74]
            _add_textbox(slide, left + Inches(0.12), top + Inches(0.34), lane_w - Inches(0.24), row_h - Inches(0.62), text[:86], size=7, color="334155")


def _render_strategy_kpi(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    labels = ["业务价值", "采用率", "交付进度", "风险控制"]
    sample_values = ["+3-5%", "70%+", "90%", "月度"]
    accents = ["1D4ED8", "047857", "C2410C", "7C3AED"]
    for idx, label in enumerate(labels):
        left = x + (idx % 2) * (w * 0.52)
        top = y + (idx // 2) * (h * 0.52)
        _add_card(slide, left, top, w * 0.46, h * 0.42, fill="FFFFFF", line="D7DEE8")
        _add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, w * 0.46, Inches(0.06), fill=accents[idx], line=accents[idx])
        _add_textbox(slide, left + Inches(0.16), top + Inches(0.18), w * 0.4, Inches(0.28), label, size=9, bold=True, color=accents[idx])
        _add_textbox(slide, left + Inches(0.16), top + Inches(0.46), Inches(0.9), Inches(0.36), sample_values[idx], size=16, bold=True, color=accents[idx])
        text = bullets[idx] if idx < len(bullets) else "定义基线、目标、负责人和复盘节奏"
        _add_textbox(slide, left + Inches(1.18), top + Inches(0.50), w * 0.28, h * 0.22, text[:105], size=7, color="334155")


def _render_strategy_current_target(slide, x, y, w, h, slide_data: dict, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    left_text = str(slide_data.get("left_content") or "")
    right_text = str(slide_data.get("right_content") or "")
    current_points = _split_bullets(left_text, limit=5) or bullets[:3] or ["数据与流程仍较割裂", "客户旅程编排有限", "价值追踪尚未嵌入治理"]
    target_points = _split_bullets(right_text, limit=5) or bullets[3:6] or ["建立统一数据底座", "形成场景化运营自动化", "用 KPI 驱动组合治理"]
    delta_points = bullets[6:9] or bullets[:3] or ["补齐数据责任机制", "优先上线高价值场景", "建立采用率和价值复盘节奏"]

    column_w = (w - Inches(0.72)) / 3
    columns = [
        ("现状", current_points, "F8FAFC", "64748B"),
        ("关键差距", delta_points, "FFF7ED", "C2410C"),
        ("目标状态", target_points, "EFF6FF", "1D4ED8"),
    ]
    for idx, (label, points, fill, accent) in enumerate(columns):
        left = x + idx * (column_w + Inches(0.36))
        _add_card(slide, left, y, column_w, h, fill=fill, line=accent)
        _add_shape(slide, MSO_SHAPE.RECTANGLE, left, y, column_w, Inches(0.08), fill=accent, line=accent)
        _add_textbox(slide, left + Inches(0.22), y + Inches(0.22), column_w - Inches(0.44), Inches(0.34), label, size=13, bold=True, color=accent)
        for point_idx, point in enumerate(points[:5]):
            top = y + Inches(0.82) + point_idx * Inches(0.72)
            _add_shape(slide, MSO_SHAPE.OVAL, left + Inches(0.24), top + Inches(0.04), Inches(0.22), Inches(0.22), fill=accent, line=accent)
            _add_textbox(slide, left + Inches(0.58), top - Inches(0.02), column_w - Inches(0.82), Inches(0.42), point[:86], size=8, color="334155")
        if idx < 2:
            _add_shape(slide, MSO_SHAPE.CHEVRON, left + column_w + Inches(0.08), y + h * 0.46, Inches(0.22), Inches(0.34), fill="CBD5E1", line="CBD5E1")


def _render_strategy_risk(slide, x, y, w, h, bullets: list[str]):
    from pptx.util import Inches

    risk_points = bullets[:4] or ["遗留系统迁移风险", "数据口径和责任风险", "一线采用率风险", "供应商锁定风险"]
    mitigations = ["架构护栏和分阶段迁移", "指定数据负责人和质量 SLA", "变革冠军和角色化培训", "退出标准和能力转移"]
    _add_textbox(slide, x, y - Inches(0.32), w * 0.46, Inches(0.2), "关键风险", size=8, bold=True, color="991B1B")
    _add_textbox(slide, x + w * 0.52, y - Inches(0.32), w * 0.46, Inches(0.2), "缓释动作与责任", size=8, bold=True, color="065F46")
    for idx, risk in enumerate(risk_points[:4]):
        top = y + idx * (h / 4)
        _add_card(slide, x, top, w * 0.46, h / 4 - Inches(0.12), fill="FEF2F2", line="DC2626")
        _add_card(slide, x + w * 0.52, top, w * 0.46, h / 4 - Inches(0.12), fill="ECFDF5", line="047857")
        _add_textbox(slide, x + Inches(0.12), top + Inches(0.1), Inches(0.34), Inches(0.24), f"R{idx + 1}", size=8, bold=True, color="991B1B")
        _add_textbox(slide, x + Inches(0.52), top + Inches(0.1), w * 0.36, Inches(0.34), risk[:82], size=7, bold=True, color="991B1B")
        mitigation = mitigations[idx] if idx < len(mitigations) else "明确责任人和监控节奏"
        _add_textbox(slide, x + w * 0.52 + Inches(0.12), top + Inches(0.1), Inches(0.34), Inches(0.24), "✓", size=8, bold=True, color="065F46")
        _add_textbox(slide, x + w * 0.52 + Inches(0.52), top + Inches(0.1), w * 0.36, Inches(0.34), mitigation, size=7, bold=True, color="065F46")


def _render_strategy_action_plan(slide, x, y, w, h, bullets: list[str]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    steps = ["1-2 周", "3-5 周", "6-8 周", "9-12 周"]
    _add_shape(slide, MSO_SHAPE.RECTANGLE, x + Inches(0.25), y + Inches(0.28), Inches(0.05), h - Inches(0.56), fill="BFDBFE", line="BFDBFE")
    for idx, step in enumerate(steps):
        top = y + idx * (h / 4)
        _add_card(slide, x, top, w, h / 4 - Inches(0.12), fill="F8FAFC", line="D7DEE8")
        _add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.12), top + Inches(0.18), Inches(0.32), Inches(0.32), fill="1D4ED8", line="1D4ED8")
        _add_textbox(slide, x + Inches(0.52), top + Inches(0.16), Inches(0.72), Inches(0.24), step, size=8, bold=True, color="1D4ED8")
        text = bullets[idx] if idx < len(bullets) else "明确行动、负责人、输入和输出"
        _add_textbox(slide, x + Inches(1.32), top + Inches(0.1), w - Inches(1.5), h / 4 - Inches(0.3), text[:155], size=8, color="334155")


def _render_digital_strategy_layout(slide, slide_data: dict, slide_number: int, layout_key: str) -> bool:
    title = str(slide_data.get("title") or "")
    body = _combined_slide_content(slide_data)
    full_canvas_layouts = {
        "executive_summary",
        "maturity_heatmap",
        "root_cause",
        "target_blueprint",
        "capability_blueprint",
        "operating_model",
        "portfolio_matrix",
        "prioritization_matrix",
        "roadmap",
        "investment_kpi",
        "risk_register",
        "action_plan",
        "current_target",
        "strategic_context",
        "customer_journey",
        "initiative_milestones",
    }
    bullets, (x, y, w, h), used_template = _prepare_strategy_canvas(
        slide,
        title,
        body,
        slide_number,
        full_canvas=layout_key in full_canvas_layouts,
        slide_data=slide_data,
    )
    if layout_key == "executive_summary":
        _render_strategy_executive_summary(slide, x, y, w, h, bullets)
    elif layout_key == "strategic_context":
        _render_strategy_context(slide, x, y, w, h, bullets)
    elif layout_key == "target_blueprint":
        _render_strategy_capability_blueprint(slide, x, y, w, h, bullets)
    elif layout_key == "maturity_heatmap":
        _render_strategy_heatmap(slide, x, y, w, h, bullets)
    elif layout_key == "root_cause":
        _render_strategy_root_cause(slide, x, y, w, h, bullets)
    elif layout_key == "capability_blueprint":
        _render_strategy_capability_blueprint(slide, x, y, w, h, bullets)
    elif layout_key == "customer_journey":
        _render_strategy_customer_journey(slide, x, y, w, h, bullets)
    elif layout_key == "operating_model":
        _render_strategy_operating_model(slide, x, y, w, h, bullets)
    elif layout_key in {"portfolio_matrix", "prioritization_matrix"}:
        _render_strategy_portfolio_matrix(slide, x, y, w, h, bullets)
    elif layout_key == "roadmap":
        _render_strategy_roadmap(slide, x, y, w, h, bullets)
    elif layout_key == "initiative_milestones":
        _render_strategy_initiative_milestones(slide, x, y, w, h, bullets)
    elif layout_key == "investment_kpi":
        _render_strategy_kpi(slide, x, y, w, h, bullets)
    elif layout_key == "current_target":
        _render_strategy_current_target(slide, x, y, w, h, slide_data, bullets)
    elif layout_key == "risk_register":
        _render_strategy_risk(slide, x, y, w, h, bullets)
    elif layout_key == "action_plan":
        _render_strategy_action_plan(slide, x, y, w, h, bullets)
    else:
        return False
    if not used_template:
        _add_slide_footer(slide)
    return True


def _render_content_slide(slide, title: str, content: str, slide_number: int, slide_data: dict | None = None):
    from pptx.util import Inches

    _clear_generated_text_shapes(slide)
    bullets = _split_bullets(content, limit=6)
    if _add_roadmap_visual(slide, title, bullets, slide_number):
        return
    used_template = (
        _set_title_named_or_placeholder_text(slide, "aria_slide_title", title)
        and _set_body_named_or_placeholder_text(slide, "aria_slide_body", "\n".join(f"- {bullet}" for bullet in bullets))
    )
    _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, content), role="content")
    _push_body_below_lead(slide)
    if not used_template:
        _clear_text_shapes(slide)
        _add_slide_header(slide, title, slide_number)
        _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, content), role="content")
    if bullets:
        hero = bullets[0]
        if used_template:
            # Template-first mode: keep the user's slide geometry intact and
            # only write into named placeholders. Extra visuals are reserved
            # for pages that explicitly expose aria_visual_area.
            return
        else:
            _add_card(slide, Inches(0.75), Inches(1.2), Inches(11.85), Inches(1.0), fill="EFF6FF", line="BFDBFE")
            _add_textbox(slide, Inches(1.0), Inches(1.38), Inches(11.35), Inches(0.55), hero, size=18, bold=True, color="1E3A8A")

    card_y = Inches(2.45)
    card_h = Inches(0.72)
    if not used_template:
        for idx, bullet in enumerate(bullets[1:6], start=1):
            y = card_y + Inches(0.82 * (idx - 1))
            _add_card(slide, Inches(0.85), y, Inches(7.85), card_h)
            _add_textbox(slide, Inches(1.05), y + Inches(0.12), Inches(0.45), Inches(0.3), f"{idx}", size=11, bold=True, color="2563EB")
            _add_textbox(slide, Inches(1.55), y + Inches(0.08), Inches(6.85), Inches(0.45), bullet, size=13, color="334155")
    if not used_template:
        _add_consulting_visual_panel(slide, Inches(9.05), Inches(2.45), Inches(3.35), Inches(4.1), bullets)
        _add_slide_footer(slide)


def _draw_generated_two_column_content(
    slide,
    title: str,
    left_content: str,
    right_content: str,
    slide_number: int,
    slide_data: dict | None = None,
    *,
    preserve_template_frame: bool = False,
):
    from pptx.util import Inches

    if preserve_template_frame:
        if not _set_title_named_or_placeholder_text(slide, "aria_slide_title", title):
            _add_slide_header(slide, title, slide_number)
    else:
        _clear_text_shapes(slide)
        _add_slide_header(slide, title, slide_number)
    _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, f"{left_content}\n{right_content}"), role="content")
    columns = [
        ("\u73b0\u72b6 / \u57fa\u7840", left_content, "F8FAFC", "2563EB", Inches(0.85)),
        ("\u76ee\u6807 / \u89c4\u6a21\u5316", right_content, "F0FDF4", "16A34A", Inches(6.85)),
    ]
    for heading, content, fill, color, x in columns:
        _add_card(slide, x, Inches(1.65), Inches(5.55), Inches(4.9), fill=fill)
        _add_textbox(slide, x + Inches(0.28), Inches(1.88), Inches(5.0), Inches(0.38), heading, size=15, bold=True, color=color)
        bullets = _split_bullets(content, limit=6)
        for idx, bullet in enumerate(bullets):
            y = Inches(2.5 + idx * 0.62)
            _add_textbox(slide, x + Inches(0.35), y, Inches(0.25), Inches(0.24), "\u2022", size=14, bold=True, color=color)
            _add_textbox(slide, x + Inches(0.65), y - Inches(0.02), Inches(4.55), Inches(0.38), bullet, size=12, color="334155")
    if not preserve_template_frame:
        _add_slide_footer(slide)


def _two_column_template_body_is_usable(slide) -> bool:
    left_body_shape = _shape_by_name_or_placeholder(slide, "aria_left_body")
    right_body_shape = _shape_by_name_or_placeholder(slide, "aria_right_body")
    return _has_usable_text_bounds(left_body_shape) and _has_usable_text_bounds(right_body_shape)


def _render_two_column_slide(slide, title: str, left_content: str, right_content: str, slide_number: int, slide_data: dict | None = None):
    from pptx.util import Inches

    _clear_generated_text_shapes(slide)
    if not _two_column_template_body_is_usable(slide):
        _remove_shape_by_name(slide, "aria_left_body")
        _remove_shape_by_name(slide, "aria_right_body")
        _draw_generated_two_column_content(
            slide,
            title,
            left_content,
            right_content,
            slide_number,
            slide_data,
            preserve_template_frame=True,
        )
        return
    title_set = _set_title_named_or_placeholder_text(slide, "aria_slide_title", title)
    exact_left_body = _shape_by_name(slide, "aria_left_body")
    exact_right_body = _shape_by_name(slide, "aria_right_body")
    has_bad_template_columns = (
        (exact_left_body is not None and not _has_usable_text_bounds(exact_left_body))
        or (exact_right_body is not None and not _has_usable_text_bounds(exact_right_body))
    )
    if has_bad_template_columns:
        _remove_shape_by_name(slide, "aria_left_body")
        _remove_shape_by_name(slide, "aria_right_body")
        _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, f"{left_content}\n{right_content}"), role="content")
        columns = [
            ("现状 / 约束", left_content, "F8FAFC", "2563EB", Inches(0.85)),
            ("目标 / 动作", right_content, "F0FDF4", "16A34A", Inches(6.85)),
        ]
        for heading, content, fill, color, x in columns:
            _add_card(slide, x, Inches(1.65), Inches(5.55), Inches(4.92), fill=fill)
            _add_textbox(slide, x + Inches(0.28), Inches(1.88), Inches(5.0), Inches(0.38), heading, size=15, bold=True, color=color)
            bullets = _split_bullets(content, limit=6)
            for idx, bullet in enumerate(bullets):
                y = Inches(2.5 + idx * 0.62)
                _add_textbox(slide, x + Inches(0.35), y, Inches(0.25), Inches(0.24), "•", size=14, bold=True, color=color)
                _add_textbox(slide, x + Inches(0.65), y - Inches(0.02), Inches(4.55), Inches(0.38), bullet, size=12, color="334155")
        return
    if (
        has_bad_template_columns
    ):
        template_columns_usable = False
    else:
        left_body_shape = _shape_by_name_or_placeholder(slide, "aria_left_body")
        right_body_shape = _shape_by_name_or_placeholder(slide, "aria_right_body")
        template_columns_usable = _has_usable_text_bounds(left_body_shape) and _has_usable_text_bounds(right_body_shape)
    used_template = False
    if template_columns_usable:
        used_template = (
            title_set
            and _set_body_named_or_placeholder_text(slide, "aria_left_body", left_content, default_size=13)
            and _set_body_named_or_placeholder_text(slide, "aria_right_body", right_content, default_size=13)
        )
        if used_template and (
            not _has_usable_text_bounds(_shape_by_name(slide, "aria_left_body"))
            or not _has_usable_text_bounds(_shape_by_name(slide, "aria_right_body"))
        ):
            used_template = False
    _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, f"{left_content}\n{right_content}"), role="content")
    _push_body_below_lead(slide, ("aria_left_body", "aria_right_body"))
    if not used_template:
        _clear_text_shapes(slide)
        _add_slide_header(slide, title, slide_number)
        _add_generated_slide_lead(slide, _slide_lead_text(slide_data, title, f"{left_content}\n{right_content}"), role="content")
    columns = [
        ("现状 / 基础", left_content, "F8FAFC", "2563EB", Inches(0.85)),
        ("目标 / 规模化", right_content, "F0FDF4", "16A34A", Inches(6.85)),
    ]
    if used_template and _two_column_template_body_is_usable(slide):
        # Keep the user's two-column template untouched apart from named text.
        return
    if used_template:
        _remove_shape_by_name(slide, "aria_left_body")
        _remove_shape_by_name(slide, "aria_right_body")
        _draw_generated_two_column_content(
            slide,
            title,
            left_content,
            right_content,
            slide_number,
            slide_data,
            preserve_template_frame=True,
        )
        return
    for heading, content, fill, color, x in columns:
        _add_card(slide, x, Inches(1.35), Inches(5.55), Inches(5.2), fill=fill)
        _add_textbox(slide, x + Inches(0.28), Inches(1.58), Inches(5.0), Inches(0.38), heading, size=15, bold=True, color=color)
        bullets = _split_bullets(content, limit=6)
        for idx, bullet in enumerate(bullets):
            y = Inches(2.22 + idx * 0.62)
            _add_textbox(slide, x + Inches(0.35), y, Inches(0.25), Inches(0.24), "•", size=14, bold=True, color=color)
            _add_textbox(slide, x + Inches(0.65), y - Inches(0.02), Inches(4.55), Inches(0.38), bullet, size=12, color="334155")
    if not used_template:
        _add_slide_footer(slide)


def _render_back_cover(slide, title: str):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    _clear_text_shapes(slide)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor.from_string("0F172A")
    background.line.fill.background()
    _add_textbox(slide, Inches(0.9), Inches(2.35), Inches(11.5), Inches(0.7), "谢谢", size=34, bold=True, color="FFFFFF")
    _add_textbox(slide, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.45), title, size=16, color="CBD5E1")
    _add_textbox(slide, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.3), "由 AriaAI 生成", size=10, color="94A3B8")


def _append_unique_items(existing: Any, additions: list[str], *, limit: int = 8) -> list[str]:
    items = [str(item).strip() for item in existing if str(item).strip()] if isinstance(existing, list) else []
    seen = {item.lower() for item in items}
    for addition in additions:
        item = str(addition or "").strip()
        if not item or item.lower() in seen:
            continue
        items.append(item)
        seen.add(item.lower())
        if len(items) >= limit:
            break
    return items


def _append_unique_bullet_lines(text: str, additions: list[str], *, limit: int = 10) -> str:
    lines = [line.rstrip() for line in str(text or "").splitlines() if line.strip()]
    normalized = {re.sub(r"^[-*\u2022]\s*", "", line).strip().lower() for line in lines}
    bullet_count = sum(1 for line in lines if re.match(r"^\s*[-*\u2022]", line))
    for addition in additions:
        item = str(addition or "").strip()
        if not item or item.lower() in normalized or bullet_count >= limit:
            continue
        lines.append(f"- {item}")
        normalized.add(item.lower())
        bullet_count += 1
    return "\n".join(lines).strip()


def _digital_strategy_depth_defaults(layout_key: str) -> tuple[list[str], list[str]]:
    defaults: dict[str, tuple[list[str], list[str]]] = {
        "executive_summary": (
            [
                "量化假设：首年聚焦 3-5 个高价值场景，每个场景设置基线、目标、采用率和财务影响口径",
                "价值口径：同时追踪增长、效率、风险、体验四类结果，避免只用项目进度衡量转型",
                "组合节奏：每月检查交付进度，每季度复盘价值兑现并调整投资组合",
            ],
            [
                "管理层需一次性确认转型范围、资金池、业务 Owner、首批试点和阶段门标准",
                "任何没有业务 KPI、采用率目标和责任人的举措不进入首批组合",
            ],
        ),
        "strategic_context": (
            [
                "外部证据：客户响应速度、服务透明度和个性化能力正在成为行业竞争的基础门槛",
                "内部证据：系统割裂、人工交接和指标口径不一致会直接削弱规模化复制速度",
                "技术窗口：AI、流程自动化和云平台降低了跨部门流程重构的实施门槛",
            ],
            [
                "高层需要判断数字化是防守型效率工程，还是增长与经营模式升级工程",
                "建议把行业压力转译成 2-3 个必须在本年度启动的业务议题",
            ],
        ),
        "maturity_heatmap": (
            [
                "评分口径：战略、客户、运营、组织、数据、技术六个维度按 L1-L5 评估",
                "典型短板：数据责任、跨部门流程、业务产品 Owner 和平台复用能力通常低于工具建设成熟度",
                "证据来源：高层访谈、流程样本、系统清单、数据质量报告、项目组合和 KPI 基线",
            ],
            [
                "优先处理会阻断多个场景复制的共性短板，而不是逐个项目补洞",
                "成熟度评分必须绑定责任部门和下一阶段改进动作",
            ],
        ),
        "root_cause": (
            [
                "症状层：响应慢、重复录入、报表口径不一致、项目价值不可追踪",
                "结构层：流程 Owner 缺失、主数据薄弱、系统点对点集成、激励与采用率脱节",
                "动作层：重画端到端流程、指定数据 Owner、建立价值复盘和例外升级机制",
            ],
            [
                "不要把所有痛点归因于系统缺失，需要区分流程、数据、组织和治理根因",
                "根因页面应直接导向治理动作和优先级，而不是停留在问题罗列",
            ],
        ),
        "customer_journey": (
            [
                "旅程证据：从获客、销售、交付、服务到续约逐段识别断点和责任边界",
                "价值泄漏：转化率下降、响应周期拉长、服务成本上升和满意度下降需要量化",
                "场景机会：下一最佳行动、客户 360、自动派单、服务预警和续约预测可作为首批候选",
            ],
            [
                "每个旅程断点都要明确业务 Owner、数据来源和流程改造动作",
                "优先选择客户可感知、财务可量化、数据可获得的场景",
            ],
        ),
        "target_blueprint": (
            [
                "目标状态：把经营决策、客户运营和流程执行连接到同一套数据与指标体系",
                "能力边界：统一主数据、指标口径、权限、安全和集成标准，保留业务创新空间",
                "北极星指标：收入提升、周期缩短、成本改善、风险下降和采用率提升",
            ],
            [
                "目标蓝图需要明确哪些能力企业统一建设，哪些能力由业务域自主迭代",
                "不要把蓝图画成技术架构图，应说明业务结果、能力和责任机制",
            ],
        ),
        "capability_blueprint": (
            [
                "能力包：客户智能、数字运营、数据基础、AI 决策、平台架构和变革运营需成体系设计",
                "复用标准：每个能力包包含流程模板、数据产品、API/模型服务、运营 SOP 和 KPI",
                "成熟路径：先支持 3-5 个场景验证，再沉淀为可复制的企业能力目录",
            ],
            [
                "平台团队负责共性能力，业务 Owner 负责价值场景和采用率",
                "能力蓝图必须说明依赖关系，避免场景先行但底座缺失",
            ],
        ),
        "operating_model": (
            [
                "角色设计：业务 Owner 管价值和采用，数据 Owner 管口径和质量，技术团队管平台和安全",
                "节奏设计：周度项目推进、月度 PMO 组合看板、季度高层价值复盘",
                "决策设计：范围、资金、跨部门取舍和风险升级进入指导委员会",
            ],
            [
                "需要把 RACI 写进治理机制，避免所有问题最终回到 IT 部门",
                "采用率和价值兑现应成为业务部门 KPI，而不是项目团队自评指标",
            ],
        ),
        "portfolio_matrix": (
            [
                "组合原则：同时配置快赢、基础能力、战略差异化和暂缓事项，避免只追逐低价值自动化",
                "排序维度：价值池、可行性、数据可得性、依赖复杂度、赞助强度和变革准备度",
                "候选场景：增长、效率、风险、员工赋能四类场景都应进入组合池",
            ],
            [
                "首批场景不宜过多，建议 3-5 个足以验证价值和机制",
                "基础能力类举措即使短期收益不显著，也要用依赖关系证明投资必要性",
            ],
        ),
        "prioritization_matrix": (
            [
                "快赢标准：90-180 天可见、业务赞助明确、数据可获得、流程边界可控",
                "基础标准：能释放多个场景，解决数据、平台、集成或治理共性瓶颈",
                "暂缓标准：价值不清、数据缺失、缺少 Owner 或变革阻力高",
            ],
            [
                "排序结果要进入资金分配和季度发布计划，而不是只停留在讨论页",
                "每个优先级都要有进入、退出和复盘标准",
            ],
        ),
        "roadmap": (
            [
                "阶段一：夯实数据和治理基础，启动首批试点并建立价值看板",
                "阶段二：复制验证场景，平台化共性组件，扩展到多业务单元",
                "阶段三：形成 AI 原生运营、生态协同和持续优化机制",
            ],
            [
                "阶段门必须绑定采用率、业务 KPI 和数据质量，不只看项目交付完成",
                "路线图需要体现依赖关系，避免底座未稳就大规模复制",
            ],
        ),
        "initiative_milestones": (
            [
                "举措字段：每个举措至少包含价值 KPI、Owner、用户群、数据依赖、里程碑和风险",
                "里程碑：用月度发布和季度价值复盘管理，不用一次性大项目交付逻辑",
                "证据包：流程图、数据字典、系统接口清单和用户采用数据应持续沉淀",
            ],
            [
                "PMO 要有权停止低采用率或价值不达标的举措",
                "财务和业务共同确认收益口径，防止价值重复计算",
            ],
        ),
        "investment_kpi": (
            [
                "投资结构：技术、数据、人才/变革、生态伙伴四类预算必须同时覆盖",
                "测算口径：设置保守、基准、进取三档收益假设，并标注验证数据来源",
                "KPI 闭环：领先指标看采用率、数据质量和流程周期，滞后指标看收入、成本和风险",
            ],
            [
                "资金采用阶段门机制，下一阶段资金释放取决于价值证据和采用率",
                "投资页必须说明不投的风险，例如效率差距扩大、客户流失或合规压力上升",
            ],
        ),
        "risk_register": (
            [
                "遗留风险：定制化系统、停机窗口、接口债务和供应商锁定需要提前识别",
                "数据风险：口径不一致、质量 SLA 缺失、权限过宽和隐私合规不可后置",
                "采用风险：一线低使用率、培训疲劳和激励错配会直接影响价值兑现",
            ],
            [
                "每项风险都要配触发条件、责任人、缓释动作和监控节奏",
                "高风险依赖必须进入指导委员会升级清单",
            ],
        ),
        "action_plan": (
            [
                "前 30 天：确认范围、赞助人、基线假设和访谈/数据收集计划",
                "第 31-60 天：完成成熟度诊断、场景池、价值测算和首批试点设计",
                "第 61-90 天：定稿蓝图、路线图、资金机制、治理 RACI 和启动材料",
            ],
            [
                "90 天计划的输出物必须可进入立项和预算流程",
                "每两周同步一次事实发现、风险和需要高层拍板的问题",
            ],
        ),
        "current_target": (
            [
                "对比维度：数据、流程、系统、组织、治理和 KPI 六个维度都要有前后变化",
                "差距判断：标明哪些差距是短期补齐，哪些需要进入年度能力建设",
                "验证标准：目标状态必须对应可观测指标和责任人",
            ],
            [
                "管理层需要确认目标状态的投入强度和组织承载能力",
                "当前到目标的差距应直接转化为路线图举措",
            ],
        ),
    }
    return defaults.get(layout_key, ([], []))


def _add_digital_strategy_depth(enriched: dict, layout_key: str) -> None:
    if str(enriched.get("type") or "") in {"title", "section"}:
        return
    data_points, implications = _digital_strategy_depth_defaults(layout_key)
    if data_points:
        enriched["data_points"] = _append_unique_items(enriched.get("data_points", []), data_points, limit=8)
    if implications:
        enriched["management_implications"] = _append_unique_items(
            enriched.get("management_implications", []),
            implications,
            limit=6,
        )
    if not layout_key and enriched.get("content"):
        enriched["content"] = _append_unique_bullet_lines(
            str(enriched.get("content") or ""),
            [
                "关键判断：本页结论需要绑定业务价值、责任人和下一步管理动作",
                "证据要求：补充访谈事实、系统数据、KPI 基线或外部 benchmark 作为支撑",
            ],
            limit=8,
        )
    elif not layout_key and str(enriched.get("type") or "") == "two_column":
        enriched["left_content"] = _append_unique_bullet_lines(
            str(enriched.get("left_content") or ""),
            [
                "判断依据：明确哪些标准会影响多个业务单元、多个系统或多个数据口径",
                "风险提示：过度分散会造成重复建设、指标不一致和后续集成成本上升",
            ],
            limit=6,
        )
        enriched["right_content"] = _append_unique_bullet_lines(
            str(enriched.get("right_content") or ""),
            [
                "落地动作：为业务创新设置试点边界、复盘节奏和可复制标准",
                "管理要求：统一底座由集团/平台负责，场景创新由业务 Owner 负责价值和采用率",
            ],
            limit=6,
        )


def _enrich_digital_strategy_slide(slide: dict, page_number: int | None = None) -> dict:
    enriched = dict(slide)
    slide_type = str(enriched.get("type") or "content")
    layout_key = _resolve_digital_strategy_layout(enriched)
    if layout_key:
        enriched.setdefault("layout_key", layout_key)
    if slide_type in {"title", "section"}:
        enriched.setdefault("page_rhythm", "breathing")
    elif layout_key in {"maturity_heatmap", "prioritization_matrix", "roadmap", "investment_kpi", "risk_register"}:
        enriched.setdefault("page_rhythm", "dense")
    elif layout_key in {"executive_summary", "target_blueprint", "capability_blueprint", "operating_model", "portfolio_matrix"}:
        enriched.setdefault("page_rhythm", "anchor")
    else:
        enriched.setdefault("page_rhythm", "dense")
    if layout_key and not enriched.get("visualization_type"):
        enriched["visualization_type"] = layout_key
    if page_number is not None:
        enriched.setdefault("page_number", page_number)
    _add_digital_strategy_depth(enriched, layout_key)
    return enriched


def _digital_strategy_slide_density(slide: dict) -> int:
    text = _combined_slide_content(slide)
    bullets = _split_bullets(text, limit=12)
    return len(text.strip()) + len(bullets) * 25


def _is_sparse_digital_strategy_deck(slides: list[dict]) -> bool:
    if not slides:
        return True
    business_slides = [slide for slide in slides if str(slide.get("type") or "") not in {"title", "section"}]
    if not business_slides:
        return True
    avg_density = sum(_digital_strategy_slide_density(slide) for slide in business_slides) / len(business_slides)
    layout_keys = {str(slide.get("layout_key") or "") for slide in business_slides if slide.get("layout_key")}
    visual_types = {str(slide.get("type") or "") for slide in business_slides}
    return avg_density < 260 or len(layout_keys) < 6 or len(visual_types) < 4


def _merge_digital_strategy_plan(existing: list[dict], plan: list[dict]) -> list[dict]:
    if not _is_sparse_digital_strategy_deck(existing) and len(existing) >= 20:
        return existing

    source_slides = [
        slide for slide in existing
        if _digital_strategy_slide_density(slide) >= 160 and str(slide.get("type") or "") not in {"title", "section"}
    ]

    merged: list[dict] = []
    seen: set[str] = set()
    for slide in plan:
        title = str(slide.get("title") or "").strip()
        if not title:
            continue
        key = title.lower()
        if key in seen:
            continue
        seen.add(key)
        current = dict(slide)
        source = _find_matching_digital_strategy_source_slide(current, source_slides)
        if source:
            current = _blend_digital_strategy_source_slide(current, source)
        merged.append(_enrich_digital_strategy_slide(current, len(merged) + 1))
        if len(merged) >= 40:
            break
    return merged


def _find_matching_digital_strategy_source_slide(target: dict, sources: list[dict]) -> dict | None:
    layout_key = str(target.get("layout_key") or "")
    target_text = f"{target.get('title', '')} {layout_key}".lower()
    keyword_map = {
        "executive_summary": ("执行摘要", "战略方案", "digital-strategy", "转型论点"),
        "maturity_heatmap": ("成熟度", "现状", "评估", "rating", "l1", "l2", "l3"),
        "root_cause": ("痛点", "根因", "瓶颈", "不足", "差距"),
        "target_blueprint": ("愿景", "目标", "2029", "north", "北极星"),
        "capability_blueprint": ("能力", "cdp", "画像", "ai", "数据", "平台"),
        "operating_model": ("运营模式", "治理", "pmo", "owner", "责任"),
        "portfolio_matrix": ("场景", "用例", "营销", "复购", "自动化"),
        "prioritization_matrix": ("优先级", "p1", "p2", "价值", "可行"),
        "roadmap": ("路线图", "阶段", "year", "年", "里程碑"),
        "investment_kpi": ("投资", "it 投入", "kpi", "指标", "roi", "营收"),
        "risk_register": ("风险", "遗留", "供应商", "采用", "安全"),
        "action_plan": ("90", "行动", "周", "启动", "下一步"),
    }
    tokens = keyword_map.get(layout_key) or tuple(part for part in target_text.split() if len(part) > 2)
    best: tuple[int, dict] | None = None
    for source in sources:
        source_text = f"{source.get('title', '')}\n{_combined_slide_content(source)}".lower()
        score = sum(1 for token in tokens if token and token.lower() in source_text)
        score += min(4, _digital_strategy_slide_density(source) // 180)
        if score and (best is None or score > best[0]):
            best = (score, source)
    return best[1] if best else None


def _blend_digital_strategy_source_slide(target: dict, source: dict) -> dict:
    blended = dict(target)
    source_bullets = _split_bullets(_combined_slide_content(source), limit=5)
    if not source_bullets:
        return blended
    source_text = "\n".join(f"- {bullet[:150]}" for bullet in source_bullets[:4])
    if str(blended.get("type") or "") == "two_column":
        right = str(blended.get("right_content") or "").strip()
        blended["right_content"] = f"{right}\n\n客户事实补充\n{source_text}".strip()
    else:
        content = str(blended.get("content") or "").strip()
        blended["content"] = f"{content}\n- 客户事实：{source_bullets[0][:145]}".strip()
        if len(source_bullets) > 1:
            data_points = list(blended.get("data_points") or [])
            data_points.extend(source_bullets[1:4])
            blended["data_points"] = data_points[:6]
    return blended


def _normalize_digital_strategy_slides(slides: list[dict]) -> list[dict]:
    normalized = [_enrich_digital_strategy_slide(dict(slide), idx + 1) for idx, slide in enumerate(slides) if slide.get("title")]

    existing = {str(slide.get("title", "")).strip().lower() for slide in normalized}
    plan = [
        {
            "type": "title",
            "title": "高层共识",
            "insight": "本章节先回答管理层最关心的方向、价值、边界和决策问题，避免数字化停留在技术项目清单。",
        },
        {
            "type": "content",
            "layout_key": "executive_summary",
            "title": "执行摘要：把数字化作为业务价值组合来管理",
            "insight": "数字化转型应从零散项目升级为由业务价值、数据能力和治理节奏共同驱动的组合管理机制。",
            "content": "- 战略判断：数字化窗口期已从试点验证转向规模复制，需要用组合管理方式统一优先级、资金和责任\n- 价值目标：围绕增长、效率、风险、体验和决策速度设置可验证 KPI，而不是只追踪项目上线进度\n- 优先动作：先补齐数据、流程和平台底座，再推进高价值场景、流程重构和 AI 赋能\n- 高层决策：确认转型范围、资金边界、业务 Owner、数据 Owner、PMO 权限和第一批试点清单\n- 成功条件：每个举措都绑定业务 KPI、采用率目标、基线数据、复盘节奏和退出标准\n- 立即请示：批准 90 天诊断与首批试点设计，建立月度 PMO 和季度指导委员会价值复盘",
        },
        {
            "type": "content",
            "layout_key": "strategic_context",
            "title": "战略背景：数字化窗口期已经从试点转向规模化",
            "insight": "竞争优势正在从单点系统上线转向数据驱动的运营速度、客户响应和组织协同能力。",
            "content": "- 外部压力：客户期望更快响应、更透明服务和更个性化体验，传统人工链路难以持续满足\n- 竞争压力：数字化成熟企业通过数据沉淀、流程自动化和 AI 决策辅助持续放大经营效率\n- 内部约束：系统割裂、人工交接、指标口径不一致和项目价值不可追踪会削弱规模化复制能力\n- 技术窗口：云、低代码、流程自动化、数据湖仓和大模型让跨部门流程重构具备现实可行性\n- 管理含义：应投资可复用能力和治理机制，而不是继续建设互不连接的单点数字项目\n- 不行动风险：效率差距扩大、客户体验落后、数据合规压力上升，并形成更高遗留系统改造成本",
        },
        {
            "type": "content",
            "layout_key": "strategic_context",
            "title": "战略目标拆解：把业务战略翻译为数字化议题",
            "insight": "数字化战略必须先承接业务增长、效率、风险和体验目标，再反推能力、数据和平台建设重点。",
            "content": "- 增长目标：识别获客、转化、复购、定价、渠道协同和客户经营中的价值池\n- 效率目标：定位人工交接、重复录入、审批周期、服务响应和资源调度的改善空间\n- 风险目标：明确合规、数据安全、运营韧性、供应链风险和权限控制的数字化管理点\n- 体验目标：把客户、员工和伙伴体验拆解为旅程断点、响应时效、满意度和一次解决率\n- 能力反推：从业务议题反推所需数据产品、流程 Owner、应用能力、平台组件和组织机制\n- 管理要求：每个数字化议题都必须对应业务 KPI、负责人、验证周期和阶段门退出标准",
        },
        {
            "type": "two_column",
            "layout_key": "operating_model",
            "title": "战略设计原则：哪些事情必须统一，哪些事情允许业务创新",
            "insight": "顶层设计的关键是建立统一底座与业务灵活性之间的边界，避免一管就死或一放就乱。",
            "left_content": "- 必须统一：主数据、客户/项目/合同编码、身份权限、集成标准、安全规范和指标口径\n- 必须统一：架构原则、项目分级、投资评审、价值复盘、阶段门和供应商准入机制\n- 必须统一：跨部门流程的责任边界、例外升级、数据质量 SLA 和审计追踪要求\n- 必须统一：代码/配置资产、模型服务、知识沉淀、复用组件和平台运维标准\n- 管理动作：由指导委员会批准红线清单，并由 PMO 按季度检查执行偏差",
            "right_content": "- 允许创新：面向不同客群、区域和业务线的经营场景、触点运营和服务体验设计\n- 允许创新：一线流程体验、业务规则优化、快速试点和小规模 A/B 验证\n- 允许创新：新技术探索、生态合作和业务域内的产品路线图节奏\n- 允许创新：在统一数据和安全边界下，业务团队自主选择场景排序和运营打法\n- 管理动作：为创新设置试点边界、复盘标准和可复制条件，避免试点成功后无法推广",
        },
        {
            "type": "title",
            "title": "现状诊断",
            "insight": "本章节从成熟度、客户体验、流程、系统和数据底座识别数字化规模化的关键约束。",
        },
        {
            "type": "content",
            "layout_key": "maturity_heatmap",
            "title": "现状诊断：从六个维度识别成熟度短板",
            "insight": "成熟度诊断要同时看战略、客户、运营、组织、数据和技术，找出限制规模化复制的共性瓶颈。",
            "content": "- 战略维度：检查数字化优先级是否直接挂钩增长、成本、风险和客户体验目标\n- 客户维度：评估渠道协同、旅程编排、客户数据完整度和服务闭环能力\n- 运营维度：识别人工交接、流程瓶颈、审批周期、例外处理和自动化机会\n- 组织维度：评估决策权、产品 Owner、数字人才、变革能力和一线采用机制\n- 数据维度：检查主数据、数据质量、指标口径、权限分级、数据产品和治理节奏\n- 技术维度：评估 API 成熟度、云就绪度、架构债务、遗留系统风险和安全控制",
        },
        {
            "type": "matrix",
            "layout_key": "maturity_heatmap",
            "title": "成熟度热力图：优势基础与制约因素",
            "insight": "热力图应区分可复用优势和必须高层介入的短板，避免把所有问题都归因于系统缺失。",
            "content": "- 优势基础：已有系统、数据资产、团队能力、业务赞助和早期项目经验可作为复制起点\n- 关键短板：数据责任、流程 Owner、平台复用、采用率机制和价值追踪通常是规模化瓶颈\n- 高层介入：跨部门流程边界、资金分配、系统整合和数据治理需要指导委员会拍板\n- 证据来源：高层访谈、流程样本、系统清单、数据质量报告、项目组合和 KPI 基线\n- 输出要求：每个低成熟度领域必须绑定责任部门、改进行动、完成时间和验证指标",
        },
        {
            "type": "content",
            "layout_key": "root_cause",
            "title": "痛点根因：区分症状、结构性原因和管理动作",
            "insight": "痛点根因分析要把业务症状、结构性原因和管理动作拆开，避免只形成问题清单。",
            "content": "- 流程症状：响应慢、重复录入、审批周期长、例外处理靠人工协调，根因通常是流程 Owner 和授权边界不清\n- 数据症状：报表口径不一致、客户/项目视图不完整，根因通常是主数据薄弱和数据责任缺失\n- 技术症状：接口复杂、系统切换多、上线慢，根因通常是点对点集成、过度定制和平台复用不足\n- 组织症状：一线不用、业务不认账、项目价值不清，根因通常是激励、培训和价值复盘机制缺位\n- 管理动作：重画端到端流程、指定数据 Owner、建立价值复盘、定义阶段门和例外升级机制\n- 输出要求：每条根因都要对应一个治理动作、一个责任人和一个可观测改善指标",
        },
        {
            "type": "two_column",
            "layout_key": "customer_journey",
            "title": "客户与一线声音：把体验问题转化为可改造的流程断点",
            "insight": "客户体验问题往往不是单一触点问题，而是跨销售、交付、服务和数据链路的连续性问题。",
            "left_content": "- 客户侧断点：响应慢、信息不一致、服务过程不可视、问题重复解释和跨渠道体验割裂\n- 一线侧断点：多系统切换、重复录入、无法获得客户全貌、授权链条长和知识检索困难\n- 管理侧断点：指标滞后、问题归因困难、跨部门责任边界模糊和服务成本不可追踪\n- 证据要求：用客户投诉、服务工单、销售漏斗、满意度和流程时长验证真实影响",
            "right_content": "- 改造方向：统一客户视图、旅程状态和关键事件触发规则\n- 改造方向：建立自动派单、服务预警、下一最佳行动和续约风险提示\n- 改造方向：把体验指标、响应时效和一次解决率纳入业务负责人 KPI\n- 改造方向：用流程数据支撑持续优化，并按月复盘体验断点变化",
        },
        {
            "type": "content",
            "layout_key": "strategic_context",
            "title": "系统与数据底座诊断：找出规模化复制前必须补齐的短板",
            "insight": "如果系统、数据和集成底座不先补齐，高价值场景很容易停留在试点而无法复制。",
            "content": "- 系统层：盘点核心系统、定制化程度、接口方式、运维风险、生命周期和可替换性\n- 数据层：识别主数据、交易数据、行为数据、外部数据的口径、质量、权限和血缘问题\n- 集成层：评估 API、消息、批处理、手工导入导出和接口债务对复制速度的影响\n- 安全层：检查身份、权限、审计、脱敏、外部访问控制和大模型调用合规边界\n- 交付层：评估开发测试环境、发布流程、配置管理、供应商依赖和知识转移机制\n- 管理动作：形成系统债务清单、数据治理优先级和平台化改造路线，而不是只列技术问题",
        },
        {
            "type": "content",
            "layout_key": "strategic_context",
            "title": "行业标杆与差距启示：用外部参照校准转型目标",
            "insight": "行业标杆的价值不是照搬方案，而是帮助管理层判断目标强度、投资节奏和能力缺口是否合理。",
            "content": "- 标杆维度：客户响应速度、线上化覆盖率、自动化处理率、数据治理成熟度和 AI 场景采用率\n- 对标对象：选择 2-3 类参照，包括行业领先者、相邻行业数字化成熟企业和本地可落地标杆\n- 差距判断：区分必须追平的基础能力、可以差异化领先的经营场景和暂不投入的低价值领域\n- 投资启示：将标杆差距转化为三年投资强度、年度能力建设重点和首批试点选择\n- 风险提示：对标数据只能作为假设，需要通过访谈、系统清单和 KPI 基线进行本企业校准\n- 管理输出：形成目标强度区间、优先能力清单和需要高层确认的取舍问题",
        },
        {
            "type": "title",
            "title": "目标蓝图",
            "insight": "本章节把诊断发现转化为目标状态、能力体系、数据 AI 架构和责任清晰的运营模式。",
        },
        {
            "type": "content",
            "layout_key": "customer_journey",
            "title": "客户与增长体验缺口：优先找到价值泄漏点",
            "insight": "增长体验缺口要从客户旅程出发，找到转化、留存、服务和销售效率中的价值泄漏点。",
            "content": "- 旅程拆解：从获客、销售、签约、交付、服务到续约逐段识别客户和一线断点\n- 价值量化：用线索转化率、响应时长、续约风险、服务成本和客户满意度量化影响\n- 场景定义：优先设计精准营销、销售助手、客户 360、服务预警和下一最佳行动\n- 数据要求：每个场景明确客户主数据、行为数据、交易数据、服务工单和外部数据需求\n- 责任机制：业务 Owner 对价值和采用率负责，数据 Owner 对口径和质量负责\n- 验证节奏：首批场景用 90-180 天验证基线改善和用户采纳，而不是直接大范围铺开",
        },
        {
            "type": "content",
            "layout_key": "target_blueprint",
            "title": "数字化愿景与目标状态：形成数据驱动的运营体系",
            "insight": "目标状态应明确企业希望怎样经营、怎样决策、怎样服务客户，而不是只描述系统建设目标。",
            "content": "- 愿景表述：建设数据驱动企业，让决策、运营和客户互动能够持续优化并形成可复制能力\n- 北极星指标：收入提升、毛利改善、周期缩短、风险事件下降、客户体验提升和采用率增长\n- 运营原则：业务牵引、数据治理、平台赋能、敏捷交付、采用率衡量和价值复盘\n- 能力目标：把重点领域从机会型试点推进到可管理、可度量、可复用的企业能力\n- 设计边界：核心平台和数据标准统一，业务场景和运营打法允许差异化创新\n- 管理含义：目标状态必须进入组织职责、预算机制和 KPI，而不是停留在蓝图图形上",
        },
        {
            "type": "content",
            "layout_key": "capability_blueprint",
            "title": "能力蓝图：围绕客户、运营、数据、AI 和平台建设",
            "insight": "能力蓝图要把客户、运营、数据、AI 和平台能力组织成可复用的能力包，而不是功能清单。",
            "content": "- 客户智能能力：统一画像、客户分群、旅程触发、服务个性化和客户价值运营\n- 数字运营能力：工作流自动化、流程挖掘、异常管理、SLA 可视化和运营指挥看板\n- 数据基础能力：主数据、质量规则、数据产品、访问控制、数据 Owner 和指标口径管理\n- AI 决策能力：预测预警、推荐决策、知识检索、智能写作和辅助执行工作台\n- 平台架构能力：API 层、事件流、云服务、安全控制、模型服务和可复用集成组件\n- 变革运营能力：业务产品 Owner、训练营、采用率追踪、反馈机制和能力转移计划",
        },
        {
            "type": "content",
            "layout_key": "capability_blueprint",
            "title": "数据与 AI 架构蓝图：从报表数据走向可运营的数据产品",
            "insight": "数据能力的目标不是多做报表，而是形成可复用、可治理、可嵌入流程的数据产品和 AI 服务。",
            "content": "- 数据产品：围绕客户、项目、合同、交付、财务、风险和知识资产建立主题数据产品\n- 数据治理：明确数据 Owner、口径规则、质量 SLA、权限分级、数据血缘和生命周期\n- AI 服务：优先建设知识检索、预测预警、推荐决策、自动生成和智能问答五类能力\n- 技术架构：通过 API、事件流、数据湖仓、向量检索和模型服务支撑业务场景调用\n- 安全合规：把脱敏、审计、访问控制、模型输出可追溯和敏感信息防护嵌入架构设计\n- 运营机制：数据产品按月复盘使用量、质量问题、场景价值和下一批需求",
        },
        {
            "type": "two_column",
            "layout_key": "operating_model",
            "title": "目标运营模式：把价值、数据、技术和变革责任拆清",
            "insight": "目标运营模式的关键不是新增层级，而是把价值、数据、技术和变革责任绑定到同一套管理节奏。",
            "left_content": "- 业务产品 Owner：负责价值 KPI、采用率、需求优先级、用户反馈和场景路线图\n- 数据 Owner：负责口径、质量、权限、数据产品和生命周期管理\n- 技术平台团队：负责可复用平台、集成组件、安全护栏和模型服务能力\n- 转型 PMO：负责组合节奏、依赖协调、收益追踪、风险升级和阶段门管理",
            "right_content": "- 指导委员会：解决范围、资金、跨部门取舍、重大风险和资源冲突\n- 领域小队：通过敏捷发布交付场景，并把复用组件沉淀到能力目录\n- 变革冠军：推动一线采纳、培训、反馈和行为习惯改变\n- 财务伙伴：验证收益口径、基线、归因和阶段门资金释放条件",
        },
        {
            "type": "title",
            "title": "场景组合与能力落地",
            "insight": "本章节把目标蓝图拆成可排序、可验证、可复制的业务场景和能力建设包。",
        },
        {
            "type": "content",
            "layout_key": "portfolio_matrix",
            "title": "场景组合：平衡快赢、基础能力和战略差异化",
            "insight": "场景组合要同时覆盖快赢、基础能力和战略差异化，避免只做低价值自动化。",
            "content": "- 增长场景：线索评分、精准营销、客户流失预警、定价优化和下一最佳行动\n- 效率场景：自动报表、流程路由、需求预测、服务运营和智能知识助手\n- 风险场景：合规监测、异常识别、权限治理、项目风险预警和经营早期预警看板\n- 员工赋能：文档生成、知识检索、培训推荐、专家匹配和项目交付助手\n- 基础能力：主数据治理、数据产品、API 集成、统一身份、安全审计和模型服务\n- 组合规则：首批选择 3-5 个高价值场景，同时配置能够释放多个场景的基础能力",
        },
        {
            "type": "matrix",
            "layout_key": "prioritization_matrix",
            "title": "场景优先级逻辑：按价值、可行性、依赖和变革准备度排序",
            "insight": "优先级排序要进入资金分配和季度发布计划，而不是只停留在讨论页。",
            "content": "- 横轴建议：实施可行性，综合数据可得性、流程边界、技术复杂度和组织承载能力\n- 纵轴建议：业务价值，综合收入提升、效率改善、风险降低、体验提升和战略差异化\n- 快赢象限：90-180 天可见、业务赞助明确、数据可获得、用户群清晰的场景优先启动\n- 基础象限：能释放多个场景的主数据、平台、集成和治理能力需要进入年度投资\n- 差异化象限：数据和责任机制成熟后推进，形成客户、成本或生态优势\n- 暂缓象限：价值不清、缺少 Owner、数据缺失或变革阻力过高的事项暂不进入首批",
        },
        {
            "type": "content",
            "layout_key": "initiative_milestones",
            "title": "首批试点设计：用 3-5 个场景验证价值、数据和组织机制",
            "insight": "首批试点不应只证明工具可用，而要同时验证业务价值、数据质量、责任机制和一线采纳。",
            "content": "- 选择原则：业务价值明确、数据可获得、负责人有动力、流程边界可控、上线周期可控\n- 试点包 1：客户经营或销售效率场景，验证增长、转化、跟进效率和客户响应指标\n- 试点包 2：运营流程自动化场景，验证周期、成本、质量、返工率和员工体验改善\n- 试点包 3：风险预警或经营看板场景，验证管理透明度、异常发现速度和决策效率\n- 试点包 4：知识助手或文档生成场景，验证知识复用、交付效率和专家经验沉淀\n- 验证标准：每个试点设置基线、目标、用户群、采纳率、收益归因和复盘节奏",
        },
        {
            "type": "two_column",
            "layout_key": "capability_blueprint",
            "title": "能力建设包：把单点场景沉淀为可复制的企业能力",
            "insight": "场景成功后要沉淀为能力包，否则每个项目都会重新搭建数据、流程和技术组件。",
            "left_content": "- 业务组件：流程模板、规则库、角色职责、运营 SOP、例外处理和服务承诺\n- 数据组件：指标口径、数据模型、质量规则、权限策略、数据血缘和数据产品说明\n- 技术组件：API、自动化流程、模型服务、监控能力、日志审计和配置资产\n- 变革组件：培训材料、采纳指标、激励机制、反馈渠道和能力转移清单",
            "right_content": "- 复用方式：形成标准能力目录、调用说明、适用边界和版本管理机制\n- 复用方式：按业务域建立产品负责人、路线图、需求池和季度优先级复盘\n- 复用方式：用平台团队维护共性组件，业务团队负责场景价值和用户采用\n- 复用方式：通过季度复盘决定扩展、停用、重构或沉淀为企业级标准",
        },
        {
            "type": "title",
            "title": "路线图与投资",
            "insight": "本章节把目标蓝图转化为分阶段路径、举措组合、投资机制、收益闭环和 KPI 复盘。",
        },
        {
            "type": "two_column",
            "layout_key": "current_target",
            "title": "现状与目标状态对比",
            "insight": "当前到目标的差距应直接转化为路线图举措，而不是停留在概念性蓝图。",
            "left_content": "- 当前数据：主数据割裂、指标口径不一、数据质量责任不清、权限与审计分散\n- 当前流程：人工交接多、审批周期长、例外处理靠协调、服务链路不可视\n- 当前系统：点对点集成、重复建设、定制化高、平台复用不足、供应商依赖高\n- 当前治理：价值追踪弱、业务 Owner 不清、采用率不进 KPI、项目停启缺少阶段门",
            "right_content": "- 目标数据：统一主数据、指标口径、数据产品、质量 SLA、权限分级和审计追踪\n- 目标流程：事件触发、自动派单、流程可视、异常闭环和端到端责任清晰\n- 目标系统：平台化能力、API 集成、模型服务、安全护栏和可复用组件目录\n- 目标治理：业务价值 KPI、采用率、阶段门、季度价值复盘和 PMO 组合管理",
        },
        {
            "type": "matrix",
            "layout_key": "prioritization_matrix",
            "title": "差距优先级矩阵：先做高价值、低复杂度和可见成果",
            "insight": "差距优先级矩阵帮助管理层决定先投哪里、缓投哪里，以及哪些能力必须先于场景建设。",
            "content": "- 快赢类：高价值、低复杂度、90-180 天可见，适合建立信心和验证价值口径\n- 基础类：数据、架构、治理、身份权限和集成能力，短期价值不显著但释放多个场景\n- 差异化类：客户经营、AI 决策、生态协同等能力，在底座成熟后形成竞争优势\n- 暂缓类：价值不清、缺少赞助、数据缺失或技术实验属性过强的事项暂不投入\n- 排序口径：价值、可行性、依赖、风险、变革准备度和与战略目标的相关性\n- 管理输出：形成季度发布包、年度投资包和需要指导委员会拍板的跨部门事项",
        },
        {
            "type": "roadmap",
            "layout_key": "roadmap",
            "title": "三阶段路线图：夯实基础、规模复制、领先优化",
            "insight": "路线图必须体现依赖关系和阶段门，避免底座未稳就大规模复制。",
            "left_content": "阶段一：夯实基础\n- 完成成熟度基线和系统/数据诊断\n- 启动主数据、指标口径和权限治理\n- 选择 3-5 个高价值试点并建立价值看板\n- 建立 PMO、业务 Owner 和阶段门机制",
            "content": "阶段二：规模复制\n- 将验证场景复制到更多业务单元\n- 平台化共性组件、API 和数据产品\n- 扩展流程自动化、客户经营和风险预警\n- 建立产品化路线图和季度价值复盘",
            "right_content": "阶段三：领先优化\n- 建设 AI 原生运营和知识资产体系\n- 形成生态协同和外部数据连接能力\n- 持续优化模型、流程和运营策略\n- 将数字化能力嵌入年度经营和预算机制",
        },
        {
            "type": "content",
            "layout_key": "initiative_milestones",
            "title": "举措组合与里程碑：每个项目都要有价值、负责人和闸口",
            "insight": "每个举措都必须有价值、负责人、数据依赖、里程碑和阶段门，否则无法进入组合管理。",
            "content": "- 举措字段：每个举措定义业务 Owner、价值 KPI、用户群、数据依赖、系统依赖、预算和风险\n- 第一年：完成成熟度基线、数据治理启动、3-5 个试点、首版价值看板和组织机制试运行\n- 第二年：推进平台集成、流程规模化、业务单元复制、人才训练营和场景产品化运营\n- 第三年：建设 AI 运营模式、生态协同、知识资产体系和持续优化机制\n- 里程碑节奏：月度发布、季度价值复盘、半年度路线图刷新和年度投资评审\n- 治理检查点：按采用率和价值结果决定停止、扩大、重设或并入平台能力",
        },
        {
            "type": "content",
            "layout_key": "investment_kpi",
            "title": "收益实现路径：把业务价值拆成可跟踪的领先指标和滞后指标",
            "insight": "收益不是项目结束时才验证，而要在路线图中嵌入基线、领先指标、责任人和复盘机制。",
            "content": "- 增长收益：线索转化率、复购率、客户留存率、客单价、销售周期和渠道协同收入\n- 效率收益：人均处理量、自动化率、审批周期、返工率、服务成本和知识复用效率\n- 风险收益：异常发现时间、合规事件数量、权限违规率、关键流程中断时间和审计追踪完整度\n- 体验收益：客户满意度、一次解决率、响应时效、服务透明度和员工使用满意度\n- 领先指标：采用率、活跃用户、数据质量、流程命中率、模型准确率和一线反馈闭环\n- 财务闭环：由业务、财务和 PMO 共同确认基线、目标、归因、复盘周期和收益归属",
        },
        {
            "type": "two_column",
            "layout_key": "investment_kpi",
            "title": "投资测算与资金机制：用阶段门把投入和价值绑定",
            "insight": "投资机制要把技术、数据、人才和变革成本纳入同一预算，并用阶段门绑定价值证据。",
            "left_content": "- 技术投资：平台、集成、安全、自动化工具、模型服务和监控能力\n- 数据投资：主数据、数据治理、质量规则、数据产品和数据目录\n- 人才/变革投资：产品 Owner、数据 Owner、训练营、采用推广和能力学院\n- 生态投资：精选伙伴、试点加速、能力转移、外部数据和联合创新\n- 投资边界：区分一次性建设、持续运营、伙伴服务和内部团队能力建设",
            "right_content": "- 资金机制：阶段门资金释放与价值证据、采用率和数据质量绑定\n- 收益假设：设置保守、基准、进取三档，并说明关键驱动因素\n- 责任分工：业务和财务共同确认 KPI、基线、收益归因和价值复盘\n- 复盘机制：季度决定停止、扩大、重设或追加投资\n- 风险控制：明确不投的风险、沉没成本上限和供应商锁定防范要求",
        },
        {
            "type": "kpi",
            "layout_key": "investment_kpi",
            "title": "投资、KPI 与风险控制：建立可追踪的价值闭环",
            "insight": "KPI 仪表盘必须连接业务结果、采用率、数据质量、交付里程碑和风险控制。",
            "content": "- 投资口径：覆盖技术、数据、人才、变革、伙伴支持和持续运营成本\n- 建议结构：40% 技术平台、20% 数据治理、25% 人才/变革、15% 生态伙伴与试点加速\n- 业务 KPI：收入提升、成本改善、周期缩短、风险下降、满意度提升和决策速度改善\n- 采用 KPI：活跃用户、关键流程覆盖率、一线使用频次、功能留存率和培训完成率\n- 数据 KPI：主数据完整率、质量问题关闭率、口径一致率、权限合规率和数据产品使用量\n- 控制节奏：月度 PMO 看板、季度高层价值复盘和半年度路线图刷新",
        },
        {
            "type": "content",
            "layout_key": "investment_kpi",
            "title": "数据治理优先级：先解决会阻断多个场景复制的共性问题",
            "insight": "数据治理不应被做成长期后台工程，而要优先服务首批场景和规模复制所需的共性数据。",
            "content": "- 优先数据域：客户、项目、合同、服务、财务、风险和组织人员应优先建立主数据责任\n- 优先规则：先统一关键指标口径、编码规则、权限分级、质量规则和数据血缘\n- 优先场景：围绕首批试点所需数据建立数据产品，避免治理与业务价值脱节\n- 责任机制：每个数据域指定业务数据 Owner、数据 Steward、IT 支撑人和质量 SLA\n- 运行节奏：月度质量看板、问题闭环、例外审批和季度数据产品复盘\n- 成功标准：数据能被业务场景稳定调用，并能解释价值、风险和客户体验变化",
        },
        {
            "type": "title",
            "title": "治理与动员",
            "insight": "本章节把路线图转化为治理节奏、组织能力、风险控制和 90 天启动动作。",
        },
        {
            "type": "content",
            "layout_key": "operating_model",
            "title": "治理与运营模式：用节奏、责任和指标保证落地",
            "insight": "治理机制要用节奏、责任和指标保证落地，而不是只建立会议和汇报模板。",
            "content": "- 指导委员会：负责优先级、资金取舍、跨职能升级、风险裁决和阶段门审批\n- 转型 PMO：管理组合节奏、收益追踪、依赖协调、问题升级和季度价值复盘\n- 产品 Owner：把业务痛点转化为场景路线图、需求优先级、用户采纳和价值目标\n- 数据 Owner：治理数据定义、质量、访问、生命周期、数据产品和口径一致性\n- 技术团队：提供可复用平台、架构标准、模型服务、安全护栏和供应商能力转移\n- 财务伙伴：共同确认收益基线、归因方式、投资释放和价值兑现口径",
        },
        {
            "type": "content",
            "layout_key": "operating_model",
            "title": "组织与人才机制：把数字化能力嵌入业务岗位和管理节奏",
            "insight": "数字化组织能力的核心不是设一个部门，而是让业务、数据、技术和变革角色进入同一套管理节奏。",
            "content": "- 角色配置：业务产品负责人、数据 Owner、架构负责人、变革负责人、价值 Owner 和一线冠军\n- 能力建设：建立产品管理、数据治理、AI 应用、敏捷交付、变革管理和价值测算训练营\n- 激励机制：把采纳率、价值实现、数据质量和跨部门协作纳入业务负责人考核\n- 伙伴策略：外部伙伴负责交付加速和方法输入，内部团队负责知识沉淀和能力接管\n- 组织嵌入：关键业务域配置兼职产品 Owner 和数据 Steward，避免数字化责任只留在 IT\n- 管理节奏：月度 PMO、季度价值复盘、半年度路线图刷新和年度投资评审",
        },
        {
            "type": "risk",
            "layout_key": "risk_register",
            "title": "风险登记与缓释计划：提前管理遗留、数据、采用和伙伴风险",
            "insight": "风险登记必须配触发条件、责任人、缓释动作和监控节奏，才能进入真实治理。",
            "content": "- 遗留系统风险：隐藏定制、停机窗口、接口债务和供应商锁定；通过分阶段迁移、架构护栏和灰度切换缓释\n- 数据风险：口径不一致、质量 SLA 缺失、权限过宽和隐私合规暴露；通过数据 Owner、质量看板和权限分级缓释\n- 采用风险：一线低使用率、培训疲劳和激励错配；通过变革冠军、角色化培训和采用率 KPI 缓释\n- 伙伴风险：能力转移不足、责任边界不清和服务依赖过强；通过退出标准、知识转移和联合交付缓释\n- 价值风险：收益归因不清、重复计算或价值无法兑现；通过财务共签、基线确认和季度复盘缓释\n- 管理节奏：高风险依赖进入指导委员会升级清单，PMO 每月更新风险状态和触发条件",
        },
        {
            "type": "next_steps",
            "layout_key": "action_plan",
            "title": "90 天行动计划：把共识转成可执行启动方案",
            "insight": "90 天计划要把高层共识转成可进入立项、预算和试点启动的行动方案。",
            "content": "- 第 1-2 周：确认转型目标、范围、赞助人、指导委员会、决策机制和基线假设\n- 第 3-5 周：开展高层访谈、成熟度评估、系统清单、数据质量和流程样本诊断\n- 第 6-8 周：排序场景池、估算价值、定义首批试点、确认业务 Owner 和数据依赖\n- 第 9-11 周：设计目标运营模式、投资测算、KPI 看板、路线图依赖和治理 RACI\n- 第 12 周：在指导委员会对齐启动计划、资金边界、负责人、阶段门和风险升级事项\n- 持续动作：每两周同步事实发现、假设变化、风险问题和需要高层拍板的事项",
        },
        {
            "type": "content",
            "layout_key": "action_plan",
            "title": "立即下一步：确认决策、验证基线并启动首批试点",
            "insight": "立即下一步应聚焦决策、基线、试点和治理，不再停留在继续讨论方向。",
            "content": "- 确认高层共识：转型范围、价值目标、资金池、业务 Owner 和指导委员会决策权\n- 验证成熟度基线：用访谈、KPI 数据、系统清单、流程证据和项目组合替代未验证假设\n- 选择首批场景：明确价值负责人、用户群、数据依赖、采用率目标和 90-180 天验证标准\n- 转化路线图：把三阶段路线图拆成有资金支持的季度发布包、依赖清单和风险控制项\n- 准备动员材料：形成指导委员会汇报包、业务 Owner 职责说明、试点章程和价值看板模板\n- 锁定复盘节奏：月度 PMO 看板、季度价值复盘和半年度路线图刷新机制",
        },
        {
            "type": "content",
            "layout_key": "action_plan",
            "title": "附录：评估与访谈指南",
            "insight": "附录应提供可执行的事实收集框架，帮助团队把假设替换成客户事实。",
            "content": "- 高层访谈：战略优先级、增长压力、风险容忍度、价值目标、组织阻力和决策边界\n- 业务访谈：客户旅程、流程瓶颈、采用障碍、服务成本、KPI 基线和一线反馈\n- IT/数据访谈：架构、集成、主数据、质量问题、权限安全、模型服务和交付约束\n- 证据包：流程图、系统清单、数据字典、项目组合、预算基线、工单数据和满意度数据\n- 输出方式：用事实发现替换假设，更新成熟度评分、场景优先级、投资测算和路线图依赖\n- 下一版材料：形成访谈纪要、问题清单、假设验证表和指导委员会决策清单",
        },
    ]

    localized_plan = [_enrich_digital_strategy_slide(slide, index + 1) for index, slide in enumerate(_localize_builtin_slides(plan))]
    if normalized:
        return _merge_digital_strategy_plan(normalized, localized_plan)

    for slide in localized_plan:
        key = slide["title"].lower()
        if key in existing:
            continue
        normalized.append(_enrich_digital_strategy_slide(slide, len(normalized) + 1))
        existing.add(key)
        if len(normalized) >= 40:
            break
    return normalized


PRESENTATION_BUILDER_PRESETS: dict[str, list[dict]] = {
    "strategy": [
        {
            "type": "title",
            "title": "Strategic Direction",
        },
        {
            "type": "content",
            "title": "Executive Answer",
            "content": "- Core recommendation and decision required\n- Value ambition and expected management impact\n- Priority moves for the next planning horizon\n- Key assumptions and evidence to validate\n- Immediate leadership asks",
        },
        {
            "type": "content",
            "title": "Strategic Context",
            "content": "- Market, customer, competitor and internal pressure points\n- Why the topic matters now\n- Current constraints and opportunity window\n- Business implications if no action is taken\n- Scope boundaries for the recommendation",
        },
        {
            "type": "two_column",
            "title": "Current State vs Target State",
            "left_content": "Current state\n- Fragmented priorities\n- Manual decision loops\n- Inconsistent ownership\n- Limited value tracking",
            "right_content": "Target state\n- Clear strategic choices\n- Measurable initiatives\n- Accountable owners\n- Quarterly value review",
        },
        {
            "type": "matrix",
            "title": "Strategic Options",
            "content": "- Option A: conservative path with lower execution risk\n- Option B: focused acceleration around priority value pools\n- Option C: bold transformation with broader operating model change\n- Trade-offs across value, feasibility, risk and speed\n- Recommended option and rationale",
            "labels": ["Quick wins", "Foundations", "Differentiators", "Defer"],
        },
        {
            "type": "title",
            "title": "Roadmap and Governance",
        },
        {
            "type": "roadmap",
            "title": "Roadmap and Investment Logic",
            "left_content": "Phase 1\n- Prove value and establish foundations\n- Confirm owners and operating cadence",
            "content": "Phase 2\n- Scale validated initiatives\n- Expand across teams or business units",
            "right_content": "Phase 3\n- Institutionalize governance\n- Optimize continuous value realization",
        },
        {
            "type": "kpi",
            "title": "Governance, KPI and Next Steps",
            "content": "- Decision forum and escalation route\n- KPI dashboard linked to business outcomes\n- Owner model for delivery and adoption\n- Top risks and mitigations\n- 30/60/90-day actions",
        },
    ],
    "proposal": [
        {
            "type": "title",
            "title": "Client Need",
        },
        {
            "type": "content",
            "title": "Client Situation and Need",
            "content": "- Client context and triggering business issue\n- What is at stake for leadership\n- Current pain points and constraints\n- Why external support is valuable now\n- Desired outcomes for the engagement",
        },
        {
            "type": "content",
            "title": "Our Understanding of the Challenge",
            "content": "- Business questions to answer\n- Stakeholder priorities and concerns\n- Data, process or organizational unknowns\n- Success criteria for the work\n- Key assumptions to validate in kickoff",
        },
        {
            "type": "title",
            "title": "Proposed Solution",
        },
        {
            "type": "roadmap",
            "title": "Proposed Approach",
            "left_content": "Phase 1\n- Diagnose current state and value pools\n- Align on business questions",
            "content": "Phase 2\n- Design recommendations and target model\n- Build roadmap and business case",
            "right_content": "Phase 3\n- Align stakeholders\n- Prepare mobilization and governance",
        },
        {
            "type": "two_column",
            "title": "Scope and Deliverables",
            "left_content": "In scope\n- Executive interviews\n- Current-state analysis\n- Option design\n- Roadmap and business case\n- Steering materials",
            "right_content": "Deliverables\n- Findings summary\n- Recommendation deck\n- Initiative backlog\n- Implementation roadmap\n- Governance playbook",
        },
        {
            "type": "title",
            "title": "Mobilization",
        },
        {
            "type": "roadmap",
            "title": "Team, Timeline and Ways of Working",
            "left_content": "Team\n- Consulting team roles\n- Client participation model",
            "content": "Timeline\n- Weekly cadence\n- Steering committee rhythm",
            "right_content": "Ways of working\n- Required inputs\n- Escalation route",
        },
        {
            "type": "risk",
            "title": "Commercials, Risks and Next Steps",
            "content": "- Fee and effort assumptions\n- Optional modules and expansion paths\n- Key risks and mitigation actions\n- Immediate next meeting agenda\n- Decision required to launch",
        },
    ],
    "project-update": [
        {
            "type": "title",
            "title": "Status Snapshot",
        },
        {
            "type": "content",
            "title": "Executive Status",
            "content": "- Overall status and confidence level\n- Progress since last update\n- Key decisions or escalations required\n- Risks that may affect timeline, budget or value\n- Next milestone and owner",
        },
        {
            "type": "two_column",
            "title": "Progress vs Plan",
            "left_content": "Planned\n- Milestones\n- Workstream outputs\n- Decisions expected\n- Dependencies",
            "right_content": "Actual\n- Completed work\n- Variances\n- Open decisions\n- Dependency status",
        },
        {
            "type": "content",
            "title": "Workstream Highlights",
            "content": "- Workstream 1: progress, blocker and next action\n- Workstream 2: progress, blocker and next action\n- Workstream 3: progress, blocker and next action\n- Cross-workstream dependencies\n- Support needed from sponsors",
        },
        {
            "type": "title",
            "title": "Risks and Decisions",
        },
        {
            "type": "risk",
            "title": "Risks, Issues and Decisions",
            "content": "- Top risks ranked by impact and likelihood\n- Active issues and resolution owner\n- Decisions needed this cycle\n- Mitigation actions and deadlines\n- Items to monitor before next update",
        },
        {
            "type": "kpi",
            "title": "Value and Adoption Signals",
            "content": "- Benefits delivered or leading indicators\n- User adoption and stakeholder feedback\n- KPI movement against baseline\n- Evidence collected this period\n- Gaps to address before scaling",
        },
        {
            "type": "next_steps",
            "title": "Next Steps",
            "content": "- Actions for the next two weeks\n- Owners and deadlines\n- Upcoming workshops or steering meetings\n- Inputs needed from client or leadership\n- Decision log updates",
        },
    ],
}

PRESENTATION_BUILDER_COMMON_SLIDES: list[dict] = [
    {
        "type": "content",
        "title": "Key Assumptions",
        "content": "- Business context and audience assumptions\n- Data and evidence currently available\n- Constraints that shape the recommendation\n- Areas requiring validation\n- Implications for the next working session",
    },
    {
        "type": "content",
        "title": "Stakeholder Implications",
        "content": "- Primary stakeholders affected by the recommendation\n- Expected benefits and concerns by stakeholder group\n- Communication messages to reinforce\n- Likely objections and response logic\n- Sponsor actions required",
    },
    {
        "type": "content",
        "title": "Decision and Action Log",
        "content": "- Decisions required from leadership\n- Actions already agreed\n- Open actions and owners\n- Due dates and dependencies\n- Escalations for the next governance meeting",
    },
    {
        "type": "risk",
        "title": "Risks and Mitigations",
        "content": "- Execution risks and likely triggers\n- Stakeholder or adoption risks\n- Data, technology or operational dependencies\n- Mitigation actions and owners\n- Monitoring cadence",
    },
    {
        "type": "kpi",
        "title": "Success Metrics",
        "content": "- Business outcome KPIs\n- Adoption and usage indicators\n- Delivery milestone metrics\n- Quality and risk indicators\n- Review cadence and accountability",
    },
    {
        "type": "content",
        "title": "Appendix: Supporting Detail",
        "content": "- Source materials and evidence pack\n- Interview or workshop notes to collect\n- Additional analysis required\n- Optional modules or future work\n- Reference data for the next version",
    },
]


CHINESE_BUILTIN_SLIDES: dict[str, dict[str, str]] = {
    "Strategic Direction": {"title": "战略方向"},
    "Executive Answer": {
        "title": "管理层答案：明确建议、价值目标和决策请求",
        "content": "- 核心建议和需要管理层拍板的事项\n- 价值目标及其对经营管理的影响\n- 下一规划周期的优先动作\n- 需要验证的关键假设和证据\n- 立即需要高层支持的事项",
    },
    "Strategic Context": {
        "title": "战略背景：说明为什么现在必须行动",
        "content": "- 市场、客户、竞争和内部压力点\n- 该议题为什么在当前阶段重要\n- 现有约束和机会窗口\n- 如果不行动，对业务的影响\n- 本次建议的范围边界",
    },
    "Strategic Options": {
        "title": "战略选项矩阵：比较价值、可行性、风险和速度",
        "content": "- 选项 A：稳健路径，执行风险较低\n- 选项 B：围绕优先价值池集中加速\n- 选项 C：更大力度的运营模式变革\n- 从价值、可行性、风险和速度比较取舍\n- 推荐选项及其理由",
    },
    "Roadmap and Governance": {"title": "路线图与治理"},
    "Roadmap and Investment Logic": {
        "title": "路线图与投资逻辑：从证明价值到规模化复制",
        "content": "阶段一\n- 证明价值并建立基础能力\n- 确认负责人和运营节奏\n阶段二\n- 复制已经验证的举措\n- 扩展到团队或业务单元\n阶段三\n- 固化治理机制\n- 持续优化价值实现",
    },
    "Governance, KPI and Next Steps": {
        "title": "治理、KPI 与下一步：把决策转成行动闭环",
        "content": "- 决策论坛和升级路径\n- 连接业务结果的 KPI 看板\n- 交付和采用的负责人模型\n- 主要风险和缓释动作\n- 30/60/90 天行动安排",
    },
    "Client Need": {"title": "客户需求"},
    "Client Situation and Need": {
        "title": "客户情境与需求：明确问题、约束和期望成果",
        "content": "- 客户背景和触发议题\n- 管理层当前最关注的风险和机会\n- 现有痛点和约束条件\n- 为什么需要外部支持\n- 期望通过项目达成的结果",
    },
    "Our Understanding of the Challenge": {
        "title": "我们对挑战的理解：把业务问题拆成可回答的问题",
        "content": "- 本项目需要回答的关键业务问题\n- 主要干系人的优先级和顾虑\n- 数据、流程或组织方面的未知项\n- 项目成功标准\n- 启动阶段需要验证的关键假设",
    },
    "Proposed Solution": {"title": "建议方案"},
    "Proposed Approach": {
        "title": "建议方法：诊断、设计、对齐与动员",
        "content": "阶段一\n- 诊断现状和价值池\n- 对齐关键业务问题\n阶段二\n- 设计建议和目标模式\n- 形成路线图和商业测算\n阶段三\n- 对齐关键干系人\n- 准备动员和治理机制",
    },
    "Scope and Deliverables": {
        "title": "范围与交付物：明确边界、产出和客户投入",
        "left_content": "范围内\n- 高层访谈\n- 现状分析\n- 选项设计\n- 路线图和商业测算\n- 指导委员会材料",
        "right_content": "交付物\n- 发现总结\n- 建议方案 PPT\n- 举措清单\n- 实施路线图\n- 治理手册",
    },
    "Mobilization": {"title": "项目动员"},
    "Team, Timeline and Ways of Working": {
        "title": "团队、时间表与工作方式：确保启动后快速进入节奏",
        "content": "团队\n- 咨询团队角色\n- 客户参与机制\n时间表\n- 周度工作节奏\n- 指导委员会节奏\n工作方式\n- 所需输入\n- 升级路径",
    },
    "Commercials, Risks and Next Steps": {
        "title": "商务、风险与下一步：明确启动条件",
        "content": "- 费用和投入假设\n- 可选模块和扩展路径\n- 关键风险和缓释动作\n- 下一次会议议程\n- 启动项目所需决策",
    },
    "Status Snapshot": {"title": "状态快照"},
    "Executive Status": {
        "title": "高层状态：进展、风险和所需决策",
        "content": "- 整体状态和信心等级\n- 相比上次更新的进展\n- 需要决策或升级的事项\n- 可能影响时间、预算或价值的风险\n- 下一里程碑和负责人",
    },
    "Progress vs Plan": {
        "title": "计划与实际：识别偏差、依赖和纠偏动作",
        "left_content": "计划\n- 里程碑\n- 工作流产出\n- 预期决策\n- 关键依赖",
        "right_content": "实际\n- 已完成工作\n- 偏差说明\n- 待定决策\n- 依赖状态",
    },
    "Workstream Highlights": {
        "title": "工作流亮点：按模块呈现进展、阻碍和下一步",
        "content": "- 工作流 1：进展、阻碍和下一步\n- 工作流 2：进展、阻碍和下一步\n- 工作流 3：进展、阻碍和下一步\n- 跨工作流依赖\n- 需要赞助人支持的事项",
    },
    "Risks and Decisions": {"title": "风险与决策"},
    "Risks, Issues and Decisions": {
        "title": "风险、问题与决策：聚焦本周期必须处理的事项",
        "content": "- 按影响和可能性排序的主要风险\n- 当前问题及解决负责人\n- 本周期需要的决策\n- 缓释动作和截止时间\n- 下次更新前需要监控的事项",
    },
    "Value and Adoption Signals": {
        "title": "价值与采用信号：用领先指标判断是否可规模化",
        "content": "- 已实现收益或领先指标\n- 用户采用情况和干系人反馈\n- KPI 相比基线的变化\n- 本周期收集的证据\n- 规模化前需要补齐的缺口",
    },
    "Next Steps": {
        "title": "下一步行动：明确负责人、时间和输入",
        "content": "- 未来两周行动\n- 负责人和截止时间\n- 即将开展的工作坊或指导委员会\n- 需要客户或管理层提供的输入\n- 决策日志更新",
    },
    "Key Assumptions": {
        "title": "关键假设：说明材料当前依据和待验证项",
        "content": "- 业务背景和受众假设\n- 当前可用的数据和证据\n- 影响建议的约束条件\n- 需要验证的领域\n- 对下一次工作会的含义",
    },
    "Stakeholder Implications": {
        "title": "干系人影响：提前管理收益、顾虑和沟通口径",
        "content": "- 受建议影响的主要干系人\n- 各类干系人的预期收益和顾虑\n- 需要强化的沟通信息\n- 可能出现的反对意见和回应逻辑\n- 需要赞助人采取的行动",
    },
    "Decision and Action Log": {
        "title": "决策与行动日志：把讨论转化为可追踪事项",
        "content": "- 需要管理层决策的事项\n- 已达成一致的行动\n- 未完成行动和负责人\n- 截止日期和依赖关系\n- 下次治理会议需要升级的问题",
    },
    "Risks and Mitigations": {
        "title": "风险与缓释：建立触发条件、责任人和监控节奏",
        "content": "- 执行风险及可能触发条件\n- 干系人或采用风险\n- 数据、技术或运营依赖\n- 缓释动作和负责人\n- 监控节奏",
    },
    "Success Metrics": {
        "title": "成功指标：把业务结果、采用和交付质量连起来",
        "content": "- 业务结果 KPI\n- 采用率和使用指标\n- 交付里程碑指标\n- 质量和风险指标\n- 复盘节奏和责任机制",
    },
    "Appendix: Supporting Detail": {
        "title": "附录：支撑材料和后续分析",
        "content": "- 来源材料和证据包\n- 待收集的访谈或工作坊记录\n- 需要补充的分析\n- 可选模块或未来工作\n- 下一版材料所需参考数据",
    },
    "Executive Alignment": {"title": "高层共识"},
    "Executive Summary": {
        "title": "执行摘要：把数字化作为业务价值组合来管理",
        "content": "- 转型论点：数字化不是单纯技术换新，而是围绕增长、效率、风险和体验的价值组合\n- 价值目标：明确收入提升、成本改善、响应速度和风险控制的量化目标\n- 优先动作：先补齐数据基础，选择高价值场景试点，并建立组合治理机制\n- 高层决策：确认转型范围、资金边界、责任人机制和首批试点清单\n- 成功条件：每个举措必须绑定业务 KPI、采用率目标和明确负责人\n- 立即请示：批准诊断范围、首批试点和治理例会节奏",
    },
    "Strategic Context and Transformation Thesis": {
        "title": "战略背景：数字化窗口期已经从试点转向规模化",
        "content": "- 市场变化：客户期望更快响应、更高透明度和更个性化的服务\n- 竞争变化：数字化成熟企业通过数据和运营速度持续放大优势\n- 内部约束：系统割裂和人工流程降低执行效率，也削弱责任追踪\n- 机会窗口：AI、自动化和云平台让跨职能流程重构具备现实可行性\n- 管理含义：应投资可复用能力，而不是继续建设互不连接的单点项目",
    },
    "Current Digital Maturity Diagnosis": {
        "title": "现状诊断：从六个维度识别成熟度短板",
        "content": "- 战略：检验数字化优先级是否与增长、成本和风险目标直接挂钩\n- 客户：评估渠道协同、旅程编排和客户数据完整度\n- 运营：识别人工交接、流程瓶颈和可自动化环节\n- 组织：评估决策权、产品负责人、数字人才和变革承载能力\n- 数据/技术：评估主数据、数据治理、API 成熟度、云就绪度和遗留系统风险",
    },
    "Maturity Heatmap: Strengths vs Constraints": {
        "title": "成熟度热力图：优势基础与制约因素",
        "left_content": "- 已具备动能的强势领域\n- 可复用的系统、数据资产或团队能力\n- 具备赞助和采用条件的业务单元\n- 能快速建立信心的早期证明点",
        "right_content": "- 阻碍规模化的薄弱领域\n- 遗留系统、数据责任、人才或流程约束\n- 需要高层介入的关键决策\n- 不应只靠工具解决的能力缺口",
    },
    "Pain Point Root Causes": {
        "title": "痛点根因：区分症状、结构性原因和管理动作",
        "content": "- 流程痛点通常来自责任边界不清，而不只是缺少系统\n- 数据痛点通常来自主数据薄弱、口径不一致和责任机制不足\n- 技术痛点多来自点对点集成、过度定制和遗留平台债务\n- 采用痛点往往来自激励和培训缺口，而不只是工具可用性\n- 根因分析要把业务症状、结构性原因和管理动作拆开",
    },
    "Target Blueprint": {"title": "目标蓝图"},
    "Customer and Growth Experience Gaps": {
        "title": "客户与增长体验缺口：优先找到价值泄漏点",
        "content": "- 梳理从获客、销售、上线、服务到留存的优先客户旅程\n- 定位断点：重复录入、响应慢、渠道不一致、个性化不足\n- 用转化流失、留存风险、服务成本和满意度缺口量化影响\n- 定义增长场景：精准营销、销售效率提升和下一最佳行动\n- 将每个场景连接到数据资产、流程变化和责任人机制",
    },
    "Digital Vision and Target State": {
        "title": "数字化愿景与目标状态：形成数据驱动的运营体系",
        "content": "- 愿景：建设数据驱动企业，让决策、运营和客户互动持续优化\n- 北极星指标：收入提升、利润改善、周期缩短和风险事件下降\n- 运营原则：业务牵引、数据治理、平台赋能、采用率可衡量\n- 能力目标：把重点领域从机会型试点推进到可管理的企业能力\n- 设计边界：核心平台标准化，同时保留业务创新空间",
    },
    "Capability Blueprint": {
        "title": "能力蓝图：围绕客户、运营、数据、AI 和平台建设",
        "content": "- 客户智能：统一画像、分群、旅程触发和服务个性化\n- 数字运营：工作流自动化、流程挖掘、异常管理和 SLA 可视化\n- 数据基础：主数据、质量规则、数据产品、访问控制和责任模型\n- AI 决策支持：预测、推荐、知识检索和辅助执行\n- 平台架构：API 层、云服务、安全控制和可复用集成组件",
    },
    "Target Operating Model Blueprint": {
        "title": "目标运营模式：把价值、数据、技术和变革责任拆清",
        "left_content": "- 业务产品负责人负责价值、采用率和需求优先级\n- 数据负责人负责口径、访问、质量和生命周期\n- 技术团队提供可复用平台和安全护栏\n- 转型 PMO 管理组合节奏和收益追踪",
        "right_content": "- 指导委员会解决范围、资金和跨部门取舍\n- 领域小队通过敏捷发布交付场景\n- 变革冠军推动一线采用和培训\n- 财务验证价值实现并支持阶段门资金",
    },
    "Use-Case Portfolio": {
        "title": "场景组合：平衡快赢、基础能力和战略差异化",
        "content": "- 增长场景：线索评分、精准营销、流失预测和定价优化\n- 效率场景：自动报表、流程路由、需求计划和服务运营\n- 风险场景：合规监测、异常识别、权限治理和预警看板\n- 员工场景：知识助手、文档生成、培训推荐和专家匹配\n- 组合原则：同时配置快赢、基础使能和战略差异化场景",
    },
    "Use-Case Prioritization Logic": {
        "title": "场景优先级逻辑：按价值、可行性、依赖和变革准备度排序",
        "left_content": "- 价值池规模和收益确定性\n- 赞助强度和业务负责人准备度\n- 数据可用性和数据质量\n- 交付复杂度和依赖关系",
        "right_content": "- 先推进可见快赢，建立组织信心\n- 投资能释放多个场景的数据和平台基础\n- 在数据和责任机制成熟后推进差异化能力\n- 暂缓缺少业务赞助的低价值自动化",
    },
    "Roadmap and Investment": {"title": "路线图与投资"},
    "Current State vs Target State": {
        "title": "现状与目标状态对比",
        "left_content": "现状\n- 数据割裂，流程依赖人工\n- 客户旅程编排能力有限\n- 价值追踪尚未嵌入治理",
        "right_content": "目标状态\n- 建立客户和业务数据底座\n- 场景化运营自动化\n- 基于 KPI 的组合治理机制",
    },
    "Gap Prioritization Matrix": {
        "title": "差距优先级矩阵：先做高价值、低复杂度和可见成果",
        "content": "- 快赢：高价值、低复杂度，并能在 90-180 天内可见\n- 基础能力：数据、架构和治理等规模化前置条件\n- 差异化能力：形成客户、成本或生态优势的能力\n- 暂缓项：缺少业务赞助、价值不清的低价值自动化或技术实验\n- 决策规则：按价值、可行性、依赖和变革准备度排序",
    },
    "Three-Horizon Roadmap": {
        "title": "三阶段路线图：夯实基础、规模复制、领先优化",
        "content": "- 阶段一：夯实数据基础，启动试点，建立治理并证明价值\n- 阶段二：把验证过的场景复制到业务单元并整合平台\n- 阶段三：建设 AI 原生运营、生态协同和持续创新机制\n- 路线依赖：数据责任机制未运转前，不宜大规模复制高级分析\n- 复盘节奏：季度价值复盘，半年度路线图刷新\n- 管理闸口：只有在采用率和业务 KPI 变化可见后才进入规模化",
    },
    "Initiative Portfolio and Milestones": {
        "title": "举措组合与里程碑：每个项目都要有价值、负责人和闸口",
        "content": "- 每个举措明确负责人、价值 KPI、用户群体、数据依赖和里程碑\n- 第一年：完成成熟度基线、启动数据治理、3-5 个试点和首版价值看板\n- 第二年：推进平台集成、流程规模化、业务单元复制和人才学院\n- 第三年：建设 AI 运营模式、生态协同和持续优化机制\n- 治理检查点：按采用率和价值结果决定停止、扩大或重设举措",
    },
    "Investment Case and Funding Model": {
        "title": "投资测算与资金机制：用阶段门把投入和价值绑定",
        "left_content": "- 技术：平台、集成、安全和自动化工具\n- 数据：主数据、治理、质量和数据产品\n- 人才/变革：产品负责人、培训、采用和能力学院\n- 生态：精选伙伴、试点和能力转移",
        "right_content": "- 阶段门资金与价值证明绑定\n- 设置基础、乐观和保守收益假设\n- KPI 责任由业务和财务共同承担\n- 季度复盘决定停止、扩大或重设",
    },
    "Investment, KPI and Risk Controls": {
        "title": "投资、KPI 与风险控制：建立可追踪的价值闭环",
        "content": "- 投资范围覆盖技术、数据、人才、变革和伙伴支持\n- 建议结构：40% 技术、30% 人才/变革、20% 数据、10% 生态实验\n- KPI 看板连接业务结果、采用率、数据质量和交付里程碑\n- 关键风险：遗留复杂度、数据责任缺口、低采用率、供应商锁定和安全暴露\n- 控制节奏：月度 PMO 看板和季度高层价值复盘",
    },
    "Governance and Mobilization": {"title": "治理与动员"},
    "Governance and Operating Model": {
        "title": "治理与运营模式：用节奏、责任和指标保证落地",
        "content": "- 指导委员会负责优先级、资金取舍和跨职能升级\n- 转型 PMO 管理组合节奏、收益追踪和依赖协调\n- 产品负责人把业务痛点转化为路线图和采用计划\n- 数据负责人治理口径、质量、访问和生命周期\n- 技术团队提供可复用平台、标准和安全护栏",
    },
    "Risk Register and Mitigation Plan": {
        "title": "风险登记与缓释计划：提前管理遗留、数据、采用和伙伴风险",
        "left_content": "- 遗留风险：隐藏定制、停机和集成债务\n- 数据风险：口径不一致、责任弱和隐私暴露\n- 采用风险：一线使用率低、培训疲劳和激励不足\n- 伙伴风险：锁定、责任不清和能力转移不足",
        "right_content": "- 采用分阶段迁移和架构护栏\n- 指定数据负责人和质量 SLA\n- 建立变革冠军和分角色赋能\n- 定义伙伴退出标准和内部能力转移",
    },
    "90-Day Action Plan": {
        "title": "90 天行动计划：把共识转成可执行启动方案",
        "content": "- 第 1-2 周：确认目标、范围、赞助人、决策机制和基线假设\n- 第 3-5 周：开展高层访谈、成熟度评估和数据/平台诊断\n- 第 6-8 周：排序场景、估算收益并定义首批试点\n- 第 9-11 周：设计运营模式、投资测算、KPI 看板和路线依赖\n- 第 12 周：在指导委员会对齐启动计划、资金和负责人",
    },
    "Immediate Next Steps": {
        "title": "立即下一步：确认决策、验证基线并启动首批试点",
        "content": "- 确认高层共识和转型组合的决策权\n- 用访谈、KPI 数据、系统清单和流程证据验证成熟度基线\n- 选择首批场景，明确价值负责人和采用率目标\n- 把路线图转成有资金支持的季度发布计划\n- 准备指导委员会材料，推动范围、资金和动员审批",
    },
    "Appendix: Assessment and Interview Guide": {
        "title": "附录：评估与访谈指南",
        "content": "- 高层访谈：战略优先级、痛点、风险偏好和价值目标\n- 业务访谈：旅程摩擦、流程瓶颈、采用障碍和 KPI 基线\n- IT/数据访谈：架构、集成、数据质量、安全和交付约束\n- 证据包：流程图、系统清单、数据字典、项目组合和预算基线\n- 用调研结果替换假设，并持续打磨下一版材料",
    },
}


def _localize_builtin_slides(slides: list[dict]) -> list[dict]:
    localized: list[dict] = []
    for slide in slides:
        item = dict(slide)
        override = CHINESE_BUILTIN_SLIDES.get(str(item.get("title") or ""))
        if override:
            item.update(override)
        localized.append(item)
    return localized


def _normalize_presentation_deck_type(deck_type: str | None) -> str:
    normalized = (deck_type or "strategy").strip().lower().replace("_", "-")
    aliases = {
        "strategic": "strategy",
        "strategy-report": "strategy",
        "executive": "strategy",
        "executive-briefing": "strategy",
        "proposal-deck": "proposal",
        "client-proposal": "proposal",
        "update": "project-update",
        "project": "project-update",
        "status": "project-update",
        "project-status": "project-update",
    }
    return aliases.get(normalized, normalized if normalized in PRESENTATION_BUILDER_PRESETS else "strategy")


def _normalize_presentation_builder_slides(slides: list[dict], deck_type: str | None = None) -> list[dict]:
    normalized = [dict(slide) for slide in slides if slide.get("title")]
    preset_key = _normalize_presentation_deck_type(deck_type)
    minimum_slide_count = 12 if preset_key in {"proposal", "project-update"} else 14
    if len(normalized) >= minimum_slide_count:
        return normalized

    existing = {str(slide.get("title", "")).strip().lower() for slide in normalized}
    for slide in _localize_builtin_slides(PRESENTATION_BUILDER_PRESETS[preset_key]):
        key = slide["title"].lower()
        if key in existing:
            continue
        normalized.append(dict(slide))
        existing.add(key)
        if len(normalized) >= minimum_slide_count:
            break
    for slide in _localize_builtin_slides(PRESENTATION_BUILDER_COMMON_SLIDES):
        if len(normalized) >= minimum_slide_count:
            break
        key = slide["title"].lower()
        if key in existing:
            continue
        normalized.append(dict(slide))
        existing.add(key)
    return normalized


def _normalize_presentation_slide_type(slide_type: Any) -> str:
    normalized = str(slide_type or "content").strip().lower().replace("-", "_").replace(" ", "_")
    aliases = {
        "section": "title",
        "divider": "title",
        "section_divider": "title",
        "compare": "two_column",
        "current_target": "two_column",
        "plan_actual": "two_column",
        "timeline": "roadmap",
        "milestones": "roadmap",
        "scorecard": "kpi",
        "metrics": "kpi",
        "risks": "risk",
        "risk_mitigation": "risk",
        "nextsteps": "next_steps",
        "actions": "next_steps",
    }
    return aliases.get(normalized, normalized)


def _normalize_presentation_builder_slide_format(slides: list[dict]) -> list[dict]:
    formatted: list[dict] = []
    for slide in slides:
        item = dict(slide)
        slide_type = _normalize_presentation_slide_type(item.get("type"))
        item["type"] = slide_type
        if slide_type == "roadmap":
            roadmap_parts = [
                str(item.get("left_content") or item.get("phase_1") or "").strip(),
                str(item.get("content") or item.get("phase_2") or "").strip(),
                str(item.get("right_content") or item.get("phase_3") or "").strip(),
            ]
            item["content"] = "\n".join(part for part in roadmap_parts if part)
        if slide_type == "matrix":
            labels = item.get("labels") or item.get("card_titles") or []
            label_lines = [f"- {label}" for label in labels if str(label).strip()]
            body = str(item.get("content") or "").strip()
            item["content"] = "\n".join([*label_lines, body]).strip()
        if slide_type in {"roadmap", "matrix", "kpi", "risk", "next_steps"}:
            item.setdefault("content", item.get("content") or "")
        if item["type"] in {"content", "roadmap", "matrix", "kpi", "risk", "next_steps"} and item.get("content"):
            bullets = _split_bullets(str(item.get("content") or ""), limit=6)
            item["content"] = "\n".join(f"- {bullet}" for bullet in bullets)
        if item["type"] == "two_column":
            left = _split_bullets(str(item.get("left_content") or ""), limit=5)
            right = _split_bullets(str(item.get("right_content") or ""), limit=5)
            item["left_content"] = "\n".join(f"- {bullet}" for bullet in left)
            item["right_content"] = "\n".join(f"- {bullet}" for bullet in right)
        formatted.append(item)
    return formatted


async def generate_ppt(
    title: str,
    slides: list[dict],
    subtitle: str = "",
    template_path: str = ""
) -> dict[str, Any]:
    """Generate a PowerPoint presentation."""
    title = _clean_ppt_text(title)
    subtitle = _clean_ppt_text(subtitle)
    slides = _prepare_ppt_slide_text(slides)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx not installed. Run: pip install python-pptx"
        }

    template_path_obj = Path(template_path) if template_path else None
    using_template = bool(template_path_obj and template_path_obj.is_file())
    template_name = template_path_obj.name.lower() if using_template else ""
    is_graphic_library_template = template_name == "graphic library.pptx"
    if using_template:
        prs = Presentation(template_path)
        if is_graphic_library_template:
            aria_cover_prototype = _find_named_prototype_slide(
                prs,
                {"aria_cover_title", "aria_cover_subtitle"},
                preferred_layout_name="Cover style 1-D",
            )
            aria_section_prototype = _find_named_prototype_slide(
                prs,
                {"aria_section_title", "aria_section_number"},
                preferred_layout_name="Divider with image",
            )
            aria_content_prototype = _find_named_prototype_slide(
                prs,
                {"aria_slide_title", "aria_slide_body"},
                preferred_layout_name="ONE COLUMN TEXT",
            )
            aria_two_col_prototype = _find_named_prototype_slide(
                prs,
                {"aria_slide_title", "aria_left_body", "aria_right_body"},
                preferred_layout_name="TWO COLUMN TEXT",
            )
            aria_title_prototype = _find_named_prototype_slide(
                prs,
                {"aria_slide_title", "aria_visual_area"},
                preferred_layout_name="Title Only",
            )

            if all((aria_cover_prototype, aria_section_prototype, aria_content_prototype, aria_two_col_prototype, aria_title_prototype)):
                original_refs = list(prs.slides._sldIdLst)
                _set_template_cover_text(aria_cover_prototype, title, subtitle)

                section_counter = 0
                for slide_index, slide_data in enumerate(slides):
                    slide_type = slide_data.get("type", "content")
                    slide_title = slide_data.get("title", "")

                    if slide_type == "section":
                        section_counter += 1
                        slide = _clone_slide_from_prototype(prs, aria_section_prototype)
                        _render_section_slide(slide, slide_title, section_counter, slide_data)
                    elif slide_type == "title":
                        section_counter += 1
                        slide = _clone_slide_from_prototype(prs, aria_title_prototype)
                        _render_graphic_library_title_slide(slide, slide_title, slide_data)
                    elif slide_type == "two_column":
                        slide = _clone_slide_from_prototype(prs, aria_two_col_prototype)
                        _render_two_column_slide(
                            slide,
                            slide_title,
                            slide_data.get("left_content", ""),
                            slide_data.get("right_content", ""),
                            slide_index + 1,
                            slide_data,
                        )
                    else:
                        slide = _clone_slide_from_prototype(prs, aria_content_prototype)
                        _render_content_slide(
                            slide,
                            slide_title,
                            slide_data.get("content", ""),
                            slide_index + 1,
                            slide_data,
                        )

                keep_ref = None
                for index, slide_ref in enumerate(original_refs):
                    if prs.slides[index] is aria_cover_prototype:
                        keep_ref = slide_ref
                        break
                for slide_ref in list(original_refs):
                    if slide_ref is keep_ref:
                        continue
                    _remove_slide_ref(prs, slide_ref)

                filename = _generate_filename("pptx")
                filepath = GENERATED_DIR / filename
                _ensure_body_min_font_sizes(prs, min_size=12)
                prs.save(filepath)
                return {
                    "success": True,
                    "file_type": "pptx",
                    "file_name": filename,
                    "file_path": str(filepath.relative_to(UPLOADS_DIR)),
                    "full_path": str(filepath),
                    "template_path": str(template_path),
                    "template_name": Path(template_path).name,
                    "template_applied": True,
                    "template_mode": "graphic_library_named_prototypes",
                    "slide_count": len(prs.slides),
                }

        if is_graphic_library_template and len(prs.slides) >= 4:
            original_refs = list(prs.slides._sldIdLst)
            _set_template_cover_text(prs.slides[0], title, subtitle)

            agenda_ref = _slide_ref(prs, 1)
            divider_ref = _slide_ref(prs, 2)
            content_ref = _slide_ref(prs, 3)
            content_prototype = prs.slides[3]
            divider_prototype = prs.slides[2]
            two_col_layout = next(
                (layout for layout in prs.slide_layouts if layout.name == "TWO COLUMN TEXT"),
                _safe_layout(prs, 1),
            )

            section_counter = 0
            for slide_index, slide_data in enumerate(slides):
                slide_type = slide_data.get("type", "content")
                slide_title = slide_data.get("title", "")

                if slide_type == "title":
                    section_counter += 1
                    slide = _clone_slide_from_prototype(prs, divider_prototype)
                    _render_section_slide(slide, slide_title, section_counter, slide_data)
                elif slide_type == "two_column":
                    slide = prs.slides.add_slide(two_col_layout)
                    _render_two_column_slide(
                        slide,
                        slide_title,
                        slide_data.get("left_content", ""),
                        slide_data.get("right_content", ""),
                        slide_index + 1,
                        slide_data,
                    )
                else:
                    slide = _clone_slide_from_prototype(prs, content_prototype)
                    _render_graphic_library_content_slide(
                        slide,
                        slide_title,
                        slide_data.get("content", ""),
                        slide_index + 1,
                        slide_data,
                    )

            for slide_ref in list(original_refs)[1:]:
                _remove_slide_ref(prs, slide_ref)

            filename = _generate_filename("pptx")
            filepath = GENERATED_DIR / filename
            _ensure_body_min_font_sizes(prs, min_size=12)
            prs.save(filepath)
            return {
                "success": True,
                "file_type": "pptx",
                "file_name": filename,
                "file_path": str(filepath.relative_to(UPLOADS_DIR)),
                "full_path": str(filepath),
                "template_path": str(template_path),
                "template_name": Path(template_path).name,
                "template_applied": True,
                "template_mode": "graphic_library_prototypes",
                "slide_count": len(prs.slides),
            }
        if len(prs.slides) >= 5:
            _set_template_cover_text(prs.slides[0], title, subtitle)

            if len(prs.slides) >= 6:
                section_prototype = prs.slides[1]
                content_prototype = prs.slides[2]
                two_col_prototype = prs.slides[3]
                visual_prototype = prs.slides[4]
            else:
                section_prototype = prs.slides[3]
                content_prototype = _find_named_prototype_slide(prs, {"aria_slide_title", "aria_slide_body"}) or prs.slides[1]
                two_col_prototype = _find_named_prototype_slide(prs, {"aria_slide_title", "aria_left_body", "aria_right_body"}) or content_prototype
                visual_prototype = _find_named_prototype_slide(prs, {"aria_slide_title", "aria_visual_area"}) or content_prototype
            back_cover_ref = _slide_ref(prs, len(prs.slides) - 1)
            prototype_refs = {_slide_ref(prs, index) for index in range(1, len(prs.slides))}

            section_counter = 0
            for slide_index, slide_data in enumerate(slides):
                slide_type = slide_data.get("type", "content")
                slide_title = slide_data.get("title", "")
                content = slide_data.get("content", "")
                strategy_layout = _resolve_digital_strategy_layout(slide_data)
                use_visual = slide_type in {"roadmap", "matrix", "kpi", "risk", "next_steps"} or (
                    slide_type == "content" and _wants_visual_slide(slide_title, content)
                )

                if slide_type in {"title", "section"}:
                    slide = _clone_slide_from_prototype(prs, section_prototype)
                elif strategy_layout or use_visual:
                    slide = _clone_slide_from_prototype(prs, visual_prototype)
                elif slide_type == "two_column":
                    slide = _clone_slide_from_prototype(prs, two_col_prototype)
                else:
                    slide = _clone_slide_from_prototype(prs, content_prototype)

                if slide_type in {"title", "section"}:
                    section_counter += 1
                    _render_section_slide(slide, slide_title, section_counter, slide_data)
                elif strategy_layout and _render_digital_strategy_layout(slide, slide_data, slide_index + 1, strategy_layout):
                    pass
                elif use_visual:
                    _render_visual_slide(slide, slide_title, content, slide_index + 1, slide_type, slide_data)
                elif slide_type == "two_column":
                    _render_two_column_slide(
                        slide,
                        slide_title,
                        slide_data.get("left_content", ""),
                        slide_data.get("right_content", ""),
                        slide_index + 1,
                        slide_data,
                    )
                elif "content" in slide_data:
                    _render_content_slide(slide, slide_title, content, slide_index + 1, slide_data)

            for slide_ref in prototype_refs:
                if slide_ref is not back_cover_ref:
                    _remove_slide_ref(prs, slide_ref)
            _move_slide_ref_to_end(prs, back_cover_ref)

            filename = _generate_filename("pptx")
            filepath = GENERATED_DIR / filename
            _ensure_body_min_font_sizes(prs, min_size=12)
            prs.save(filepath)
            return {
                "success": True,
                "file_type": "pptx",
                "file_name": filename,
                "file_path": str(filepath.relative_to(UPLOADS_DIR)),
                "full_path": str(filepath),
                "template_path": str(template_path),
                "template_name": Path(template_path).name,
                "template_applied": True,
                "template_mode": "cloned_prototype_slides",
                "slide_count": len(prs.slides),
            }

    if using_template:
        prs = Presentation(template_path)

        # ── Cover slide (template slide 0) ──────────────────────────────────
        cover_slide = prs.slides[0]
        _set_template_cover_text(cover_slide, title, subtitle)
        has_template_content_prototype = len(prs.slides) > 2

        # Layout helpers for content slides
        content_layout = _safe_layout(prs, 1)   # aria_slide_title + aria_slide_body
        two_col_layout = _safe_layout(prs, 2)   # aria_slide_title + aria_left_body + aria_right_body

    else:
        prs = Presentation()
        has_template_content_prototype = False

        # Blank template: add a simple cover slide
        cover_layout = _safe_layout(prs, 0)
        slide = prs.slides.add_slide(cover_layout)
        if slide.shapes.title:
            _write_title_preserving_style(slide.shapes.title, title)
        if subtitle:
            body = _find_body_placeholder(slide)
            if body is not None:
                body.text_frame.text = subtitle

        content_layout = _safe_layout(prs, 1)
        two_col_layout = _safe_layout(prs, 1)

    # ── Content slides ───────────────────────────────────────────────────────
    for slide_index, slide_data in enumerate(slides):
        slide_type  = slide_data.get("type", "content")
        slide_title = slide_data.get("title", "")

        layout = two_col_layout if slide_type == "two_column" else content_layout
        if using_template and has_template_content_prototype and slide_index == 0:
            slide = prs.slides[1]
        else:
            slide = prs.slides.add_slide(layout)

        if slide_type == "content" and "content" in slide_data:
            _render_content_slide(slide, slide_title, slide_data["content"], slide_index + 1, slide_data)

        elif slide_type == "two_column":
            _render_two_column_slide(
                slide,
                slide_title,
                slide_data.get("left_content", ""),
                slide_data.get("right_content", ""),
                slide_index + 1,
                slide_data,
            )

    # ── Move back cover to the end ────────────────────────────────────────────
    # Template order after adding slides:
    # - With prototype: [cover, content1(prototype), back_cover, content2, ...]
    # - Without:        [cover, back_cover, content1, ...]
    # Desired order:    [cover, content1, content2, ..., back_cover]
    if using_template and len(prs.slides) >= 2:
        sldIdLst = prs.slides._sldIdLst
        back_cover_index = 2 if has_template_content_prototype and len(prs.slides) > 2 else 1
        back_cover_ref = list(sldIdLst)[back_cover_index]
        sldIdLst.remove(back_cover_ref)
        sldIdLst.append(back_cover_ref)

    # ── Save ─────────────────────────────────────────────────────────────────
    filename = _generate_filename("pptx")
    filepath = GENERATED_DIR / filename
    _ensure_body_min_font_sizes(prs, min_size=12)
    prs.save(filepath)

    return {
        "success": True,
        "file_type": "pptx",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
        **({
            "template_path": str(template_path),
            "template_name": Path(template_path).name,
            "template_applied": True,
            "template_mode": "layout_fallback",
        } if using_template else {"template_applied": False}),
        "slide_count": len(slides) + 2,  # cover + content + back cover
    }


@registry.register(
    name="generate_ppt_from_skill",
    description="Generate PowerPoint using a Skill's template. Looks for KPMG-Template.pptx or template.pptx in the skill's assets folder.",
    input_schema={
        "type": "object",
        "properties": {
            "skill_name": {
                "type": "string",
                "description": "Skill folder name (e.g., 'ai-strategy-report')"
            },
            "title": {
                "type": "string",
                "description": "Presentation title"
            },
            "subtitle": {
                "type": "string",
                "description": "Optional subtitle"
            },
            "slides": {
                "type": "array",
                "description": "Array of slide objects",
                "items": {
                    "type": "object",
                    "properties": {
                        "type": {
                            "type": "string",
                            "enum": [
                                "title",
                                "section",
                                "content",
                                "two_column",
                                "roadmap",
                                "matrix",
                                "kpi",
                                "risk",
                                "next_steps",
                            ],
                        },
                        "title": {"type": "string"},
                        "content": {"type": "string"},
                        "left_content": {"type": "string"},
                        "right_content": {"type": "string"}
                    },
                    "required": ["type", "title"]
                }
            },
            "deck_type": {
                "type": "string",
                "description": "Optional deck preset for presentation-builder: strategy, proposal, or project-update"
            },
            "template_key": {
                "type": "string",
                "description": "Optional template hint. presentation-builder can use digital-strategy as its base template."
            }
        },
        "required": ["skill_name", "title", "slides"]
    }
)
async def generate_ppt_from_skill(
    skill_name: str,
    title: str,
    slides: list[dict],
    subtitle: str = "",
    deck_type: str = "",
    template_key: str = "",
) -> dict[str, Any]:
    """Generate PPT using a skill's template."""
    title = _clean_ppt_text(title)
    subtitle = _clean_ppt_text(subtitle)
    if skill_name == "digital-strategy":
        slides = _normalize_digital_strategy_slides(slides)
    elif skill_name == "presentation-builder":
        slides = _normalize_presentation_builder_slides(slides, deck_type or template_key)
        slides = _normalize_presentation_builder_slide_format(slides)
    slides = _prepare_ppt_slide_text(slides)

    strict_template_skills = {"digital-strategy", "presentation-builder"}

    # Search for template in assets/ then references/ (both locations are valid)
    template_path = None
    searched_paths: list[str] = []
    template_skill_names = [skill_name]
    if skill_name == "presentation-builder":
        template_skill_names.append("digital-strategy")
    for template_skill_name in template_skill_names:
        for folder in ("assets", "references"):
            for filename in ("KPMG-Template.pptx", "Template.pptx", "template.pptx", "Graphic library.pptx"):
                candidate = SKILLS_DIR / template_skill_name / folder / filename
                searched_paths.append(str(candidate))
                if candidate.is_file():
                    template_path = candidate
                    break
            if template_path:
                break
        if template_path:
            break

    if not template_path:
        if skill_name in strict_template_skills:
            return {
                "success": False,
                "error": f"{skill_name} template not found; refusing to generate a blank deck.",
                "searched_paths": searched_paths,
                "template_applied": False,
            }
        # Fallback to default generation
        return await generate_ppt(title, slides, subtitle)

    return await generate_ppt(title, slides, subtitle, str(template_path))


@registry.register(
    name="generate_docx",
    description="Generate a Word (.docx) document with formatted content",
    input_schema={
        "type": "object",
        "properties": {
            "title": {
                "type": "string",
                "description": "Document title"
            },
            "sections": {
                "type": "array",
                "description": "Array of document sections",
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {
                            "type": "string",
                            "description": "Section heading"
                        },
                        "content": {
                            "type": "string",
                            "description": "Section content"
                        },
                        "level": {
                            "type": "integer",
                            "enum": [1, 2, 3],
                            "description": "Heading level"
                        }
                    },
                    "required": ["heading", "content"]
                }
            }
        },
        "required": ["title", "sections"]
    }
)
async def generate_docx(
    title: str,
    sections: list[dict]
) -> dict[str, Any]:
    """Generate a Word document."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return {
            "success": False,
            "error": "python-docx not installed. Run: pip install python-docx"
        }
    
    doc = Document()
    
    # Title
    title_para = doc.add_paragraph()
    title_run = title_para.add_run(title)
    title_run.bold = True
    title_run.font.size = Pt(18)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph()  # Spacing
    
    # Sections
    for section in sections:
        heading = section.get("heading", "")
        content = section.get("content", "")
        level = section.get("level", 1)
        
        # Add heading
        heading_para = doc.add_heading(heading, level=level)
        
        # Add content
        for line in content.split("\n"):
            line = line.strip()
            if line:
                # Check for bold (**text**)
                if line.startswith("**") and line.endswith("**"):
                    p = doc.add_paragraph()
                    run = p.add_run(line[2:-2])
                    run.bold = True
                elif line.startswith("- ") or line.startswith("* "):
                    doc.add_paragraph(line[2:], style='List Bullet')
                else:
                    doc.add_paragraph(line)
    
    # Save file
    filename = _generate_filename("docx")
    filepath = GENERATED_DIR / filename
    doc.save(filepath)
    
    return {
        "success": True,
        "file_type": "docx",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
        "section_count": len(sections),
    }


@registry.register(
    name="generate_xlsx",
    description="Generate an Excel (.xlsx) spreadsheet with data tables",
    input_schema={
        "type": "object",
        "properties": {
            "sheets": {
                "type": "array",
                "description": "Array of sheet objects",
                "items": {
                    "type": "object",
                    "properties": {
                        "name": {
                            "type": "string",
                            "description": "Sheet name"
                        },
                        "headers": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Column headers"
                        },
                        "data": {
                            "type": "array",
                            "description": "Array of row arrays",
                            "items": {
                                "type": "array",
                                "items": {"type": ["string", "number", "boolean"]}
                            }
                        }
                    },
                    "required": ["name", "headers", "data"]
                }
            }
        },
        "required": ["sheets"]
    }
)
async def generate_xlsx(
    sheets: list[dict]
) -> dict[str, Any]:
    """Generate an Excel spreadsheet."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill
    except ImportError:
        return {
            "success": False,
            "error": "openpyxl not installed. Run: pip install openpyxl"
        }
    
    wb = Workbook()
    
    for idx, sheet_data in enumerate(sheets):
        sheet_name = sheet_data.get("name", f"Sheet{idx+1}")
        headers = sheet_data.get("headers", [])
        data = sheet_data.get("data", [])
        
        if idx == 0:
            ws = wb.active
            ws.title = sheet_name
        else:
            ws = wb.create_sheet(title=sheet_name)
        
        # Add headers
        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="B4C7DC", end_color="B4C7DC", fill_type="solid")
        
        # Add data
        for row_idx, row_data in enumerate(data, 2):
            for col_idx, value in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=value)
    
    # Save file
    filename = _generate_filename("xlsx")
    filepath = GENERATED_DIR / filename
    wb.save(filepath)
    
    return {
        "success": True,
        "file_type": "xlsx",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
        "sheet_count": len(sheets),
    }


@registry.register(
    name="generate_pdf",
    description="Generate a PDF document from markdown or HTML content",
    input_schema={
        "type": "object",
        "properties": {
            "title": {
                "type": "string",
                "description": "Document title"
            },
            "content": {
                "type": "string",
                "description": "Document content (Markdown format supported)"
            },
            "orientation": {
                "type": "string",
                "enum": ["portrait", "landscape"],
                "description": "Page orientation"
            }
        },
        "required": ["title", "content"]
    }
)
async def generate_pdf(
    title: str,
    content: str,
    orientation: str = "portrait"
) -> dict[str, Any]:
    """Generate a PDF document."""
    try:
        # Try to use reportlab if available
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
    except ImportError:
        # Fallback: generate markdown file with note
        filename = _generate_filename("md")
        filepath = GENERATED_DIR / filename
        
        md_content = f"# {title}\n\n{content}\n\n---\n*Note: PDF generation requires reportlab. Install with: pip install reportlab*"
        filepath.write_text(md_content, encoding="utf-8")
        
        return {
            "success": True,
            "file_type": "md",
            "file_name": filename,
            "file_path": str(filepath.relative_to(UPLOADS_DIR)),
            "full_path": str(filepath),
            "note": "Generated as Markdown. Install reportlab for true PDF support."
        }
    
    # Generate actual PDF
    filename = _generate_filename("pdf")
    filepath = GENERATED_DIR / filename
    
    page_size = landscape(A4) if orientation == "landscape" else A4
    doc = SimpleDocTemplate(str(filepath), pagesize=page_size)
    
    styles = getSampleStyleSheet()
    story = []
    
    # Title
    story.append(Paragraph(title, styles['Title']))
    story.append(Spacer(1, 12))
    
    # Content
    for line in content.split("\n"):
        if line.strip():
            story.append(Paragraph(line, styles['Normal']))
            story.append(Spacer(1, 6))
    
    doc.build(story)
    
    return {
        "success": True,
        "file_type": "pdf",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
    }


@registry.register(
    name="save_json",
    description="Save structured data as a JSON file",
    input_schema={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Filename without extension"
            },
            "data": {
                "type": "object",
                "description": "JSON-serializable data object"
            }
        },
        "required": ["filename", "data"]
    }
)
async def save_json(
    filename: str,
    data: dict
) -> dict[str, Any]:
    """Save data as JSON file."""
    if not filename.endswith(".json"):
        filename += ".json"
    
    filepath = GENERATED_DIR / filename
    
    with open(filepath, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    
    return {
        "success": True,
        "file_type": "json",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
    }


@registry.register(
    name="save_text",
    description="Save plain text content as a text file",
    input_schema={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Filename without extension"
            },
            "content": {
                "type": "string",
                "description": "Text content"
            },
            "extension": {
                "type": "string",
                "enum": ["txt", "md", "csv", "html"],
                "description": "File extension"
            }
        },
        "required": ["filename", "content"]
    }
)
async def save_text(
    filename: str,
    content: str,
    extension: str = "txt"
) -> dict[str, Any]:
    """Save text content to file."""
    if not filename.endswith(f".{extension}"):
        filename += f".{extension}"
    
    filepath = GENERATED_DIR / filename
    filepath.write_text(content, encoding="utf-8")
    
    return {
        "success": True,
        "file_type": extension,
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
    }
