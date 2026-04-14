"""Projects router — CRUD for projects, milestones, file uploads."""
from __future__ import annotations

import shutil
import uuid
from datetime import datetime
from pathlib import Path
from typing import List, Optional

from fastapi import APIRouter, BackgroundTasks, Depends, Form, HTTPException, UploadFile, File
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlmodel import Session, select

import json
from app.config import UPLOADS_DIR
from app.database import get_session
from app.models.db import Conversation, Message, Project, Milestone, ProjectFile, ProjectFolder, ProjectPayment, ProjectTodo, ProjectMember, User
from app.services import claude as _claude_svc, openai_compat as _kimi_svc
from app.models.db import Setting as _Setting
from app.services.cache import projects_cache
from app.services.provider_selector import get_selected_model, resolve_provider_from_model, _load_provider_module
from app.routers.auth import get_current_user
from app.routers.chat_export import build_markdown_export_content

_PROJECTS_TTL = 120.0


def _bust_project(project_id: int) -> None:
    """Invalidate all caches that reference this project."""
    projects_cache.delete(f"detail:{project_id}")
    projects_cache.delete_prefix("list:")


async def _complete(messages: list[dict], max_tokens: int = 4000) -> str:
    """Call the active LLM provider using the user's selected model."""
    from app.database import engine
    from sqlmodel import Session as _S
    with _S(engine) as s:
        model = get_selected_model(s)
        provider = resolve_provider_from_model(model)
    llm = _load_provider_module(provider)
    return await llm.complete(messages, model=model, max_tokens=max_tokens)


async def _stream(messages: list[dict], max_tokens: int = 4000):
    """Stream response chunks from the active LLM provider."""
    from app.database import engine
    from sqlmodel import Session as _S
    with _S(engine) as s:
        model = get_selected_model(s)
        provider = resolve_provider_from_model(model)
    llm = _load_provider_module(provider)
    async for chunk in llm.stream_response(messages, model=model, max_tokens=max_tokens):
        yield chunk


DEFAULT_FOLDER_NAMES = ["项目需求", "方案和报价", "项目交付文档", "项目归档信息"]

# ── Shared file text extraction ────────────────────────────────────────────────

LEGACY_PRESALES_FOLDER_TARGETS = {
    "项目需求": "02_需求与方案",
    "方案和报价": "02_需求与方案",
    "项目交付文档": "03_会议与推进",
    "项目归档信息": "03_会议与推进",
}

PRESALES_TEMPLATE_FOLDERS = [
    "00_项目总览",
    "01_客户与关系",
    "02_需求与方案",
    "03_会议与推进",
]

PRESALES_TEMPLATE_FILES = [
    {
        "folder": "00_项目总览",
        "name": "00_项目概览.md",
        "summary": "咨询售前项目概览",
        "content": """# 项目概览

## 项目名称

## 客户名称

## 当前阶段
- 线索
- 初步沟通
- 需求澄清
- 方案准备
- 商务谈判
- 赢单 / 失单

## 客户想解决的核心问题

## 为什么这个机会值得跟进

## 当前最大不确定性

## 下一步关键动作

## 最新更新时间
""",
    },
    {
        "folder": "01_客户与关系",
        "name": "01_客户与关系.md",
        "summary": "客户背景和关键关系图谱",
        "content": """# 客户与关系

## 客户公司背景

## 项目发起背景

## Stakeholder Map
| 姓名 | 职位 | 角色 | 态度 | 影响力 |
|---|---|---|---|---|

## 决策链路
- 谁提出需求
- 谁影响预算
- 谁参与选型
- 谁最终拍板

## 当前关系判断
- 支持者
- 中立者
- 阻力方

## 关系推进策略
""",
    },
    {
        "folder": "02_需求与方案",
        "name": "02_需求与问题定义.md",
        "summary": "客户需求和问题定义",
        "content": """# 需求与问题定义

## 客户表面需求

## 客户真实需求

## 当前痛点

## 业务目标

## 成功标准
- 客户如何判断项目成功
- 可量化指标
- 不能触碰的约束

## 已确认事实

## 待确认问题

## 我们的判断与假设
""",
    },
    {
        "folder": "02_需求与方案",
        "name": "03_方案思路.md",
        "summary": "咨询方案初步思路",
        "content": """# 方案思路

## 方案目标

## 核心价值主张

## 工作范围

## 项目路径与方法

## 关键交付物

## 差异化亮点

## 客户可能追问的问题

## 下一轮需要补齐的材料
""",
    },
    {
        "folder": "02_需求与方案",
        "name": "04_竞争与替代方案.md",
        "summary": "竞争格局与替代方案判断",
        "content": """# 竞争与替代方案

## 客户当前可能的替代方案
- 内部自己做
- 继续沿用现状
- 交给其他咨询公司
- 暂缓不做

## 潜在竞争对手

## 我们的优势

## 我们的短板

## 当前竞争风险

## 应对策略
""",
    },
    {
        "folder": "02_需求与方案",
        "name": "05_商务推进.md",
        "summary": "预算、采购和商务推进记录",
        "content": """# 商务推进

## 预算情况
- 已明确
- 模糊
- 暂未确认

## 采购方式
- 直接签约
- 比价
- 招标
- 框架协议
- 其他

## 商务约束
- 价格
- 周期
- 合同条款
- 付款方式
- 验收方式

## 当前报价思路

## 客户反馈

## 商务下一步
""",
    },
    {
        "folder": "03_会议与推进",
        "name": "06_立项沟通纪要.md",
        "summary": "立项或首次沟通纪要",
        "content": """# 立项沟通纪要

## 基本信息
- 时间：
- 参会人：
- 会议目的：

## 客户表达的重点

## 我们确认到的信息

## 待补充材料

## 会后动作
- 客户侧：
- 我方：
""",
    },
    {
        "folder": "03_会议与推进",
        "name": "07_需求澄清会议纪要.md",
        "summary": "需求澄清会议纪要",
        "content": """# 需求澄清会议纪要

## 基本信息
- 时间：
- 参会人：
- 会议目的：

## 客户重点关注

## 关键问题与回答

## 未决事项

## 会后动作
- 客户侧：
- 我方：
""",
    },
    {
        "folder": "03_会议与推进",
        "name": "08_行动清单.md",
        "summary": "咨询售前行动清单",
        "content": """# 行动清单

## 本周重点动作
- [ ] 
- [ ] 
- [ ] 

## 客户待反馈
- [ ] 
- [ ] 

## 我方待准备
- [ ] 方案
- [ ] 报价
- [ ] 案例
- [ ] 演示
- [ ] 合同条款

## 时间节点
| 事项 | 截止时间 | 负责人 | 状态 |
|---|---|---|---|
""",
    },
    {
        "folder": "03_会议与推进",
        "name": "09_下一步与赢单判断.md",
        "summary": "赢单信号与下一步推进判断",
        "content": """# 下一步与赢单判断

## 下一步最关键的一步

## 未来 7 天推进计划

## 赢单信号
- 决策人开始参与
- 客户愿意给更多内部资料
- 客户明确时间表
- 客户主动讨论报价或合同
- 客户要求定制化方案

## 风险信号
- 一直在沟通但没有推进
- 需求频繁变化
- 预算始终不明确
- 决策人迟迟不出现
- 客户只是做信息收集

## 当前赢单判断
- 高
- 中
- 低

## 判断依据
""",
    },
]

try:
    import pdfplumber as _pdfplumber
    _HAS_PDF = True
except ImportError:
    _HAS_PDF = False

try:
    from docx import Document as _DocxDocument
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

try:
    from pptx import Presentation as _Presentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

try:
    import openpyxl as _openpyxl
    _HAS_XLSX = True
except ImportError:
    _HAS_XLSX = False


def _extract_file_text(path: Path, file_type: str, max_chars: int = 4000) -> str:
    """Extract plain text from a project file for AI context injection."""
    if not path.exists():
        return "[File not found]"
    try:
        ft = file_type.lower()
        if ft == "pdf" and _HAS_PDF:
            with _pdfplumber.open(path) as pdf:
                pages = [p.extract_text() or "" for p in pdf.pages[:15]]
            text = "\n".join(pages)
        elif ft == "docx" and _HAS_DOCX:
            doc = _DocxDocument(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ft == "pptx" and _HAS_PPTX:
            prs = _Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
            text = "\n\n".join(parts)
        elif ft in ("xlsx", "xls") and _HAS_XLSX:
            wb = _openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                rows = []
                for row in sheet.iter_rows(max_row=200, values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        rows.append("\t".join(cells))
                if rows:
                    parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
            wb.close()
            text = "\n\n".join(parts)
        elif ft in ("txt", "md", "csv", "json"):
            text = path.read_text(encoding="utf-8", errors="replace")
        else:
            return ""  # Binary — skip auto-summary
        text = text.strip()
        if len(text) > max_chars:
            return text[:max_chars] + "\n…[truncated]"
        return text if text else ""
    except Exception:
        return ""

def _init_default_folders(project_id: int, session: Session) -> list[ProjectFolder]:
    """Create the 4 default folders for a project. Safe to call if they already exist."""
    existing = session.exec(
        select(ProjectFolder).where(ProjectFolder.project_id == project_id)
    ).all()
    if existing:
        return existing
    folders = [
        ProjectFolder(project_id=project_id, name=name, sort_order=i)
        for i, name in enumerate(DEFAULT_FOLDER_NAMES)
    ]
    for f in folders:
        session.add(f)
    session.commit()
    for f in folders:
        session.refresh(f)
    return folders

router = APIRouter(prefix="/projects", tags=["projects"])


# ── Schemas ──────────────────────────────────────────────────────────────────

class ProjectCreate(BaseModel):
    name: str
    client: str
    description: str = ""
    status: str = "lead"
    contract_amount: float = 0.0
    notes: str = ""
    md_notes: str = ""


class ProjectUpdate(BaseModel):
    name: Optional[str] = None
    client: Optional[str] = None
    description: Optional[str] = None
    status: Optional[str] = None
    context_freshness: Optional[float] = None
    contract_amount: Optional[float] = None
    context_summary: Optional[str] = None
    notes: Optional[str] = None
    md_notes: Optional[str] = None


class TodoCreate(BaseModel):
    content: str
    is_done: bool = False
    due_date: Optional[str] = None
    assigned_to_user_id: Optional[int] = None


class TodoUpdate(BaseModel):
    content: Optional[str] = None
    is_done: Optional[bool] = None
    due_date: Optional[str] = None
    assigned_to_user_id: Optional[int] = None


class NoteBody(BaseModel):
    content: str
    append: bool = True   # True = append to existing notes; False = overwrite


class NotePolishBody(BaseModel):
    draft: str


class MemberCreate(BaseModel):
    user_id: int


class MemberUserOut(BaseModel):
    id: int
    display_name: str


class MemberOut(BaseModel):
    id: int
    project_id: int
    user_id: int
    user: MemberUserOut
    created_at: datetime


class PaymentCreate(BaseModel):
    amount: float
    payment_date: str               # YYYY-MM-DD
    note: str = ""
    payment_type: str = "received"  # received | expense | milestone_payment


class MilestoneCreate(BaseModel):
    title: str
    priority: str = "medium"
    due_date: Optional[str] = None


class MilestoneUpdate(BaseModel):
    title: Optional[str] = None
    is_done: Optional[bool] = None
    priority: Optional[str] = None
    due_date: Optional[str] = None


class FolderCreate(BaseModel):
    name: str
    sort_order: int = 0


class SaveConversationMarkdownRequest(BaseModel):
    action: str = "new"  # merge | new
    folder_id: Optional[int] = None
    file_id: Optional[int] = None
    file_name: Optional[str] = None


class SaveMessageToDocumentRequest(BaseModel):
    action: str  # merge | new
    file_id: Optional[int] = None
    file_name: Optional[str] = None
    folder_id: Optional[int] = None
    prepend_header: bool = True


class ProjectDocumentCreate(BaseModel):
    folder_id: Optional[int] = None
    name: str
    content: str = ""


class ProjectDocumentUpdate(BaseModel):
    content: Optional[str] = None
    name: Optional[str] = None
    folder_id: Optional[int] = None


class InitPresalesTemplateRequest(BaseModel):
    overwrite: bool = False


# ── Projects ──────────────────────────────────────────────────────────────────

@router.get("")
def list_projects(
    status: Optional[str] = None,
    member_user_id: Optional[int] = None,
    session: Session = Depends(get_session),
):
    cache_key = f"list:{status or ''}:member:{member_user_id or ''}"
    cached = projects_cache.get(cache_key)
    if cached is not None:
        return cached
    stmt = select(Project).order_by(Project.updated_at.desc())
    if status:
        stmt = stmt.where(Project.status == status)
    if member_user_id is not None:
        from sqlalchemy.orm import joinedload
        # join with members and filter; distinct to avoid duplicate projects
        stmt = (
            stmt.join(ProjectMember, ProjectMember.project_id == Project.id)
            .where(ProjectMember.user_id == member_user_id)
            .distinct()
        )
    result = session.exec(stmt).all()
    projects_cache.set(cache_key, result, _PROJECTS_TTL)
    return result


@router.post("", status_code=201)
def create_project(data: ProjectCreate, session: Session = Depends(get_session)):
    project = Project(**data.model_dump())
    session.add(project)
    session.commit()
    session.refresh(project)
    _init_default_folders(project.id, session)
    session.refresh(project)  # re-refresh: _init_default_folders commits, expiring project
    projects_cache.delete_prefix("list:")   # no detail key yet — project just created
    return project


@router.get("/{project_id}")
def get_project(project_id: int, session: Session = Depends(get_session)):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")
    return project


@router.get("/{project_id}/detail")
def get_project_detail(project_id: int, session: Session = Depends(get_session)):
    """Single-request combined endpoint: project + files + milestones + folders + financials.

    Reduces 4-5 round trips to Supabase down to 1 HTTP call with 5 fast local queries.
    """
    cache_key = f"detail:{project_id}"
    cached = projects_cache.get(cache_key)
    if cached is not None:
        return cached

    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    files = session.exec(
        select(ProjectFile).where(ProjectFile.project_id == project_id)
    ).all()
    milestones = session.exec(
        select(Milestone).where(Milestone.project_id == project_id)
    ).all()
    folders = session.exec(
        select(ProjectFolder)
        .where(ProjectFolder.project_id == project_id)
        .order_by(ProjectFolder.sort_order)
    ).all()
    if not folders:
        folders = _init_default_folders(project_id, session)

    payments = session.exec(
        select(ProjectPayment)
        .where(ProjectPayment.project_id == project_id)
        .order_by(ProjectPayment.payment_date)
    ).all()
    received = sum(p.amount for p in payments if p.payment_type in ("received", "milestone_payment"))
    expenses = sum(abs(p.amount) for p in payments if p.payment_type == "expense")
    invoiced = sum(p.amount for p in payments if p.payment_type == "invoiced")
    contract = project.contract_amount or 0.0

    todos = session.exec(
        select(ProjectTodo)
        .where(ProjectTodo.project_id == project_id)
        .order_by(ProjectTodo.is_done, ProjectTodo.updated_at.desc())
    ).all()

    members = session.exec(
        select(ProjectMember).where(ProjectMember.project_id == project_id)
    ).all()

    def _todo_dict(t: ProjectTodo) -> dict:
        return {
            "id": t.id,
            "project_id": t.project_id,
            "content": t.content,
            "is_done": t.is_done,
            "due_date": t.due_date,
            "assigned_to_user_id": t.assigned_to_user_id,
            "assigned_user": (
                {"id": t.assigned_user.id, "display_name": t.assigned_user.display_name}
                if t.assigned_user else None
            ),
            "created_at": t.created_at,
            "updated_at": t.updated_at,
        }

    result = {
        "project": project,
        "files": files,
        "milestones": milestones,
        "folders": folders,
        "md_notes": project.md_notes or "",
        "todos": [_todo_dict(t) for t in todos],
        "members": [
            {
                "id": m.id,
                "project_id": m.project_id,
                "user_id": m.user_id,
                "user": {"id": m.user.id, "display_name": m.user.display_name} if m.user else None,
                "created_at": m.created_at,
            }
            for m in members
        ],
        "financials": {
            "contract_amount": contract,
            "total_received": received,
            "total_expense": expenses,
            "total_invoiced": invoiced,
            "uncollected": invoiced - received,
            "remaining": contract - received,
            "payments": payments,
        },
    }
    projects_cache.set(cache_key, result, _PROJECTS_TTL)
    return result


@router.patch("/{project_id}")
def update_project(project_id: int, data: ProjectUpdate, session: Session = Depends(get_session)):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")
    for k, v in data.model_dump(exclude_none=True).items():
        setattr(project, k, v)
    project.updated_at = datetime.utcnow()
    session.add(project)
    session.commit()
    session.refresh(project)
    _bust_project(project_id)
    return project


@router.delete("/{project_id}")
def delete_project(project_id: int, session: Session = Depends(get_session)):
    from app.models.db import Conversation, Message
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    # Delete in dependency order to avoid FK violations
    # 1. Messages in project conversations
    convs = session.exec(select(Conversation).where(Conversation.project_id == project_id)).all()
    for conv in convs:
        session.exec(select(Message).where(Message.conversation_id == conv.id)).all()
        for msg in session.exec(select(Message).where(Message.conversation_id == conv.id)).all():
            session.delete(msg)
    session.flush()
    for conv in convs:
        session.delete(conv)
    session.flush()

    # 2. Files (must come before folders)
    for f in session.exec(select(ProjectFile).where(ProjectFile.project_id == project_id)).all():
        session.delete(f)
    session.flush()

    # 3. Folders
    for folder in session.exec(select(ProjectFolder).where(ProjectFolder.project_id == project_id)).all():
        session.delete(folder)
    session.flush()

    # 4. Milestones
    for ms in session.exec(select(Milestone).where(Milestone.project_id == project_id)).all():
        session.delete(ms)
    session.flush()

    # 5. Payments
    for p in session.exec(select(ProjectPayment).where(ProjectPayment.project_id == project_id)).all():
        session.delete(p)
    session.flush()

    # 6. Todos
    for t in session.exec(select(ProjectTodo).where(ProjectTodo.project_id == project_id)).all():
        session.delete(t)
    session.flush()

    # 7. Members
    for m in session.exec(select(ProjectMember).where(ProjectMember.project_id == project_id)).all():
        session.delete(m)
    session.flush()

    # 8. Project itself
    session.delete(project)
    session.commit()
    _bust_project(project_id)
    return {"ok": True}


# ── Milestones ────────────────────────────────────────────────────────────────

@router.get("/{project_id}/milestones")
def list_milestones(project_id: int, session: Session = Depends(get_session)):
    return session.exec(select(Milestone).where(Milestone.project_id == project_id)).all()


@router.post("/{project_id}/milestones", status_code=201)
def create_milestone(project_id: int, data: MilestoneCreate, session: Session = Depends(get_session)):
    ms = Milestone(project_id=project_id, **data.model_dump())
    session.add(ms)
    session.commit()
    session.refresh(ms)
    _bust_project(project_id)
    return ms


@router.patch("/{project_id}/milestones/{ms_id}")
def update_milestone(project_id: int, ms_id: int, data: MilestoneUpdate, session: Session = Depends(get_session)):
    ms = session.get(Milestone, ms_id)
    if not ms or ms.project_id != project_id:
        raise HTTPException(404, "Milestone not found")
    for k, v in data.model_dump(exclude_none=True).items():
        setattr(ms, k, v)
    session.add(ms)
    session.commit()
    session.refresh(ms)
    _bust_project(project_id)
    return ms


@router.delete("/{project_id}/milestones/{ms_id}")
def delete_milestone(project_id: int, ms_id: int, session: Session = Depends(get_session)):
    ms = session.get(Milestone, ms_id)
    if not ms or ms.project_id != project_id:
        raise HTTPException(404, "Milestone not found")
    session.delete(ms)
    session.commit()
    _bust_project(project_id)
    return {"ok": True}


# ── Files ─────────────────────────────────────────────────────────────────────

@router.get("/{project_id}/files")
def list_files(project_id: int, session: Session = Depends(get_session)):
    return session.exec(select(ProjectFile).where(ProjectFile.project_id == project_id)).all()


def _sanitize_markdown_filename(name: str) -> str:
    sanitized = "".join(c if c.isalnum() or c in (" ", "-", "_") else "_" for c in (name or "").strip())
    sanitized = "_".join(part for part in sanitized.split())
    return sanitized[:80] or "conversation"


def _resolve_project_folder(
    session: Session,
    project_id: int,
    preferred_folder_id: Optional[int] = None,
) -> Optional[ProjectFolder]:
    if preferred_folder_id is not None:
        folder = session.get(ProjectFolder, preferred_folder_id)
        if not folder or folder.project_id != project_id:
            raise HTTPException(404, "Folder not found")
        return folder

    folders = session.exec(
        select(ProjectFolder)
        .where(ProjectFolder.project_id == project_id)
        .order_by(ProjectFolder.sort_order, ProjectFolder.id)
    ).all()
    if not folders:
        folders = _init_default_folders(project_id, session)

    for folder in folders:
        if folder.sort_order == 2:
            return folder
    return folders[0] if folders else None


def _project_documents_dir(project_id: int) -> Path:
    dest_dir = UPLOADS_DIR / "projects" / str(project_id)
    dest_dir.mkdir(parents=True, exist_ok=True)
    return dest_dir


def _create_markdown_project_file(
    session: Session,
    project_id: int,
    name: str,
    content: str,
    folder_id: Optional[int] = None,
    summary: str = "",
) -> ProjectFile:
    safe_name = name if name.lower().endswith(".md") else f"{name}.md"
    dest_dir = _project_documents_dir(project_id)
    dest_file = dest_dir / f"{uuid.uuid4().hex}_{safe_name}"
    dest_file.write_text(content, encoding="utf-8")

    project_file = ProjectFile(
        project_id=project_id,
        folder_id=folder_id,
        name=safe_name,
        file_type="md",
        path=str(dest_file.relative_to(UPLOADS_DIR)),
        size_bytes=dest_file.stat().st_size,
        summary=summary,
    )
    session.add(project_file)
    session.commit()
    session.refresh(project_file)
    _bust_project(project_id)
    return project_file


def _get_project_file_or_404(session: Session, project_id: int, file_id: int) -> ProjectFile:
    project_file = session.get(ProjectFile, file_id)
    if not project_file or project_file.project_id != project_id:
        raise HTTPException(404, "File not found")
    return project_file


def _read_project_file_content(project_file: ProjectFile) -> str:
    full_path = UPLOADS_DIR / project_file.path
    if not full_path.exists():
        raise HTTPException(404, "File not found on disk")
    return full_path.read_text(encoding="utf-8", errors="replace")


def _cleanup_legacy_presales_folders(
    session: Session,
    project_id: int,
    folders_by_name: dict[str, ProjectFolder],
) -> int:
    cleaned_count = 0
    for legacy_name, target_name in LEGACY_PRESALES_FOLDER_TARGETS.items():
        legacy_folder = folders_by_name.get(legacy_name)
        target_folder = folders_by_name.get(target_name)
        if not legacy_folder or not target_folder or legacy_folder.id == target_folder.id:
            continue

        files = session.exec(select(ProjectFile).where(ProjectFile.folder_id == legacy_folder.id)).all()
        for project_file in files:
            project_file.folder_id = target_folder.id
            session.add(project_file)

        session.flush()
        session.delete(legacy_folder)
        session.commit()
        cleaned_count += 1
        folders_by_name.pop(legacy_name, None)

    if cleaned_count:
        _bust_project(project_id)
    return cleaned_count


@router.post("/{project_id}/notes/templates/presales", status_code=201)
def init_presales_notes_template(
    project_id: int,
    body: InitPresalesTemplateRequest,
    session: Session = Depends(get_session),
):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    existing_folders = session.exec(
        select(ProjectFolder)
        .where(ProjectFolder.project_id == project_id)
        .order_by(ProjectFolder.sort_order, ProjectFolder.id)
    ).all()
    folders_by_name = {folder.name: folder for folder in existing_folders}

    for index, folder_name in enumerate(PRESALES_TEMPLATE_FOLDERS):
        folder = folders_by_name.get(folder_name)
        if folder:
            continue
        folder = ProjectFolder(project_id=project_id, name=folder_name, sort_order=index)
        session.add(folder)
        session.commit()
        session.refresh(folder)
        folders_by_name[folder_name] = folder

    cleaned_folders = _cleanup_legacy_presales_folders(session, project_id, folders_by_name)

    existing_docs = session.exec(
        select(ProjectFile).where(ProjectFile.project_id == project_id, ProjectFile.file_type == "md")
    ).all()
    existing_by_key = {(doc.folder_id, doc.name): doc for doc in existing_docs}

    created_files: list[ProjectFile] = []
    updated_files: list[ProjectFile] = []
    for template in PRESALES_TEMPLATE_FILES:
        folder = folders_by_name[template["folder"]]
        existing = existing_by_key.get((folder.id, template["name"]))
        if existing and not body.overwrite:
            continue
        if existing:
            full_path = UPLOADS_DIR / existing.path
            full_path.write_text(template["content"], encoding="utf-8")
            existing.summary = template["summary"]
            existing.size_bytes = full_path.stat().st_size
            session.add(existing)
            session.commit()
            session.refresh(existing)
            updated_files.append(existing)
            continue

        created_files.append(
            _create_markdown_project_file(
                session=session,
                project_id=project_id,
                folder_id=folder.id,
                name=template["name"],
                content=template["content"],
                summary=template["summary"],
            )
        )

    return {
        "ok": True,
        "created_count": len(created_files),
        "updated_count": len(updated_files),
        "cleaned_folder_count": cleaned_folders,
        "folders": list(folders_by_name.values()),
        "files": created_files + updated_files,
    }


@router.post("/{project_id}/documents", status_code=201)
def create_project_document(
    project_id: int,
    data: ProjectDocumentCreate,
    session: Session = Depends(get_session),
):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    folder = _resolve_project_folder(session, project_id, data.folder_id) if data.folder_id is not None else None
    filename = _sanitize_markdown_filename(data.name)
    if not filename.lower().endswith(".md"):
        filename = f"{filename}.md"

    return _create_markdown_project_file(
        session=session,
        project_id=project_id,
        folder_id=folder.id if folder else None,
        name=filename,
        content=data.content,
        summary="Project note document",
    )


@router.get("/{project_id}/documents/{file_id}")
def get_project_document(project_id: int, file_id: int, session: Session = Depends(get_session)):
    project_file = _get_project_file_or_404(session, project_id, file_id)
    if project_file.file_type.lower() != "md":
        raise HTTPException(400, "Only markdown documents are supported")
    return {
        "id": project_file.id,
        "project_id": project_file.project_id,
        "folder_id": project_file.folder_id,
        "name": project_file.name,
        "content": _read_project_file_content(project_file),
        "summary": project_file.summary,
        "uploaded_at": project_file.uploaded_at,
    }


@router.patch("/{project_id}/documents/{file_id}")
def update_project_document(
    project_id: int,
    file_id: int,
    data: ProjectDocumentUpdate,
    session: Session = Depends(get_session),
):
    project_file = _get_project_file_or_404(session, project_id, file_id)
    if project_file.file_type.lower() != "md":
        raise HTTPException(400, "Only markdown documents are supported")

    full_path = UPLOADS_DIR / project_file.path
    if not full_path.exists():
        raise HTTPException(404, "File not found on disk")

    if data.content is not None:
        full_path.write_text(data.content, encoding="utf-8")
        project_file.size_bytes = full_path.stat().st_size

    if data.name is not None:
        next_name = _sanitize_markdown_filename(data.name)
        if not next_name.lower().endswith(".md"):
            next_name = f"{next_name}.md"
        project_file.name = next_name

    if data.folder_id is not None:
        folder = _resolve_project_folder(session, project_id, data.folder_id)
        project_file.folder_id = folder.id if folder else None

    session.add(project_file)
    session.commit()
    session.refresh(project_file)
    _bust_project(project_id)
    return {
        "ok": True,
        "id": project_file.id,
        "name": project_file.name,
        "folder_id": project_file.folder_id,
        "size_bytes": project_file.size_bytes,
    }


@router.post("/{project_id}/conversations/{conv_id}/save-markdown", status_code=201)
def save_conversation_markdown(
    project_id: int,
    conv_id: int,
    data: SaveConversationMarkdownRequest,
    session: Session = Depends(get_session),
):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    conv = session.get(Conversation, conv_id)
    if not conv or conv.project_id != project_id:
        raise HTTPException(404, "Conversation not found")

    messages = session.exec(
        select(Message)
        .where(Message.conversation_id == conv_id)
        .order_by(Message.created_at)
    ).all()
    if not messages:
        raise HTTPException(400, "Conversation has no messages")

    markdown_content = build_markdown_export_content(conv, messages)

    if data.action == "merge":
        if not data.file_id:
            raise HTTPException(400, "file_id is required for merge action")
        project_file = _get_project_file_or_404(session, project_id, data.file_id)
        if project_file.file_type.lower() != "md":
            raise HTTPException(400, "Only markdown documents can be merged")

        full_path = UPLOADS_DIR / project_file.path
        if not full_path.exists():
            raise HTTPException(404, "File not found on disk")

        existing = full_path.read_text(encoding="utf-8", errors="replace")
        timestamp = datetime.utcnow().strftime("%Y-%m-%d %H:%M")
        header = f"\n\n---\n\n> From project conversation | {timestamp}\n\n"
        new_content = existing + header + markdown_content
        full_path.write_text(new_content, encoding="utf-8")
        project_file.size_bytes = full_path.stat().st_size
        session.add(project_file)
        session.commit()
        session.refresh(project_file)
        _bust_project(project_id)
        return {
            "ok": True,
            "action": "merge",
            "id": project_file.id,
            "name": project_file.name,
            "folder_id": project_file.folder_id,
            "size_bytes": project_file.size_bytes,
        }

    # action == "new"
    target_folder = _resolve_project_folder(session, project_id, data.folder_id) if data.folder_id is not None else None
    base_name = _sanitize_markdown_filename(data.file_name or conv.title or "conversation")
    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    filename = f"{base_name}_{timestamp}.md"

    new_file = _create_markdown_project_file(
        session=session,
        project_id=project_id,
        folder_id=target_folder.id if target_folder else None,
        name=filename,
        content=markdown_content,
        summary=f"Saved from conversation: {conv.title or 'Untitled Conversation'}",
    )
    return {
        "ok": True,
        "action": "new",
        "id": new_file.id,
        "name": new_file.name,
        "folder_id": new_file.folder_id,
        "size_bytes": new_file.size_bytes,
    }


@router.post("/{project_id}/messages/{message_id}/save-to-document", status_code=201)
def save_message_to_document(
    project_id: int,
    message_id: int,
    data: SaveMessageToDocumentRequest,
    session: Session = Depends(get_session),
):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    message = session.get(Message, message_id)
    if not message:
        raise HTTPException(404, "Message not found")

    conv = session.get(Conversation, message.conversation_id)
    if not conv or conv.project_id != project_id:
        raise HTTPException(404, "Message does not belong to this project")

    timestamp = datetime.utcnow().strftime("%Y-%m-%d %H:%M")
    header = f"\n\n---\n\n> From project conversation | {timestamp}\n\n"
    content_block = header + message.content if data.prepend_header else message.content

    if data.action == "merge":
        if not data.file_id:
            raise HTTPException(400, "file_id is required for merge action")
        project_file = _get_project_file_or_404(session, project_id, data.file_id)
        if project_file.file_type.lower() != "md":
            raise HTTPException(400, "Only markdown documents can be merged")

        full_path = UPLOADS_DIR / project_file.path
        if not full_path.exists():
            raise HTTPException(404, "File not found on disk")

        existing = full_path.read_text(encoding="utf-8", errors="replace")
        new_content = existing + content_block
        full_path.write_text(new_content, encoding="utf-8")
        project_file.size_bytes = full_path.stat().st_size
        session.add(project_file)
        session.commit()
        session.refresh(project_file)
        _bust_project(project_id)
        return {
            "ok": True,
            "action": "merge",
            "id": project_file.id,
            "name": project_file.name,
            "size_bytes": project_file.size_bytes,
        }

    # action == "new"
    target_folder = _resolve_project_folder(session, project_id, data.folder_id) if data.folder_id is not None else None
    base_name = _sanitize_markdown_filename(data.file_name or f"message_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}")
    if not base_name.lower().endswith(".md"):
        base_name = f"{base_name}.md"

    new_file = _create_markdown_project_file(
        session=session,
        project_id=project_id,
        folder_id=target_folder.id if target_folder else None,
        name=base_name,
        content=message.content,
        summary=f"Saved from conversation: {conv.title or 'Untitled Conversation'}",
    )
    return {
        "ok": True,
        "action": "new",
        "id": new_file.id,
        "name": new_file.name,
        "folder_id": new_file.folder_id,
        "size_bytes": new_file.size_bytes,
    }


@router.post("/{project_id}/files", status_code=201)
async def upload_file(
    project_id: int,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    folder_id: Optional[int] = Form(None),
    session: Session = Depends(get_session),
):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    suffix = Path(file.filename or "file").suffix.lower()
    dest_name = f"{uuid.uuid4().hex}{suffix}"
    dest_path = UPLOADS_DIR / "projects" / str(project_id)
    dest_path.mkdir(parents=True, exist_ok=True)
    dest_file = dest_path / dest_name

    with dest_file.open("wb") as f:
        shutil.copyfileobj(file.file, f)

    pf = ProjectFile(
        project_id=project_id,
        folder_id=folder_id,
        name=file.filename or dest_name,
        file_type=suffix.lstrip("."),
        path=str(dest_file.relative_to(UPLOADS_DIR)),
        size_bytes=dest_file.stat().st_size,
    )
    session.add(pf)
    session.commit()
    session.refresh(pf)

    # Auto-generate file summary in the background
    background_tasks.add_task(_auto_summarize_file, pf.id, str(dest_file), suffix.lstrip("."))

    _bust_project(project_id)
    return pf


@router.delete("/{project_id}/files/{file_id}")
def delete_file(project_id: int, file_id: int, session: Session = Depends(get_session)):
    pf = session.get(ProjectFile, file_id)
    if not pf or pf.project_id != project_id:
        raise HTTPException(404, "File not found")
    full_path = UPLOADS_DIR / pf.path
    if full_path.exists():
        full_path.unlink()
    session.delete(pf)
    session.commit()
    _bust_project(project_id)
    return {"ok": True}


@router.get("/{project_id}/files/{file_id}/download")
def download_file(project_id: int, file_id: int, session: Session = Depends(get_session)):
    """Download a project file."""
    from fastapi.responses import FileResponse
    
    pf = session.get(ProjectFile, file_id)
    if not pf or pf.project_id != project_id:
        raise HTTPException(404, "File not found")
    
    full_path = UPLOADS_DIR / pf.path
    if not full_path.exists():
        raise HTTPException(404, "File not found on disk")
    
    return FileResponse(
        path=str(full_path),
        filename=pf.name,
        media_type="application/octet-stream"
    )


# ── Folders ───────────────────────────────────────────────────────────────────

@router.get("/{project_id}/folders")
def list_folders(project_id: int, session: Session = Depends(get_session)):
    folders = session.exec(
        select(ProjectFolder)
        .where(ProjectFolder.project_id == project_id)
        .order_by(ProjectFolder.sort_order)
    ).all()
    if not folders:
        folders = _init_default_folders(project_id, session)
    return folders


@router.post("/{project_id}/folders", status_code=201)
def create_folder(project_id: int, data: FolderCreate, session: Session = Depends(get_session)):
    folder = ProjectFolder(project_id=project_id, **data.model_dump())
    session.add(folder)
    session.commit()
    session.refresh(folder)
    _bust_project(project_id)
    return folder


@router.delete("/{project_id}/folders/{folder_id}")
def delete_folder(project_id: int, folder_id: int, session: Session = Depends(get_session)):
    folder = session.get(ProjectFolder, folder_id)
    if not folder or folder.project_id != project_id:
        raise HTTPException(404, "Folder not found")
    # Unlink files from this folder before deleting
    files = session.exec(select(ProjectFile).where(ProjectFile.folder_id == folder_id)).all()
    for f in files:
        f.folder_id = None
        session.add(f)
    session.delete(folder)
    session.commit()
    _bust_project(project_id)
    return {"ok": True}


# ── Financials ────────────────────────────────────────────────────────────────

@router.get("/{project_id}/financials")
def get_financials(project_id: int, session: Session = Depends(get_session)):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")
    payments = session.exec(
        select(ProjectPayment)
        .where(ProjectPayment.project_id == project_id)
        .order_by(ProjectPayment.payment_date)
    ).all()
    received = sum(p.amount for p in payments if p.payment_type in ("received", "milestone_payment"))
    expenses = sum(abs(p.amount) for p in payments if p.payment_type == "expense")
    invoiced = sum(p.amount for p in payments if p.payment_type == "invoiced")
    return {
        "contract_amount": project.contract_amount,
        "total_received": received,
        "total_expense": expenses,
        "total_invoiced": invoiced,
        "uncollected": invoiced - received,   # 已开票未收款
        "remaining": project.contract_amount - received,
        "payments": payments,
    }


@router.post("/{project_id}/financials", status_code=201)
def add_payment(project_id: int, data: PaymentCreate, session: Session = Depends(get_session)):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")
    amount = -abs(data.amount) if data.payment_type == "expense" else abs(data.amount)
    payment = ProjectPayment(project_id=project_id, amount=amount,
                             payment_date=data.payment_date, note=data.note,
                             payment_type=data.payment_type)
    session.add(payment)
    session.commit()
    session.refresh(payment)
    _bust_project(project_id)
    return payment


@router.delete("/{project_id}/financials/{payment_id}")
def delete_payment(project_id: int, payment_id: int, session: Session = Depends(get_session)):
    payment = session.get(ProjectPayment, payment_id)
    if not payment or payment.project_id != project_id:
        raise HTTPException(404, "Payment not found")
    session.delete(payment)
    session.commit()
    _bust_project(project_id)
    return {"ok": True}


# ── AI Suggest ────────────────────────────────────────────────────────────────

class ProjectAISuggestQuery(BaseModel):
    query: str             # rough idea the user typed
    client_name: str = ""
    client_industry: str = ""


class ProjectAISuggestion(BaseModel):
    name: str
    description: str


@router.post("/ai-suggest", response_model=list[ProjectAISuggestion])
async def ai_suggest_project(body: ProjectAISuggestQuery):
    """Ask Claude to propose 1-3 consulting project names + descriptions."""
    client_context = ""
    if body.client_name:
        client_context = f"Client: {body.client_name}"
        if body.client_industry:
            client_context += f" ({body.client_industry})"

    prompt = f"""You are a senior consultant at a top-tier consulting firm.
{f"The project is for: {client_context}" if client_context else ""}
The user described the project as: "{body.query}"

Generate 1 to 3 consulting project name and description suggestions.
- If the idea is specific, return 1 suggestion.
- If the idea is broad or ambiguous, return up to 3 distinct angle variations.

Return ONLY a valid JSON array (no markdown, no extra text):
[
  {{
    "name": "Crisp, professional project title (5-8 words max)",
    "description": "2-3 sentence scope statement: objectives, key workstreams, and expected deliverable"
  }}
]

Rules:
- name: concise, consulting-style (e.g. "China Market Entry Strategy", "Digital Transformation Roadmap")
- description: professional, specific, actionable — no filler phrases
- Return pure JSON array only"""

    try:
        raw = await _complete(
            messages=[{"role": "user", "content": prompt}],
            max_tokens=4000,
        )
        text = raw.strip()
        if text.startswith("```"):
            text = text.split("```")[1]
            if text.startswith("json"):
                text = text[4:]
        suggestions = json.loads(text)
        return [ProjectAISuggestion(**s) for s in suggestions[:3]]
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI suggestion failed: {e}")


# ── Background: auto-summarize uploaded file ──────────────────────────────────

async def _auto_summarize_file(file_id: int, file_path: str, file_type: str) -> None:
    """Generate a 2-3 sentence summary for an uploaded project file and persist it."""
    from app.database import engine
    from sqlmodel import Session as _Session

    text = _extract_file_text(Path(file_path), file_type, max_chars=3000)
    if not text or text.startswith("["):
        return  # Nothing to summarize

    prompt = (
        "You are a professional consultant analyst. "
        "Read the following document excerpt and write a concise 2-3 sentence summary "
        "covering: what this document is, its main purpose, and the most important information it contains. "
        "Be specific and professional. Return ONLY the summary, no preamble.\n\n"
        f"Document excerpt:\n{text}"
    )
    try:
        summary = await _complete(
            messages=[{"role": "user", "content": prompt}],
            max_tokens=2000,
        )
        summary = summary.strip()
        if summary:
            with _Session(engine) as session:
                pf = session.get(ProjectFile, file_id)
                if pf:
                    pf.summary = summary
                    session.add(pf)
                    session.commit()
    except Exception:
        pass  # Best-effort — don't break the upload flow


# ── Generate project context summary ──────────────────────────────────────────

@router.post("/{project_id}/generate-context")
async def generate_project_context(project_id: int, session: Session = Depends(get_session)):
    """Ask LLM to generate a structured context summary (SSE streaming)."""
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    milestones = session.exec(select(Milestone).where(Milestone.project_id == project_id)).all()
    files = session.exec(select(ProjectFile).where(ProjectFile.project_id == project_id)).all()

    # Build project data block
    lines = [
        f"Project: {project.name}",
        f"Client: {project.client}",
        f"Status: {project.status}",
    ]
    if project.description:
        lines.append(f"Description: {project.description}")
    if milestones:
        lines.append(f"Milestones ({len(milestones)} total, {sum(1 for m in milestones if m.is_done)} completed):")
        for m in milestones:
            status = "✓" if m.is_done else "○"
            lines.append(f"  {status} {m.title}" + (f" [{m.priority}]" if m.priority == "high" else ""))
    if files:
        lines.append(f"Uploaded files ({len(files)}):")
        for f in files:
            lines.append(f"  - {f.name}" + (f": {f.summary[:120]}" if f.summary else ""))

    project_data = "\n".join(lines)

    prompt = (
        "You are an AI consultant assistant. Based on the project data below, "
        "treat the current project as the only source of truth. "
        "Do not blend in facts, progress, or risks from other projects under the same client unless explicitly stated in the project data below. "
        "If some information appears ambiguous, stay conservative and note the uncertainty rather than borrowing context from elsewhere. "
        "generate a concise context summary of 3-5 bullet points that capture: "
        "the project's core objective, current stage, key risks or open questions, "
        "critical milestones, and important context a consultant should always remember. "
        "Each bullet should be specific and actionable, not generic. "
        "Use **bold** for key terms or milestones within each bullet. "
        "Return ONLY the bullet points, one per line, starting with '•'. "
        "Write in the same language as the project name (Chinese if Chinese, English if English).\n\n"
        f"Project data:\n{project_data}"
    )

    messages = [{"role": "user", "content": prompt}]

    async def event_stream():
        accumulated: list[str] = []
        try:
            async for chunk in _stream(messages, max_tokens=4000):
                # Skip tool_use JSON blobs and TOOL_START markers
                if chunk.startswith('{"type": "tool_use"') or chunk.startswith("[TOOL_START:"):
                    continue
                accumulated.append(chunk)
                yield f"data: {json.dumps({'type': 'text', 'content': chunk}, ensure_ascii=False)}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"
            return

        summary = "".join(accumulated).strip()

        # Save to DB with a fresh session
        from app.database import engine as _engine
        from sqlmodel import Session as _S
        with _S(_engine) as write_session:
            p = write_session.get(Project, project_id)
            if p:
                p.context_summary = summary
                p.updated_at = datetime.utcnow()
                write_session.add(p)
                write_session.commit()
                _bust_project(project_id)

        yield f"data: {json.dumps({'type': 'done', 'context_summary': summary}, ensure_ascii=False)}\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")


# ── Project notes (沉淀到项目) ─────────────────────────────────────────────────

@router.post("/{project_id}/notes")
def save_project_note(project_id: int, body: NoteBody, session: Session = Depends(get_session)):
    """Append or overwrite project notes."""
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    if body.append and project.notes:
        timestamp = datetime.utcnow().strftime("%Y-%m-%d %H:%M")
        project.notes = f"{project.notes}\n\n---\n[{timestamp}]\n{body.content}"
    else:
        project.notes = body.content

    project.updated_at = datetime.utcnow()
    session.add(project)
    session.commit()
    session.refresh(project)
    _bust_project(project_id)
    return {"notes": project.notes}

# ── Project Todos ────────────────────────────────────────────────────────────

def _serialize_todo(todo: ProjectTodo) -> dict:
    return {
        "id": todo.id,
        "project_id": todo.project_id,
        "content": todo.content,
        "is_done": todo.is_done,
        "due_date": todo.due_date,
        "assigned_to_user_id": todo.assigned_to_user_id,
        "assigned_user": (
            {"id": todo.assigned_user.id, "display_name": todo.assigned_user.display_name}
            if todo.assigned_user else None
        ),
        "created_at": todo.created_at,
        "updated_at": todo.updated_at,
    }


@router.get("/{project_id}/todos")
def list_todos(project_id: int, session: Session = Depends(get_session)):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")
    todos = session.exec(
        select(ProjectTodo)
        .where(ProjectTodo.project_id == project_id)
        .order_by(ProjectTodo.is_done, ProjectTodo.updated_at.desc())
    ).all()
    return [_serialize_todo(t) for t in todos]


@router.post("/{project_id}/todos", status_code=201)
def create_todo(project_id: int, body: TodoCreate, session: Session = Depends(get_session)):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")
    todo = ProjectTodo(
        project_id=project_id,
        content=body.content,
        is_done=body.is_done,
        due_date=body.due_date,
        assigned_to_user_id=body.assigned_to_user_id,
    )
    session.add(todo)
    session.commit()
    session.refresh(todo)
    _bust_project(project_id)
    return _serialize_todo(todo)


@router.patch("/{project_id}/todos/{todo_id}")
def update_todo(project_id: int, todo_id: int, body: TodoUpdate, session: Session = Depends(get_session)):
    todo = session.exec(
        select(ProjectTodo).where(ProjectTodo.id == todo_id, ProjectTodo.project_id == project_id)
    ).first()
    if not todo:
        raise HTTPException(404, "Todo not found")
    for k, v in body.model_dump(exclude_none=True).items():
        setattr(todo, k, v)
    todo.updated_at = datetime.utcnow()
    session.add(todo)
    session.commit()
    session.refresh(todo)
    _bust_project(project_id)
    return _serialize_todo(todo)


@router.delete("/{project_id}/todos/{todo_id}")
def delete_todo(project_id: int, todo_id: int, session: Session = Depends(get_session)):
    todo = session.exec(
        select(ProjectTodo).where(ProjectTodo.id == todo_id, ProjectTodo.project_id == project_id)
    ).first()
    if not todo:
        raise HTTPException(404, "Todo not found")
    session.delete(todo)
    session.commit()
    _bust_project(project_id)
    return {"ok": True}


# ── Project Members ───────────────────────────────────────────────────────────

@router.get("/{project_id}/members", response_model=list[MemberOut])
def list_members(project_id: int, session: Session = Depends(get_session)):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")
    members = session.exec(
        select(ProjectMember).where(ProjectMember.project_id == project_id)
    ).all()
    return [
        MemberOut(
            id=m.id,
            project_id=m.project_id,
            user_id=m.user_id,
            user=MemberUserOut(id=m.user.id, display_name=m.user.display_name),
            created_at=m.created_at,
        )
        for m in members if m.user
    ]


@router.post("/{project_id}/members", status_code=201, response_model=MemberOut)
def add_member(project_id: int, body: MemberCreate, session: Session = Depends(get_session)):
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")
    user = session.get(User, body.user_id)
    if not user:
        raise HTTPException(404, "User not found")
    existing = session.exec(
        select(ProjectMember).where(
            ProjectMember.project_id == project_id,
            ProjectMember.user_id == body.user_id,
        )
    ).first()
    if existing:
        raise HTTPException(409, "User is already a member of this project")
    member = ProjectMember(project_id=project_id, user_id=body.user_id)
    session.add(member)
    session.commit()
    session.refresh(member)
    _bust_project(project_id)
    return MemberOut(
        id=member.id,
        project_id=member.project_id,
        user_id=member.user_id,
        user=MemberUserOut(id=user.id, display_name=user.display_name),
        created_at=member.created_at,
    )


@router.delete("/{project_id}/members/{user_id}")
def remove_member(project_id: int, user_id: int, session: Session = Depends(get_session)):
    member = session.exec(
        select(ProjectMember).where(
            ProjectMember.project_id == project_id,
            ProjectMember.user_id == user_id,
        )
    ).first()
    if not member:
        raise HTTPException(404, "Member not found")
    session.delete(member)
    session.commit()
    _bust_project(project_id)
    return {"ok": True}


# ── AI Polish for Project Notes ──────────────────────────────────────────────

@router.post("/{project_id}/notes/ai-polish")
async def ai_polish_project_notes(project_id: int, body: NotePolishBody, session: Session = Depends(get_session)):
    """Use the active LLM to polish a rough draft into structured Markdown project notes."""
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    system_prompt = (
        "You are a helpful assistant that turns rough drafts into well-structured Markdown project notes. "
        "Keep the user's original meaning, organize content with headings, bullet points, and checklists where appropriate, "
        "and output clean Markdown without wrapping it in code blocks."
    )
    user_prompt = f"""Please polish the following rough draft into well-structured Markdown project notes.

Project name: {project.name}
Client: {project.client}

Draft:
{body.draft}
"""
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]
    result = await _complete(messages, max_tokens=4000)
    return {"result": result}


@router.post("/{project_id}/notes/ai-polish-stream")
async def ai_polish_project_notes_stream(project_id: int, body: NotePolishBody, session: Session = Depends(get_session)):
    """Stream the active LLM polishing a rough draft into structured Markdown project notes."""
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    system_prompt = (
        "You are a helpful assistant that turns rough drafts into well-structured Markdown project notes. "
        "Keep the user's original meaning, organize content with headings, bullet points, and checklists where appropriate, "
        "and output clean Markdown without wrapping it in code blocks."
    )
    user_prompt = f"""Please polish the following rough draft into well-structured Markdown project notes.

Project name: {project.name}
Client: {project.client}

Draft:
{body.draft}
"""
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]

    async def event_stream():
        try:
            async for chunk in _stream(messages, max_tokens=4000):
                if chunk.startswith('{"type": "tool_use"') or chunk.startswith("[TOOL_START:"):
                    continue
                yield f"data: {json.dumps({'type': 'text', 'content': chunk}, ensure_ascii=False)}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"
            return
        yield f"data: {json.dumps({'type': 'done'}, ensure_ascii=False)}\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")


@router.get("/todos/my")
def list_my_todos(
    current_user: User = Depends(get_current_user),
    session: Session = Depends(get_session),
):
    """Return pending todos assigned to the current user across all projects."""
    rows = session.exec(
        select(ProjectTodo, Project)
        .join(Project, ProjectTodo.project_id == Project.id)
        .where(
            ProjectTodo.assigned_to_user_id == current_user.id,
            ProjectTodo.is_done == False,
        )
        .order_by(ProjectTodo.updated_at.desc())
    ).all()
    return [
        {
            "id": t.id,
            "project_id": t.project_id,
            "project_name": p.name,
            "content": t.content,
            "due_date": t.due_date,
            "created_at": t.created_at,
            "updated_at": t.updated_at,
        }
        for t, p in rows
    ]
