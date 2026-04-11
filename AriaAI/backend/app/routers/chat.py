"""Chat router — SSE streaming, conversation history, RAG injection."""
from __future__ import annotations

import asyncio
import json
from datetime import datetime
from pathlib import Path
from typing import List, Optional

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlmodel import Session, select

from app.config import UPLOADS_DIR
from app.database import get_session
from app.models.db import Conversation, Message, Milestone, Project, ProjectFile, ProjectPayment, Skill
from app.services import claude, rag, openai_compat
from app.services.cache import conversations_cache

_CONV_TTL = 20.0
from app.models.db import Setting as _Setting
from app.services.tool_executor import format_tools_for_claude
from app.tools import registry

try:
    import pdfplumber as _pdfplumber
    _HAS_PDF = True
except ImportError:
    _HAS_PDF = False

try:
    from docx import Document as _DocxDocument
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

try:
    from pptx import Presentation as _Presentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

try:
    import openpyxl as _openpyxl
    _HAS_XLSX = True
except ImportError:
    _HAS_XLSX = False


def _extract_file_text(path: Path, file_type: str, max_chars: int = 4000) -> str:
    """Extract plain text from a project file for AI context injection."""
    if not path.exists():
        return "[File not found]"
    try:
        ft = file_type.lower()
        if ft == "pdf" and _HAS_PDF:
            with _pdfplumber.open(path) as pdf:
                pages = [p.extract_text() or "" for p in pdf.pages[:15]]
            text = "\n".join(pages)
        elif ft == "docx" and _HAS_DOCX:
            doc = _DocxDocument(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ft == "pptx" and _HAS_PPTX:
            prs = _Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
            text = "\n\n".join(parts)
        elif ft in ("xlsx", "xls") and _HAS_XLSX:
            wb = _openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                rows = []
                for row in sheet.iter_rows(max_row=200, values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        rows.append("\t".join(cells))
                if rows:
                    parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
            wb.close()
            text = "\n\n".join(parts)
        elif ft in ("txt", "md", "csv", "json"):
            text = path.read_text(encoding="utf-8", errors="replace")
        else:
            return "[Binary file — text extraction not supported for this format]"
        text = text.strip()
        if len(text) > max_chars:
            return text[:max_chars] + "\n…[truncated]"
        return text if text else "[File appears to be empty]"
    except Exception as exc:
        return f"[Could not extract text: {exc}]"

router = APIRouter(prefix="/chat", tags=["chat"])


async def _generate_title_bg(conv_id: int, user_content: str, bind, complete_fn) -> None:
    """Generate conversation title in the background — called after SSE done is sent."""
    from sqlmodel import Session as _Sess
    try:
        raw = await complete_fn(
            messages=[{"role": "user", "content": (
                f"Write a short title for this conversation (max 12 Chinese characters "
                f"or 6 English words, no quotes, no punctuation at end).\n"
                f"User said: {user_content[:200]}\n"
                f"Return ONLY the title."
            )}],
            max_tokens=20,
        )
        title = raw.strip().strip('"').strip("'")[:60] or user_content[:40]
    except Exception:
        return
    try:
        with _Sess(bind) as s:
            c = s.get(Conversation, conv_id)
            if c:
                c.title = title
                s.add(c)
                s.commit()
                conversations_cache.delete_prefix("list:")
    except Exception:
        pass


def _get_llm(session: Session):
    """Return the active LLM service module based on the llm_provider setting."""
    setting = session.get(_Setting, "llm_provider")
    provider = (setting.value if setting and setting.value else "claude").lower().strip()
    if provider == "kimi":
        return openai_compat
    return claude


# Old short aliases → correct full model IDs
_MODEL_ALIASES: dict[str, str] = {
    "claude-opus-4":   "claude-opus-4-6",
    "claude-sonnet-4": "claude-sonnet-4-6",
    "claude-haiku-4":  "claude-haiku-4-5-20251001",
}


def _get_selected_model(session: Session, provider: str) -> str:
    """Get the selected model for the given provider."""
    setting = session.get(_Setting, "selected_model")
    if setting and setting.value:
        model = setting.value.strip()
        return _MODEL_ALIASES.get(model, model)
    # Return default model based on provider
    if provider == "kimi":
        return "moonshot-v1-32k"
    return "claude-sonnet-4-6"


def _get_setting_value(session: Session, key: str, default: str = "") -> str:
    """Get a setting value from database."""
    setting = session.get(_Setting, key)
    return setting.value if setting and setting.value else default


def _get_float_setting(session: Session, key: str, default: float = 0.0) -> float:
    """Get a float setting value from database."""
    value = _get_setting_value(session, key)
    if value:
        try:
            return float(value)
        except ValueError:
            pass
    return default


def _get_int_setting(session: Session, key: str, default: int = 0) -> int:
    """Get an int setting value from database."""
    value = _get_setting_value(session, key)
    if value:
        try:
            return int(value)
        except ValueError:
            pass
    return default


# ── Schemas ──────────────────────────────────────────────────────────────────

class SendMessageRequest(BaseModel):
    conversation_id: Optional[int] = None
    content: str
    project_id: Optional[int] = None
    skill_id: Optional[int] = None
    rag_doc_ids: List[int] = []
    file_ids: List[int] = []


class ConversationOut(BaseModel):
    id: int
    title: str
    project_id: Optional[int]
    skill_id: Optional[int]
    created_at: datetime
    updated_at: datetime


class MessageOut(BaseModel):
    id: int
    conversation_id: int
    role: str
    content: str
    metadata_json: str = "{}"  # references: skill_id, doc_ids, etc.
    created_at: datetime
    
    @property
    def metadata(self) -> dict:
        try:
            import json
            return json.loads(self.metadata_json)
        except:
            return {}


# ── Endpoints ─────────────────────────────────────────────────────────────────

@router.get("/conversations", response_model=List[ConversationOut])
def list_conversations(
    project_id: Optional[int] = None,
    standalone: bool = False,
    session: Session = Depends(get_session),
):
    cache_key = f"list:{project_id or ''}:{'s' if standalone else ''}"
    cached = conversations_cache.get(cache_key)
    if cached is not None:
        return cached
    stmt = select(Conversation).order_by(Conversation.updated_at.desc()).limit(100)
    if project_id:
        stmt = stmt.where(Conversation.project_id == project_id)
    elif standalone:
        # Only return conversations not associated with any project
        stmt = stmt.where(Conversation.project_id == None)  # noqa: E711
    result = session.exec(stmt).all()
    conversations_cache.set(cache_key, result, _CONV_TTL)
    return result


@router.get("/conversations/{conv_id}", response_model=ConversationOut)
def get_conversation(conv_id: int, session: Session = Depends(get_session)):
    conv = session.get(Conversation, conv_id)
    if not conv:
        raise HTTPException(404, "Conversation not found")
    return conv


class CreateConversationRequest(BaseModel):
    project_id: Optional[int] = None
    skill_id: Optional[int] = None
    title: Optional[str] = None


@router.post("/conversations", response_model=ConversationOut)
def create_conversation(
    req: CreateConversationRequest,
    session: Session = Depends(get_session),
):
    conv = Conversation(project_id=req.project_id, skill_id=req.skill_id, title=req.title or "")
    session.add(conv)
    session.commit()
    session.refresh(conv)
    conversations_cache.delete_prefix("list:")
    return conv


@router.get("/conversations/{conv_id}/messages", response_model=List[MessageOut])
def get_messages(
    conv_id: int,
    limit: int = 30,
    before_id: Optional[int] = None,
    session: Session = Depends(get_session),
):
    conv = session.get(Conversation, conv_id)
    if not conv:
        raise HTTPException(404, "Conversation not found")
    stmt = (
        select(Message)
        .where(Message.conversation_id == conv_id)
        .order_by(Message.created_at.desc())
        .limit(limit)
    )
    if before_id is not None:
        stmt = stmt.where(Message.id < before_id)
    msgs = session.exec(stmt).all()
    msgs.reverse()  # restore chronological order for the client
    return msgs


@router.post("/send")
async def send_message(req: SendMessageRequest, session: Session = Depends(get_session)):
    """Stream Claude response via SSE. Creates conversation if needed."""

    # Get or create conversation
    if req.conversation_id:
        conv = session.get(Conversation, req.conversation_id)
        if not conv:
            raise HTTPException(404, "Conversation not found")
    else:
        conv = Conversation(project_id=req.project_id, skill_id=req.skill_id)
        session.add(conv)
        session.commit()
        session.refresh(conv)

    # Persist user message with metadata (references)
    metadata = {}
    if req.skill_id:
        metadata["skill_id"] = req.skill_id
    if req.rag_doc_ids:
        metadata["doc_ids"] = req.rag_doc_ids
    if req.file_ids:
        metadata["file_ids"] = req.file_ids
    if req.project_id:
        metadata["project_id"] = req.project_id
    
    user_msg = Message(
        conversation_id=conv.id, 
        role="user", 
        content=req.content,
        metadata_json=json.dumps(metadata, ensure_ascii=False) if metadata else "{}"
    )
    session.add(user_msg)
    session.commit()

    # Build system prompt and load tools
    skill_prompt = ""
    tools = None
    # Read max_tokens and temperature from DB settings (fallback to defaults)
    max_tokens = _get_int_setting(session, "max_tokens", 4096) or 4096
    temperature = _get_float_setting(session, "temperature", 0.7) or 0.7
    if req.skill_id:
        skill = session.get(Skill, req.skill_id)
        if skill:
            skill_prompt = skill.system_prompt
            max_tokens = skill.max_tokens or max_tokens
            # Load tools from skill
            if skill.tools_definition_json and skill.tools_definition_json.strip():
                try:
                    tools = format_tools_for_claude(json.loads(skill.tools_definition_json))
                except json.JSONDecodeError:
                    pass

    project_context = ""
    if not req.project_id:
        # ── Global workspace context ──────────────────────────────────────────
        # Inject a compact overview of all non-archived projects so the AI can
        # answer portfolio-level questions (progress, milestones, risks, etc.)
        all_projects = session.exec(
            select(Project).where(Project.status != "archived").order_by(Project.updated_at.desc())
        ).all()
        if all_projects:
            today_str = datetime.utcnow().strftime("%Y-%m-%d")
            ws_lines = [
                f"# 工作台全局数据（截至 {today_str}）",
                f"当前共有 {len(all_projects)} 个活跃项目。以下是每个项目的详细信息：\n",
            ]
            for p in all_projects:
                ws_lines.append(f"## 项目：{p.name}")
                ws_lines.append(f"- 客户：{p.client}")
                ws_lines.append(f"- 阶段：{p.status}")
                if p.contract_amount:
                    ws_lines.append(f"- 合同金额：¥{p.contract_amount:,.0f}")
                if p.description:
                    ws_lines.append(f"- 简介：{p.description[:200]}")
                if p.context_summary:
                    ws_lines.append(f"- AI 摘要：{p.context_summary[:300]}")
                if p.notes:
                    ws_lines.append(f"- 项目笔记：{p.notes[:200]}")
                # Milestones
                milestones = session.exec(
                    select(Milestone).where(Milestone.project_id == p.id)
                ).all()
                if milestones:
                    done = sum(1 for m in milestones if m.is_done)
                    overdue = [m for m in milestones if not m.is_done and m.due_date and m.due_date < today_str]
                    ws_lines.append(f"- 里程碑：{done}/{len(milestones)} 已完成" + (f"，{len(overdue)} 个已逾期" if overdue else ""))
                    for m in milestones:
                        status_icon = "✓" if m.is_done else ("⚠ 逾期" if m.due_date and m.due_date < today_str else "○")
                        due = f"（截止 {m.due_date}）" if m.due_date else ""
                        priority = "【高优先级】" if m.priority == "high" else ""
                        ws_lines.append(f"  {status_icon} {m.title}{priority}{due}")
                # Financials
                payments = session.exec(
                    select(ProjectPayment).where(ProjectPayment.project_id == p.id)
                ).all()
                if payments:
                    received = sum(pay.amount for pay in payments if pay.payment_type == "received")
                    expense = sum(pay.amount for pay in payments if pay.payment_type == "expense")
                    if received or expense:
                        ws_lines.append(f"- 财务：已收款 ¥{received:,.0f}，支出 ¥{abs(expense):,.0f}")
                ws_lines.append("")
            project_context = "\n".join(ws_lines)

    if req.project_id:
        project = session.get(Project, req.project_id)
        if project:
            lines = [
                f"**Project Name:** {project.name}",
                f"**Client:** {project.client}",
                f"**Status:** {project.status}",
            ]
            if project.description:
                lines.append(f"**Description:** {project.description}")
            if project.contract_amount:
                lines.append(f"**Contract Amount:** ¥{project.contract_amount:,.0f}")

            # AI-generated context summary (V1.1)
            if project.context_summary:
                lines.append(f"\n**Project Context Summary:**\n{project.context_summary}")

            # Accumulated project notes (V1.1)
            if project.notes:
                lines.append(f"\n**Project Notes:**\n{project.notes}")

            milestones = session.exec(
                select(Milestone).where(Milestone.project_id == project.id)
            ).all()
            if milestones:
                lines.append("\n**Milestones:**")
                for m in milestones:
                    status_icon = "✓" if m.is_done else "○"
                    due = f" (due: {m.due_date})" if m.due_date else ""
                    priority = f" [{m.priority} priority]" if m.priority == "high" else ""
                    lines.append(f"  {status_icon} {m.title}{priority}{due}")

            files = session.exec(
                select(ProjectFile).where(ProjectFile.project_id == project.id)
            ).all()
            file_content_sections = []
            if files:
                lines.append("\n**Uploaded Documents (full content auto-injected below):**")
                total_chars = 0
                MAX_TOTAL_CHARS = 40000  # cap total injected content to ~10k tokens
                for f in files:
                    summary_hint = f" — {f.summary[:80]}" if f.summary else ""
                    lines.append(f"  - {f.name} ({f.file_type.upper()}){summary_hint}")
                    # Auto-inject readable file content (skip if budget exceeded)
                    if f.file_type.lower() in ("pdf", "docx", "pptx", "xlsx", "xls", "txt", "md", "csv", "json") and total_chars < MAX_TOTAL_CHARS:
                        full_path = UPLOADS_DIR / f.path
                        text = _extract_file_text(full_path, f.file_type, max_chars=min(8000, MAX_TOTAL_CHARS - total_chars))
                        if text and not text.startswith("["):
                            file_content_sections.append(f"### {f.name}\n{text}")
                            total_chars += len(text)

            payments = session.exec(
                select(ProjectPayment).where(ProjectPayment.project_id == project.id)
            ).all()
            if payments:
                received = sum(p.amount for p in payments if p.payment_type == "received")
                expense = sum(p.amount for p in payments if p.payment_type == "expense")
                if received or expense:
                    lines.append(f"\n**Financials:** Received ¥{received:,.0f} | Expenses ¥{abs(expense):,.0f}")

            project_context = "\n".join(lines)

            # Append auto-injected file contents after the summary block
            if file_content_sections:
                project_context += "\n\n## Project File Contents\n" + "\n\n---\n\n".join(file_content_sections)

    rag_context = ""
    if req.rag_doc_ids:
        rag_context = rag.retrieve(req.content, session, req.rag_doc_ids)
    elif "#doc" in req.content:
        rag_context = rag.retrieve(req.content, session)

    # Inject selected project file contents
    print(f"[chat] file_ids received: {req.file_ids}", flush=True)
    if req.file_ids:
        file_sections = []
        for fid in req.file_ids:
            pf = session.get(ProjectFile, fid)
            if pf:
                full_path = UPLOADS_DIR / pf.path
                text = _extract_file_text(full_path, pf.file_type)
                file_sections.append(f"### {pf.name}\n{text}")
        if file_sections:
            attachment_block = "\n\n---\n\n".join(file_sections)
            project_context = (project_context + "\n\n## Attached Files\n" + attachment_block
                               if project_context else "## Attached Files\n" + attachment_block)

    llm = _get_llm(session)
    provider = "kimi" if llm == openai_compat else "claude"
    selected_model = _get_selected_model(session, provider)
    system = llm.build_system_prompt(skill_prompt, rag_context, project_context)

    # Build message history — skip empty assistant messages (from prior failures)
    history = session.exec(
        select(Message)
        .where(Message.conversation_id == conv.id)
        .order_by(Message.created_at)
    ).all()
    api_messages = [
        {"role": m.role, "content": m.content}
        for m in history
        if m.content.strip()
    ]

    async def event_stream():
        conv_id = conv.id
        yield f"data: {json.dumps({'type': 'conversation_id', 'id': conv_id})}\n\n"

        try:
            # ── Phase 1: stream first Claude turn ────────────────────────────
            # claude.stream_response yields either:
            #   • plain text chunks  (stream to user immediately)
            #   • complete tool_use JSON strings  (intercept, do NOT stream)
            text_buffer = ""          # user-visible text from this turn
            tool_use_blocks = []      # collected tool_use dicts
            reasoning_content = ""   # kimi-k2.5 reasoning (needed for multi-turn tool calls)

            print(f"[P1] starting stream, tools={[t.get('name') for t in (tools or [])]}", flush=True)

            async for chunk in llm.stream_response(
                api_messages, system=system, model=selected_model, tools=tools, max_tokens=max_tokens, temperature=temperature
            ):
                stripped = chunk.strip()
                # 检查是否是提前通知前端工具正在生成的 marker
                if stripped.startswith("[TOOL_START:") and stripped.endswith("]"):
                    tool_name = stripped[12:-1]
                    progress_msg = f"Generating 15 slides... (this may take 1-2 minutes)"
                    yield f"data: {json.dumps({'type': 'tool_executing', 'tool_name': tool_name, 'message': progress_msg})}\n\n"
                    continue

                # Detect complete tool_use JSON emitted by claude.py / openai_compat.py
                if (
                    stripped.startswith("{")
                    and stripped.endswith("}")
                    and '"type"' in stripped
                ):
                    try:
                        block = json.loads(stripped)
                        if block.get("type") == "tool_use":
                            print(f"[P1] tool_use detected: {block.get('name')}, id={block.get('id')}, input_keys={list(block.get('input', {}).keys())}", flush=True)
                            tool_use_blocks.append(block)
                            continue  # do NOT yield to frontend
                        if block.get("type") == "reasoning_content":
                            reasoning_content = block.get("content", "")
                            continue  # internal only, not sent to frontend
                    except json.JSONDecodeError:
                        pass  # not valid JSON, treat as text

                # Regular text chunk
                text_buffer += chunk
                yield f"data: {json.dumps({'type': 'text', 'content': chunk})}\n\n"

            print(f"[P1] done. text_len={len(text_buffer)}, tool_use_count={len(tool_use_blocks)}", flush=True)

            # ── Phase 2: execute tools if any ────────────────────────────────
            executed_results = []   # raw registry.execute() outputs
            tool_result_blocks = [] # Anthropic-format tool_result blocks

            for tool_data in tool_use_blocks:
                tool_name  = tool_data.get("name", "")
                tool_input = tool_data.get("input", {})
                tool_id    = tool_data.get("id", "")

                if not tool_name or not isinstance(tool_input, dict):
                    continue

                # Progress notification to frontend
                if tool_name in ("generate_ppt", "generate_ppt_from_skill"):
                    slides     = tool_input.get("slides", [])
                    slide_count = len(slides)
                    title      = tool_input.get("title", "Untitled")
                    progress   = {
                        "message": f"Generating \"{title}\" ({slide_count} slides)…",
                        "total": slide_count, "current": 0,
                    }
                elif tool_name == "generate_docx":
                    progress = {"message": f"Generating document \"{tool_input.get('title', 'Untitled')}\"…"}
                elif tool_name == "generate_xlsx":
                    progress = {"message": f"Generating spreadsheet ({len(tool_input.get('sheets', []))} sheets)…"}
                elif tool_name == "generate_pdf":
                    progress = {"message": f"Generating PDF \"{tool_input.get('title', 'Untitled')}\"…"}
                else:
                    progress = {"message": f"Executing {tool_name}…"}

                yield f"data: {json.dumps({'type': 'tool_executing', 'tool_name': tool_name, **progress})}\n\n"

                print(f"[P2] executing tool: {tool_name}, input_keys={list(tool_input.keys())}", flush=True)
                try:
                    result = await registry.execute(tool_name, tool_input)
                except Exception as exc:
                    result = {"type": "tool_result", "tool_name": tool_name,
                              "status": "error", "error": str(exc)}

                print(f"[P2] tool result: status={result.get('status')}, keys={list(result.keys())}", flush=True)
                executed_results.append(result)
                yield f"data: {json.dumps({'type': 'tool_result', 'result': result})}\n\n"

                # Build Anthropic-format tool_result block
                output = result.get("output", result)
                tool_result_blocks.append({
                    "type": "tool_result",
                    "tool_use_id": tool_id,
                    "content": json.dumps(output, ensure_ascii=False),
                })

            print(f"[P2] done. tool_result_blocks={len(tool_result_blocks)}", flush=True)

            # ── Phase 3: optional follow-up turn after tool execution ─────────
            follow_up_text = ""
            if tool_use_blocks and tool_result_blocks:
                # Build proper Anthropic multi-turn messages
                # Assistant turn: text (if any) + tool_use blocks
                assistant_content: list = []
                if text_buffer.strip():
                    assistant_content.append({"type": "text", "text": text_buffer.strip()})
                for tb in tool_use_blocks:
                    assistant_content.append({
                        "type": "tool_use",
                        "id":    tb["id"],
                        "name":  tb["name"],
                        "input": tb.get("input", {}),
                    })

                continuation_messages = api_messages + [
                    {"role": "assistant", "content": assistant_content,
                     **({"reasoning_content": reasoning_content} if reasoning_content else {})},
                    {"role": "user",      "content": tool_result_blocks},
                ]

                print(f"[P3] starting follow-up. continuation_messages={len(continuation_messages)}", flush=True)
                # Stream follow-up response (no tools needed)
                async for chunk in llm.stream_response(
                    continuation_messages, system=system, model=selected_model,
                    tools=None, max_tokens=max_tokens, temperature=temperature
                ):
                    follow_up_text += chunk
                    yield f"data: {json.dumps({'type': 'text', 'content': chunk})}\n\n"

                print(f"[P3] done. follow_up_text_len={len(follow_up_text)}", flush=True)

            # ── Phase 4: persist ─────────────────────────────────────────────
            # Save user-visible text only (no tool_use JSON blobs)
            full_text = text_buffer.strip()
            if follow_up_text.strip():
                full_text = (full_text + "\n\n" + follow_up_text.strip()).strip()

            print(f"[P4] persisting. full_text_len={len(full_text)}", flush=True)
            need_title = False
            if full_text:
                with Session(session.get_bind()) as new_session:
                    asst_msg = Message(
                        conversation_id=conv_id,
                        role="assistant",
                        content=full_text,
                    )
                    new_session.add(asst_msg)
                    c = new_session.get(Conversation, conv_id)
                    if c:
                        c.updated_at = datetime.utcnow()
                        if c.title == "New Workstream":
                            # Placeholder — real title generated in background
                            c.title = req.content[:40] + ("…" if len(req.content) > 40 else "")
                            need_title = True
                        new_session.add(c)
                    new_session.commit()

        except Exception as e:
            import traceback
            print(f"[event_stream error] {e}\n{traceback.format_exc()}", flush=True)
            
            # Provide user-friendly error messages
            error_msg = str(e)
            user_friendly_msg = error_msg
            
            # Check for specific error patterns
            if "429" in error_msg or "engine_overloaded" in error_msg:
                user_friendly_msg = "AI 服务当前繁忙，请稍后重试。这是临时状况，几秒钟后再试即可。"
            elif "Kimi 服务当前繁忙" in error_msg:
                user_friendly_msg = error_msg  # Already user-friendly
            elif "No Kimi API key" in error_msg or "No Claude API key" in error_msg:
                user_friendly_msg = "请先配置 API Key。前往「设置」页面添加您的 API Key。"
            elif "timeout" in error_msg.lower() or "Connection refused" in error_msg:
                user_friendly_msg = "连接超时，请检查网络或稍后重试。"
            elif "rate limit" in error_msg.lower():
                user_friendly_msg = "请求频率过高，请稍等片刻后重试。"
            
            yield f"data: {json.dumps({'type': 'error', 'message': user_friendly_msg})}\n\n"
            return

        # Send done immediately — don't wait for title generation
        yield f"data: {json.dumps({'type': 'done'}, ensure_ascii=False)}\n\n"

        # Generate title in background (after done is already sent to client)
        if need_title and full_text:
            asyncio.ensure_future(_generate_title_bg(
                conv_id=conv_id,
                user_content=req.content,
                bind=session.get_bind(),
                complete_fn=llm.complete,
            ))

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@router.delete("/conversations/{conv_id}")
def delete_conversation(conv_id: int, session: Session = Depends(get_session)):
    conv = session.get(Conversation, conv_id)
    if not conv:
        raise HTTPException(404, "Conversation not found")
    for m in session.exec(select(Message).where(Message.conversation_id == conv_id)).all():
        session.delete(m)
    session.delete(conv)
    session.commit()
    conversations_cache.delete_prefix("list:")
    return {"ok": True}


# ---------------------------------------------------------------------------
# Test Connection Endpoints
# ---------------------------------------------------------------------------

class TestConnectionRequest(BaseModel):
    provider: str
    model: Optional[str] = None


@router.post("/test-connection")
async def test_connection(req: TestConnectionRequest):
    """Test API key connectivity for a provider."""
    from app.core.security import get_api_key, get_kimi_api_key
    
    provider = req.provider
    
    # Only support anthropic and moonshot for now
    if provider not in ["anthropic", "moonshot"]:
        return {"success": False, "message": f"Provider not supported: {provider}"}
    
    try:
        if provider == "anthropic":
            api_key = get_api_key()
            if not api_key:
                return {"success": False, "message": "No API key configured"}
            # Test with a simple request
            import httpx
            async with httpx.AsyncClient() as client:
                resp = await client.post(
                    "https://api.anthropic.com/v1/messages",
                    headers={
                        "x-api-key": api_key,
                        "anthropic-version": "2023-06-01",
                        "content-type": "application/json",
                    },
                    json={
                        "model": req.model or "claude-3-5-haiku-20241022",
                        "max_tokens": 10,
                        "messages": [{"role": "user", "content": "Hi"}],
                    },
                    timeout=30.0,
                )
                if resp.status_code == 200:
                    return {"success": True, "message": "Connection successful"}
                else:
                    return {"success": False, "message": f"API error: {resp.status_code}"}
                    
        elif provider == "moonshot":
            api_key = get_kimi_api_key()
            if not api_key:
                return {"success": False, "message": "No API key configured"}
            import httpx
            async with httpx.AsyncClient() as client:
                resp = await client.post(
                    "https://api.moonshot.cn/v1/chat/completions",
                    headers={"Authorization": f"Bearer {api_key}"},
                    json={
                        "model": req.model or "kimi-k2-0711-preview",
                        "messages": [{"role": "user", "content": "Hi"}],
                        "max_tokens": 10,
                    },
                    timeout=30.0,
                )
                if resp.status_code == 200:
                    return {"success": True, "message": "Connection successful"}
                else:
                    return {"success": False, "message": f"API error: {resp.status_code}"}
        else:
            return {"success": False, "message": f"Unknown provider: {provider}"}
            
    except Exception as e:
        return {"success": False, "message": f"Connection failed: {str(e)}"}


class TestModelRequest(BaseModel):
    message: str
    model: str
    temperature: float = 0.7
    max_tokens: int = 100


@router.post("/test-model")
async def test_model(req: TestModelRequest):
    """Test a model with a simple message."""
    try:
        # Determine provider from model
        provider = "anthropic"
        if req.model.startswith("moonshot-"):
            provider = "moonshot"
        elif req.model.startswith("claude-"):
            provider = "anthropic"
        else:
            return {"success": False, "message": f"Model not supported: {req.model}"}
        
        # Only support anthropic and moonshot for now
        if provider not in ["anthropic", "moonshot"]:
            return {"success": False, "message": f"Provider not supported: {provider}"}
        
        # Get the appropriate LLM client
        if provider == "anthropic":
            from app.services import claude as llm
        elif provider == "moonshot":
            from app.services import openai_compat as llm
        else:
            return {"success": False, "message": f"Unsupported provider: {provider}"}
        
        # Make a simple completion
        messages = [{"role": "user", "content": req.message}]
        response = await llm.complete(
            messages=messages,
            model=req.model,
            temperature=req.temperature,
            max_tokens=req.max_tokens,
        )
        
        return {
            "success": True,
            "message": "Model test successful",
            "response": response[:200] + "..." if len(response) > 200 else response,
        }
    except Exception as e:
        return {"success": False, "message": f"Model test failed: {str(e)}"}
