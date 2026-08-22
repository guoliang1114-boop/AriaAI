from __future__ import annotations

import asyncio
import shutil
import uuid
from pathlib import Path
from typing import Any, Literal

from fastapi import HTTPException
from sqlmodel import Session, select

from app.config import UPLOADS_DIR
from app.database import engine
from app.models.db import Project, ProjectFile, ProjectFolder
from app.services.cache import projects_cache
from app.services.document_text import extract_text_from_file
from app.services.project_contexts import mark_project_memory_stale
from app.services.project_core import init_default_project_folders
from app.services.project_documents import infer_project_folder
from app.services.project_files import active_project_files_stmt, archive_project_file
from app.services.tool_descriptions import tool_description
from app.tools import registry
from app.tools.file_generators import generate_docx, generate_pdf, generate_ppt, generate_xlsx

READ_PROJECT_FILE_TOOL_NAME = "read_project_file"
WRITE_PROJECT_OFFICE_DOCUMENT_TOOL_NAME = "write_project_office_document"
EDIT_PROJECT_OFFICE_DOCUMENT_TOOL_NAME = "edit_project_office_document"
MANAGE_PROJECT_FOLDERS_TOOL_NAME = "manage_project_folders"
MANAGE_PROJECT_FILES_TOOL_NAME = "manage_project_files"

_READABLE_TYPES = {"pdf", "docx", "pptx", "xlsx", "xls", "md", "txt", "csv", "json"}
_WRITABLE_TYPES = {"docx", "xlsx", "pptx", "pdf"}
_READ_MAX_CHARS = 20000


def _bust_project_cache(project_id: int) -> None:
    projects_cache.delete(f"detail:{project_id}")
    projects_cache.delete_prefix("list:")


def _safe_name(name: str, extension: str) -> str:
    stem = "_".join((name or "document").strip().split())
    for char in '/\\:*?"<>|':
        stem = stem.replace(char, "_")
    stem = stem.strip("._")[:96] or "document"
    suffix = f".{extension.lower()}"
    if not stem.lower().endswith(suffix):
        stem = f"{stem}{suffix}"
    return stem


def _file_path(project_file: ProjectFile) -> Path:
    full_path = UPLOADS_DIR / Path(project_file.path)
    try:
        full_path.resolve().relative_to(UPLOADS_DIR.resolve())
    except ValueError as exc:
        raise HTTPException(400, "Invalid project file path") from exc
    if not full_path.is_file():
        raise HTTPException(404, "File not found on disk")
    return full_path


def _find_project_file(session: Session, project_id: int, file_id: int | None, file_name: str | None) -> ProjectFile:
    if file_id is not None:
        project_file = session.get(ProjectFile, file_id)
        if not project_file or project_file.project_id != project_id or project_file.deleted_at is not None:
            raise HTTPException(404, "File not found")
        return project_file

    normalized = (file_name or "").strip().lower()
    if not normalized:
        raise HTTPException(400, "Provide file_id or file_name")

    files = session.exec(active_project_files_stmt(project_id)).all()
    for project_file in files:
        if project_file.name.strip().lower() == normalized:
            return project_file
    raise HTTPException(404, "File not found")


def _list_files(session: Session, project_id: int, file_types: list[str] | None) -> dict:
    normalized_types = {item.lower().lstrip(".") for item in file_types or [] if item}
    stmt = active_project_files_stmt(project_id)
    if normalized_types:
        stmt = stmt.where(ProjectFile.file_type.in_(normalized_types))
    files = session.exec(stmt.order_by(ProjectFile.uploaded_at.desc(), ProjectFile.id.desc())).all()

    folder_names: dict[int, str] = {}
    folder_ids = {item.folder_id for item in files if item.folder_id is not None}
    if folder_ids:
        folders = session.exec(select(ProjectFolder).where(ProjectFolder.id.in_(folder_ids))).all()
        folder_names = {folder.id: folder.name for folder in folders}

    return {
        "ok": True,
        "count": len(files),
        "files": [
            {
                "id": item.id,
                "name": item.name,
                "file_type": item.file_type,
                "folder_id": item.folder_id,
                "folder": folder_names.get(item.folder_id, "") if item.folder_id else "",
                "summary": item.summary or "",
                "origin": item.origin,
                "size_bytes": item.size_bytes,
            }
            for item in files
        ],
    }


def _normalize_file_ids(file_ids: list[int] | None, file_id: int | None) -> list[int]:
    normalized: list[int] = []
    for value in file_ids or []:
        try:
            parsed = int(value)
        except (TypeError, ValueError):
            continue
        if parsed > 0 and parsed not in normalized:
            normalized.append(parsed)
    if file_id is not None:
        try:
            parsed_file_id = int(file_id)
        except (TypeError, ValueError):
            parsed_file_id = 0
        if parsed_file_id > 0 and parsed_file_id not in normalized:
            normalized.append(parsed_file_id)
    return normalized


def _list_folders(session: Session, project_id: int) -> list[ProjectFolder]:
    folders = session.exec(
        select(ProjectFolder)
        .where(ProjectFolder.project_id == project_id)
        .order_by(ProjectFolder.sort_order, ProjectFolder.id)
    ).all()
    if not folders:
        folders = init_default_project_folders(session, project_id)
    return folders


def _folder_payload(session: Session, project_id: int) -> dict:
    folders = _list_folders(session, project_id)
    counts: dict[int, int] = {folder.id: 0 for folder in folders if folder.id is not None}
    files = session.exec(active_project_files_stmt(project_id)).all()
    unassigned_count = 0
    for project_file in files:
        if project_file.folder_id in counts:
            counts[project_file.folder_id] += 1
        else:
            unassigned_count += 1
    return {
        "ok": True,
        "folders": [
            {
                "id": folder.id,
                "name": folder.name,
                "sort_order": folder.sort_order,
                "file_count": counts.get(folder.id, 0),
            }
            for folder in folders
        ],
        "unassigned_file_count": unassigned_count,
    }


def _find_folder_by_name(session: Session, project_id: int, folder_name: str | None) -> ProjectFolder | None:
    normalized = (folder_name or "").strip().lower()
    if not normalized:
        return None
    folders = _list_folders(session, project_id)
    for folder in folders:
        if folder.name.strip().lower() == normalized:
            return folder
    return None


def _resolve_folder_target(
    session: Session,
    project_id: int,
    *,
    folder_id: int | None = None,
    folder_name: str | None = None,
) -> ProjectFolder:
    if folder_id is not None:
        folder = session.get(ProjectFolder, folder_id)
        if not folder or folder.project_id != project_id:
            raise HTTPException(404, "Folder not found")
        return folder
    folder = _find_folder_by_name(session, project_id, folder_name)
    if folder:
        return folder
    raise HTTPException(404, "Folder not found")


def _next_folder_sort_order(session: Session, project_id: int) -> int:
    folders = _list_folders(session, project_id)
    if not folders:
        return 0
    return max(folder.sort_order for folder in folders) + 1


@registry.register(
    name=MANAGE_PROJECT_FOLDERS_TOOL_NAME,
    description=tool_description(
        MANAGE_PROJECT_FOLDERS_TOOL_NAME,
        "Manage project space folders and file placement. "
        "Use it to list folders, create or rename folders, move generated files into the right folder, or delete empty folders."
    ),
    input_schema={
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["list", "create", "rename", "move_file", "delete"],
                "description": "Folder operation to perform.",
            },
            "folder_id": {"type": "integer", "description": "Target folder id."},
            "folder_name": {"type": "string", "description": "Target folder name when id is not known."},
            "new_name": {"type": "string", "description": "New name for create or rename."},
            "sort_order": {"type": "integer", "description": "Optional folder order for create/rename."},
            "file_id": {"type": "integer", "description": "Project file id for move_file."},
            "file_name": {"type": "string", "description": "Project file name for move_file when file_id is not known."},
        },
        "required": ["action"],
    },
)
async def manage_project_folders(
    *,
    project_id: int,
    action: Literal["list", "create", "rename", "move_file", "delete"],
    folder_id: int | None = None,
    folder_name: str | None = None,
    new_name: str | None = None,
    sort_order: int | None = None,
    file_id: int | None = None,
    file_name: str | None = None,
) -> dict:
    if not project_id:
        raise HTTPException(400, "Project id is required")

    with Session(engine) as session:
        project = session.get(Project, project_id)
        if not project:
            raise HTTPException(404, "Project not found")

        if action == "list":
            return _folder_payload(session, project_id)

        if action == "create":
            name = (new_name or folder_name or "").strip()
            if not name:
                raise HTTPException(400, "Folder name is required")
            existing = _find_folder_by_name(session, project_id, name)
            if existing:
                return {"ok": True, "action": "exists", "folder": {"id": existing.id, "name": existing.name, "sort_order": existing.sort_order}}
            folder = ProjectFolder(
                project_id=project_id,
                name=name,
                sort_order=sort_order if sort_order is not None else _next_folder_sort_order(session, project_id),
            )
            session.add(folder)
            session.commit()
            session.refresh(folder)
            mark_project_memory_stale(session, project_id, trigger="folder_tool_create")
            _bust_project_cache(project_id)
            return {"ok": True, "action": "created", "folder": {"id": folder.id, "name": folder.name, "sort_order": folder.sort_order}}

        if action == "rename":
            folder = _resolve_folder_target(session, project_id, folder_id=folder_id, folder_name=folder_name)
            name = (new_name or "").strip()
            if not name:
                raise HTTPException(400, "new_name is required")
            folder.name = name
            if sort_order is not None:
                folder.sort_order = sort_order
            session.add(folder)
            session.commit()
            session.refresh(folder)
            mark_project_memory_stale(session, project_id, trigger="folder_tool_rename")
            _bust_project_cache(project_id)
            return {"ok": True, "action": "renamed", "folder": {"id": folder.id, "name": folder.name, "sort_order": folder.sort_order}}

        if action == "move_file":
            project_file = _find_project_file(session, project_id, file_id, file_name)
            folder = _resolve_folder_target(session, project_id, folder_id=folder_id, folder_name=folder_name)
            project_file.folder_id = folder.id
            session.add(project_file)
            session.commit()
            session.refresh(project_file)
            mark_project_memory_stale(session, project_id, trigger="folder_tool_move_file")
            _bust_project_cache(project_id)
            return {
                "ok": True,
                "action": "moved",
                "file": {"id": project_file.id, "name": project_file.name, "folder_id": project_file.folder_id},
                "folder": {"id": folder.id, "name": folder.name},
            }

        # action == "delete"
        folder = _resolve_folder_target(session, project_id, folder_id=folder_id, folder_name=folder_name)
        file_count = session.exec(
            select(ProjectFile).where(ProjectFile.folder_id == folder.id, ProjectFile.deleted_at.is_(None))
        ).all()
        if file_count:
            raise HTTPException(400, "Folder is not empty. Move files before deleting it.")
        session.delete(folder)
        session.commit()
        mark_project_memory_stale(session, project_id, trigger="folder_tool_delete")
        _bust_project_cache(project_id)
        return {"ok": True, "action": "deleted", "folder_id": folder.id}


@registry.register(
    name=MANAGE_PROJECT_FILES_TOOL_NAME,
    description=tool_description(
        MANAGE_PROJECT_FILES_TOOL_NAME,
        "List project-space files and delete explicitly selected project files after user confirmation. "
        "Use action='list' to inspect space clutter. Use action='delete' only with concrete file_ids."
    ),
    input_schema={
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["list", "delete"],
                "description": "list returns project files; delete removes selected files by id.",
            },
            "file_id": {
                "type": "integer",
                "description": "Single project file id to delete. Prefer file_ids for multiple files.",
            },
            "file_ids": {
                "type": "array",
                "items": {"type": "integer"},
                "description": "Project file ids to delete. Required for bulk cleanup.",
            },
            "file_types": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Optional extensions for list, for example ['pdf', 'docx', 'xlsx'].",
            },
            "reason": {
                "type": "string",
                "description": "Short deletion rationale shown in trace and confirmation UI.",
            },
        },
        "required": ["action"],
    },
)
async def manage_project_files(
    *,
    project_id: int,
    action: Literal["list", "delete"],
    file_id: int | None = None,
    file_ids: list[int] | None = None,
    file_types: list[str] | None = None,
    reason: str | None = None,
) -> dict:
    if not project_id:
        raise HTTPException(400, "Project id is required")

    with Session(engine) as session:
        project = session.get(Project, project_id)
        if not project:
            raise HTTPException(404, "Project not found")

        if action == "list":
            return _list_files(session, project_id, file_types)

        ids = _normalize_file_ids(file_ids, file_id)
        if not ids:
            raise HTTPException(400, "Provide file_ids to delete project files")
        if len(ids) > 50:
            raise HTTPException(400, "At most 50 files can be deleted at once")

        files = session.exec(
            select(ProjectFile).where(
                ProjectFile.project_id == project_id,
                ProjectFile.id.in_(ids),
                ProjectFile.deleted_at.is_(None),
            )
        ).all()
        found_by_id = {item.id: item for item in files if item.id is not None}
        missing_ids = [item for item in ids if item not in found_by_id]
        if missing_ids:
            raise HTTPException(404, f"Files not found: {missing_ids}")

        deleted_files = [
            {
                "id": item.id,
                "name": item.name,
                "file_type": item.file_type,
                "size_bytes": item.size_bytes,
            }
            for item in files
        ]
        delete_batch_id = uuid.uuid4().hex
        for target_id in ids:
            archive_project_file(
                session,
                project_id,
                target_id,
                reason=reason or "HITAS approved project file cleanup",
                batch_id=delete_batch_id,
            )

        mark_project_memory_stale(session, project_id, trigger="file_tool_delete")
        _bust_project_cache(project_id)
        return {
            "ok": True,
            "action": "archived",
            "deleted_count": len(deleted_files),
            "deleted_files": deleted_files,
            "trash": True,
            "reason": reason or "",
        }


def _read_project_file_sync(
    *,
    project_id: int,
    action: Literal["list", "read"],
    file_id: int | None = None,
    file_name: str | None = None,
    file_types: list[str] | None = None,
    max_chars: int = _READ_MAX_CHARS,
) -> dict:
    if not project_id:
        raise HTTPException(400, "Project id is required")

    with Session(engine) as session:
        if action == "list":
            return _list_files(session, project_id, file_types)

        project_file = _find_project_file(session, project_id, file_id, file_name)
        file_type = project_file.file_type.lower().lstrip(".")
        if file_type not in _READABLE_TYPES:
            raise HTTPException(400, f"Unsupported readable file type: {project_file.file_type}")

        content = extract_text_from_file(
            _file_path(project_file),
            file_type,
            max_chars=max(1000, min(max_chars or _READ_MAX_CHARS, 60000)),
            empty_placeholder="[No text extracted]",
            unsupported_placeholder="[Unsupported file type]",
            error_prefix="Extract failed: ",
        )
        return {
            "ok": True,
            "id": project_file.id,
            "name": project_file.name,
            "file_type": project_file.file_type,
            "size_bytes": project_file.size_bytes,
            "truncated": content.endswith("\n…[truncated]"),
            "content": content,
        }


@registry.register(
    name=READ_PROJECT_FILE_TOOL_NAME,
    description=tool_description(
        READ_PROJECT_FILE_TOOL_NAME,
        "List or read files in the current project. Supports PDF, DOCX, PPTX, XLSX/XLS, Markdown, text, CSV, and JSON. "
        "Use action='list' to discover files. Use action='read' with file_id or file_name to extract text for analysis."
    ),
    input_schema={
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["list", "read"],
                "description": "list returns project files; read extracts text from one file.",
            },
            "file_id": {
                "type": "integer",
                "description": "Project file id to read. Prefer this when available.",
            },
            "file_name": {
                "type": "string",
                "description": "Project file name to read when file_id is not known.",
            },
            "file_types": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Optional extensions for list, for example ['pdf', 'docx', 'xlsx'].",
            },
            "max_chars": {
                "type": "integer",
                "description": "Maximum extracted characters for read. Default 20000.",
            },
        },
        "required": ["action"],
    },
)
async def read_project_file(
    *,
    project_id: int,
    action: Literal["list", "read"],
    file_id: int | None = None,
    file_name: str | None = None,
    file_types: list[str] | None = None,
    max_chars: int = _READ_MAX_CHARS,
) -> dict:
    """Read project files off the event loop for safe parallel batches."""

    return await asyncio.to_thread(
        _read_project_file_sync,
        project_id=project_id,
        action=action,
        file_id=file_id,
        file_name=file_name,
        file_types=file_types,
        max_chars=max_chars,
    )


def _content_preview(file_type: str, *, title: str, content: str, sections: list[dict] | None, sheets: list[dict] | None) -> str:
    if content:
        return content
    if sections:
        return "\n\n".join(f"{item.get('heading', '')}\n{item.get('content', '')}" for item in sections)
    if sheets:
        names = ", ".join(str(item.get("name", "")) for item in sheets)
        return f"{title}\nSheets: {names}"
    return file_type


def _register_generated_project_file(
    session: Session,
    project_id: int,
    *,
    source_path: Path,
    file_name: str,
    file_type: str,
    folder_id: int | None,
    summary: str,
    preview_text: str,
) -> ProjectFile:
    project = session.get(Project, project_id)
    if not project:
        raise HTTPException(404, "Project not found")

    folder = infer_project_folder(
        session,
        project_id,
        init_default_folders=init_default_project_folders,
        preferred_folder_id=folder_id,
        name=file_name,
        content=preview_text,
        summary=summary,
    )

    dest_dir = UPLOADS_DIR / "projects" / str(project_id)
    dest_dir.mkdir(parents=True, exist_ok=True)
    safe_name = _safe_name(file_name, file_type)
    dest_path = dest_dir / f"{uuid.uuid4().hex}_{safe_name}"
    shutil.copyfile(source_path, dest_path)

    project_file = ProjectFile(
        project_id=project_id,
        folder_id=folder.id if folder else None,
        name=safe_name,
        file_type=file_type,
        path=str(dest_path.relative_to(UPLOADS_DIR)),
        size_bytes=dest_path.stat().st_size,
        summary=summary,
        origin="ai_generated",
    )
    session.add(project_file)
    session.commit()
    session.refresh(project_file)
    return project_file


@registry.register(
    name=WRITE_PROJECT_OFFICE_DOCUMENT_TOOL_NAME,
    description=tool_description(
        WRITE_PROJECT_OFFICE_DOCUMENT_TOOL_NAME,
        "Create a Word, Excel, PowerPoint, or PDF file and save it into the current project space. "
        "Use only when the user explicitly asks to create, generate, export, or save a DOCX/XLSX/PPTX/PDF deliverable. "
        "Do not use for analysis-only chat answers. Use DOCX/PDF for narrative documents, XLSX for tables, and PPTX for presentation slides."
    ),
    input_schema={
        "type": "object",
        "properties": {
            "file_type": {
                "type": "string",
                "enum": ["docx", "xlsx", "pptx", "pdf"],
                "description": "Output file type.",
            },
            "file_name": {"type": "string", "description": "Output file name."},
            "title": {"type": "string", "description": "Document or presentation title."},
            "content": {"type": "string", "description": "Markdown/text content for PDF or simple DOCX."},
            "sections": {
                "type": "array",
                "description": "DOCX sections: heading/content/level.",
                "items": {"type": "object"},
            },
            "sheets": {
                "type": "array",
                "description": "XLSX sheets: name/headers/data.",
                "items": {"type": "object"},
            },
            "slides": {
                "type": "array",
                "description": (
                    "PPTX slides. Supported type values include title, section, content, two_column, "
                    "table, quote, process, roadmap, matrix, kpi, risk, and next_steps. "
                    "Use columns/rows for tables, items/steps for process slides, quote/source for quote slides."
                ),
                "items": {"type": "object"},
            },
            "folder_id": {"type": "integer", "description": "Optional project folder id."},
            "summary": {"type": "string", "description": "Short file summary for project space."},
        },
        "required": ["file_type", "file_name"],
    },
)
async def write_project_office_document(
    *,
    project_id: int,
    file_type: Literal["docx", "xlsx", "pptx", "pdf"],
    file_name: str,
    title: str = "",
    content: str = "",
    sections: list[dict] | None = None,
    sheets: list[dict] | None = None,
    slides: list[dict] | None = None,
    folder_id: int | None = None,
    summary: str = "",
) -> dict[str, Any]:
    if not project_id:
        raise HTTPException(400, "Project id is required")
    file_type = file_type.lower()
    if file_type not in _WRITABLE_TYPES:
        raise HTTPException(400, f"Unsupported writable file type: {file_type}")

    title = title.strip() or Path(file_name).stem or "Document"
    if file_type == "docx":
        doc_sections = sections or [{"heading": title, "content": content or title, "level": 1}]
        result = await generate_docx(title=title, sections=doc_sections)
    elif file_type == "xlsx":
        if not sheets:
            raise HTTPException(400, "sheets is required for xlsx")
        result = await generate_xlsx(sheets=sheets)
    elif file_type == "pptx":
        if not slides:
            raise HTTPException(400, "slides is required for pptx")
        result = await generate_ppt(title=title, slides=slides)
    else:
        if not content.strip():
            raise HTTPException(400, "content is required for pdf")
        result = await generate_pdf(title=title, content=content)

    if not result.get("success"):
        raise HTTPException(500, result.get("error") or "Failed to generate document")

    source_path = Path(str(result.get("full_path") or ""))
    if not source_path.is_file():
        raise HTTPException(500, "Generated file not found")

    preview_text = _content_preview(file_type, title=title, content=content, sections=sections, sheets=sheets)
    with Session(engine) as session:
        project_file = _register_generated_project_file(
            session,
            project_id,
            source_path=source_path,
            file_name=file_name,
            file_type=file_type,
            folder_id=folder_id,
            summary=summary or f"AI generated {file_type.upper()} document",
            preview_text=preview_text,
        )
        mark_project_memory_stale(session, project_id, trigger=f"{file_type}_tool_create")
        _bust_project_cache(project_id)
        output = {
            "ok": True,
            "id": project_file.id,
            "name": project_file.name,
            "file_type": project_file.file_type,
            "folder_id": project_file.folder_id,
            "size_bytes": project_file.size_bytes,
            "path": project_file.path,
            "message": f"Created {project_file.name}",
        }

    return output


# ---------------------------------------------------------------------------
# Edit existing Office document
# ---------------------------------------------------------------------------

def _edit_pptx(file_path: Path, edits: list[dict]) -> dict:
    """Edit an existing PPTX file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {"success": False, "error": "python-pptx not installed"}

    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        return {"success": False, "error": f"Failed to open PPTX: {e}"}

    slide_count = len(prs.slides)
    changes = []

    for edit in edits:
        action = edit.get("action", "update_slide")
        slide_index = edit.get("slide_index")

        if action == "update_slide":
            if slide_index is None or slide_index < 0 or slide_index >= slide_count:
                changes.append(f"Skipped update_slide: invalid slide_index {slide_index}")
                continue
            slide = prs.slides[slide_index]
            new_title = edit.get("title")
            new_content = edit.get("content")

            title_updated = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                        if new_title:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    run.text = ""
                                if para.runs:
                                    para.runs[0].text = new_title
                                else:
                                    para.text = new_title
                            title_updated = True
                    elif new_content and not title_updated:
                        pass

            if new_content:
                body_shapes = [
                    s for s in slide.shapes
                    if s.has_text_frame and (not slide.shapes.title or s.shape_id != slide.shapes.title.shape_id)
                ]
                if body_shapes:
                    tf = body_shapes[0].text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.text = ""
                    lines = new_content.split("\n")
                    for i, line in enumerate(lines):
                        if i == 0:
                            if tf.paragraphs:
                                tf.paragraphs[0].text = line.lstrip("- ").lstrip("• ")
                        else:
                            p = tf.add_paragraph()
                            p.text = line.lstrip("- ").lstrip("• ")
                            p.level = 0

            changes.append(f"Updated slide {slide_index}")

        elif action == "update_text":
            if slide_index is None or slide_index < 0 or slide_index >= slide_count:
                changes.append(f"Skipped update_text: invalid slide_index {slide_index}")
                continue
            slide = prs.slides[slide_index]
            old_text = edit.get("old_text", "")
            new_text = edit.get("new_text", "")
            if not old_text:
                changes.append(f"Skipped update_text: missing old_text")
                continue
            replaced = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        full_text = para.text
                        if old_text in full_text:
                            for run in para.runs:
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)
                                    replaced += 1
            changes.append(f"Replaced {replaced} occurrences of text on slide {slide_index}")

        elif action == "add_slide":
            layout_index = edit.get("layout_index", 1)
            title = edit.get("title", "")
            content = edit.get("content", "")
            try:
                layout = prs.slide_layouts[min(layout_index, len(prs.slide_layouts) - 1)]
                slide = prs.slides.add_slide(layout)
                if title and slide.shapes.title:
                    slide.shapes.title.text = title
                if content:
                    body = [
                        s for s in slide.shapes
                        if s.has_text_frame and (not slide.shapes.title or s.shape_id != slide.shapes.title.shape_id)
                    ]
                    if body:
                        body[0].text_frame.text = content
                changes.append(f"Added new slide at index {len(prs.slides) - 1}")
            except Exception as e:
                changes.append(f"Failed to add slide: {e}")

        elif action == "delete_slide":
            if slide_index is None or slide_index < 0 or slide_index >= slide_count:
                changes.append(f"Skipped delete_slide: invalid slide_index {slide_index}")
                continue
            try:
                rId = prs.slides._sldIdLst[slide_index].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[slide_index])
                changes.append(f"Deleted slide {slide_index}")
            except Exception as e:
                changes.append(f"Failed to delete slide {slide_index}: {e}")

    try:
        prs.save(str(file_path))
    except Exception as e:
        return {"success": False, "error": f"Failed to save PPTX: {e}", "changes": changes}

    return {"success": True, "changes": changes}


def _edit_docx(file_path: Path, edits: list[dict]) -> dict:
    """Edit an existing DOCX file."""
    try:
        from docx import Document
    except ImportError:
        return {"success": False, "error": "python-docx not installed"}

    try:
        doc = Document(str(file_path))
    except Exception as e:
        return {"success": False, "error": f"Failed to open DOCX: {e}"}

    changes = []

    for edit in edits:
        action = edit.get("action", "update_text")

        if action == "update_text":
            old_text = edit.get("old_text", "")
            new_text = edit.get("new_text", "")
            if not old_text:
                changes.append("Skipped update_text: missing old_text")
                continue
            replaced = 0
            for para in doc.paragraphs:
                if old_text in para.text:
                    for run in para.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            replaced += 1
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if old_text in cell.text:
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    if old_text in run.text:
                                        run.text = run.text.replace(old_text, new_text)
                                        replaced += 1
            changes.append(f"Replaced {replaced} occurrences of '{old_text[:30]}'")

        elif action == "replace_paragraph":
            heading_text = edit.get("heading", "")
            new_content = edit.get("content", "")
            if not heading_text:
                changes.append("Skipped replace_paragraph: missing heading")
                continue
            found = False
            for i, para in enumerate(doc.paragraphs):
                if para.style.name.startswith("Heading") and heading_text.lower() in para.text.lower():
                    j = i + 1
                    while j < len(doc.paragraphs) and not doc.paragraphs[j].style.name.startswith("Heading"):
                        j += 1
                    for k in range(j - 1, i, -1):
                        p = doc.paragraphs[k]._element
                        p.getparent().remove(p)
                    if new_content:
                        for line in new_content.split("\n"):
                            new_para = doc.paragraphs[i].add_paragraph(line.lstrip("- ").lstrip("• "))
                            new_para.style = doc.styles["Body Text"] if "Body Text" in [s.name for s in doc.styles] else doc.paragraphs[i].style
                    found = True
                    changes.append(f"Replaced content under heading '{heading_text}'")
                    break
            if not found:
                changes.append(f"Heading '{heading_text}' not found")

        elif action == "add_paragraph":
            after_heading = edit.get("after_heading", "")
            content = edit.get("content", "")
            if not content:
                changes.append("Skipped add_paragraph: missing content")
                continue
            if after_heading:
                found = False
                for i, para in enumerate(doc.paragraphs):
                    if para.style.name.startswith("Heading") and after_heading.lower() in para.text.lower():
                        j = i + 1
                        while j < len(doc.paragraphs) and not doc.paragraphs[j].style.name.startswith("Heading"):
                            j += 1
                        for line in content.split("\n"):
                            new_para = doc.add_paragraph(line.lstrip("- ").lstrip("• "))
                            new_para.style = doc.styles["Body Text"] if "Body Text" in [s.name for s in doc.styles] else doc.paragraphs[0].style
                            new_para._element.addprevious(doc.paragraphs[j - 1]._element)
                        found = True
                        changes.append(f"Added paragraph after heading '{after_heading}'")
                        break
                if not found:
                    for line in content.split("\n"):
                        doc.add_paragraph(line.lstrip("- ").lstrip("• "))
                    changes.append(f"Heading '{after_heading}' not found, appended to end")
            else:
                for line in content.split("\n"):
                    doc.add_paragraph(line.lstrip("- ").lstrip("• "))
                changes.append("Appended paragraph to end")

    try:
        doc.save(str(file_path))
    except Exception as e:
        return {"success": False, "error": f"Failed to save DOCX: {e}", "changes": changes}

    return {"success": True, "changes": changes}


def _edit_xlsx(file_path: Path, edits: list[dict]) -> dict:
    """Edit an existing XLSX file."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return {"success": False, "error": "openpyxl not installed"}

    try:
        wb = load_workbook(str(file_path))
    except Exception as e:
        return {"success": False, "error": f"Failed to open XLSX: {e}"}

    changes = []

    for edit in edits:
        action = edit.get("action", "update_cell")

        if action == "update_cell":
            sheet_name = edit.get("sheet")
            cell_ref = edit.get("cell")
            value = edit.get("value")
            if not sheet_name or not cell_ref:
                changes.append("Skipped update_cell: missing sheet or cell")
                continue
            if sheet_name not in wb.sheetnames:
                changes.append(f"Skipped update_cell: sheet '{sheet_name}' not found")
                continue
            ws = wb[sheet_name]
            try:
                ws[cell_ref] = value
                changes.append(f"Updated {sheet_name}!{cell_ref} = {value}")
            except Exception as e:
                changes.append(f"Failed to update {cell_ref}: {e}")

        elif action == "update_cells":
            sheet_name = edit.get("sheet")
            updates = edit.get("updates", {})
            if not sheet_name:
                changes.append("Skipped update_cells: missing sheet")
                continue
            if sheet_name not in wb.sheetnames:
                changes.append(f"Skipped update_cells: sheet '{sheet_name}' not found")
                continue
            ws = wb[sheet_name]
            count = 0
            for cell_ref, value in updates.items():
                try:
                    ws[cell_ref] = value
                    count += 1
                except Exception:
                    pass
            changes.append(f"Updated {count} cells in {sheet_name}")

        elif action == "update_row":
            sheet_name = edit.get("sheet")
            row_num = edit.get("row")
            values = edit.get("values", [])
            if not sheet_name or row_num is None:
                changes.append("Skipped update_row: missing sheet or row")
                continue
            if sheet_name not in wb.sheetnames:
                changes.append(f"Skipped update_row: sheet '{sheet_name}' not found")
                continue
            ws = wb[sheet_name]
            for col_idx, value in enumerate(values, 1):
                ws.cell(row=row_num, column=col_idx, value=value)
            changes.append(f"Updated row {row_num} in {sheet_name} with {len(values)} values")

        elif action == "add_row":
            sheet_name = edit.get("sheet")
            values = edit.get("values", [])
            after_row = edit.get("after_row")
            if not sheet_name:
                changes.append("Skipped add_row: missing sheet")
                continue
            if sheet_name not in wb.sheetnames:
                changes.append(f"Skipped add_row: sheet '{sheet_name}' not found")
                continue
            ws = wb[sheet_name]
            if after_row:
                ws.insert_rows(after_row + 1)
                for col_idx, value in enumerate(values, 1):
                    ws.cell(row=after_row + 1, column=col_idx, value=value)
                changes.append(f"Inserted row after row {after_row} in {sheet_name}")
            else:
                next_row = ws.max_row + 1
                for col_idx, value in enumerate(values, 1):
                    ws.cell(row=next_row, column=col_idx, value=value)
                changes.append(f"Appended row to {sheet_name}")

        elif action == "delete_row":
            sheet_name = edit.get("sheet")
            row_num = edit.get("row")
            if not sheet_name or row_num is None:
                changes.append("Skipped delete_row: missing sheet or row")
                continue
            if sheet_name not in wb.sheetnames:
                changes.append(f"Skipped delete_row: sheet '{sheet_name}' not found")
                continue
            ws = wb[sheet_name]
            ws.delete_rows(row_num)
            changes.append(f"Deleted row {row_num} in {sheet_name}")

        elif action == "add_sheet":
            sheet_name = edit.get("name", "Sheet")
            try:
                wb.create_sheet(title=sheet_name)
                changes.append(f"Added sheet '{sheet_name}'")
            except Exception as e:
                changes.append(f"Failed to add sheet: {e}")

    try:
        wb.save(str(file_path))
    except Exception as e:
        return {"success": False, "error": f"Failed to save XLSX: {e}", "changes": changes}

    return {"success": True, "changes": changes}


EDITABLE_TYPES = {"docx", "xlsx", "pptx"}

EDIT_PROJECT_OFFICE_DOCUMENT_DESCRIPTION = (
    "Edit an existing Office document (PPTX, DOCX, XLSX) in the project space. "
    "Supports: PPT — update slide title/content, add/delete slides, replace text. "
    "DOCX — replace text, update section content, add paragraphs. "
    "XLSX — update cells/rows, add/delete rows, add sheets. "
    "Use read_project_file first to understand the file structure before editing."
)


@registry.register(
    name=EDIT_PROJECT_OFFICE_DOCUMENT_TOOL_NAME,
    description=tool_description(EDIT_PROJECT_OFFICE_DOCUMENT_TOOL_NAME, EDIT_PROJECT_OFFICE_DOCUMENT_DESCRIPTION),
    input_schema={
        "type": "object",
        "properties": {
            "file_id": {"type": "integer", "description": "Project file ID of the existing document."},
            "file_name": {"type": "string", "description": "Project file name (alternative to file_id)."},
            "file_type": {
                "type": "string",
                "enum": ["docx", "xlsx", "pptx"],
                "description": "File type. Auto-detected from file extension if not provided.",
            },
            "edits": {
                "type": "array",
                "description": "List of edit operations to apply.",
                "items": {"type": "object"},
            },
            "output_name": {
                "type": "string",
                "description": "Output file name. If provided, saves as a new file instead of overwriting.",
            },
        },
        "required": ["file_id", "edits"],
    },
)
async def edit_project_office_document(
    *,
    project_id: int,
    file_id: int | None = None,
    file_name: str | None = None,
    file_type: str | None = None,
    edits: list[dict] | None = None,
    output_name: str | None = None,
) -> dict[str, Any]:
    if not project_id:
        raise HTTPException(400, "Project id is required")
    if not edits:
        raise HTTPException(400, "edits is required")

    with Session(engine) as session:
        project_file = _find_project_file(session, project_id, file_id, file_name)
        source_path = _file_path(project_file)

        detected_type = file_type or project_file.file_type or source_path.suffix.lstrip(".").lower()
        if detected_type not in EDITABLE_TYPES:
            raise HTTPException(400, f"Cannot edit file type '{detected_type}'. Supported: {', '.join(sorted(EDITABLE_TYPES))}")

        if output_name:
            out_name = _safe_name(output_name, detected_type)
            out_path = source_path.parent / out_name
            shutil.copy2(source_path, out_path)
            target_path = out_path
        else:
            target_path = source_path

    if detected_type == "pptx":
        result = _edit_pptx(target_path, edits)
    elif detected_type == "docx":
        result = _edit_docx(target_path, edits)
    elif detected_type == "xlsx":
        result = _edit_xlsx(target_path, edits)
    else:
        result = {"success": False, "error": f"Unsupported type: {detected_type}"}

    if not result.get("success"):
        if output_name and target_path.is_file():
            target_path.unlink()
        raise HTTPException(500, result.get("error") or "Failed to edit document")

    with Session(engine) as session:
        if output_name:
            project_file = _register_generated_project_file(
                session,
                project_id,
                source_path=target_path,
                file_name=output_name,
                file_type=detected_type,
                folder_id=project_file.folder_id,
                summary=f"Edited copy of {project_file.name}",
                preview_text=f"Edited: {', '.join(result.get('changes', []))}",
            )
            output = {
                "ok": True,
                "id": project_file.id,
                "name": project_file.name,
                "file_type": project_file.file_type,
                "changes": result.get("changes", []),
                "message": f"Created edited copy: {project_file.name}",
            }
        else:
            project_file.size_bytes = target_path.stat().st_size
            session.add(project_file)
            session.commit()
            mark_project_memory_stale(session, project_id, trigger=f"{detected_type}_tool_edit")
            _bust_project_cache(project_id)
            output = {
                "ok": True,
                "id": project_file.id,
                "name": project_file.name,
                "file_type": project_file.file_type,
                "changes": result.get("changes", []),
                "message": f"Edited {project_file.name}",
            }

    return output
