"""Document text extraction — PDF, DOCX, PPTX, XLSX, and plain text."""
from __future__ import annotations

from pathlib import Path
from typing import Union


def extract_text(path: Union[str, Path]) -> str:
    p = Path(path)
    suffix = p.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(p)
    elif suffix in (".docx", ".doc"):
        return _extract_docx(p)
    elif suffix == ".pptx":
        return _extract_pptx(p)
    elif suffix in (".xlsx", ".xls"):
        return _extract_xlsx(p)
    elif suffix in (".txt", ".md", ".csv", ".json"):
        return p.read_text(encoding="utf-8", errors="replace")
    else:
        return ""


def _extract_pdf(path: Path) -> str:
    import pdfplumber
    texts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
    return "\n\n".join(texts)


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    rows = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            text = "\t".join(str(c) for c in row if c is not None)
            if text.strip():
                rows.append(text)
    return "\n".join(rows)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            texts.extend(_extract_pptx_shape_text(shape))
        if texts:
            slides.append(f"[Slide {index + 1}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_pptx_shape_text(shape) -> list[str]:
    texts: list[str] = []
    text = getattr(shape, "text", "")
    if text and text.strip():
        texts.append(text.strip())

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                texts.append("\t".join(cells))

    for child in getattr(shape, "shapes", []):
        texts.extend(_extract_pptx_shape_text(child))

    return texts
