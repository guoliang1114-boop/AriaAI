"""Tests for document_text service — text extraction from various file formats."""
import unittest
import tempfile
import os
from pathlib import Path

from app.services.document_text import extract_text_from_file


class ExtractTextFromTxtTestCase(unittest.TestCase):
    def test_extracts_plain_text(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8') as f:
            f.write("Hello world\nSecond line")
            path = Path(f.name)
        try:
            result = extract_text_from_file(path, 'txt')
            self.assertIn('Hello world', result)
            self.assertIn('Second line', result)
        finally:
            os.unlink(path)

    def test_truncates_at_max_chars(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8') as f:
            f.write("A" * 5000)
            path = Path(f.name)
        try:
            result = extract_text_from_file(path, 'txt', max_chars=100)
            self.assertLessEqual(len(result), 150)
            self.assertIn('truncated', result)
        finally:
            os.unlink(path)

    def test_returns_empty_placeholder_for_empty_file(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8') as f:
            f.write("")
            path = Path(f.name)
        try:
            result = extract_text_from_file(path, 'txt', empty_placeholder='[EMPTY]')
            self.assertEqual(result, '[EMPTY]')
        finally:
            os.unlink(path)

    def test_nonexistent_file_returns_not_found(self):
        result = extract_text_from_file(Path('/nonexistent/file.txt'), 'txt')
        self.assertIn('not found', result.lower())

    def test_unsupported_type_returns_placeholder(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.xyz', delete=False) as f:
            f.write("data")
            path = Path(f.name)
        try:
            result = extract_text_from_file(path, 'xyz', unsupported_placeholder='[UNSUPPORTED]')
            self.assertEqual(result, '[UNSUPPORTED]')
        finally:
            os.unlink(path)

    def test_error_prefix_on_exception(self):
        result = extract_text_from_file(Path('/nonexistent'), 'txt', error_prefix='[ERR] ')
        self.assertTrue(result.startswith('[ERR]') or 'not found' in result.lower())


class ExtractTextFromJsonTestCase(unittest.TestCase):
    def test_extracts_json(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8') as f:
            f.write('{"key": "value", "number": 42}')
            path = Path(f.name)
        try:
            result = extract_text_from_file(path, 'json')
            self.assertIn('key', result)
            self.assertIn('value', result)
        finally:
            os.unlink(path)


class ExtractTextFromCsvTestCase(unittest.TestCase):
    def test_extracts_csv(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.csv', delete=False, encoding='utf-8') as f:
            f.write('name,age\nAlice,30\nBob,25')
            path = Path(f.name)
        try:
            result = extract_text_from_file(path, 'csv')
            self.assertIn('Alice', result)
            self.assertIn('Bob', result)
        finally:
            os.unlink(path)


class ExtractTextFromMarkdownTestCase(unittest.TestCase):
    def test_extracts_markdown(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as f:
            f.write('# Title\n\nSome **bold** text')
            path = Path(f.name)
        try:
            result = extract_text_from_file(path, 'md')
            self.assertIn('Title', result)
            self.assertIn('bold', result)
        finally:
            os.unlink(path)


class ExtractTextFromDocxTestCase(unittest.TestCase):
    def test_extracts_docx(self):
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            path = Path(f.name)
        try:
            doc = Document()
            doc.add_paragraph("First paragraph")
            doc.add_paragraph("Second paragraph")
            doc.save(str(path))
            result = extract_text_from_file(path, 'docx')
            self.assertIn('First paragraph', result)
            self.assertIn('Second paragraph', result)
        finally:
            os.unlink(path)


class ExtractTextFromXlsxTestCase(unittest.TestCase):
    def test_extracts_xlsx(self):
        import openpyxl
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
            path = Path(f.name)
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "TestData"
            ws.append(["Name", "Age"])
            ws.append(["Alice", 30])
            wb.save(str(path))
            wb.close()
            result = extract_text_from_file(path, 'xlsx')
            self.assertIn('TestData', result)
            self.assertIn('Alice', result)
        finally:
            os.unlink(path)


class ExtractTextFromPptxTestCase(unittest.TestCase):
    def test_extracts_pptx(self):
        from pptx import Presentation
        from pptx.util import Inches
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            path = Path(f.name)
        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Slide Title"
            prs.save(str(path))
            result = extract_text_from_file(path, 'pptx')
            self.assertIn('Slide Title', result)
        finally:
            os.unlink(path)


class ExtractTextErrorHandlingTestCase(unittest.TestCase):
    def test_error_returns_empty_without_prefix(self):
        result = extract_text_from_file(Path('/nonexistent/file.txt'), 'txt')
        self.assertIn('not found', result.lower())

    def test_error_returns_with_prefix(self):
        result = extract_text_from_file(Path('/nonexistent/file.txt'), 'txt', error_prefix='ERR: ')
        self.assertTrue('not found' in result.lower() or result.startswith('ERR'))


if __name__ == "__main__":
    unittest.main()
