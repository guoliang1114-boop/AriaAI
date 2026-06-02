"""Tests for parser service — extract_text from various file types."""
import unittest
import tempfile
import os
from pathlib import Path

from app.services.parser import extract_text, _extract_docx, _extract_pptx, _extract_xlsx


class ExtractTextTestCase(unittest.TestCase):
    def test_txt_file(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8') as f:
            f.write("Hello world\nSecond line")
            path = f.name
        try:
            result = extract_text(path)
            self.assertIn('Hello world', result)
            self.assertIn('Second line', result)
        finally:
            os.unlink(path)

    def test_md_file_extracts_plain_text(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as f:
            f.write("# Title\n\nParagraph text")
            path = f.name
        try:
            result = extract_text(path)
            self.assertIn("Title", result)
            self.assertIn("Paragraph text", result)
        finally:
            os.unlink(path)

    def test_csv_file_extracts_plain_text(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.csv', delete=False, encoding='utf-8') as f:
            f.write('name,age\nAlice,30\nBob,25')
            path = f.name
        try:
            result = extract_text(path)
            self.assertIn("Alice", result)
        finally:
            os.unlink(path)

    def test_json_file_extracts_plain_text(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8') as f:
            f.write('{"key": "value"}')
            path = f.name
        try:
            result = extract_text(path)
            self.assertIn('"key"', result)
        finally:
            os.unlink(path)

    def test_unsupported_returns_empty(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.xyz', delete=False) as f:
            f.write("data")
            path = f.name
        try:
            result = extract_text(path)
            self.assertEqual(result, "")
        finally:
            os.unlink(path)

    def test_path_object(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8') as f:
            f.write("Path object test")
            path = Path(f.name)
        try:
            result = extract_text(path)
            self.assertIn('Path object test', result)
        finally:
            os.unlink(path)


class ExtractDocxTestCase(unittest.TestCase):
    def test_extract_docx(self):
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            path = f.name
        try:
            doc = Document()
            doc.add_paragraph("First paragraph")
            doc.add_paragraph("Second paragraph")
            doc.add_paragraph("")
            doc.save(path)
            result = _extract_docx(Path(path))
            self.assertIn("First paragraph", result)
            self.assertIn("Second paragraph", result)
        finally:
            os.unlink(path)

    def test_extract_docx_via_extract_text(self):
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            path = f.name
        try:
            doc = Document()
            doc.add_paragraph("Hello docx")
            doc.save(path)
            result = extract_text(path)
            self.assertIn("Hello docx", result)
        finally:
            os.unlink(path)

    def test_extract_doc_suffix(self):
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as f:
            path = f.name
        try:
            doc = Document()
            doc.add_paragraph("Doc format test")
            doc.save(path)
            result = extract_text(path)
            self.assertIn("Doc format test", result)
        finally:
            os.unlink(path)


class ExtractXlsxTestCase(unittest.TestCase):
    def test_extract_xlsx(self):
        import openpyxl
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
            path = f.name
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.append(["Name", "Age"])
            ws.append(["Alice", 30])
            ws.append(["Bob", 25])
            wb.save(path)
            wb.close()
            result = _extract_xlsx(Path(path))
            self.assertIn("Name", result)
            self.assertIn("Alice", result)
            self.assertIn("30", result)
        finally:
            os.unlink(path)

    def test_extract_xlsx_via_extract_text(self):
        import openpyxl
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
            path = f.name
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.append(["Col1", "Col2"])
            ws.append(["Val1", "Val2"])
            wb.save(path)
            wb.close()
            result = extract_text(path)
            self.assertIn("Val1", result)
        finally:
            os.unlink(path)

    def test_xlsx_skips_empty_rows(self):
        import openpyxl
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
            path = f.name
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.append(["A", "B"])
            ws.append([None, None])
            ws.append(["C", "D"])
            wb.save(path)
            wb.close()
            result = _extract_xlsx(Path(path))
            self.assertIn("A", result)
            self.assertIn("C", result)
        finally:
            os.unlink(path)


class ExtractPptxTestCase(unittest.TestCase):
    def test_extract_pptx(self):
        from pptx import Presentation
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            path = f.name
        try:
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Strategy title"
            slide.placeholders[1].text = "PPT body text"
            presentation.save(path)

            result = _extract_pptx(Path(path))
            self.assertIn("Strategy title", result)
            self.assertIn("PPT body text", result)
        finally:
            os.unlink(path)

    def test_extract_pptx_via_extract_text(self):
        from pptx import Presentation
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            path = f.name
        try:
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "Knowledge deck"
            presentation.save(path)

            result = extract_text(path)
            self.assertIn("Knowledge deck", result)
        finally:
            os.unlink(path)

    def test_extract_pptx_table_text(self):
        from pptx import Presentation
        from pptx.util import Inches
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            path = f.name
        try:
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            table = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(6), Inches(1)).table
            table.cell(0, 0).text = "Client"
            table.cell(0, 1).text = "IBM strategy"
            presentation.save(path)

            result = _extract_pptx(Path(path))
            self.assertIn("Client", result)
            self.assertIn("IBM strategy", result)
        finally:
            os.unlink(path)
