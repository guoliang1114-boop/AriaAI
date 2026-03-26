"""File generation tools for Claude Function Calling.

Provides tools to generate PPT, Word, Excel, and PDF files.
"""
from __future__ import annotations

import json
import os
from datetime import datetime
from pathlib import Path
from typing import Any

from app.config import UPLOADS_DIR
from app.tools import registry

# Ensure output directory exists
GENERATED_DIR = UPLOADS_DIR / "generated"
GENERATED_DIR.mkdir(parents=True, exist_ok=True)

# Skills directory for templates
SKILLS_DIR = Path(__file__).parent.parent.parent / "skills"


def _generate_filename(extension: str) -> str:
    """Generate a unique filename with timestamp."""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"generated_{timestamp}.{extension}"


@registry.register(
    name="generate_ppt",
    description="Generate a PowerPoint (.pptx) presentation with structured slides. Supports title slides, content slides, and two-column layouts.",
    input_schema={
        "type": "object",
        "properties": {
            "title": {
                "type": "string",
                "description": "Presentation title"
            },
            "subtitle": {
                "type": "string",
                "description": "Optional subtitle for title slide"
            },
            "slides": {
                "type": "array",
                "description": "Array of slide objects",
                "items": {
                    "type": "object",
                    "properties": {
                        "type": {
                            "type": "string",
                            "enum": ["title", "content", "two_column"],
                            "description": "Slide layout type"
                        },
                        "title": {
                            "type": "string",
                            "description": "Slide title"
                        },
                        "content": {
                            "type": "string",
                            "description": "Main content (bullet points supported with - or *)"
                        },
                        "left_content": {
                            "type": "string",
                            "description": "Left column content (for two_column type)"
                        },
                        "right_content": {
                            "type": "string",
                            "description": "Right column content (for two_column type)"
                        }
                    },
                    "required": ["type", "title"]
                }
            },
            "template_path": {
                "type": "string",
                "description": "Optional path to a .pptx template file. If not provided, uses default blank template."
            }
        },
        "required": ["title", "slides"]
    }
)
async def generate_ppt(
    title: str,
    slides: list[dict],
    subtitle: str = "",
    template_path: str = ""
) -> dict[str, Any]:
    """Generate a PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx not installed. Run: pip install python-pptx"
        }
    
    # Load template if provided, otherwise create blank
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = subtitle
        except KeyError:
            pass
    
    # Content slides
    for slide_data in slides:
        slide_type = slide_data.get("type", "content")
        slide_title = slide_data.get("title", "")
        
        if slide_type == "title":
            layout = prs.slide_layouts[0]
        elif slide_type == "two_column":
            layout = prs.slide_layouts[5]  # Blank for custom layout
        else:
            layout = prs.slide_layouts[1]  # Title and content
        
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_title
        
        if slide_type == "content" and "content" in slide_data:
            try:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = slide_data["content"]
            except KeyError:
                pass
        
        elif slide_type == "two_column":
            # Add text boxes for two-column layout
            left = Inches(1)
            top = Inches(2)
            width = Inches(4)
            height = Inches(5)
            
            left_box = slide.shapes.add_textbox(left, top, width, height)
            left_box.text_frame.text = slide_data.get("left_content", "")
            
            right_box = slide.shapes.add_textbox(Inches(5.5), top, width, height)
            right_box.text_frame.text = slide_data.get("right_content", "")
    
    # Save file
    filename = _generate_filename("pptx")
    filepath = GENERATED_DIR / filename
    prs.save(filepath)
    
    return {
        "success": True,
        "file_type": "pptx",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
        "slide_count": len(slides) + 1,  # +1 for title slide
    }


@registry.register(
    name="generate_ppt_from_skill",
    description="Generate PowerPoint using a Skill's template. Looks for KPMG-Template.pptx or template.pptx in the skill's assets folder.",
    input_schema={
        "type": "object",
        "properties": {
            "skill_name": {
                "type": "string",
                "description": "Skill folder name (e.g., 'ai-strategy-report')"
            },
            "title": {
                "type": "string",
                "description": "Presentation title"
            },
            "subtitle": {
                "type": "string",
                "description": "Optional subtitle"
            },
            "slides": {
                "type": "array",
                "description": "Array of slide objects",
                "items": {
                    "type": "object",
                    "properties": {
                        "type": {"type": "string", "enum": ["title", "content", "two_column"]},
                        "title": {"type": "string"},
                        "content": {"type": "string"},
                        "left_content": {"type": "string"},
                        "right_content": {"type": "string"}
                    },
                    "required": ["type", "title"]
                }
            }
        },
        "required": ["skill_name", "title", "slides"]
    }
)
async def generate_ppt_from_skill(
    skill_name: str,
    title: str,
    slides: list[dict],
    subtitle: str = ""
) -> dict[str, Any]:
    """Generate PPT using a skill's template."""
    # Look for template in skill's assets folder (try KPMG-Template.pptx first, then template.pptx)
    skill_assets = SKILLS_DIR / skill_name / "assets"
    template_path = skill_assets / "KPMG-Template.pptx"
    
    if not template_path.exists():
        template_path = skill_assets / "template.pptx"
    
    if not template_path.exists():
        # Fallback to default generation
        return await generate_ppt(title, slides, subtitle)
    
    # Use skill template
    return await generate_ppt(title, slides, subtitle, str(template_path))


@registry.register(
    name="generate_docx",
    description="Generate a Word (.docx) document with formatted content",
    input_schema={
        "type": "object",
        "properties": {
            "title": {
                "type": "string",
                "description": "Document title"
            },
            "sections": {
                "type": "array",
                "description": "Array of document sections",
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {
                            "type": "string",
                            "description": "Section heading"
                        },
                        "content": {
                            "type": "string",
                            "description": "Section content"
                        },
                        "level": {
                            "type": "integer",
                            "enum": [1, 2, 3],
                            "description": "Heading level"
                        }
                    },
                    "required": ["heading", "content"]
                }
            }
        },
        "required": ["title", "sections"]
    }
)
async def generate_docx(
    title: str,
    sections: list[dict]
) -> dict[str, Any]:
    """Generate a Word document."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return {
            "success": False,
            "error": "python-docx not installed. Run: pip install python-docx"
        }
    
    doc = Document()
    
    # Title
    title_para = doc.add_paragraph()
    title_run = title_para.add_run(title)
    title_run.bold = True
    title_run.font.size = Pt(18)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph()  # Spacing
    
    # Sections
    for section in sections:
        heading = section.get("heading", "")
        content = section.get("content", "")
        level = section.get("level", 1)
        
        # Add heading
        heading_para = doc.add_heading(heading, level=level)
        
        # Add content
        for line in content.split("\n"):
            line = line.strip()
            if line:
                # Check for bold (**text**)
                if line.startswith("**") and line.endswith("**"):
                    p = doc.add_paragraph()
                    run = p.add_run(line[2:-2])
                    run.bold = True
                elif line.startswith("- ") or line.startswith("* "):
                    doc.add_paragraph(line[2:], style='List Bullet')
                else:
                    doc.add_paragraph(line)
    
    # Save file
    filename = _generate_filename("docx")
    filepath = GENERATED_DIR / filename
    doc.save(filepath)
    
    return {
        "success": True,
        "file_type": "docx",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
        "section_count": len(sections),
    }


@registry.register(
    name="generate_xlsx",
    description="Generate an Excel (.xlsx) spreadsheet with data tables",
    input_schema={
        "type": "object",
        "properties": {
            "sheets": {
                "type": "array",
                "description": "Array of sheet objects",
                "items": {
                    "type": "object",
                    "properties": {
                        "name": {
                            "type": "string",
                            "description": "Sheet name"
                        },
                        "headers": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Column headers"
                        },
                        "data": {
                            "type": "array",
                            "description": "Array of row arrays",
                            "items": {
                                "type": "array",
                                "items": {"type": ["string", "number", "boolean"]}
                            }
                        }
                    },
                    "required": ["name", "headers", "data"]
                }
            }
        },
        "required": ["sheets"]
    }
)
async def generate_xlsx(
    sheets: list[dict]
) -> dict[str, Any]:
    """Generate an Excel spreadsheet."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill
    except ImportError:
        return {
            "success": False,
            "error": "openpyxl not installed. Run: pip install openpyxl"
        }
    
    wb = Workbook()
    
    for idx, sheet_data in enumerate(sheets):
        sheet_name = sheet_data.get("name", f"Sheet{idx+1}")
        headers = sheet_data.get("headers", [])
        data = sheet_data.get("data", [])
        
        if idx == 0:
            ws = wb.active
            ws.title = sheet_name
        else:
            ws = wb.create_sheet(title=sheet_name)
        
        # Add headers
        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="B4C7DC", end_color="B4C7DC", fill_type="solid")
        
        # Add data
        for row_idx, row_data in enumerate(data, 2):
            for col_idx, value in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=value)
    
    # Save file
    filename = _generate_filename("xlsx")
    filepath = GENERATED_DIR / filename
    wb.save(filepath)
    
    return {
        "success": True,
        "file_type": "xlsx",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
        "sheet_count": len(sheets),
    }


@registry.register(
    name="generate_pdf",
    description="Generate a PDF document from markdown or HTML content",
    input_schema={
        "type": "object",
        "properties": {
            "title": {
                "type": "string",
                "description": "Document title"
            },
            "content": {
                "type": "string",
                "description": "Document content (Markdown format supported)"
            },
            "orientation": {
                "type": "string",
                "enum": ["portrait", "landscape"],
                "description": "Page orientation"
            }
        },
        "required": ["title", "content"]
    }
)
async def generate_pdf(
    title: str,
    content: str,
    orientation: str = "portrait"
) -> dict[str, Any]:
    """Generate a PDF document."""
    try:
        # Try to use reportlab if available
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
    except ImportError:
        # Fallback: generate markdown file with note
        filename = _generate_filename("md")
        filepath = GENERATED_DIR / filename
        
        md_content = f"# {title}\n\n{content}\n\n---\n*Note: PDF generation requires reportlab. Install with: pip install reportlab*"
        filepath.write_text(md_content, encoding="utf-8")
        
        return {
            "success": True,
            "file_type": "md",
            "file_name": filename,
            "file_path": str(filepath.relative_to(UPLOADS_DIR)),
            "full_path": str(filepath),
            "note": "Generated as Markdown. Install reportlab for true PDF support."
        }
    
    # Generate actual PDF
    filename = _generate_filename("pdf")
    filepath = GENERATED_DIR / filename
    
    page_size = landscape(A4) if orientation == "landscape" else A4
    doc = SimpleDocTemplate(str(filepath), pagesize=page_size)
    
    styles = getSampleStyleSheet()
    story = []
    
    # Title
    story.append(Paragraph(title, styles['Title']))
    story.append(Spacer(1, 12))
    
    # Content
    for line in content.split("\n"):
        if line.strip():
            story.append(Paragraph(line, styles['Normal']))
            story.append(Spacer(1, 6))
    
    doc.build(story)
    
    return {
        "success": True,
        "file_type": "pdf",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
    }


@registry.register(
    name="save_json",
    description="Save structured data as a JSON file",
    input_schema={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Filename without extension"
            },
            "data": {
                "type": "object",
                "description": "JSON-serializable data object"
            }
        },
        "required": ["filename", "data"]
    }
)
async def save_json(
    filename: str,
    data: dict
) -> dict[str, Any]:
    """Save data as JSON file."""
    if not filename.endswith(".json"):
        filename += ".json"
    
    filepath = GENERATED_DIR / filename
    
    with open(filepath, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    
    return {
        "success": True,
        "file_type": "json",
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
    }


@registry.register(
    name="save_text",
    description="Save plain text content as a text file",
    input_schema={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Filename without extension"
            },
            "content": {
                "type": "string",
                "description": "Text content"
            },
            "extension": {
                "type": "string",
                "enum": ["txt", "md", "csv", "html"],
                "description": "File extension"
            }
        },
        "required": ["filename", "content"]
    }
)
async def save_text(
    filename: str,
    content: str,
    extension: str = "txt"
) -> dict[str, Any]:
    """Save text content to file."""
    if not filename.endswith(f".{extension}"):
        filename += f".{extension}"
    
    filepath = GENERATED_DIR / filename
    filepath.write_text(content, encoding="utf-8")
    
    return {
        "success": True,
        "file_type": extension,
        "file_name": filename,
        "file_path": str(filepath.relative_to(UPLOADS_DIR)),
        "full_path": str(filepath),
    }
